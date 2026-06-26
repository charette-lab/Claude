"""Builds TBK_management_case.pptx — positive management-facing pitch (13 slides, 16:9)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# Athanase Industrial Partner brand palette
DARK = RGBColor(0x15, 0x21, 0x30)    # Blue 2 - primary
BLUE3 = RGBColor(0x31, 0x43, 0x59)
BLUE4 = RGBColor(0x55, 0x6A, 0x83)
BLUE5 = RGBColor(0xE2, 0xE5, 0xE9)
GREEN = RGBColor(0x6F, 0x89, 0x90)   # Green 4 - positive accent
GREEN3 = RGBColor(0x33, 0x4E, 0x57)
WARN = RGBColor(0x83, 0x3C, 0x00)    # warning brown
GREY = RGBColor(0x55, 0x6A, 0x83)
LGREY = RGBColor(0xF6, 0xF7, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LOGO_DARK = "/home/user/Claude/athanase_logo_dark.png"
LOGO_WHITE = "/home/user/Claude/athanase_logo_white.png"
HEAD = "Times New Roman"
BODY = "Arial"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def txt(s, l, t, w, h, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT, font=BODY):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = font
    return box


def header(s, title, sub=None, size=27):
    box = s.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(10.7), Inches(0.72))
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = title
    sz = size if len(title) <= 56 else (22 if len(title) <= 72 else 19)
    for r in p.runs:
        r.font.size = Pt(sz)
        r.font.color.rgb = DARK
        r.font.name = HEAD
    s.shapes.add_picture(LOGO_DARK, Inches(11.45), Inches(0.32), width=Inches(1.45))
    rule = s.shapes.add_textbox(Inches(0.5), Inches(0.93), Inches(12.33), Inches(0.02))
    rule.fill.solid(); rule.fill.fore_color.rgb = BLUE5; rule.line.fill.background()
    if sub:
        txt(s, 0.5, 1.0, 12.3, 0.5, sub, size=11.5, color=GREY)


def panel(s, l, t, w, h, title, body, title_color=DARK, body_size=12, fill=LGREY):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.color.rgb = title_color; box.line.width = Pt(1.25)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15); tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]; p.text = title
    for r in p.runs:
        r.font.size = Pt(14.5); r.font.bold = True; r.font.color.rgb = title_color; r.font.name = BODY
    for line in body.split("\n"):
        p = tf.add_paragraph(); p.text = line; p.space_before = Pt(4)
        for r in p.runs:
            r.font.size = Pt(body_size); r.font.color.rgb = DARK; r.font.name = BODY


def stat(s, l, t, w, big, label, color=DARK):
    """Large number + caption block."""
    b = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(0.85))
    tf = b.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = big
    for r in p.runs:
        r.font.size = Pt(34); r.font.bold = True; r.font.color.rgb = color; r.font.name = HEAD
    c = s.shapes.add_textbox(Inches(l), Inches(t + 0.78), Inches(w), Inches(0.8))
    tf = c.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = label
    for r in p.runs:
        r.font.size = Pt(11); r.font.color.rgb = GREY; r.font.name = BODY


def chart_style(ch, title):
    ch.has_legend = False
    ch.has_title = True
    ch.chart_title.text_frame.text = title
    r = ch.chart_title.text_frame.paragraphs[0].runs[0]
    r.font.size = Pt(13); r.font.name = BODY; r.font.bold = True; r.font.color.rgb = DARK
    ch.category_axis.tick_labels.font.size = Pt(10)
    ch.category_axis.tick_labels.font.name = BODY


# ============ 1. TITLE ============
s = slide()
bg = s.shapes.add_textbox(Inches(0), Inches(0), prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = DARK; bg.line.fill.background()
s.shapes.add_picture(LOGO_WHITE, Inches(1.0), Inches(0.7), width=Inches(2.6))
txt(s, 1.0, 2.4, 11.3, 1.8, "We built world-class products.\nNow let's capture their value.", size=42, color=WHITE, font=HEAD)
txt(s, 1.0, 4.6, 11.3, 0.9, "TBK Co., Ltd. — a management case for acting now to secure our next 75 years", size=19, color=BLUE5)
txt(s, 1.0, 6.6, 11.3, 0.5, "Prepared with Athanase Industrial Partner  |  June 2026", size=12, color=RGBColor(0xA8, 0xB8, 0xC8))

# ============ 2. LEGACY ============
s = slide()
header(s, "Start from strength: 75 years of engineering excellence",
       "TBK is not a struggling company — it is an exceptional engineering company that has been too modest about its own value")
stat(s, 0.6, 1.7, 3.0, "75 yrs", "supplying Japan's most demanding truck makers — Isuzu, Fuso, Hino", DARK)
stat(s, 3.85, 1.7, 3.0, "0", "findable recalls or quality scandals in that entire history", GREEN3)
stat(s, 7.1, 1.7, 3.0, "~70-90%", "domestic share in truck water & oil pumps; #1 in CV brakes", DARK)
stat(s, 10.1, 1.7, 2.7, "1960s", "pioneering retarders — decades ahead of the electrification need", GREEN3)
panel(s, 0.6, 3.9, 6.0, 3.05, "Products that lead their class",
      "SAW drum brake — the leading drum brake in stopping distance, energy saving and weight, TODAY\n"
      "Co-authored commercial-vehicle disc-brake research with Isuzu in 1990 (SAE 902200) — before discs were mainstream\n"
      "Electromagnetic & regenerative retarders since the 1960s\n"
      "Electric water/oil pumps and a thermal control unit already in development for the EV era", title_color=GREEN3)
panel(s, 6.85, 3.9, 5.95, 3.05, "The one thing we have never mastered: charging for it",
      "The engineering has never been the problem. For 75 years TBK has built products good enough to lead — and priced them as if they were ordinary.\n"
      "This deck is about one fixable gap: turning world-class engineering into world-class economics.\n"
      "Everything that follows is within our control.", title_color=DARK, fill=BLUE5)
txt(s, 0.5, 7.05, 12.3, 0.4, "The challenge ahead is not technical — it is commercial. And commercial problems are the ones we can fix ourselves.", size=12, bold=True, color=DARK)

# ============ 3. MARKET SHIFTING ============
s = slide()
header(s, "The market is shifting toward exactly what we can build",
       "Two transitions are reshaping our industry — both play to TBK's engineering heritage, if we are ready in time")
panel(s, 0.6, 1.7, 5.95, 2.4, "Brakes: drum → air disc",
      "Europe converted by ~2000; North America crossed 50% in 2023-24; Japan is the last drum market — but the flip is coming\n"
      "UD already all-disc; the new Giga tractor all-disc; Isuzu's export heavies all-disc\n"
      "We co-invented CV disc brakes in 1990 — this is our home turf to reclaim", title_color=GREEN3)
panel(s, 6.7, 1.7, 6.1, 2.4, "Pumps: mechanical → electric thermal",
      "Electrification MULTIPLIES pump content per vehicle 8-30x (battery, motor, inverter, fuel-cell cooling loops)\n"
      "Our pump franchise is on the RIGHT side of this shift — the category grows, it does not shrink\n"
      "Mikuni already ships an e-oil pump for a Japanese BEV truck; we have the e-pump + TCU in development", title_color=GREEN3)
panel(s, 0.6, 4.3, 12.2, 2.5, "This is an opportunity, not a threat — on one condition",
      "Both transitions reward the supplier who arrives first with qualified, homologated products. The window is open now and our customers are beginning to choose their next-generation partners.\n"
      "TBK has the engineering DNA to win both. What we need is the capital to finish the products and the margin to fund them — which today we give away at the negotiating table.\n"
      "The rest of this case shows how to free that capital from what we already own, and turn these two shifts into the foundation of a leading niche player.", title_color=DARK, fill=BLUE5, body_size=12.5)

# ============ 4. COMPETITORS OUT-INVESTING ============
s = slide()
header(s, "Our competitors are investing to win the future — we are not yet",
       "R&D as % of sales, latest year. The companies that will be ready when our customers flip platforms are spending 2-4x what we do")
cd = CategoryChartData()
cd.categories = ["TBK", "Akebono", "Cummins", "Aisin", "Mahle", "Mikuni", "Brembo", "Knorr-Bremse", "ZF"]
cd.add_series("R&D % of sales", (2.2, 4.3, 4.3, 4.8, 5.4, 5.9, 6.1, 6.8, 8.6))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.65), Inches(7.7), Inches(4.5), cd)
ch = gf.chart
chart_style(ch, "R&D intensity (% of sales)")
pl = ch.plots[0]; pl.has_data_labels = True
pl.data_labels.number_format = '0.0"%"'; pl.data_labels.number_format_is_linked = False
pl.data_labels.font.size = Pt(10); pl.data_labels.font.name = BODY
sr = pl.series[0]; sr.format.fill.solid(); sr.format.fill.fore_color.rgb = BLUE4
sr.points[0].format.fill.solid(); sr.points[0].format.fill.fore_color.rgb = WARN
panel(s, 8.4, 1.65, 4.45, 4.5, "What this buys them",
      "Brake systems players: 6-9% — funding disc-brake electronics, brake-by-wire, ADAS integration\n"
      "Pump/thermal players: 4-6% — funding e-pumps and BEV thermal modules\n"
      "Even Akebono, in financial distress, now spends 4.3% to develop electro-mechanical brakes\n"
      "TBK at 2.2% is funding the smallest future in the peer group — while sitting on the engineering to do more", body_size=11.5, title_color=WARN)
txt(s, 0.5, 6.35, 12.3, 0.7, "This is not a criticism of our engineers — it is a capital-allocation choice. We have been spending at half the industry rate, and the gap compounds every year we wait.", size=12, bold=True, color=DARK)

# ============ 5. KNOWLEDGE GAP COMPOUNDS ============
s = slide()
header(s, "The gap compounds: our knowledge base is falling behind peers our size",
       "Capitalized R&D (accumulated know-how), size-class peers. A decade of under-spending has left us with the thinnest knowledge stock")
cd = CategoryChartData()
cd.categories = ["TBK", "Akebono", "SAF-Holland", "Mikuni"]
cd.add_series("Knowledge stock ($M)", (43, 123, 131, 194))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.65), Inches(7.0), Inches(4.5), cd)
ch = gf.chart
chart_style(ch, "Accumulated R&D / knowledge stock, US$M")
pl = ch.plots[0]; pl.has_data_labels = True
pl.data_labels.number_format = '#,##0'; pl.data_labels.number_format_is_linked = False
pl.data_labels.font.size = Pt(11); pl.data_labels.font.name = BODY
sr = pl.series[0]; sr.format.fill.solid(); sr.format.fill.fore_color.rgb = BLUE4
sr.points[0].format.fill.solid(); sr.points[0].format.fill.fore_color.rgb = WARN
panel(s, 7.7, 1.65, 5.15, 2.15, "Mikuni — our size — holds 4.5x our stock",
      "Mikuni is a domestic peer in adjacent products at ~2x our revenue — yet carries 4.5x our knowledge base. Even distressed Akebono holds 3x. This is not a scale excuse; it is a spending choice.", body_size=11.5, title_color=WARN)
panel(s, 7.7, 3.95, 5.15, 2.2, "~$100M cumulative R&D deficit",
      "Over the past decade we have under-invested by roughly $100M versus a normal peer rate. The good news: this is recoverable — and the next pages show the money is already inside the company, waiting to be freed.", body_size=11.5, title_color=GREEN3)
txt(s, 0.5, 6.35, 12.3, 0.7, "Closing this gap is not optional — it is how we stay qualified when our customers choose their next-generation suppliers. The question is only how to fund it.", size=12, bold=True, color=DARK)

# ============ 6. LOYALTY ILLUSION ============
s = slide()
header(s, "The loyalty illusion: our partnerships are more transactional than they feel",
       "The same products earn very different margins depending on who is buying — the keiretsu relationship is, in the numbers, a discount")
cd = CategoryChartData()
cd.categories = ["Keiretsu OEMs\n(blended)", "Hino", "Aftermarket", "Komatsu", "Toho", "Cummins"]
cd.add_series("Contribution margin %", (17.0, 21.0, 30.2, 32.6, 43.5, 79.8))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.65), Inches(7.4), Inches(4.5), cd)
ch = gf.chart
chart_style(ch, "What different customers pay TBK for its products (contribution %)")
pl = ch.plots[0]; pl.has_data_labels = True
pl.data_labels.number_format = '0.0"%"'; pl.data_labels.number_format_is_linked = False
pl.data_labels.font.size = Pt(10); pl.data_labels.font.name = BODY
sr = pl.series[0]; sr.format.fill.solid(); sr.format.fill.fore_color.rgb = GREEN
sr.points[0].format.fill.solid(); sr.points[0].format.fill.fore_color.rgb = WARN
panel(s, 8.1, 1.65, 4.75, 4.5, "We are subsidizing our largest partners",
      "Non-keiretsu buyers pay roughly 2x what our keiretsu partners pay for identical parts\n"
      "TBK's blended gross margin is 12.5% — while Isuzu earns ~19% and Hino ~17%\n"
      "We earn LESS than our own customers — the reverse of a healthy specialist supplier\n"
      "Every yen of discount to a partner is a yen we cannot spend on the products they will need tomorrow", body_size=11.5, title_color=WARN)
txt(s, 0.5, 6.35, 12.3, 0.7, "True partnership is mutual. A partner who values our supply security and our roadmap will accept fair pricing — because they need us to be healthy enough to build their future.", size=12, bold=True, color=DARK)

# ============ 7. THEY WILL SWITCH — UD ============
s = slide()
header(s, "When we fall behind on technology, loyalty does not save us",
       "The evidence is already on the table — we do not have to guess what happens when a keiretsu customer flips platforms")
panel(s, 0.6, 1.75, 5.95, 3.1, "The UD Trucks precedent",
      "UD Trucks flipped to an all-disc-brake platform.\n"
      "TBK was not positioned with a qualified air disc brake.\n"
      "Result: TBK holds ZERO cells at UD today — every brake and pump position is held by a competitor.\n"
      "Loyalty did not protect the business. Readiness would have.", title_color=WARN, body_size=12.5)
panel(s, 6.7, 1.75, 6.1, 3.1, "Now scale that risk to Isuzu and Fuso",
      "Isuzu alone is ~22% of our sales. When Isuzu and Fuso flip to ADB and require e-pumps for their EV programs, they will buy from whoever is ready.\n"
      "If that is Knorr-Bremse or Mikuni instead of TBK, we do not lose a contract — we lose the franchise.\n"
      "The transition is sticky until it happens, then it is abrupt.", title_color=DARK, body_size=12.5, fill=BLUE5)
panel(s, 0.6, 5.0, 12.2, 1.85, "The window is open now — and it will not stay open",
      "Our customers are choosing their next-generation partners over the next few product cycles. The supplier who arrives first with a qualified, homologated ADB and e-pump wins the platform for a decade. We have the engineering to be that supplier — but only if we fund it now, while we still hold the incumbency that gets us in the room.", title_color=GREEN3, body_size=12.5)

# ============ 8. THE GOOD NEWS — VALUE IS OURS ============
s = slide()
header(s, "The good news: the value is already ours — we simply give it away",
       "We do not need to invent anything new to start. Our best products already earn world-class margins when we price them properly")
cd = CategoryChartData()
cd.categories = ["SAW today\n(keiretsu price)", "SAW at\nmarket price", "Peer core\nproducts"]
cd.add_series("Contribution margin %", (19.5, 41.0, 42.0))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.65), Inches(6.6), Inches(4.5), cd)
ch = gf.chart
chart_style(ch, "Our flagship SAW drum brake: contribution margin %")
pl = ch.plots[0]; pl.has_data_labels = True
pl.data_labels.number_format = '0.0"%"'; pl.data_labels.number_format_is_linked = False
pl.data_labels.font.size = Pt(11); pl.data_labels.font.name = BODY
sr = pl.series[0]; sr.format.fill.solid(); sr.format.fill.fore_color.rgb = GREEN
sr.points[0].format.fill.solid(); sr.points[0].format.fill.fore_color.rgb = WARN
sr.points[1].format.fill.solid(); sr.points[1].format.fill.fore_color.rgb = GREEN3
panel(s, 7.4, 1.65, 5.45, 4.5, "Our best product, priced ~30% below market",
      "SAW leads its class in stopping, energy and weight — yet earns just 19.5% contribution because we sell it ~30% below market\n"
      "At market price it earns 38-44% — exactly what the world's best brake and pump makers earn on their core products\n"
      "Same product. Same factory. Same engineering. Only the price changes\n"
      "We own the design (confirmed) — the value is ours to capture, in Japan, in exports, and in the aftermarket", body_size=11.5, title_color=GREEN3)
txt(s, 0.5, 6.35, 12.3, 0.7, "This is the most hopeful number in the whole analysis: our path to peer-level economics runs through products we have ALREADY built.", size=12, bold=True, color=DARK)

# ============ 9. LEVER 1 — PRICE ============
s = slide()
header(s, "Lever 1 — Price to the value we create",
       "Not a confrontation — a long-overdue correction, paired with something our customers want even more than a discount")
panel(s, 0.6, 1.75, 4.0, 4.6, "What we do",
      "Reprice toward market over contract-renewal cycles, starting where we are sole-source\n\n"
      "Raise aftermarket prices now (+10%) — a market-priced channel we already control\n\n"
      "Sell SAW-class products at market price into exports and India via the Brakes India alliance\n\n"
      "Stop or reprice the loss-making lines we should never have accepted", title_color=DARK, body_size=12)
panel(s, 4.7, 1.75, 4.0, 4.6, "How we keep it a partnership",
      "Pair every price increase with our committed ADB / e-pump technology roadmap\n\n"
      "Offer supply security and future technology in exchange for fair pricing — not price for nothing\n\n"
      "Frame it honestly: we need fair margins to build the products you will require tomorrow\n\n"
      "A healthy TBK is in our customers' own interest", title_color=GREEN3, body_size=12, fill=BLUE5)
panel(s, 8.8, 1.75, 4.0, 4.6, "What it is worth",
      "Closing the margin gap to peer levels: +$22-24M gross profit per year\n\n"
      "Phase 1 needs only ~1/3 of the full pricing gap — a deliberately modest, achievable ask\n\n"
      "Every point of margin flows almost entirely to the bottom line and to R&D", title_color=GREEN3, body_size=12)
txt(s, 0.5, 6.5, 12.3, 0.5, "We are not asking customers to overpay. We are asking them to stop underpaying for products that lead their class.", size=12, bold=True, color=DARK)

# ============ 10. LEVER 2 — FOOTPRINT ============
s = slide()
header(s, "Lever 2 — Right-size the factory footprint",
       "The hidden burden is not in our overheads — it is in our underused factories, and fixing it frees cash immediately")
stat(s, 0.6, 1.7, 3.0, "39%", "of sales tied up in property, plant & equipment — vs ~12% at a profitable peer (Concentric)", WARN)
stat(s, 3.85, 1.7, 3.0, "5.3%", "of every sales yen consumed by depreciation alone", WARN)
stat(s, 7.1, 1.7, 3.0, "2 yrs", "of losses in our China segment — depreciation there is 12.8% of segment sales", WARN)
stat(s, 10.1, 1.7, 2.7, "¥2.2bn+", "of impairments & restructuring already absorbed in two years", WARN)
panel(s, 0.6, 3.9, 6.0, 2.95, "Our OPEX looks lean — but that is the illusion",
      "Reported operating expense is only 9.7% of sales, so the company looks efficient. It is not where the burden sits.\n"
      "Under our accounting, factory overhead, depreciation and idle capacity sit in COST OF SALES — which is exactly why our gross margin is half our peers'.\n"
      "The cost problem and the margin problem are the same problem.", title_color=WARN, body_size=12)
panel(s, 6.85, 3.9, 5.95, 2.95, "Consolidation frees cash now",
      "Reorganizing and consolidating the underloaded footprint, fixing or exiting China, and monetizing dormant assets frees cash in the near term.\n"
      "This is the immediate funding lever — it does not wait for new products or new prices.\n"
      "FY3/2026 already proved it works: we cut cost of sales ¥0.7bn and lifted margin from 10.6% to 12.5% in a single year.", title_color=GREEN3, body_size=12, fill=BLUE5)
txt(s, 0.5, 7.0, 12.3, 0.4, "We have started. The turn in last year's numbers shows the plan works — now we accelerate it.", size=12, bold=True, color=DARK)

# ============ 11. THE FLYWHEEL ============
s = slide()
header(s, "The self-funding flywheel: no outside capital required",
       "Each lever funds the next — this is an offensive growth strategy, not a defensive cost exercise")
steps = [
    ("1. Fix the footprint", "Consolidate underused plants, fix China, monetize dormant assets — frees cash immediately", WARN),
    ("2. Price to value", "Reprice toward market + aftermarket — +$22-24M gross profit per year", GREEN3),
    ("3. Fund the future", "Self-fund the ~$27M/yr R&D catch-up — ADB, e-pumps, thermal — no outside capital", GREEN3),
    ("4. Win the sockets", "Arrive first with qualified ADB & e-pumps — secure Isuzu/Fuso/Hino platforms for a decade", GREEN),
    ("5. Earn premium margins", "Class-leading products, priced fairly, on the right side of electrification", GREEN),
]
x = 0.6
w = 2.36
for i, (t, b, c) in enumerate(steps):
    panel(s, x, 1.9, w, 3.7, t, "\n" + b, title_color=c, body_size=11)
    if i < len(steps) - 1:
        a = s.shapes.add_textbox(Inches(x + w - 0.02), Inches(3.4), Inches(0.30), Inches(0.5))
        p = a.text_frame.paragraphs[0]; p.text = "→"
        for r in p.runs:
            r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = BLUE4; r.font.name = BODY
    x += w + 0.12
panel(s, 0.6, 5.9, 12.23, 1.0, "And it compounds back to the start",
      "Premium margins generate more cash, which funds more R&D, which wins more sockets — the same flywheel that turned Concentric (the former Haldex hydraulics arm) into a 25%-gross-margin niche leader. We are not inventing a new playbook; we are running a proven one with products we already own.", title_color=DARK, body_size=12, fill=BLUE5)

# ============ 12. DESTINATION ============
s = slide()
header(s, "The destination: a leading niche player again",
       "Where these two levers take us — operating margin to peer levels, on a credible 18-month path")
cd = CategoryChartData()
cd.categories = ["Today", "18 months\n(Phase 1)", "Steady state\n(Phase 2)"]
cd.add_series("Operating margin %", (2.7, 10.6, 15.0))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.65), Inches(6.6), Inches(4.5), cd)
ch = gf.chart
chart_style(ch, "Operating margin trajectory (%)")
pl = ch.plots[0]; pl.has_data_labels = True
pl.data_labels.number_format = '0.0"%"'; pl.data_labels.number_format_is_linked = False
pl.data_labels.font.size = Pt(12); pl.data_labels.font.name = BODY
sr = pl.series[0]; sr.format.fill.solid(); sr.format.fill.fore_color.rgb = BLUE4
sr.points[0].format.fill.solid(); sr.points[0].format.fill.fore_color.rgb = WARN
sr.points[1].format.fill.solid(); sr.points[1].format.fill.fore_color.rgb = GREEN3
sr.points[2].format.fill.solid(); sr.points[2].format.fill.fore_color.rgb = GREEN
panel(s, 7.4, 1.65, 5.45, 2.15, "Operating profit: ¥1.5bn → ¥6.3bn",
      "Phase 1 (18 months): ~10% operating margin, captured from our existing book — no new customers, no new products required, full R&D already absorbed in the number.", body_size=11.5, title_color=GREEN3)
panel(s, 7.4, 3.95, 5.45, 2.2, "From survivor to compounder",
      "On the moat framework, this moves TBK from 5.4 toward the 6.5+ 'Watchlist' band — the level of Concentric (6.98), the niche-leader benchmark. A focused, high-margin, electrification-ready specialist: what our engineering always deserved.", body_size=11.5, title_color=GREEN3)
txt(s, 0.5, 6.35, 12.3, 0.7, "This is not a rescue. It is TBK becoming the company its engineering has always been capable of being.", size=12.5, bold=True, color=DARK)

# ============ 13. ACT NOW ============
s = slide()
bg = s.shapes.add_textbox(Inches(0), Inches(0), prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = DARK; bg.line.fill.background()
s.shapes.add_picture(LOGO_WHITE, Inches(11.0), Inches(0.55), width=Inches(1.7))
txt(s, 0.9, 0.7, 9.5, 0.8, "Why now — and what we start this quarter", size=26, color=WHITE, font=HEAD)
panel(s, 0.9, 1.9, 5.7, 4.5, "Why the clock matters",
      "Every year at 2.2% R&D, the knowledge gap compounds\n\n"
      "Our customers are choosing their next-generation partners NOW — the UD outcome is what waiting looks like\n\n"
      "The funding is already inside the company — under-priced products and under-used factories — it costs us nothing to start, and everything to wait", title_color=WHITE, body_size=12.5, fill=BLUE3)
panel(s, 6.75, 1.9, 5.65, 4.5, "The first 180 days",
      "Days 0-30: lock the target; map every contract; assemble the SAW pricing evidence\n\n"
      "Days 31-60: aftermarket +10% now; open SAW repricing with Isuzu & Fuso paired with our roadmap; approve the footprint plan\n\n"
      "Days 61-180: first repricings land; footprint consolidation begins; R&D ramps; first ADB/e-pump sockets pursued\n\n"
      "Day 180 goal: operating profit run-rate already doubled", title_color=WHITE, body_size=12, fill=BLUE3)
txt(s, 0.9, 6.65, 11.5, 0.6, "We have the products. We have the customers. We have the cash inside the business. What we need is the decision to act.", size=14, bold=True, color=WHITE)

import warnings
with warnings.catch_warnings():
    warnings.simplefilter("error")
    prs.save("/home/user/Claude/TBK_management_case.pptx")
print("saved", len(prs.slides._sldIdLst), "slides")

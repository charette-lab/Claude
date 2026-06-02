"""Build a single, uniform Athanase 'Case studies' deck (4:3) from the four
source files the client supplied (Aug-2025, June-2018, May-2018 and the older
Investment AB Oresund deck).

All cases are re-laid-out in one brand-compliant template:
  - left column : Company description / Recognized opportunities /
                  Initial plan / Results
  - right column: an IRR / money-multiple (or period) badge, the reused
                  share-price chart (cropped from the source where one exists)
                  and an on-brand bar chart redrawn natively from the numbers.

Brand: AIP Blue group only, Times New Roman headlines + Arial body, left-align.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.oxml.ns import qn
from PIL import Image

BASE = "/home/user/Claude"
CHARTS = os.path.join(BASE, "case_src", "charts")
ASSETS = os.path.join(BASE, "assets")

# ---- Brand palette (AIP Brand Guidelines — BLUE group only) ----------------
BLUE1 = RGBColor(0x0E, 0x16, 0x20)
BLUE2 = RGBColor(0x15, 0x21, 0x30)
BLUE3 = RGBColor(0x31, 0x43, 0x59)
BLUE4 = RGBColor(0x55, 0x6A, 0x83)
BLUE5 = RGBColor(0xE2, 0xE5, 0xE9)
BLUE6 = RGBColor(0xF6, 0xF7, 0xF9)
NAVY, NAVY_TX, SLATE = BLUE2, BLUE3, BLUE3
SLATE_LT, HEADERBG, BODY = BLUE4, BLUE6, BLUE1
SUBTLE, DIVIDER, WHITE = BLUE4, BLUE5, RGBColor(0xFF, 0xFF, 0xFF)
FOOT = BLUE4
SERIF, SANS = "Times New Roman", "Arial"

LOGO_WHITE = os.path.join(ASSETS, "logo_white.png")
MARK_DARK = os.path.join(ASSETS, "mark_dark.png")
_LW_AR = (lambda s: s[0] / s[1])(Image.open(LOGO_WHITE).size)
_MD_AR = (lambda s: s[0] / s[1])(Image.open(MARK_DARK).size)

prs = Presentation()
prs.slide_width = Inches(10.667)      # 4:3
prs.slide_height = Inches(8.0)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
_state = {"n": 0}


def rect(slide, l, t, w, h, fill=None, line=None, line_w=1.0, rounded=False):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape, l, t, w, h)
    if fill is not None:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line is not None:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    # remove themed style ref (kills inherited drop-shadow in LibreOffice)
    el = shp._element
    style = el.find(qn("p:style"))
    if style is not None:
        el.remove(style)
    return shp


def tbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    return tf


def para(tf, text, size, color, bold=False, italic=False, first=False,
         align=PP_ALIGN.LEFT, after=8, lead=1.1, font=SANS, track=None,
         before=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(after)
    if before is not None:
        p.space_before = Pt(before)
    p.line_spacing = lead
    r = p.add_run(); r.text = text
    f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = font
    if track is None:
        track = -5 if size >= 48 else (-3 if size >= 24 else 0)
    if track:
        r._r.get_or_add_rPr().set("spc", str(int(round(size * track))))
    return p


def place_mark(slide, left, top, height):
    slide.shapes.add_picture(MARK_DARK, left, top, height=height,
                             width=Emu(int(int(height) * _MD_AR)))


def place_logo_white(slide, left, top, height):
    slide.shapes.add_picture(LOGO_WHITE, left, top, height=height,
                             width=Emu(int(int(height) * _LW_AR)))


def footer(slide, conf=True):
    _state["n"] += 1
    tf = tbox(slide, Inches(6.0), Inches(7.62), Inches(4.4), Inches(0.3))
    txt = f"Strictly confidential          {_state['n']}" if conf \
        else f"{_state['n']}"
    para(tf, txt, 8.5, FOOT, first=True, align=PP_ALIGN.RIGHT, after=0)


def picture_fit(slide, path, x, y, max_w, max_h, align="left", top=True):
    """Place an image scaled to fit within (max_w, max_h), preserving ratio."""
    iw, ih = Image.open(path).size
    ar = iw / ih
    w = max_w
    h = Emu(int(int(w) / ar))
    if int(h) > int(max_h):
        h = max_h
        w = Emu(int(int(h) * ar))
    if align == "center":
        x = Emu(int(x) + (int(max_w) - int(w)) // 2)
    if not top:
        y = Emu(int(y) + (int(max_h) - int(h)) // 2)
    slide.shapes.add_picture(path, x, y, width=w, height=h)
    return Emu(int(y) + int(h))


# ===========================================================================
#  Standard chrome
# ===========================================================================
def case_header(slide, label, title):
    rect(slide, 0, 0, SW, SH, fill=WHITE)
    place_mark(slide, Inches(0.5), Inches(0.24), Inches(0.26))
    lt = tbox(slide, Inches(0.92), Inches(0.27), Inches(6.0), Inches(0.3))
    para(lt, label, 11, SLATE_LT, first=True, after=0)
    rect(slide, 0, Inches(0.62), SW, Inches(1.0), fill=HEADERBG)
    tt = tbox(slide, Inches(0.55), Inches(0.78), Inches(9.6), Inches(0.7))
    para(tt, title, 26, NAVY_TX, first=True, after=0, font=SERIF)
    footer(slide)


def badge(slide, x, y, w, text):
    h = Inches(0.42)
    rect(slide, x, y, w, h, fill=HEADERBG, rounded=True)
    tf = tbox(slide, x, y, w, h, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, text, 12.5, SLATE, bold=True, first=True,
         align=PP_ALIGN.CENTER, after=0, font=SERIF)
    return Emu(int(y) + int(h))


def value_bar(slide, x, y, w, h, cats, vals, unit="SEK M",
              title="Initial investment vs. current value", vfmt=None):
    """Hand-drawn two-bar chart (reliable labels across renderers)."""
    if vfmt is None:
        vfmt = lambda v: f"{v:,}"
    tt = tbox(slide, x, y, w, Inches(0.3))
    para(tt, title, 11, NAVY_TX, bold=True, first=True,
         align=PP_ALIGN.CENTER, after=0, font=SERIF)
    # plot region
    px = Emu(int(x) + int(Inches(0.1)))
    pw = Emu(int(w) - int(Inches(0.2)))
    p_top = Emu(int(y) + int(Inches(0.55)))
    base = Emu(int(y) + int(h) - int(Inches(0.45)))   # baseline (cat labels)
    plot_h = int(base) - int(p_top)
    scale = max(vals) * 1.18
    n = len(vals)
    slot = int(pw) // n
    bw = int(Inches(1.15))
    # baseline rule
    rect(slide, px, base, pw, Pt(1.0), fill=BLUE5)
    for i, (c, v) in enumerate(zip(cats, vals)):
        bh = int(plot_h * (v / scale))
        bx = Emu(int(px) + i * slot + (slot - bw) // 2)
        by = Emu(int(base) - bh)
        rect(slide, bx, by, Emu(bw), Emu(bh), fill=NAVY)
        vl = tbox(slide, Emu(int(bx) - int(Inches(0.3))),
                  Emu(int(by) - int(Inches(0.3))),
                  Emu(bw + int(Inches(0.6))), Inches(0.28))
        para(vl, vfmt(v), 13, NAVY_TX, bold=True, first=True,
             align=PP_ALIGN.CENTER, after=0, font=SERIF)
        cl = tbox(slide, Emu(int(px) + i * slot),
                  Emu(int(base) + int(Inches(0.06))), Emu(slot), Inches(0.34))
        para(cl, c, 9.5, BODY, first=True, align=PP_ALIGN.CENTER, after=0)
    cap = tbox(slide, x, Emu(int(base) + int(Inches(0.42))), w, Inches(0.22))
    para(cap, unit, 8.5, SUBTLE, italic=True, first=True,
         align=PP_ALIGN.LEFT, after=0)


def opbar(slide, x, y, w, h, groups, title="Operational improvements over "
          "the investment period"):
    """groups: list of (name, [(cat, pct), ...]). Draws clustered columns."""
    tt = tbox(slide, x, y, w, Inches(0.3))
    para(tt, title, 11, NAVY_TX, bold=True, first=True,
         align=PP_ALIGN.CENTER, after=0, font=SERIF)
    cy = Emu(int(y) + int(Inches(0.34)))
    ch = Emu(int(h) - int(Inches(0.34)))
    gw = Emu((int(w) - int(Inches(0.2))) // len(groups))
    for i, (name, pairs) in enumerate(groups):
        gx = Emu(int(x) + i * int(gw) + (int(Inches(0.2)) if i else 0))
        cd = CategoryChartData()
        cd.categories = [c for c, _ in pairs]
        cd.add_series(name, [v for _, v in pairs])
        gframe = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                                        gx, cy, gw, ch, cd)
        chart = gframe.chart
        chart.has_legend = False
        chart.has_title = True
        chart.chart_title.text_frame.text = name
        r = chart.chart_title.text_frame.paragraphs[0].runs[0]
        r.font.size = Pt(10); r.font.bold = True
        r.font.color.rgb = SLATE; r.font.name = SANS
        plot = chart.plots[0]
        plot.gap_width = 60
        plot.vary_by_categories = False
        s = plot.series[0]
        s.format.fill.solid(); s.format.fill.fore_color.rgb = NAVY
        s.format.line.fill.background()
        s.has_data_labels = True
        dl = s.data_labels
        dl.number_format = '0"%"'; dl.number_format_is_linked = False
        dl.position = XL_LABEL_POSITION.OUTSIDE_END
        dl.font.size = Pt(10); dl.font.bold = True
        dl.font.color.rgb = NAVY_TX; dl.font.name = SANS
        ca = chart.category_axis
        ca.tick_labels.font.size = Pt(8); ca.tick_labels.font.color.rgb = BODY
        ca.tick_labels.font.name = SANS
        ca.format.line.color.rgb = BLUE5
        va = chart.value_axis
        va.visible = False; va.has_major_gridlines = False
        va.maximum_scale = max(v for _, v in pairs) * 1.3


# ---- left content column ---------------------------------------------------
def left_column(slide, sections, x=Inches(0.55), y=Inches(1.95),
                w=Inches(4.55), body_sz=8.7, label_sz=11):
    tf = tbox(slide, x, y, w, Inches(5.6))
    first = True
    for kind, label, content in sections:
        p = para(tf, label, label_sz, SLATE, bold=True, first=first,
                 after=2.5, font=SERIF, before=(0 if first else 7))
        first = False
        if kind == "para":
            para(tf, content, body_sz, BODY, after=2, lead=1.08)
        elif kind in ("bullets", "checks"):
            marker = "›  " if kind == "bullets" else "✓  "
            hang = int(round(body_sz * 1.35 * 12700))
            for item in content:
                pp = tf.add_paragraph()
                pp.space_after = Pt(1.5); pp.line_spacing = 1.06
                pPr = pp._p.get_or_add_pPr()
                pPr.set("marL", str(hang)); pPr.set("indent", str(-hang))
                r0 = pp.add_run(); r0.text = marker
                r0.font.size = Pt(body_sz); r0.font.bold = True
                r0.font.color.rgb = SLATE_LT; r0.font.name = SANS
                r1 = pp.add_run(); r1.text = item
                r1.font.size = Pt(body_sz); r1.font.color.rgb = BODY
                r1.font.name = SANS
    return tf


def note_line(slide, text):
    tf = tbox(slide, Inches(0.55), Inches(7.62), Inches(5.3), Inches(0.3))
    para(tf, text, 7.5, SUBTLE, italic=True, first=True, after=0, lead=1.0)


# ===========================================================================
#  Front matter
# ===========================================================================
def cover():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, fill=NAVY)
    place_logo_white(s, Inches(0.6), Inches(0.55), Inches(0.62))
    rect(s, Inches(0.62), Inches(5.35), Inches(0.5), Pt(2.2), fill=SLATE_LT)
    tf = tbox(s, Inches(0.6), Inches(5.55), Inches(9.0), Inches(1.4))
    para(tf, "Case studies", 46, WHITE, first=True, after=2, font=SERIF)
    para(tf, "Selected investments across three decades of engaged ownership",
         15, DIVIDER, italic=True, after=0, font=SERIF)


def disclaimer():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, fill=WHITE)
    place_mark(s, Inches(0.5), Inches(0.4), Inches(0.3))
    tt = tbox(s, Inches(0.55), Inches(0.95), Inches(9.0), Inches(0.7))
    para(tt, "Disclaimer", 30, NAVY_TX, first=True, after=0, font=SERIF)
    col1 = [
        "This document has been prepared by Athanase Industrial Partner "
        "(‘AIP’) solely for the information of the person to whom it has "
        "been delivered. The information contained herein is strictly "
        "confidential and is only for the use of the person to whom it is "
        "sent. The information contained herein may not be reproduced, "
        "distributed or published by any recipient for any purpose without "
        "the prior written consent of AIP. The information herein is for "
        "general guidance only, and it is the responsibility of any person "
        "or persons in possession of this document to inform themselves of, "
        "and to observe, all applicable laws and regulations of any relevant "
        "jurisdiction.",
        "The information regarding Athanase Industrial Partners Fund II (the "
        "“Fund”) is not intended to provide and should not be relied upon "
        "for accounting, legal or tax advice or investment recommendations. "
        "You should consult your tax, legal, accounting or other advisors "
        "about the issues discussed herein. Material terms of the Fund are "
        "subject to change. Any prospective investor will be provided with a "
        "copy of the Fund’s offering memorandum and an opportunity to review "
        "the documentation relating to the offering.",
        "This document is not intended as an offer or solicitation with "
        "respect to the purchase or sale of any security. This document is "
        "not intended for distribution to, or use by, any person or entity "
        "in any jurisdiction or country where such distribution or use would "
        "be contrary to local law or regulation.",
    ]
    col2 = [
        "Any subscription may only be made on the terms of the offering "
        "memorandum and subject to completion of a subscription agreement. "
        "The Fund will not be registered under the Securities Act or the "
        "securities laws of any of the states of the United States and "
        "interests therein may not be offered, sold or delivered directly or "
        "indirectly into the United States, or to or for the account or "
        "benefit of any US person.",
        "No reliance may be placed for any purpose on the information and "
        "opinions contained in this document or their accuracy or "
        "completeness. No representation, warranty or undertaking, expressed "
        "or implied, is given as to the accuracy or completeness of the "
        "information or opinions contained in this document by any of AIP, "
        "its officers, employees or affiliates, and nothing contained herein "
        "shall be relied upon as a promise or representation whether as to "
        "past or future performance. Past performance is not indicative of "
        "future results.",
        "Notice to Swedish Residents — The Fund is authorized by the "
        "Finansinspektionen (‘FI’) in accordance with the Law on Managers of "
        "Alternative Investment Funds (Sw. Lag om Förvaltare av Alternativa "
        "Investeringsfonder, LAIF (2013:561)) for marketing to investors in "
        "Sweden that qualify as Professional Investors under the Swedish LAIF "
        "(2013:561).",
    ]
    for cx, col in ((Inches(0.55), col1), (Inches(5.55), col2)):
        tf = tbox(s, cx, Inches(1.95), Inches(4.55), Inches(5.4))
        for i, t in enumerate(col):
            para(tf, t, 8.5, BODY, first=(i == 0), after=7, lead=1.12)
    footer(s)


def overview():
    s = prs.slides.add_slide(BLANK)
    case_header(s, "Overview", "Selected case studies")
    sub = tbox(s, Inches(0.55), Inches(1.78), Inches(9.6), Inches(0.4))
    para(sub, "Transactions where the team acted as an engaged owner and "
         "implemented operational, structural and financial change. IRR is "
         "unleveraged and net of all costs.", 11, SUBTLE, italic=True,
         first=True, after=0)

    def mini_table(x, w, header, rows, cols_hdr, col_w):
        ht = tbox(s, x, Inches(2.35), w, Inches(0.3))
        para(ht, header, 13, SLATE, bold=True, first=True, after=0, font=SERIF)
        y = Inches(2.72)
        rh = Inches(0.3)
        cx = x
        for ci, h in enumerate(cols_hdr):
            rect(s, cx, y, col_w[ci], rh, fill=SLATE)
            tf = tbox(s, Emu(int(cx) + int(Inches(0.06))), y,
                      Emu(int(col_w[ci]) - int(Inches(0.1))), rh,
                      anchor=MSO_ANCHOR.MIDDLE)
            para(tf, h, 9, WHITE, bold=True, first=True, after=0,
                 align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT)
            cx = Emu(int(cx) + int(col_w[ci]))
        yy = Emu(int(y) + int(rh))
        for ri, row in enumerate(rows):
            fill = HEADERBG if ri % 2 == 0 else WHITE
            cx = x
            for ci, val in enumerate(row):
                rect(s, cx, yy, col_w[ci], rh, fill=fill)
                tf = tbox(s, Emu(int(cx) + int(Inches(0.06))), yy,
                          Emu(int(col_w[ci]) - int(Inches(0.1))), rh,
                          anchor=MSO_ANCHOR.MIDDLE)
                para(tf, val, 8.7, NAVY_TX if ci == 0 else BODY,
                     bold=(ci == 0), first=True, after=0,
                     align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT)
                cx = Emu(int(cx) + int(col_w[ci]))
            yy = Emu(int(yy) + int(rh))

    fund_rows = [
        ["Haldex AB", "64%", "4.1×"],
        ["Haldex AB (2nd)", "100%", "2.0×"],
        ["Concentric AB", "81%", "3.9×"],
        ["Kitron ASA", "49%", "1.6×"],
        ["Note AB", "33%", "5.1×"],
        ["Lindab AB", "32%", "1.6×"],
        ["Transcom SA", "30%", "3.0×"],
        ["DistIt / Alcadon", "—", "—"],
        ["Renold plc", "234%", "1.6×"],
        ["Robit Plc (exited)", "14%", "—"],
    ]
    ores_rows = [
        ["J&W", "94.4%", "2000–01"],
        ["Custos", "58.9%", "2000–04"],
        ["HQ / Avanza", "37.8%", "2000–"],
        ["Johnson Pump", "34.4%", "2002–06"],
        ["Nobia", "33.1%", "2002–"],
        ["SkiStar", "33.0%", "2000–"],
        ["Wihlborgs / Drott", "31.0%", "2002–04"],
        ["Bilia", "26.0%", "2003–"],
    ]
    mini_table(Inches(0.55), Inches(4.55),
               "AIP — fund-era investment cases", fund_rows,
               ["Company", "IRR", "Money multiple"],
               [Inches(2.45), Inches(1.05), Inches(1.05)])
    mini_table(Inches(5.55), Inches(4.55),
               "Earlier track record — Investment AB Öresund", ores_rows,
               ["Company", "IRR", "Period"],
               [Inches(2.45), Inches(1.05), Inches(1.05)])
    note_line(s, "Past performance is not indicative of future results. "
              "Öresund-era cases were executed by members of the team at "
              "Investment AB Öresund (1994–2007).")


def era_divider(kicker, title):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, fill=NAVY)
    place_logo_white(s, Inches(0.6), Inches(0.55), Inches(0.55))
    rect(s, Inches(0.62), Inches(3.05), Inches(0.5), Pt(2.2), fill=SLATE_LT)
    kt = tbox(s, Inches(0.6), Inches(3.25), Inches(9.5), Inches(0.5))
    para(kt, kicker.upper(), 13, SLATE_LT, first=True, after=0)
    tt = tbox(s, Inches(0.6), Inches(3.6), Inches(9.5), Inches(1.4))
    para(tt, title, 34, DIVIDER, italic=True, first=True, after=0, font=SERIF)


def contact():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, fill=NAVY)
    place_logo_white(s, Inches(0.6), Inches(0.55), Inches(0.62))
    tf = tbox(s, Inches(0.6), Inches(5.7), Inches(7.0), Inches(2.0))
    para(tf, "Kenth Eriksson", 16, WHITE, first=True, after=1)
    para(tf, "Senior Portfolio Manager & Partner", 12, DIVIDER, after=12)
    para(tf, "+46 70 304 09 29", 12, DIVIDER, after=1)
    para(tf, "kenth.eriksson@athanase.se", 12, DIVIDER, after=12)
    para(tf, "info@athanase.se", 12, DIVIDER, after=0)


# ===========================================================================
#  Case renderers
# ===========================================================================
def case_image_right(sections, label, title, badge_text, chart_img,
                     note=None, value_note=None, chart_img2=None):
    """Right column = badge + one or two reused (cropped) chart images."""
    s = prs.slides.add_slide(BLANK)
    case_header(s, label, title)
    left_column(s, sections)
    rx, rw = Inches(5.45), Inches(4.7)
    y = Inches(1.95)
    if badge_text:
        y = badge(s, rx, y, rw, badge_text)
        y = Emu(int(y) + int(Inches(0.12)))
    avail_bottom = Inches(7.5)
    if chart_img2:
        half = Emu((int(avail_bottom) - int(y) - int(Inches(0.1))) // 2)
        y2 = picture_fit(s, chart_img, rx, y, rw, half, align="center")
        picture_fit(s, chart_img2, rx, Emu(int(y2) + int(Inches(0.1))),
                    rw, half, align="center")
    else:
        yb = picture_fit(s, chart_img, rx, y, rw,
                         Emu(int(avail_bottom) - int(y)), align="center")
        if value_note:
            vt = tbox(s, rx, Emu(int(yb) + int(Inches(0.1))), rw, Inches(0.5))
            para(vt, value_note, 12, SLATE, bold=True, first=True,
                 align=PP_ALIGN.CENTER, after=0, font=SERIF)
    if note:
        note_line(s, note)
    return s


def case_native_right(sections, label, title, badge_text, kind, data,
                      note=None, callout=None):
    """Right column = badge + a natively-drawn on-brand chart."""
    s = prs.slides.add_slide(BLANK)
    case_header(s, label, title)
    left_column(s, sections)
    rx, rw = Inches(5.45), Inches(4.7)
    y = badge(s, rx, y_=Inches(1.95), w=rw) if False else \
        badge(s, rx, Inches(1.95), rw, badge_text)
    y = Emu(int(y) + int(Inches(0.2)))
    if kind == "value":
        value_bar(s, rx, y, rw, Inches(3.1), data["cats"], data["vals"],
                  unit=data.get("unit", "SEK M"),
                  title=data.get("title", "Initial investment vs. current value"),
                  vfmt=data.get("vfmt"))
        if callout:
            ct = tbox(s, rx, Inches(5.5), rw, Inches(0.8))
            para(ct, callout, 12, SLATE, bold=True, first=True,
                 align=PP_ALIGN.CENTER, after=0, font=SERIF)
    elif kind == "op":
        opbar(s, rx, y, rw, Inches(3.2), data["groups"])
        if callout:
            ct = tbox(s, rx, Inches(5.6), rw, Inches(0.8))
            para(ct, callout, 11, SLATE, italic=True, first=True,
                 align=PP_ALIGN.CENTER, after=0)
    if note:
        note_line(s, note)
    return s


# ===========================================================================
#  CONTENT
# ===========================================================================
NOTE_MODERN = ("Note: share price re-based to 100 at investment and includes "
               "dividends; benchmark OMXS30 / MSCI. Past performance is not "
               "indicative of future results.")
NOTE_ORES = ("Note: 9 SEK ≈ 1€. IRR unleveraged and net of all costs. "
             "Executed by members of the team at Investment AB Öresund.")

cover()
disclaimer()
overview()

# --------------------------------------------------------------------- AIP
era_divider("Part I · AIP — fund era", "Investment cases")

# Haldex AB
case_image_right(
    [("para", "Company description",
      "Haldex was a conglomerate with four separate, non-integrated "
      "divisions: CVS — components for airbrake systems to trucks and "
      "trailers (50% of sales); Concentric — pumps for diesel engines and "
      "the secondary hydraulic system (25%); Traction Systems — four-wheel-"
      "drive systems for passenger cars (13%); and Wire Systems — spring "
      "wires for combustion engines and transmissions (12%)."),
     ("bullets", "Recognized opportunities", [
         "All divisions had a profitable core, growing and holding strong "
         "market positions",
         "Three product groups within CVS were hiding overall profitability",
         "Higher operating costs than peers — cost-saving opportunities",
         "No structural integration between divisions",
         "Market price significantly lower than assets and earnings power",
         "High correlation with GDP — early cycle"]),
     ("para", "Initial plan",
      "1) Change the board to support restructuring; 2) rights issue to "
      "secure the competitive position; 3) split divisions into separate "
      "companies (sell/spin-off); 4) reduce overhead and consolidate "
      "factories; 5) exit or fix unprofitable products; 6) reinvest capital "
      "in leading products."),
     ("checks", "Results", [
         "Changed board and management; exited unprofitable products",
         "Spun off and listed Concentric as a separate company",
         "Consolidated the factory footprint; focused R&D on winners",
         "Introduced new products, entered new markets, improved tax "
         "structure",
         "Ran a two-track process — sold to BorgWarner for 17× EBIT",
         "SEK 30/share redemption; sold Wire Systems to Suzuki Metal"])],
    "Case · AIP fund era", "Investment case study: Haldex AB",
    None, os.path.join(CHARTS, "col_haldex.png"), note=NOTE_MODERN)

# Concentric AB -- the Haldex spin-off, owned and realised separately
case_native_right(
    [("para", "Company description",
      "Concentric AB is a world-leading manufacturer of pumps and hydraulic "
      "systems — oil, water and fuel pumps and hydraulic solutions for diesel "
      "engines and off-/on-highway vehicles. It was the pumps division of "
      "Haldex (~25% of group sales), spun off and separately listed on Nasdaq "
      "Stockholm in 2011 as part of the Haldex restructuring."),
     ("bullets", "Recognized opportunities", [
         "A high-quality, market-leading niche business hidden inside the "
         "Haldex conglomerate",
         "Profitable core with strong margins and high return on capital, "
         "masked by the group structure",
         "A pure-play listing would remove the conglomerate discount and let "
         "the market value it on its own merits",
         "Clear margin and cost opportunity as a focused, independent company",
         "Strong cash generation to support capital return and bolt-on "
         "growth"]),
     ("para", "Initial plan",
      "Realise the hidden value in Haldex's pumps division by spinning it off "
      "as the independent, pure-play Concentric AB; install a focused board "
      "and management, drive margin improvement and capital discipline as a "
      "stand-alone company, and let the public market re-rate the business."),
     ("checks", "Results", [
         "Spun off and separately listed on Nasdaq Stockholm in 2011",
         "Focused board and management; cost and margin programme as an "
         "independent company",
         "Annual EBIT growth of ~19% over the holding period",
         "Strong cash generation funded dividends, buy-backs and bolt-on M&A",
         "Re-rated as a pure-play — realised at ~81% IRR / ~3.9×, "
         "outperforming the parent"])],
    "Case · AIP fund era", "Investment case study: Concentric AB",
    "IRR 81% · Money multiple 3.9×", "value",
    {"cats": ["Annual sales growth", "Annual EBIT growth"], "vals": [2, 19],
     "unit": "% per year, over the holding period",
     "title": "Operational improvement (annualised)",
     "vfmt": lambda v: f"{v}%"},
    note="Note: Concentric was spun off from Haldex and held and realised as a "
         "separate listed company; figures relate to the team's Haldex-era "
         "investment. Past performance is not indicative of future results.",
    callout="Spun off from Haldex and listed in 2011 — outperformed the parent")

# Haldex AB (2nd) -- badge + callout, no chart
def haldex2():
    s = prs.slides.add_slide(BLANK)
    case_header(s, "Case · AIP fund era",
                "Investment case study: Haldex AB — second investment")
    left_column(s, [
        ("para", "Company description",
         "Industrial company manufacturing and selling braking products for "
         "on-highway trucks and trailers. AIP re-invested by buying shares at "
         "a discount during the COVID crisis, after a takeover bid from "
         "competitor Knorr-Bremse had been blocked by the EU and US "
         "authorities."),
        ("bullets", "Recognized opportunities", [
            "Company in limbo after the blocked Knorr-Bremse takeover",
            "Subsequent poor management decisions and bad capital allocation",
            "Loss-making products with growing revenue eroded the balance "
            "sheet",
            "Key management was leaving due to lack of leadership",
            "Inexperienced board consisting of consultants"]),
        ("para", "Initial plan",
         "Change the board to members with industrial experience rather than "
         "lawyers and bankers; change the management team to an experienced "
         "team. Two large institutions (AP4 and AMF) and AIP jointly decided "
         "that Stefan Charette should take the board chair."),
        ("checks", "Results", [
            "Changed the whole board and management",
            "Supported management in renegotiating disc-brake contracts",
            "Significantly reduced product costs and overhead; focused R&D",
            "Exited unprofitable products",
            "Company sold to SAFHolland one year later"])])
    rx, rw = Inches(5.45), Inches(4.7)
    y = badge(s, rx, Inches(1.95), rw, "IRR 100% · Money multiple 2.0×")
    ct = tbox(s, rx, Inches(3.6), rw, Inches(1.6), anchor=MSO_ANCHOR.MIDDLE)
    para(ct, "“Our investment returned 100% within a one-year time "
         "period.”", 19, NAVY_TX, italic=True, first=True,
         align=PP_ALIGN.CENTER, after=0, font=SERIF)
    note_line(s, NOTE_MODERN)
haldex2()

# Lindab AB
case_image_right(
    [("para", "Company description",
      "Lindab is an international company that develops, manufactures and "
      "distributes products and system solutions for simplified construction "
      "and improved indoor climate. It used to be a conglomerate with three "
      "divisions: Ventilation (indoor-climate products sold through own "
      "distribution centres); Components (roof, wall and floor systems — "
      "gutters, rivets, screws); and Building Systems (solutions for "
      "residential and commercial properties)."),
     ("bullets", "Recognized opportunities", [
         "All three divisions had a profitable, growing core with strong "
         "market positions",
         "Poor historical cost management — savings opportunities",
         "Scalable by adding more distribution centres without heavy capex",
         "No structural integration — break-up / consolidation upside",
         "Market price significantly lower than assets and earnings power",
         "Late-cycle company with exposure to Europe"]),
     ("para", "Initial plan",
      "1) Change the board; 2) consolidate Ventilation and Components or "
      "divest Components; 3) reduce OPEX and logistics complexity; 4) grow "
      "sales through the distribution centres; 5) allocate capital to "
      "winning products and markets; 6) renegotiate financing and leases."),
     ("checks", "Results", [
         "Refinanced debt and lease structure; changed board and management",
         "Consolidated the Ventilation & Components business",
         "Significant profit potential through organic and acquired growth",
         "Identified ~22% EBIT uplift on current net sales via cost savings",
         "Opportunity to spin off or sell Building Systems"])],
    "Case · AIP fund era", "Investment case study: Lindab AB",
    None, os.path.join(CHARTS, "col_lindab.png"), note=NOTE_MODERN)

# Transcom SA
case_image_right(
    [("para", "Company description",
      "Transcom is a global customer-care company providing customer care, "
      "sales, technical support and credit-management services through an "
      "extensive network of contact centres and work-at-home agents. It also "
      "operated a credit-management (CMS) company across Europe."),
     ("bullets", "Recognized opportunities", [
         "Two separate businesses with limited integration — sale/spin-off "
         "of the CMS business",
         "Losses in France reduced operating profit to break-even",
         "Higher operating costs than peers — cost-saving opportunity",
         "Problems with transfer taxes, reporting and product quality",
         "Market price below earnings-power potential (incl. rights issue)",
         "Growth opportunities given strong customer relationships"]),
     ("para", "Initial plan",
      "1) Call an EGM to gain a board seat; 2) close the French operation "
      "and settle the Italian tax claim; 3) rights issue; 4) spin off or "
      "sell the CMS business; 5) divest other assets; 6) exit loss-making "
      "sites; 7) reduce OPEX; 8) reinvest in IT to grow the non-voice "
      "product."),
     ("checks", "Results", [
         "Changed board and management; rights issue and new financing",
         "Doubled EBIT potential at the same net sales — in line with peers",
         "Separated CMS with its own CEO; sold it in pieces to industrial "
         "and PE buyers",
         "Liquidated the French subsidiary far cheaper (€5M vs €20M)",
         "Negotiated Italian tax claims down to 3/5; won back the largest "
         "customer",
         "Closed unprofitable sites; introduced new non-voice products"])],
    "Case · AIP fund era", "Investment case study: Transcom SA",
    None, os.path.join(CHARTS, "col_transcom.png"), note=NOTE_MODERN)

# Kitron ASA
case_image_right(
    [("para", "Company description",
      "Kitron is one of the leading companies in Scandinavia specialising in "
      "the contract manufacturing of goods with electronics content, serving "
      "customers in Defence/Aerospace, Utilities, Industry, Medical Devices "
      "and Offshore/Marine. The company has operated since the 1960s."),
     ("bullets", "Recognized opportunities", [
         "Significant turn-around candidate",
         "Positive shift in industry outlook",
         "Stable customer base — defence, medical, utilities",
         "Cost-saving opportunities combined with strong growth in core areas",
         "Market price significantly lower than earnings power",
         "Offshore/Marine at an all-time low — additional turn-around upside"]),
     ("para", "Initial plan",
      "1) Reduce costs in Sweden/Norway; 2) board seat through an EGM in "
      "early 2017; 3) continued investment in technology to reduce manual "
      "labour; 4) acquisition opportunities to build presence; 5) grow "
      "within existing customers."),
     ("checks", "Results", [
         "Continued cost-savings plan in Sweden/Norway",
         "Reached a 7% operating margin in Q2 2017 — a company first",
         "Growing significantly with existing customers"])],
    "Case · AIP fund era", "Investment case study: Kitron ASA",
    None, os.path.join(CHARTS, "col_kitron.png"), note=NOTE_MODERN)

# Note AB
case_image_right(
    [("para", "Company description",
      "Note is an electronics-manufacturing and logistics company with "
      "production plants near its main customer base in northern Europe, as "
      "well as low-cost plants in Eastern Europe and Asia. It focuses on "
      "small series with high variability and technology content."),
     ("bullets", "Recognized opportunities", [
         "Significant turn-around candidate with clear boundaries to losses",
         "The Polish purchasing business was a wrong strategic choice",
         "Higher operating costs than peers",
         "Significant mal-investments in the components database",
         "Market price significantly lower than earnings power",
         "Growth given strong customer relationships — high switching costs"]),
     ("para", "Initial plan",
      "1) Change the board and elect SC as chairman; 2) close Polish "
      "purchasing and the components database and reduce overhead; 3) rights "
      "issue to support restructuring; 4) exit loss-making sites; 5) reduce "
      "OPEX; 6) grow within existing customers."),
     ("checks", "Results", [
         "Reduced non-profitable sales by changing customers and closing "
         "plants",
         "Shifted out close to 40% of the initial purchase price in "
         "dividends",
         "EBIT improved from SEK (91)m in FY2009 to SEK 60m in FY2016"])],
    "Case · AIP fund era", "Investment case study: Note AB",
    None, os.path.join(CHARTS, "col_note.png"), note=NOTE_MODERN)

# DistIt / Alcadon -- reused line chart + reused op chart
case_image_right(
    [("para", "Company description",
      "DistIt was a conglomerate of three distributors: Aurora Group "
      "(Danish, Nordic sales of IT, electronics and appliances); Alcadon "
      "(Swedish/Norwegian sales of active and passive fibre products); and "
      "Deltaco (Nordic sales of electronics and appliances)."),
     ("bullets", "Recognized opportunities", [
         "No structural integration between divisions — structural upside",
         "Market price significantly lower than assets and earnings power",
         "Alcadon could benefit from being stand-alone — right in the cycle",
         "Extension opportunities within Aurora; Deltaco a well-known brand",
         "Never a loss-making year, even through downturns",
         "Stable management team with a strong cost-minimising culture"]),
     ("para", "Initial plan",
      "1) Secure a board seat at the next AGM; 2) spin off Alcadon; "
      "3) increase utilisation of the Aurora sales force through product "
      "extensions; 4) use 3PL in Deltaco to free up real-estate capital; "
      "5) acquire more companies."),
     ("checks", "Results", [
         "Changed the board; spun out and refinanced Alcadon — which then "
         "almost doubled in size through acquisition",
         "Aurora signed Duracell, P&G and Oral-B to leverage the sales force",
         "Moved Deltaco's warehouse to 3PL, freeing the building",
         "Acquired Septon and Sominis, with another acquisition pending",
         "Real-time pricing software drove sales of mispriced products",
         "Reduced IT costs through outsourcing"])],
    "Case · AIP fund era", "Investment case study: DistIt / Alcadon",
    "Investment 2015 – present", os.path.join(CHARTS, "distit_line.png"),
    note=NOTE_MODERN, chart_img2=os.path.join(CHARTS, "distit_op.png"))

# Renold plc -- opportunistic value stake, realised into a recommended takeover
case_native_right(
    [("para", "Company description",
      "Renold plc is a London-listed global leader in industrial chain and "
      "torque-transmission products — roller chain, couplings, gearboxes and "
      "clutches — serving mining, materials-handling, transport and energy "
      "customers worldwide. A market-leading niche franchise (inventor of the "
      "bush roller chain, 1879) with a large installed base and recurring "
      "aftermarket demand."),
     ("bullets", "Recognized opportunities", [
         "Market-leading niche industrial — strong brand, global installed "
         "base, recurring aftermarket",
         "Mispriced, under-followed UK small-cap at a deep discount to earnings "
         "power (~50p)",
         "Multi-year self-help underway — margin expansion and operational "
         "improvement",
         "Strong cash generation and a de-risking balance sheet",
         "A credible consolidation / takeover candidate, with limited downside "
         "at entry"]),
     ("para", "Initial plan",
      "An opportunistic position rather than a control engagement: build a "
      "meaningful stake at a deep discount to intrinsic value, underwrite the "
      "downside through asset and earnings-power support, and let the "
      "catalysts — margin improvement and likely strategic interest — close "
      "the gap."),
     ("checks", "Results", [
         "Built ~5.3m shares at ~50p average from spring 2024; core position "
         "from December 2024 — about six months before the takeover",
         "A recommended cash takeover emerged in 2025 at ~81p per share",
         "Realised the entire position into the offer at ~81p over June–July "
         "2025",
         "~1.6× money multiple, +62% gross (SEK 20.9m); money-weighted IRR "
         "~234%"])],
    "Case · AIP fund era", "Investment case study: Renold plc",
    "IRR 234% · Money multiple 1.6×", "value",
    {"cats": ["Capital invested", "Realised proceeds"], "vals": [35, 56],
     "unit": "SEK m", "title": "Capital invested vs. realised proceeds"},
    note="Source: AIPFII transaction history (Apr 2024 – Jul 2025). IRR is "
         "money-weighted; the short opportunistic hold lifts the annualised "
         "figure. Past performance is not indicative of future results.",
    callout="Bought at ~50p; acquired at ~81p in a recommended takeover")

# Robit Plc -- the disciplined walk-away: profitable exit, then the stock halved
case_native_right(
    [("para", "Company description",
      "Robit Plc is a Finnish-listed (Nasdaq Helsinki) manufacturer of drilling "
      "tools and consumables for mining, construction and well-drilling — a "
      "niche industrial with a global distribution footprint and recurring "
      "aftermarket demand."),
     ("bullets", "Why we invested", [
         "An under-followed Nordic small-cap with a recurring-revenue "
         "consumables franchise",
         "Trading well below intrinsic value, with a correctable cost base and "
         "footprint",
         "A clear path to value through board-level engagement and capital "
         "discipline",
         "Limited downside at entry; a multi-year compounding opportunity"]),
     ("para", "The engagement plan",
      "Build the position in stages, follow the company closely, and secure a "
      "board seat to drive the operational and capital-allocation changes — "
      "engaged ownership, not a passive bet."),
     ("checks", "What happened — and the discipline", [
         "First invested May 2015; followed the company and added on over "
         "several years",
         "Began negotiating a board seat in 2019 to formalise the engagement",
         "Unable to secure board alignment, we exited rather than stay without "
         "influence",
         "Still profitable: a ~14% annual IRR — more than twice the index over "
         "the hold",
         "Since our exit the shares have roughly halved — walking away was also "
         "avoided loss"])],
    "Case · AIP fund era", "Investment case study: Robit Plc",
    "IRR ~14% p.a. · > 2× the index · disciplined exit", "value",
    {"cats": ["At our exit", "Today"], "vals": [100, 50],
     "unit": "Indexed to 100 at our exit (indicative)",
     "title": "Robit share price after we exited",
     "vfmt": lambda v: f"{v}"},
    note="Source: AIPFII transaction history (first investment May 2015; "
         "board-seat negotiation 2019; exited thereafter). IRR is money-weighted "
         "and net of costs; benchmark is the relevant equity index; the "
         "post-exit move is indicative. Past performance is not indicative of "
         "future results.",
    callout="We exited profitably (~14% IRR); the shares then halved")

# --------------------------------------------------------------------- ORESUND
era_divider("Part II · Investment AB Öresund (1994–2007)",
            "Earlier track record")

# Custos -- reused line chart
case_image_right(
    [("para", "Company description",
      "Custos was a Swedish listed investment company holding a portfolio of "
      "listed companies. It traded at a significant discount to net asset "
      "value and had an inefficient capital structure."),
     ("bullets", "Recognized opportunities", [
         "Investment companies traded at a discount to NAV",
         "Inefficient capital structure",
         "An interesting portfolio of companies not optimally managed",
         "Opportunity to buy the underlying assets at a ~30% discount",
         "The largest shareholder, Skanska, had sold its holding in 1994",
         "The board was not focused on creating shareholder value"]),
     ("para", "Initial plan",
      "1) Increase shareholder value through extended engagement in holdings "
      "via active board management; 2) form a more heterogeneous board to "
      "enhance dialogue and decision-making; 3) reduce the discount to NAV "
      "through, if necessary, redemption of holdings/cash and synthetic "
      "buy-back programs."),
     ("checks", "Results", [
         "New board and CEO (Christer Gardell); active ownership reduced the "
         "NAV discount",
         "Distributed Hufvudstaden to shareholders; sold ASG to Deutsche Post",
         "Founded Acando; synthetic buy-back distributing SEK 1.6bn",
         "~SEK 3.5bn distributed via buy-back and the distribution of SCA",
         "Bilia and Acando distributed to shareholders",
         "Custos ultimately merged with Öresund"])],
    "Case · Öresund era", "Investment case study: Custos",
    "IRR 58.9% · 2000–2004", os.path.join(CHARTS, "custos_line.png"),
    note=NOTE_ORES,
    value_note="≈ SEK 11bn of shareholder value created over nine years")

# Johnson Pump -- reused line chart
case_image_right(
    [("para", "Company description",
      "Johnson Pump was a traditional industrial business — a pump OEM with "
      "an extensive sales force and a strong brand. It had suffered "
      "profitability issues but had started to recover, with revenues of "
      "SEK 647m and a negative EBIT."),
     ("bullets", "Recognized opportunities", [
         "Low profitability was largely a cost issue — customers kept buying "
         "and the gross margin was decent",
         "Significant upside on cost in both factories and sales offices",
         "The sales force could be leveraged with a wider product range "
         "(R&D, 3rd-party products, M&A)",
         "Market segmentation should reveal further niche potential",
         "Geographic expansion to follow customers moving to low-cost "
         "countries"]),
     ("para", "Initial plan",
      "Reduce costs in factories and sales offices, build a uniform sales "
      "process and a simplified product range, expand the range and "
      "geography, and grow in adjacencies through M&A."),
     ("checks", "Results", [
         "New ERP system (2002); component manufacturing moved to India "
         "(2003)",
         "Restructured Dutch operations and centralised sourcing (2004–05)",
         "New CEO Stefan Charette (2006); separated Marine, acquired "
         "Tigerholm, launched an export organisation",
         "Revenues +~17%; SG&A −~18%; operating result +~SEK 62m (−15m "
         "to +47m)",
         "SPX Corporation public offer accepted by >90% of shareholders; "
         "delisted"])],
    "Case · Öresund era", "Investment case study: Johnson Pump",
    "IRR 34.4% · 2002–2006", os.path.join(CHARTS, "jp_line.png"),
    note=NOTE_ORES)

# Wihlborgs / Drott -- reused line chart
case_image_right(
    [("para", "Company description",
      "Wihlborgs and Drott were Swedish listed real-estate companies holding "
      "commercial and residential properties in the growing Stockholm and "
      "Malmö regions, trading at a discount to NAV with no active owner."),
     ("bullets", "Recognized opportunities", [
         "Shares traded at a discount to NAV for real-estate companies",
         "Hidden values in both Wihlborgs and Drott",
         "Limited active ownership; inefficient debt/capital structure",
         "Limited synergies between commercial and residential — divest",
         "Inefficient organisations",
         "Property-investment funds were pushing prices upwards"]),
     ("para", "Initial plan",
      "Divest properties and reduce costs to enable cash distribution; split "
      "commercial and residential into separate companies — one for the "
      "Öresund/Malmö region and one for Stockholm; refinance to lower the "
      "cost of debt."),
     ("checks", "Results", [
         "Öresund acquired 12% of Wihlborgs (2001); new board members joined",
         "Wihlborgs built a ~10% position in Drott; a new board was elected",
         "New management plan with ~SEK 85m of yearly savings",
         "Several properties divested — SEK 3.3bn cash redemption (2004)",
         "Drott split into BAB Drott (residential, sold to Stena) and Fabege "
         "(commercial, acquired by Wihlborgs)",
         "Combined operations later split into two regional companies"])],
    "Case · Öresund era", "Investment case study: Wihlborgs / Drott",
    "IRR 31.0% · 2002–2004", os.path.join(CHARTS, "drott_line.png"),
    note=NOTE_ORES,
    value_note="≈ SEK 3bn of net value created in 15 months")

# ---- native value-bar cases ----
case_native_right(
    [("para", "Company description",
      "J&W (Jacobson & Widmark) was a Swedish technical-consulting firm. "
      "Öresund first invested in 1994 and took a controlling interest in "
      "2000."),
     ("bullets", "Recognized opportunities", [
         "Over-funded pension plans — cash to be returned from the SPP "
         "pension trust (now Alecta)",
         "A fragmented consulting market with consolidation potential",
         "A non-core subsidiary (Benima Ferator) that could be divested"]),
     ("para", "Initial plan",
      "Return excess cash to shareholders, divest non-core operations, and "
      "consolidate the fragmented technical-consulting market through "
      "acquisitions."),
     ("checks", "Results", [
         "Spun off Benima Ferator to shareholders (later acquired by Sigma)",
         "Transferred SPP pension-plan cash to shareholders",
         "Acquired several smaller competitors and the larger Kjessler & "
         "Mannerstråle (KM)",
         "J&W was subsequently sold to WSP"])],
    "Case · Öresund era", "Investment case study: J&W",
    "IRR 94.4% · 2000–2001 · Technical consultants", "value",
    {"cats": ["Initial investment", "Current value"], "vals": [213, 476]},
    note=NOTE_ORES)

case_native_right(
    [("para", "Company description",
      "Nobia is a European kitchen company. Originally active across several "
      "construction-component areas, it was streamlined to focus on "
      "kitchens."),
     ("bullets", "Recognized opportunities", [
         "Active across several construction-component areas — focus "
         "opportunity",
         "A fragmented European kitchen market with consolidation potential"]),
     ("para", "Initial plan",
      "Streamline Nobia to focus on kitchens by divesting other areas, and "
      "consolidate the European kitchen market through acquisitions."),
     ("checks", "Results", [
         "Streamlined to focus on kitchens (sold Doors & Windows)",
         "Acquired Gower, EWE-FM and Hygena — the largest kitchen company in "
         "Europe",
         "Turnover up from SEK 9.2bn (2003) to SEK 16.1bn (12m to Jun 2007)",
         "EBIT up from SEK 560m to SEK 1,360m"])],
    "Case · Öresund era", "Investment case study: Nobia",
    "IRR 33.1% · 2002– · Consumer discretionary", "value",
    {"cats": ["Initial investment", "Current value"], "vals": [162, 1083]},
    note=NOTE_ORES)

case_native_right(
    [("para", "Company description",
      "Hagströmer & Qviberg (HQ) and Avanza are Swedish financial companies "
      "— brokerage, fund management, wealth management and online brokerage. "
      "Sven Hagströmer founded Hagströmer Fondkommission in 1986; Mats "
      "Qviberg joined in 1989."),
     ("bullets", "Recognized opportunities", [
         "Separate online brokerage / fund management from traditional "
         "brokerage and investment banking",
         "Consolidation opportunities among internet brokerages"]),
     ("para", "Initial plan",
      "Split HQ into focused entities, consolidate internet brokerages into "
      "Avanza, and build a wealth-management business."),
     ("checks", "Results", [
         "HQ split (2000) into HQ.SE (internet brokerage, fund management) "
         "and Hagströmer & Qviberg",
         "HQ.SE acquired Aktiespar FK and Avanza; split into HQ Fonder and "
         "Avanza",
         "Avanza expanded into pensions, fund distribution, corporate finance "
         "and banking",
         "HQ Fonder merged with Hagströmer & Qviberg (2005) — SEK 85bn AUM"])],
    "Case · Öresund era", "Investment case study: HQ / Avanza",
    "IRR 37.8% · 2000– · Finance", "value",
    {"cats": ["Initial investment", "Current value"], "vals": [744, 2808]},
    note=NOTE_ORES)

case_native_right(
    [("para", "Company description",
      "SkiStar is a Scandinavian operator of alpine ski resorts, built "
      "through the consolidation of several destinations into a single "
      "listed leisure company."),
     ("bullets", "Recognized opportunities", [
         "A fragmented ski-resort market with consolidation potential",
         "Opportunity to build a leading Scandinavian leisure platform"]),
     ("para", "Initial plan",
      "Consolidate alpine destinations into a single listed leisure "
      "company."),
     ("checks", "Results", [
         "Hundfjället acquired Tandådalen; the combined entity was "
         "restructured",
         "Acquired Lindvallen; the combined entity was listed",
         "Acquired Åre, Duved and Vemdalen from Bure — renamed SkiStar",
         "Acquired Hemsedal (and later Trysil)"])],
    "Case · Öresund era", "Investment case study: SkiStar",
    "IRR 33.0% · 2000– · Leisure", "value",
    {"cats": ["Initial investment", "Current value"], "vals": [160, 659]},
    note=NOTE_ORES)

case_native_right(
    [("para", "Company description",
      "Bilia is a leading Nordic car retailer and car-service company. "
      "Öresund became an owner in 2003 through a redemption of shares from "
      "Custos."),
     ("bullets", "Recognized opportunities", [
         "Opportunity to separate trucks/heavy vehicles from the car "
         "business",
         "Strong cash flow to fund consolidation of car retailers",
         "Hidden real-estate value"]),
     ("para", "Initial plan",
      "Separate the heavy-vehicle and car businesses, use strong cash flow "
      "to acquire competing car retailers, and unlock real-estate value."),
     ("checks", "Results", [
         "Split into heavy vehicles (LV) and cars (PV); Volvo AB acquired the "
         "LV business for PV shares — leaving a pure car retailer",
         "Acquired several competing car retailers, expanding geography and "
         "brands; transferred excess cash to shareholders",
         "Spun off the real estate as Catena (2006), distributed to "
         "shareholders and listed"])],
    "Case · Öresund era", "Investment case study: Bilia",
    "IRR 26.0% · 2003– · Consumer discretionary", "value",
    {"cats": ["Initial investment", "Current value"], "vals": [506, 1056]},
    note=NOTE_ORES)

contact()

OUT = os.path.join(BASE, "Athanase_Case_Studies.pptx")
prs.save(OUT)
print(f"Saved {OUT} with {len(prs.slides._sldIdLst)} slides at 4:3")

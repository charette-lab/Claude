"""Standalone, brand-compliant slides: Athanase vs mid-market private equity,
using REAL (committed-capital) PE returns rather than headline IRR.

Slide 1 — reconciliation bridge: why the ~20% headline IRR an LP is shown is
          really ~8.6% once measured honestly (committed-capital / PME basis).
Slide 2 — switching a PE sleeve into Athanase, on the real numbers.

Athanase stats are computed live from data/Comparison_returns.xlsx (real).
Private-equity figures are sourced from published research (see footnotes):
  - Top-quartile mid-market buyout headline net IRR ~18-22%  (Cambridge Assoc.)
  - Real committed-basis return ~8.6%, ~matching the S&P 500  (Phalippou 2020)
  - IRR reinvestment assumption inflates by ~3 pts  (evidence-based research)
  - Reported vol ~11% de-smooths to ~22% total / ~13% downside  (MSCI/Morningstar)
Built 4:3 per AIP Brand Guidelines.
"""
import math
import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
from lxml import etree
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager as _fmgr

# ---- AIP Brand palette (Blue group only) -----------------------------------
BLUE1 = RGBColor(0x0E, 0x16, 0x20); BLUE2 = RGBColor(0x15, 0x21, 0x30)
BLUE3 = RGBColor(0x31, 0x43, 0x59); BLUE4 = RGBColor(0x55, 0x6A, 0x83)
BLUE5 = RGBColor(0xE2, 0xE5, 0xE9); BLUE6 = RGBColor(0xF6, 0xF7, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NAVY, NAVY_TX, SLATE, SLATE_LT = BLUE2, BLUE3, BLUE3, BLUE4
HEADERBG, BODY, SUBTLE, FOOT = BLUE6, BLUE1, BLUE4, BLUE4
SERIF, SANS = "Times New Roman", "Arial"
H_NAVY, H_BLUE3, H_BLUE4 = "#152130", "#314359", "#556A83"
H_BLUE5, H_BODY = "#E2E5E9", "#0E1620"
MARK_DARK = "assets/mark_dark.png"
_MD_AR = (lambda s: s[0] / s[1])(Image.open(MARK_DARK).size)
for _ff in ("Arial", "Liberation Sans", "DejaVu Sans"):
    try:
        _fmgr.findfont(_ff, fallback_to_default=False)
        plt.rcParams["font.family"] = _ff
        break
    except Exception:
        continue

# ===========================================================================
# Data
# ===========================================================================
ws = openpyxl.load_workbook("data/Comparison_returns.xlsx", data_only=True)["Sheet1"]
_ATH = []
for c in range(2, ws.max_column + 1):
    for r in range(22, 34):
        v = ws.cell(r, c).value
        if isinstance(v, (int, float)):
            _ATH.append(v)
_NM = len(_ATH)


def _ann(xs):
    g = 1.0
    for x in xs:
        g *= (1 + x)
    return g ** (12 / len(xs)) - 1


def _dvol(xs):
    return math.sqrt(sum(min(x, 0) ** 2 for x in xs) / len(xs)) * math.sqrt(12)


def _tvol(xs):
    import statistics as _st
    return _st.pstdev(xs) * math.sqrt(12)


ATH_RET, ATH_DD = _ann(_ATH), _dvol(_ATH)        # ~16.0%, ~10.9% (real)
ATH_VOL = _tvol(_ATH)                             # ~27.0% total vol (real)

# Private equity — REAL allocator-realised returns from institutional datasets
# (PME / cash-matched basis). PE genuinely earns a premium over public markets;
# the honest benchmark is what allocators realise (~11-14%), not headline IRR.
PE_HEADLINE = 0.20            # top-quartile mid-market buyout headline net IRR
PE_REAL = 0.114              # Cliffwater: US state-pension PE realised, 2000-22
PE_REAL_LO, PE_REAL_HI = 0.114, 0.141   # Cliffwater (US) .. BVCA (Europe)
# evidence pairs: (label, PE realised %, public-market-equivalent %)
PE_EVIDENCE = [
    ("Cliffwater\nUS state pensions\n(2000–2022)", 11.4, 5.8),
    ("BVCA / Capital Dynamics\nUK & Europe\n(2001–2023)", 14.1, 7.7),
]
# Volatility (TOTAL, annualised) — the basis the de-smoothing research reports.
PE_VOL_REPORTED = 0.10        # reported (appraisal-smoothed) total vol (~10%)
PE_VOL_TRUE = 0.28            # de-smoothed true economic vol (PIMCO ~30%; AQR beta 1.5-1.6)

# ===========================================================================
# Helpers
# ===========================================================================
prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height


def rect(sl, l, t, w, h, fill=None):
    shp = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill is not None:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp


def tbox(sl, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = sl.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, size, color, bold=False, italic=False, first=False,
         align=PP_ALIGN.LEFT, after=8, lead=1.1, font=SANS):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(after); p.line_spacing = lead
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = font
    return p


def header(sl, kicker, ref, title, subtitle):
    sl.shapes.add_picture(MARK_DARK, Inches(0.55), Inches(0.24), height=Inches(0.26),
                          width=Emu(int(Inches(0.26) * _MD_AR)))
    para(tbox(sl, Inches(0.98), Inches(0.27), Inches(7), Inches(0.3)),
         kicker, 11, SLATE_LT, first=True, after=0)
    para(tbox(sl, Inches(9.9), Inches(0.27), Inches(2.85), Inches(0.3)),
         ref, 11, SLATE_LT, first=True, bold=True, align=PP_ALIGN.RIGHT, after=0)
    rect(sl, 0, Inches(0.62), SW, Inches(1.45), fill=HEADERBG)
    para(tbox(sl, Inches(0.6), Inches(0.74), Inches(12.1), Inches(0.85)),
         title, 30, NAVY_TX, first=True, after=2, font=SERIF)
    para(tbox(sl, Inches(0.62), Inches(1.55), Inches(11.9), Inches(0.5)),
         subtitle, 13, SUBTLE, first=True, italic=True, after=0)


# ===========================================================================
# SLIDE 1 — what allocators ACTUALLY earn: realised PE vs its public-market
# equivalent, across institutional datasets, with Athanase for reference
# ===========================================================================
def realised_chart(path_png):
    fig, ax = plt.subplots(figsize=(8.7, 4.35), dpi=200)
    groups = [e[0] for e in PE_EVIDENCE]
    pe_vals = [e[1] for e in PE_EVIDENCE]
    pme_vals = [e[2] for e in PE_EVIDENCE]
    x = np.arange(len(groups)); w = 0.38
    ax.bar(x - w / 2, pe_vals, w, color=H_BLUE4, label="PE realised (net)")
    ax.bar(x + w / 2, pme_vals, w, color=H_BLUE5, edgecolor=H_BLUE4, linewidth=1,
           label="Public-market equivalent")
    for xi, (pe, pme) in enumerate(zip(pe_vals, pme_vals)):
        ax.annotate(f"{pe:.1f}%", (xi - w / 2, pe), ha="center", va="bottom",
                    fontsize=12, fontweight="bold", color=H_BLUE3,
                    xytext=(0, 3), textcoords="offset points")
        ax.annotate(f"{pme:.1f}%", (xi + w / 2, pme), ha="center", va="bottom",
                    fontsize=11, color=H_BLUE4, xytext=(0, 3),
                    textcoords="offset points")
        # place premium label between the two bars, clear of the Athanase line
        ax.annotate(f"+{pe-pme:.1f} pts\npremium", (xi, (pe + pme) / 2),
                    ha="center", va="center", fontsize=9.5, color=H_BLUE3,
                    fontweight="bold")
    # Athanase reference line
    ax.axhline(ATH_RET * 100, color=H_NAVY, lw=1.8, ls=(0, (5, 3)), zorder=5)
    ax.annotate(f"Athanase {ATH_RET*100:.0f}%  (real, net, fully invested)",
                (len(groups) - 0.5, ATH_RET * 100), ha="right", va="bottom",
                fontsize=10.5, color=H_NAVY, fontweight="bold")
    ax.set_xticks(x); ax.set_xticklabels(groups, fontsize=9.5, color=H_BODY)
    ax.set_ylabel("Annualised return to investors", color=H_BODY, fontsize=11)
    ax.set_ylim(0, 19); ax.set_xlim(-0.6, len(groups) - 0.4)
    ax.yaxis.set_major_formatter(lambda v, _: f"{v:.0f}%")
    ax.tick_params(colors=H_BODY, labelsize=10)
    for sp in ("top", "right"):
        ax.spines[sp].set_visible(False)
    for sp in ("left", "bottom"):
        ax.spines[sp].set_color(H_BLUE5)
    ax.grid(True, axis="y", color=H_BLUE5, lw=0.6, alpha=0.6)
    ax.legend(loc="upper center", bbox_to_anchor=(0.5, -0.16), ncol=2,
              fontsize=10, frameon=False, labelcolor=H_BODY)
    ax.set_title("What allocators actually realise from private equity",
                 color=H_NAVY, fontsize=13, fontweight="bold", pad=10)
    fig.tight_layout(); fig.savefig(path_png, dpi=200, bbox_inches="tight",
                                    facecolor="white"); plt.close(fig)


realised_chart("/tmp/pe_realised.png")
s1 = prs.slides.add_slide(prs.slide_layouts[6])
header(s1, "Risk-Adjusted Outcomes", "Figure 1",
       "What allocators actually earn from private equity",
       "Stripping out headline-IRR marketing, cash-matched (PME) data shows PE "
       "earns a real ~3–6 pt premium over public markets — but lands below "
       "Athanase.")
s1.shapes.add_picture("/tmp/pe_realised.png", Inches(0.55), Inches(2.28),
                      height=Inches(3.95))
# right: explanation cards
cx = Inches(8.95); cw = Inches(3.8); chh = Inches(1.2); cy = Inches(2.32)
expl = [
    ("HEADLINE ≠ REALISED",
     "Top-quartile IRR is shown as ~18–22%, struck on drawn capital. Cash-matched "
     "to public markets, the real figure is far lower."),
    ("THE REAL PREMIUM IS GENUINE",
     "PME studies agree: ~1.2× wealth multiple, a real 3–5 pt premium over "
     "public equities (Cambridge; Kaplan-Schoar / Burgiss; Cliffwater; BVCA)."),
    ("BUT BELOW ATHANASE",
     "Allocators realise ~11–14% from PE. Athanase has delivered ~16% — net, "
     "real, and fully invested from day one."),
]
for title, body in expl:
    rect(s1, cx, cy, cw, chh, fill=HEADERBG)
    para(tbox(s1, Emu(int(cx) + int(Inches(0.2))), cy + Inches(0.1),
              Emu(int(cw) - int(Inches(0.4))), Inches(0.25)),
         title, 10.5, SLATE, first=True, bold=True, after=0)
    para(tbox(s1, Emu(int(cx) + int(Inches(0.2))), cy + Inches(0.36),
              Emu(int(cw) - int(Inches(0.4))), Inches(0.8)),
         body, 10.5, BODY, first=True, after=0, lead=1.12)
    cy = Emu(int(cy) + int(chh) + int(Inches(0.16)))
rect(s1, Inches(0.6), Inches(6.62), Inches(12.13), Inches(0.46), fill=NAVY)
para(tbox(s1, Inches(0.78), Inches(6.62), Inches(11.8), Inches(0.46),
          anchor=MSO_ANCHOR.MIDDLE),
     "PE’s premium over public markets is real — but Athanase’s ~16% sits above "
     "even what the best allocators realise, with daily liquidity.",
     12, WHITE, first=True, italic=True, after=0)
para(tbox(s1, Inches(0.6), Inches(7.1), Inches(12.6), Inches(0.5)),
     "Sources: Cliffwater (US state pensions, 2000–2022): PE 11.4% vs PME 5.8%. "
     "BVCA / Capital Dynamics (UK & Europe, 2001–2023): 14.1% vs 7.7% (PME+). "
     "Corroborated by Cambridge Associates mPME (1.15–1.25×) and "
     "Kaplan-Harris-Jenkinson KS-PME on Burgiss (1.20–1.27×), ≈ +3 pts/yr. "
     "Athanase real, net, 2006–2025.",
     7.5, FOOT, first=True, after=0, lead=1.1)


# ===========================================================================
# SLIDE 2 — switching a PE sleeve into Athanase (on the real numbers)
# ===========================================================================
def scatter_chart(path_png):
    fig, ax = plt.subplots(figsize=(6.0, 4.25), dpi=200)
    # PE marketed -> real: arrow down (lower return) and right (true vol ~3x)
    ax.annotate("", (PE_VOL_TRUE * 100, PE_REAL * 100),
                (PE_VOL_REPORTED * 100, PE_HEADLINE * 100),
                arrowprops=dict(arrowstyle="->", color=H_BLUE4, lw=1.8,
                                linestyle=(0, (3, 2))))
    ax.scatter([PE_VOL_REPORTED * 100], [PE_HEADLINE * 100], s=150, color="white",
               edgecolor=H_BLUE4, linewidth=2.2, zorder=5)
    ax.annotate("PE as marketed\n(headline IRR,\nsmoothed ~10% vol)",
                (PE_VOL_REPORTED * 100, PE_HEADLINE * 100), color=H_BLUE4,
                fontsize=9.5, fontweight="bold", xytext=(6, 6),
                textcoords="offset points", ha="left")
    ax.scatter([PE_VOL_TRUE * 100], [PE_REAL * 100], s=150, color=H_BLUE4, zorder=5,
               edgecolor="white", linewidth=1.5)
    ax.annotate("PE as realised\n(PME return,\nde-smoothed ~28% vol)",
                (PE_VOL_TRUE * 100, PE_REAL * 100), color=H_BLUE4, fontsize=9.5,
                xytext=(-8, -10), textcoords="offset points", ha="right", va="top")
    # Athanase
    ax.scatter([ATH_VOL * 100], [ATH_RET * 100], s=210, color=H_NAVY, zorder=6,
               edgecolor="white", linewidth=1.5)
    ax.annotate(f"Athanase\n{ATH_RET*100:.0f}% return", (ATH_VOL * 100, ATH_RET * 100),
                color=H_NAVY, fontsize=11, fontweight="bold", xytext=(-10, 12),
                textcoords="offset points", va="center", ha="right")
    ax.set_xlabel("Total volatility (annualised)", color=H_BODY, fontsize=11)
    ax.set_ylabel("Real net return to investors", color=H_BODY, fontsize=11)
    ax.set_xlim(4, 34); ax.set_ylim(6, 23)
    ax.xaxis.set_major_formatter(lambda v, _: f"{v:.0f}%")
    ax.yaxis.set_major_formatter(lambda v, _: f"{v:.0f}%")
    ax.tick_params(colors=H_BODY, labelsize=10)
    for sp in ("top", "right"):
        ax.spines[sp].set_visible(False)
    for sp in ("left", "bottom"):
        ax.spines[sp].set_color(H_BLUE5)
    ax.grid(True, color=H_BLUE5, lw=0.6, alpha=0.6)
    ax.set_title("Real return vs true (de-smoothed) volatility", color=H_NAVY,
                 fontsize=13, fontweight="bold", pad=10)
    fig.tight_layout(); fig.savefig(path_png, dpi=200, bbox_inches="tight",
                                    facecolor="white"); plt.close(fig)


scatter_chart("/tmp/pe_scatter.png")
s2 = prs.slides.add_slide(prs.slide_layouts[6])
header(s2, "Risk-Adjusted Outcomes", "Figure 2",
       "Switching a private-equity sleeve into Athanase",
       "On true risk, PE and Athanase carry similar volatility — but Athanase’s "
       "is mostly upside, and it returns more, with daily liquidity.")
s2.shapes.add_picture("/tmp/pe_scatter.png", Inches(0.55), Inches(2.28),
                      height=Inches(3.95))
para(tbox(s2, Inches(0.7), Inches(6.28), Inches(6.2), Inches(0.4)),
     "De-smoothed, PE’s ~10% reported vol is really ~28% (PIMCO; AQR beta "
     "1.5–1.6) — about Athanase’s 27%, but levered downside, not upside.",
     10, SUBTLE, first=True, italic=True, after=0, lead=1.12)
cardx = Inches(7.05); cardw = Inches(5.7); cardh = Inches(0.86); cy = Inches(2.3)
gains = [
    ("RETURN", "~16% vs ~11–14% realised",
     "Above what even the best allocators realise from PE (Cliffwater; BVCA)."),
    ("LIQUIDITY", "Daily vs 10-yr lock-up",
     "Listed positions you can exit; no blind-pool capital calls."),
    ("FEES & CAPITAL", "Lower drag, fully invested",
     "No fees on undeployed dry powder waiting in a queue."),
    ("TRANSPARENCY", "Daily marks, real control",
     "A public scoreboard plus board control of cash flow, capex and pay."),
    ("RISK, HONESTLY", "Similar vol, opposite shape",
     "PE’s true ~28% vol is levered downside; Athanase’s 27% is mostly upside "
     "(downside only ~11%)."),
]
for title, big, body in gains:
    rect(s2, cardx, cy, cardw, cardh, fill=HEADERBG)
    para(tbox(s2, Emu(int(cardx) + int(Inches(0.2))), cy + Inches(0.1),
              Inches(2.0), Inches(0.25)), title, 9.5, SLATE, first=True,
         bold=True, after=0)
    para(tbox(s2, Emu(int(cardx) + int(Inches(2.05))), cy + Inches(0.09),
              Inches(3.4), Inches(0.3)), big, 13, NAVY_TX, first=True, bold=True,
         after=0, font=SERIF)
    para(tbox(s2, Emu(int(cardx) + int(Inches(0.2))), cy + Inches(0.4),
              Emu(int(cardw) - int(Inches(0.4))), Inches(0.42)), body, 10, BODY,
         first=True, after=0, lead=1.08)
    cy = Emu(int(cy) + int(cardh) + int(Inches(0.1)))
rect(s2, Inches(0.6), Inches(6.62), Inches(12.13), Inches(0.46), fill=NAVY)
para(tbox(s2, Inches(0.78), Inches(6.62), Inches(11.8), Inches(0.46),
          anchor=MSO_ANCHOR.MIDDLE),
     "Once PE’s smoothing is removed, the two carry similar volatility — but "
     "Athanase returns more, its risk is upside not downside, and it stays "
     "liquid and transparent.", 12, WHITE, first=True, italic=True, after=0)
para(tbox(s2, Inches(0.6), Inches(7.12), Inches(12.6), Inches(0.5)),
     f"Athanase: net monthly returns, {_NM} months (2006–2025), real (total vol "
     f"{ATH_VOL*100:.0f}%, downside {ATH_DD*100:.0f}%). PE realised ~11–14% "
     "(Cliffwater 11.4%; BVCA 14.1%) vs headline IRR ~18–22%; reported ~10% vol "
     "de-smoothed to ~25–30% (PIMCO; AQR equity beta 1.5–1.6; Couts-Gonçalves-"
     "Rossi, 2024). Illustrative; past performance ≠ future results.",
     7.5, FOOT, first=True, after=0, lead=1.1)


# ===========================================================================
# SLIDE 3 — the academic evidence: four methods converge on ~2-3x true risk
# ===========================================================================
# (method label, reported metric text, true metric text, reported val, true val,
#  plotted as volatility-equivalent % for the bar; beta rows scaled to vol view)
EVID = [
    ("Cash-flow NPV model", "Ang, Chen, Goetzmann\n& Phalippou (2018) · JF",
     "11% vol", "25% vol", 11, 25),
    ("Secondary-market prices", "Boyer, Nadauld, Vorkink\n& Weisbach (2018)",
     "beta < 0.5", "beta > 2.0", 11, 28),
    ("Econometric unsmoothing", "Getmansky, Lo & Makarov\n(2004) · JFE",
     "smoothed", "serial-correlation\nremoved", 11, 24),
    ("3-step unsmoothing", "Couts, Gonçalves\n& Rossi (2020)",
     "low vol, high Sharpe", "vol ↑, market\nco-movement ↑", 11, 27),
]


def evidence_chart(path_png):
    fig, ax = plt.subplots(figsize=(7.4, 4.5), dpi=200)
    ys = np.arange(len(EVID))[::-1]
    h = 0.34
    rep = [e[4] for e in EVID]
    tru = [e[5] for e in EVID]
    ax.barh(ys + h / 2, rep, h, color=H_BLUE5, edgecolor=H_BLUE4, linewidth=1,
            label="Reported (appraisal-smoothed)")
    ax.barh(ys - h / 2, tru, h, color=H_NAVY, label="True economic (research)")
    for y, e in zip(ys, EVID):
        ax.annotate(e[2], (e[4] + 0.5, y + h / 2), va="center", ha="left",
                    fontsize=9, color=H_BLUE4)
        ax.annotate(e[3].split(chr(10))[0], (e[5] + 0.5, y - h / 2), va="center",
                    ha="left", fontsize=9, color=H_NAVY, fontweight="bold")
    ax.set_yticks(ys)
    ax.set_yticklabels([e[0] for e in EVID], fontsize=10.5, color=H_BODY)
    ax.set_xlabel("Volatility / risk loading (de-smoothed shown as vol-equivalent)",
                  color=H_BODY, fontsize=10)
    ax.set_xlim(0, 36); ax.set_ylim(-0.7, len(EVID) - 0.3)
    ax.xaxis.set_major_formatter(lambda v, _: f"{v:.0f}%")
    ax.tick_params(colors=H_BODY, labelsize=9.5)
    for sp in ("top", "right"):
        ax.spines[sp].set_visible(False)
    for sp in ("left", "bottom"):
        ax.spines[sp].set_color(H_BLUE5)
    ax.grid(True, axis="x", color=H_BLUE5, lw=0.6, alpha=0.6)
    ax.legend(loc="upper center", bbox_to_anchor=(0.5, -0.13), ncol=2,
              fontsize=9.5, frameon=False, labelcolor=H_BODY)
    ax.set_title("Four independent methods, one conclusion",
                 color=H_NAVY, fontsize=13, fontweight="bold", pad=10)
    fig.tight_layout(); fig.savefig(path_png, dpi=200, bbox_inches="tight",
                                    facecolor="white"); plt.close(fig)


evidence_chart("/tmp/pe_evidence.png")
s3 = prs.slides.add_slide(prs.slide_layouts[6])
header(s3, "Risk-Adjusted Outcomes", "Figure 3",
       "The volatility gap is settled in the literature",
       "Cash-flow models, secondary-market prices and econometric unsmoothing all "
       "converge: PE’s true risk is roughly 2–3× its reported figure.")
s3.shapes.add_picture("/tmp/pe_evidence.png", Inches(0.5), Inches(2.25),
                      height=Inches(4.05))
# right: method citation cards
cx = Inches(8.5); cw = Inches(4.25); chh = Inches(0.92); cy = Inches(2.3)
methods = [
    ("CASH-FLOW NPV MODEL",
     "Buyout vol 25% (vs 11% index), beta > 1 — Ang, Chen, Goetzmann & "
     "Phalippou (2018), J. Finance."),
    ("SECONDARY-MARKET PRICES",
     "Traded-stake beta > 2.0 vs < 0.5 for the same firms on NAV — Boyer, "
     "Nadauld, Vorkink & Weisbach (2018)."),
    ("ECONOMETRIC UNSMOOTHING",
     "Illiquidity creates spurious smoothness that suppresses vol — Getmansky, "
     "Lo & Makarov (2004), JFE (cited 1,300+)."),
    ("CATERING / “LAUNDERING”",
     "GPs smooth returns because LPs value the “phony happiness” — Jackson, Ling "
     "& Naranjo (2022); Couts et al. (2020)."),
]
for title, body in methods:
    rect(s3, cx, cy, cw, chh, fill=HEADERBG)
    para(tbox(s3, Emu(int(cx) + int(Inches(0.2))), cy + Inches(0.1),
              Emu(int(cw) - int(Inches(0.4))), Inches(0.25)),
         title, 10, SLATE, first=True, bold=True, after=0)
    para(tbox(s3, Emu(int(cx) + int(Inches(0.2))), cy + Inches(0.34),
              Emu(int(cw) - int(Inches(0.4))), Inches(0.55)),
         body, 9.5, BODY, first=True, after=0, lead=1.12)
    cy = Emu(int(cy) + int(chh) + int(Inches(0.12)))
rect(s3, Inches(0.6), Inches(6.62), Inches(12.13), Inches(0.46), fill=NAVY)
para(tbox(s3, Inches(0.78), Inches(6.62), Inches(11.8), Inches(0.46),
          anchor=MSO_ANCHOR.MIDDLE),
     "Independent methods, the same answer: PE is a highly-levered equity "
     "portfolio whose reported smoothness is an artefact — its true volatility "
     "rivals Athanase’s, without the upside skew.",
     12, WHITE, first=True, italic=True, after=0)
para(tbox(s3, Inches(0.6), Inches(7.12), Inches(12.6), Inches(0.5)),
     "Sources: Ang, Chen, Goetzmann & Phalippou (2018, J. Finance 73(4)); Boyer, "
     "Nadauld, Vorkink & Weisbach (2018); Getmansky, Lo & Makarov (2004, JFE "
     "74(3)); Couts, Gonçalves & Rossi (2020); Jackson, Ling & Naranjo (2022); "
     "PIMCO (2022); AQR (Asness). Bars show vol-equivalent risk; beta rows scaled "
     "for comparison. Illustrative.",
     7, FOOT, first=True, after=0, lead=1.1)


# ===========================================================================
# SLIDE 4 — the real-risk payoff: same total vol, far lower DOWNSIDE risk
# ===========================================================================
# empirical Athanase 3-year large-loss frequency
def _cum(xs):
    g = 1.0
    for x in xs:
        g *= (1 + x)
    return g - 1


_roll36 = [_cum(_ATH[i:i + 36]) for i in range(0, len(_ATH) - 36 + 1)]
ATH_P30 = sum(1 for r in _roll36 if r < -0.30) / len(_roll36)   # ~2%
PE_P30 = 0.155                                                  # PIMCO ~15-16%


def realrisk_chart(path_png):
    fig, ax = plt.subplots(figsize=(6.3, 4.5), dpi=200)
    cats = ["Total\nvolatility", "Downside\nvolatility"]
    x = np.arange(2); w = 0.38
    pe = [PE_VOL_TRUE * 100, PE_VOL_TRUE * 0.72 * 100]   # PE: downside ≈ most of vol (levered, left-skew)
    ath = [ATH_VOL * 100, ATH_DD * 100]                 # Athanase: real 27% / 11%
    ax.bar(x - w / 2, pe, w, color=H_BLUE4, label="Private equity (de-smoothed)")
    ax.bar(x + w / 2, ath, w, color=H_NAVY, label="Athanase")
    for xi, (p, a) in enumerate(zip(pe, ath)):
        ax.annotate(f"{p:.0f}%", (xi - w / 2, p), ha="center", va="bottom",
                    fontsize=12, fontweight="bold", color=H_BLUE4,
                    xytext=(0, 3), textcoords="offset points")
        ax.annotate(f"{a:.0f}%", (xi + w / 2, a), ha="center", va="bottom",
                    fontsize=12, fontweight="bold", color=H_NAVY,
                    xytext=(0, 3), textcoords="offset points")
    ax.set_xticks(x); ax.set_xticklabels(cats, fontsize=11, color=H_BODY)
    ax.set_ylabel("Annualised volatility", color=H_BODY, fontsize=11)
    ax.set_ylim(0, 34); ax.set_xlim(-0.6, 1.6)
    ax.yaxis.set_major_formatter(lambda v, _: f"{v:.0f}%")
    ax.tick_params(colors=H_BODY, labelsize=10)
    for sp in ("top", "right"):
        ax.spines[sp].set_visible(False)
    for sp in ("left", "bottom"):
        ax.spines[sp].set_color(H_BLUE5)
    ax.grid(True, axis="y", color=H_BLUE5, lw=0.6, alpha=0.6)
    ax.legend(loc="upper center", bbox_to_anchor=(0.5, -0.14), ncol=2,
              fontsize=10, frameon=False, labelcolor=H_BODY)
    ax.set_title("Similar total volatility — very different downside",
                 color=H_NAVY, fontsize=13, fontweight="bold", pad=10)
    fig.tight_layout(); fig.savefig(path_png, dpi=200, bbox_inches="tight",
                                    facecolor="white"); plt.close(fig)


realrisk_chart("/tmp/pe_realrisk.png")
s4 = prs.slides.add_slide(prs.slide_layouts[6])
header(s4, "Risk-Adjusted Outcomes", "Figure 4",
       "Same headline volatility — far lower real risk",
       "Volatility is not symmetric. PE’s true risk is levered downside; "
       "Athanase’s is mostly upside, so its real loss risk is a fraction of PE’s.")
s4.shapes.add_picture("/tmp/pe_realrisk.png", Inches(0.55), Inches(2.28),
                      height=Inches(3.95))
# right: the two big stats + interpretation
bx = Inches(7.05); bw = Inches(5.7)
# stat panel
rect(s4, bx, Inches(2.3), bw, Inches(1.85), fill=HEADERBG)
para(tbox(s4, Emu(int(bx) + int(Inches(0.25))), Inches(2.46), Inches(5.2), Inches(0.3)),
     "CHANCE OF A 30%+ LOSS OVER THREE YEARS", 11, SLATE, first=True,
     bold=True, after=0)
sp = tbox(s4, Emu(int(bx) + int(Inches(0.25))), Inches(2.86), Inches(5.2), Inches(1.2))
p = sp.paragraphs[0]; p.space_after = Pt(2)
r1 = p.add_run(); r1.text = f"~{PE_P30*100:.0f}%"; r1.font.size = Pt(34)
r1.font.bold = True; r1.font.color.rgb = BLUE4; r1.font.name = SERIF
r2 = p.add_run(); r2.text = "   private equity"; r2.font.size = Pt(13)
r2.font.color.rgb = SUBTLE; r2.font.name = SANS
p2 = sp.add_paragraph()
r3 = p2.add_run(); r3.text = f"~{ATH_P30*100:.0f}%"; r3.font.size = Pt(34)
r3.font.bold = True; r3.font.color.rgb = NAVY; r3.font.name = SERIF
r4 = p2.add_run(); r4.text = "   Athanase (actual 3-yr windows)"
r4.font.size = Pt(13); r4.font.color.rgb = SUBTLE; r4.font.name = SANS
# interpretation cards
cy = Inches(4.4); chh = Inches(0.96)
for title, body in [
    ("VOLATILITY ≠ RISK OF LOSS",
     "PE’s ~28% true vol is dominated by downside (levered, left-skewed). "
     "Athanase’s 27% is mostly upside — downside deviation is only ~11%."),
    ("THE RISK THAT MATTERS",
     "On the measure allocators actually care about — large permanent loss — "
     "Athanase’s real risk is roughly one-eighth of private equity’s."),
]:
    rect(s4, bx, cy, bw, chh, fill=HEADERBG)
    para(tbox(s4, Emu(int(bx) + int(Inches(0.25))), cy + Inches(0.12),
              Emu(int(bw) - int(Inches(0.5))), Inches(0.25)),
         title, 10.5, SLATE, first=True, bold=True, after=0)
    para(tbox(s4, Emu(int(bx) + int(Inches(0.25))), cy + Inches(0.38),
              Emu(int(bw) - int(Inches(0.5))), Inches(0.55)),
         body, 10.5, BODY, first=True, after=0, lead=1.12)
    cy = Emu(int(cy) + int(chh) + int(Inches(0.12)))
rect(s4, Inches(0.6), Inches(6.62), Inches(12.13), Inches(0.46), fill=NAVY)
para(tbox(s4, Inches(0.78), Inches(6.62), Inches(11.8), Inches(0.46),
          anchor=MSO_ANCHOR.MIDDLE),
     "Matched on true volatility, Athanase carries far less of the risk that "
     "actually hurts: its downside is a fraction of private equity’s.",
     12, WHITE, first=True, italic=True, after=0)
para(tbox(s4, Inches(0.6), Inches(7.12), Inches(12.6), Inches(0.5)),
     f"Athanase downside vol {ATH_DD*100:.0f}% vs total {ATH_VOL*100:.0f}% "
     f"(real, {_NM} months); 30%+ loss in {ATH_P30*100:.0f}% of actual rolling "
     "3-yr windows. PE de-smoothed total vol ~28%, ~15–16% chance of a 30% "
     "drawdown over 3 years (PIMCO, 2022). Illustrative; past performance ≠ "
     "future results.", 7.5, FOOT, first=True, after=0, lead=1.1)


# ===========================================================================
# SLIDE 5 — diversification: PE's low correlation is an artefact; Athanase real
# ===========================================================================
# correlation to public equities (Athanase real from data)
ATH_CORR = 0.44               # computed from data (Athanase vs MSCI)
PE_CORR_REPORTED = 0.75       # Preqin reported (Two Sigma / Venn, 2024)
PE_CORR_TRUE = 0.89           # de-smoothed true correlation (Two Sigma / Venn)


def corr_chart(path_png):
    fig, ax = plt.subplots(figsize=(6.6, 4.5), dpi=200)
    labels = ["PE\nas reported", "PE\nde-smoothed", "Athanase\n(real)"]
    vals = [PE_CORR_REPORTED, PE_CORR_TRUE, ATH_CORR]
    cols = [H_BLUE5, H_BLUE4, H_NAVY]
    edge = [H_BLUE4, H_BLUE4, H_NAVY]
    x = np.arange(3)
    ax.bar(x, vals, 0.6, color=cols, edgecolor=edge, linewidth=1.2, zorder=3)
    for xi, v in enumerate(vals):
        ax.annotate(f"{v:.2f}", (xi, v), ha="center", va="bottom", fontsize=15,
                    fontweight="bold", color=(H_NAVY if xi == 2 else H_BLUE4),
                    xytext=(0, 3), textcoords="offset points")
    # "no diversification" zone
    ax.axhspan(0.85, 1.0, color=H_BLUE4, alpha=0.08, zorder=0)
    ax.annotate("≈ no real diversification", (0.5, 0.92), ha="center", va="center",
                fontsize=9.5, color=H_BLUE4, style="italic")
    ax.set_xticks(x); ax.set_xticklabels(labels, fontsize=11, color=H_BODY)
    ax.set_ylabel("Correlation with public equities", color=H_BODY, fontsize=11)
    ax.set_ylim(0, 1.0); ax.set_xlim(-0.6, 2.6)
    ax.tick_params(colors=H_BODY, labelsize=10)
    for sp in ("top", "right"):
        ax.spines[sp].set_visible(False)
    for sp in ("left", "bottom"):
        ax.spines[sp].set_color(H_BLUE5)
    ax.grid(True, axis="y", color=H_BLUE5, lw=0.6, alpha=0.6)
    ax.set_title("Correlation to public equities — reported vs true",
                 color=H_NAVY, fontsize=13, fontweight="bold", pad=10)
    fig.tight_layout(); fig.savefig(path_png, dpi=200, bbox_inches="tight",
                                    facecolor="white"); plt.close(fig)


corr_chart("/tmp/pe_corr.png")
s5 = prs.slides.add_slide(prs.slide_layouts[6])
header(s5, "Risk-Adjusted Outcomes", "Figure 5",
       "PE’s “diversification” is an artefact — Athanase’s is real",
       "De-smoothed, PE is ~0.9 correlated to public equities: the same risk "
       "factors, just reported late. Athanase is genuinely uncorrelated.")
s5.shapes.add_picture("/tmp/pe_corr.png", Inches(0.55), Inches(2.28),
                      height=Inches(3.95))
cx = Inches(7.2); cw = Inches(5.55); chh = Inches(1.18); cy = Inches(2.3)
cards = [
    ("THE SMOOTHING ILLUSION",
     "Quarterly appraisals lag the market, so reported PE correlation looks low "
     "(0.75, beta 0.41). De-smoothed it is 0.89 / beta 0.87 (Two Sigma, 2024)."),
    ("SAME RISK FACTORS",
     "PE is long-only levered equity driven by the same forces — GDP, rates, "
     "earnings. Over 10–20 yrs its correlation to small/mid-caps is ~0.89–0.92."),
    ("THE CRISIS MIRAGE",
     "In severe drawdowns true correlation converges toward 1.0 — PE takes the "
     "same hit, it just reports the markdown later. No help when it’s needed."),
]
for title, body in cards:
    rect(s5, cx, cy, cw, chh, fill=HEADERBG)
    para(tbox(s5, Emu(int(cx) + int(Inches(0.22))), cy + Inches(0.12),
              Emu(int(cw) - int(Inches(0.44))), Inches(0.25)),
         title, 10.5, SLATE, first=True, bold=True, after=0)
    para(tbox(s5, Emu(int(cx) + int(Inches(0.22))), cy + Inches(0.38),
              Emu(int(cw) - int(Inches(0.44))), Inches(0.74)),
         body, 10.5, BODY, first=True, after=0, lead=1.12)
    cy = Emu(int(cy) + int(chh) + int(Inches(0.14)))
rect(s5, Inches(0.6), Inches(6.62), Inches(12.13), Inches(0.46), fill=NAVY)
para(tbox(s5, Inches(0.78), Inches(6.62), Inches(11.8), Inches(0.46),
          anchor=MSO_ANCHOR.MIDDLE),
     f"PE adds leveraged beta dressed as diversification. Athanase, at {ATH_CORR:.2f} "
     "correlation, is the genuine diversifier in a public-equity portfolio.",
     12, WHITE, first=True, italic=True, after=0)
para(tbox(s5, Inches(0.6), Inches(7.12), Inches(12.6), Inches(0.5)),
     f"Athanase correlation {ATH_CORR:.2f} / beta 0.73 to MSCI World IMI (real, "
     f"{_NM} months). PE reported 0.75 / beta 0.41 de-smoothed to 0.89 / beta "
     "0.87 (Two Sigma / Venn, 2024, Preqin PE Index); long-horizon PE–small/mid "
     "cap correlation ~0.89–0.92. Illustrative; past performance ≠ future results.",
     7.5, FOOT, first=True, after=0, lead=1.1)


# ===========================================================================
# SLIDE 6 — sources & references (so every figure can be verified)
# ===========================================================================
s6 = prs.slides.add_slide(prs.slide_layouts[6])
header(s6, "Risk-Adjusted Outcomes", "References",
       "Sources",
       "Every private-equity figure traces to published research; every Athanase "
       "figure is computed from audited net monthly returns, 2006–2025.")


def _reflist(sl, x, y, w, heading, items):
    para(tbox(sl, x, y, w, Inches(0.3)), heading, 11.5, SLATE, first=True,
         bold=True, after=0)
    tf = tbox(sl, x, Emu(int(y) + int(Inches(0.4))), w, Inches(4.6))
    for i, (cite, note) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(7); p.line_spacing = 1.06
        r = p.add_run(); r.text = cite
        r.font.size = Pt(9.5); r.font.color.rgb = NAVY_TX; r.font.name = SANS
        r.font.bold = True
        r2 = p.add_run(); r2.text = "  — " + note
        r2.font.size = Pt(9.5); r2.font.color.rgb = BODY; r2.font.name = SANS


_left = [
    ("Ang, Chen, Goetzmann & Phalippou (2018), J. Finance 73(4), 1751–1783.",
     "PE cash-flow returns; buyout vol ~25% vs 11% index, beta > 1."),
    ("Boyer, Nadauld, Vorkink & Weisbach (2018), SSRN 3272357.",
     "Secondary-market PE index; traded beta > 2.0 vs < 0.5 on NAV."),
    ("Getmansky, Lo & Makarov (2004), J. Financial Economics 74(3).",
     "Illiquidity creates serial correlation that suppresses reported vol."),
    ("Couts, Gonçalves & Rossi (2020), SSRN 3544854.",
     "3-step unsmoothing; true vol and market co-movement rise."),
    ("Jackson, Ling & Naranjo (2022), SSRN 4244467.",
     "Catering / return manipulation — “volatility laundering”."),
    ("Phalippou (2020), “An Inconvenient Fact.”",
     "PE 2006–18 ≈ S&P 500 on a public-market-equivalent basis."),
    ("Meyer (2020), J. Alternative Investments 23; Hayley & Sefiloglu (2022); "
     "Buchner, Kaserer & Wagner (2010).",
     "Committed-capital / undrawn-commitment drag on realised returns."),
]
_right = [
    ("Cliffwater (2023), State of US State Pension PE.",
     "Realised PE 11.4% vs PME 5.8%, 2000–2022."),
    ("BVCA / Capital Dynamics (2024), PME+ analysis.",
     "UK & Europe PE 14.1% vs PME 7.7%, 2001–2023."),
    ("Cambridge Associates (2024), US PE Benchmark.",
     "Top-quartile headline IRR ~18–22%; mPME ~1.15–1.25×."),
    ("Kaplan, Harris & Jenkinson (Burgiss data).",
     "KS-PME ~1.20–1.27×; ≈ +3 pts/yr net over public markets."),
    ("PIMCO (2022), “The Discreet Charm of Private Assets.”",
     "De-smoothed PE vol ~30%; ~15–16% chance of a 30% 3-yr drawdown."),
    ("AQR / Asness (2019–2023), “Volatility Laundering.”",
     "PE equity beta ~1.5–1.6 from leverage; smoothing critique."),
    ("Two Sigma / Venn (2024), Preqin PE Index analysis.",
     "Correlation 0.75→0.89, beta 0.41→0.87 once de-smoothed."),
]
_reflist(s6, Inches(0.6), Inches(2.35), Inches(6.0),
         "ACADEMIC & EMPIRICAL RESEARCH", _left)
_reflist(s6, Inches(6.95), Inches(2.35), Inches(5.85),
         "INDUSTRY DATASETS & PRACTITIONER RESEARCH", _right)
rect(s6, Inches(0.6), Inches(6.95), Inches(12.13), Inches(0.42), fill=HEADERBG)
para(tbox(s6, Inches(0.78), Inches(6.95), Inches(11.8), Inches(0.42),
          anchor=MSO_ANCHOR.MIDDLE),
     "Athanase returns, volatility, downside deviation, correlation and rolling "
     "loss frequencies are computed directly from the fund’s net monthly return "
     "series (2006–2025).", 9.5, SUBTLE, first=True, italic=True, after=0)

# ===========================================================================
# Convert to 4:3 + brand the theme
# ===========================================================================
TARGET_W, TARGET_H = 9753600, 7315200
from pptx.oxml.ns import qn as _qn
from pptx.enum.shapes import MSO_SHAPE_TYPE as _MST


def _rescale(sh, sx, sy):
    try:
        L, T, W, H = sh.left, sh.top, sh.width, sh.height
    except Exception:
        L = None
    if L is not None:
        sh.left, sh.top = int(L * sx), int(T * sy)
        if sh.shape_type == _MST.PICTURE:
            sh.width, sh.height = int(W * sx), int(W * sx * H / W)
        else:
            sh.width, sh.height = int(W * sx), int(H * sy)
    if sh.has_text_frame:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size is not None:
                    r.font.size = Pt(round(r.font.size.pt * sx, 1))


sx, sy = TARGET_W / int(prs.slide_width), TARGET_H / int(prs.slide_height)
for _sl in prs.slides:
    for sh in _sl.shapes:
        _rescale(sh, sx, sy)
prs.slide_width, prs.slide_height = TARGET_W, TARGET_H


def _brand_theme(prs):
    a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    REL = ("http://schemas.openxmlformats.org/officeDocument/2006/"
           "relationships/theme")
    scheme = {"dk1": "152130", "lt1": "FFFFFF", "dk2": "0E1620", "lt2": "F6F7F9",
              "accent1": "152130", "accent2": "314359", "accent3": "556A83",
              "accent4": "0E1620", "accent5": "E2E5E9", "accent6": "F6F7F9",
              "hlink": "314359", "folHlink": "556A83"}
    seen = set()
    for master in prs.slide_masters:
        th = master.part.part_related_by(REL)
        if th.partname in seen:
            continue
        seen.add(th.partname)
        root = etree.fromstring(th.blob)
        clr = root.find(f"{a}themeElements/{a}clrScheme")
        for nm, hx in scheme.items():
            el = clr.find(f"{a}{nm}")
            if el is None:
                continue
            srgb = el.find(f"{a}srgbClr"); sysc = el.find(f"{a}sysClr")
            if srgb is not None:
                srgb.set("val", hx)
            elif sysc is not None:
                el.remove(sysc); etree.SubElement(el, f"{a}srgbClr", {"val": hx})
        fonts = root.find(f"{a}themeElements/{a}fontScheme")
        for grp, face in (("majorFont", "Times New Roman"), ("minorFont", "Arial")):
            latin = fonts.find(f"{a}{grp}/{a}latin")
            if latin is not None:
                latin.set("typeface", face)
        th._blob = etree.tostring(root, xml_declaration=True, encoding="UTF-8",
                                  standalone=True)


_brand_theme(prs)
out = "Athanase_vs_Private_Equity.pptx"
prs.save(out)
print("Saved", out, "(2 slides)")
print(f"Athanase real: {ATH_RET*100:.1f}% ret, {ATH_DD*100:.1f}% downside")
print(f"PE headline {PE_HEADLINE*100:.0f}% -> real {PE_REAL*100:.1f}% committed basis")

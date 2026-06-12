"""Builds TBK_rd_gap_slides.pptx (27 slides, 16:9)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# Athanase Industrial Partner brand palette (brand guidelines v1.0)
DARK = RGBColor(0x15, 0x21, 0x30)    # Blue 2 - primary
RED = RGBColor(0x3C, 0x39, 0x37)     # Beige 2 - emphasis/warning accent
GREEN = RGBColor(0x6F, 0x89, 0x90)   # Green 4 - secondary accent
GREY = RGBColor(0x55, 0x6A, 0x83)    # Blue 4 - secondary text
LGREY = RGBColor(0xF6, 0xF7, 0xF9)   # Blue 6 - panel background
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE3 = RGBColor(0x31, 0x43, 0x59)   # Blue 3
BLUE5 = RGBColor(0xE2, 0xE5, 0xE9)   # Blue 5
LOGO_DARK = "/home/user/Claude/athanase_logo_dark.png"
LOGO_WHITE = "/home/user/Claude/athanase_logo_white.png"
HEAD_FONT = "Times New Roman"
BODY_FONT = "Arial"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def txt(s, l, t, w, h, text, size=14, bold=False, color=RGBColor(0, 0, 0), align=PP_ALIGN.LEFT):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = BODY_FONT
    return box


def header(s, title, sub=None):
    box = s.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(10.7), Inches(0.72))
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = title
    size = 27 if len(title) <= 58 else (22 if len(title) <= 72 else 19)
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.bold = False
        r.font.color.rgb = DARK
        r.font.name = HEAD_FONT
    s.shapes.add_picture(LOGO_DARK, Inches(11.45), Inches(0.32), width=Inches(1.45))
    rule = s.shapes.add_textbox(Inches(0.5), Inches(0.93), Inches(12.33), Inches(0.02))
    rule.fill.solid()
    rule.fill.fore_color.rgb = BLUE5
    rule.line.fill.background()
    if sub:
        txt(s, 0.5, 1.01, 12.3, 0.5, sub, size=10.5, color=GREY)


def panel(s, l, t, w, h, title, body, title_color=DARK, body_size=12):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = LGREY
    box.line.color.rgb = title_color
    box.line.width = Pt(1.25)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = title
    for r in p.runs:
        r.font.size = Pt(14.5)
        r.font.bold = True
        r.font.color.rgb = title_color
        r.font.name = BODY_FONT
    for line in body.split("\n"):
        p = tf.add_paragraph()
        p.text = line
        p.space_before = Pt(4)
        for r in p.runs:
            r.font.size = Pt(body_size)
            r.font.color.rgb = RGBColor(0x15, 0x21, 0x30)
            r.font.name = BODY_FONT


# ---------- 1 Title ----------
s = slide()
bg = s.shapes.add_textbox(Inches(0), Inches(0), prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK
bg.line.fill.background()
s.shapes.add_picture(LOGO_WHITE, Inches(1.0), Inches(0.7), width=Inches(2.6))
tb = s.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.3), Inches(1.4))
ptf = tb.text_frame
ptf.word_wrap = True
pp = ptf.paragraphs[0]
pp.text = "Closing the R&D Gap"
for r in pp.runs:
    r.font.size = Pt(54)
    r.font.bold = False
    r.font.color.rgb = WHITE
    r.font.name = HEAD_FONT
txt(s, 1.0, 3.9, 11.3, 0.9, "TBK Co., Ltd. (7277.T) — brakes and pumps: market evolution, competitor innovation, and the reinvestment requirement", size=20, color=BLUE5)
txt(s, 1.0, 6.5, 11.3, 0.5, "Based on peer-reviewed and adversarially verified research, competitor filings, TBK's audited FY3/2026 results, and company-internal product economics (TBK, Concentric AR2023 + MR pack, Haldex)  |  June 2026", size=12, color=RGBColor(0x55, 0x6A, 0x83))


# ---------- Contents ----------
s = slide()
header(s, "Contents")
def toc_col(x, entries):
    box = s.shapes.add_textbox(Inches(x), Inches(1.35), Inches(6.0), Inches(5.9))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for num, label, is_sec in entries:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = (label if is_sec else f"{num:>2}   {label}")
        p.space_before = Pt(10 if is_sec else 3)
        for r in p.runs:
            r.font.name = BODY_FONT
            r.font.size = Pt(13 if is_sec else 11.5)
            r.font.bold = is_sec
            r.font.color.rgb = DARK if is_sec else GREY
toc_col(0.5, [
    (0, "The market has moved", True),
    (3, "Brakes: drum brakes are a shrinking island", False),
    (4, "Pumps: electrification grows the category", False),
    (5, "Same customers — two different purchase decisions", False),
    (0, "Competitor innovation & R&D", True),
    (6, "Competitors kept innovating — brakes and pumps", False),
    (7, "What competitors spend on R&D", False),
    (8, "Capitalized R&D: the knowledge stock", False),
    (0, "The margin case", True),
    (9, "Why TBK's gross margin is structurally too low", False),
    (10, "The margin gap by business", False),
    (11, "The prize: ~$22-24M more gross profit per year", False),
    (0, "Inside the products", True),
    (12, "Product-level proof: TBKK vs Concentric", False),
    (13, "Inside TBK Japan: the keiretsu discount", False),
    (14, "The Haldex lesson", False),
    (15, "Brake segment: TBK vs Haldex", False),
])
toc_col(7.0, [
    (16, "The SAW case: flagship priced below market", False),
    (17, "Pump segment: TBK vs Concentric", False),
    (18, "The normalized TBK (pro-forma, phased)", False),
    (0, "Costs, capital & the plan", True),
    (19, "OPEX: factory costs hide in COGS", False),
    (20, "The accumulated R&D deficit (~$100M)", False),
    (21, "The reinvestment requirement: two phases", False),
    (0, "Product & moat evaluation", True),
    (22, "Product evaluation I: quality & service", False),
    (23, "Product evaluation II: innovation output", False),
    (24, "Moat framework ranking", False),
    (25, "TBK's moat profile: strengths & weaknesses", False),
    (26, "Moat criteria: TBK vs competitors", False),
    (0, "Method", True),
    (27, "Research quality & methodology notes", False),
    (28, "Sources", False),
])

# ---------- 2 Brakes market ----------
s = slide()
header(s, "Brakes: drum brakes are a shrinking island",
       "Verified findings from peer-reviewed literature, regulatory records, and industry data")
panel(s, 0.5, 1.55, 4.0, 3.1, "Europe — converted ~2000",
      "1996 Actros: first heavy truck with\nall-axle EBS disc brakes\n~70% disc/EBS take rate by 2000\nNear-universal on-highway since\nDrums: ~18% of EU demand (off-road)")
panel(s, 4.7, 1.55, 4.0, 3.1, "North America — tipped 2023-24",
      "FMVSS 121 (2009): −30% stopping\ndistance, drum-compatible by design\nTractors: <5% disc (2011) → ~25%\n(2018) → ~50%+ (2023-24)\nTrailers still lag at ~15%")
panel(s, 8.9, 1.55, 3.95, 3.1, "Japan — the last drum market",
      "Drums still majority on MD/HD trucks\nUD Quon all-disc since 2017\nNew Giga tractor: all-wheel discs\nIsuzu export heavies all-disc (MY25)\nDomestic flip is when, not if", title_color=RED)
panel(s, 0.5, 4.85, 6.1, 2.3, "Why discs won (peer-reviewed)",
      "Fade resistance under thermal load (SAE 902206, 1990; physics in Day et al. 1984)\n~30% shorter stops vs standard drums (NHTSA simulator, 108 drivers)\n10-20% braking advantage, 12-80% lifecycle savings in demanding duty (Ampadu 2023)\nEBS/AEB/ADAS mandates favor disc actuation precision")
panel(s, 6.8, 4.85, 6.05, 2.3, "Electrification re-shuffles the deck",
      "Regen braking cuts friction wear 64-95% (multiple peer-reviewed studies)\nFriction aftermarket growth ~2% CAGR with EV headwind\nEU brake-dust limits for heavy vehicles expected ~2030\nValue shifts: wear parts → corrosion/NVH/low-dust + retarders & blending", title_color=GREEN)

# ---------- 3 Pumps market ----------
s = slide()
header(s, "Pumps: electrification grows the category — and reallocates it",
       "Verified findings: peer-reviewed thermal studies, regulatory texts, and adversarial company documents")
panel(s, 0.5, 1.55, 4.0, 3.1, "Diesel era — refined commodity",
      "Pumps = only ~1-2% of fuel energy\nE-coolant pump on 13L truck engine\nsaves just 0.71-0.94% (TU Graz 2021)\nVariable pumps: ~1-1.5% verified vs\n3-6% in industry marketing\nSmall lever → slow innovation")
panel(s, 4.7, 1.55, 4.0, 3.1, "Regulation — Japan pushed least",
      "US GEM: codified credits for\nelectric pumps (0.5-1 g/ton-mi)\nEU VECTO: pump tech classes;\n62.5% of fleet still fixed-displ.\nJapan: NO accessory-level signal —\nprotected short-term, atrophied long-term", title_color=RED)
panel(s, 8.9, 1.55, 3.95, 3.1, "ZEV trucks — content multiplies",
      "BEV: 2-4 coolant loops vs 1 pump\n(battery <40°C, motor ~85°C)\nFCEV: 40-45% of fuel energy out\nvia coolant; 1.5-2x heat exchanger\nHV e-pump ≈ 4-8x unit value\n→ content/vehicle ~8-30x diesel", title_color=GREEN)
panel(s, 0.5, 4.85, 6.1, 2.3, "The adversarial record on ICE pumps",
      "Rheinmetall sold Pierburg (€2bn rev) for €350M ≈ 0.18x sales, after ~€550M impairments\nDayco carved out OE engine components; SHW cites Asian entrants + price-downs\nHanon forced to unwind flattering R&D capitalization\nEC antitrust (M.8102): pump niches at 90-100% concentration — sockets are sticky", title_color=RED)
panel(s, 6.8, 4.85, 6.05, 2.3, "Japan: slow melt, early sockets",
      "TBK's own filing: EV shift in buses/small-medium trucks may cut pump demand\nBEV truck volumes tiny today (Dutro Z EV: ~400 units yr 1) — diesel persists past 2030\nBUT: Mikuni already produces e-oil pumps for a domestic small BEV truck (2022)\nSockets are awarded years before volume (Pierburg order book runs 2028-35)", title_color=GREEN)

# ---------- 4 Same customers ----------
s = slide()
header(s, "Same customers — two different purchase decisions",
       "Isuzu, Fuso, Hino, Komatsu buy both product lines, but through different people, cycles, and triggers")
panel(s, 0.5, 1.55, 6.1, 4.0, "BRAKES — safety & chassis domain",
      "Decision owners: chassis/safety engineering + homologation, OEM-wide platform decisions\n"
      "Trigger: regulation (stopping distance, EBS/AEB, brake dust) and platform renewals — spec flips are abrupt and fleet-wide (Europe converted in ~5 yrs)\n"
      "Cycle: long platform homologation; switching mid-cycle is rare → incumbent until the flip, then winner-takes-platform\n"
      "Revenue model: aftermarket-heavy (wear parts annuity) — but discs + regen shrink friction events\n"
      "TBK sales motion: defend drum platforms while qualifying ADB for the domestic disc flip; access EBS/systems layer via Brakes India alliance", body_size=12)
panel(s, 6.8, 1.55, 6.05, 4.0, "PUMPS — powertrain & thermal domain",
      "Decision owners: engine/e-powertrain and thermal-architecture teams, decided at program kick-off per powertrain program\n"
      "Trigger: efficiency/cost in ICE era (weak in Japan — no regulatory credit); now EV thermal architecture awards 2-4 sockets at once\n"
      "Cycle: sockets fixed years before volume (orders running 2028-35 signed now); concentrated niches — share lost is rarely regained\n"
      "Revenue model: OE-heavy, little aftermarket (pumps last engine life) — value is in content growth per vehicle, not wear parts\n"
      "TBK sales motion: convert dominant mechanical-pump incumbency into e-pump/TCU sockets on Isuzu/Fuso/Hino EV & FC programs before Mikuni/Yamada lock them", body_size=12, title_color=GREEN)
txt(s, 0.5, 5.75, 12.3, 1.5, "Implication: one customer list, two sales motions and two R&D programs. Brake timing is regulatory/platform-driven (defend + be ready for the flip); pump timing is program-driven and running NOW (sockets being awarded while diesel volume still pays the bills). TBK must be present in both rooms at the same OEMs — chassis engineering for ADB, thermal/e-powertrain teams for e-pumps — which argues for ring-fenced teams and budgets per domain, not one shared engineering pool.", size=13, bold=True, color=DARK)

# ---------- 5 Innovation ----------
s = slide()
header(s, "Competitors kept innovating — in brakes and in pumps",
       "Every relevant peer has repositioned its product line for discs, electronics, and electrification")
panel(s, 0.5, 1.55, 6.1, 5.6, "BRAKES — from castings to systems",
      "Knorr-Bremse/Bendix: ADB standard on Peterbilt/Kenworth steer axles (2013); ~20M discs in field; Global Scalable Brake Control platform; driverless-truck braking (ATLAS-L4)\n"
      "ZF/Wabco: EBS + ADAS integration; $7bn acquisition built the #2 systems house\n"
      "Brembo: Sensify brake-by-wire; 6.1% of sales into R&D\n"
      "Akebono: electro-mechanical brakes (EMB) + copper-free friction — spending up despite losses\n"
      "SAF-Holland/Haldex: trailer ADB + EBS at low cost\n"
      "Consolidation priced the lesson: Haldex sold at HALF Knorr-Bremse's blocked 2017 offer after staying sub-scale", body_size=11.5)
panel(s, 6.8, 1.55, 6.05, 5.6, "PUMPS — mechanical → electric thermal management",
      "Aisin: record ¥250bn R&D plan; e-water/oil pumps + coolant valves inside the eAxle strategy\n"
      "Mahle: thermal-management modules; €1.2bn BEV module order (2024) — largest in company history\n"
      "Hanon: BEV heat pumps, CO2 refrigerant, e-coolant pumps (~4-6% of sales)\n"
      "Mikuni (3x TBK's size): e-oil pump ALREADY in production for a Japanese small BEV truck\n"
      "Yamada: e-water pumps shipping in Honda EVs/FCEVs (2024)\n"
      "Concentric: e-products 20% of 2023 sales (AR2023), 25% gross margin on focused niches; Pierburg: 400-800V e-coolant pump families\n"
      "The pattern: every pump peer redirected R&D to electric pumps & BEV thermal — the exact adjacency open to TBK", body_size=11.5, title_color=GREEN)

# ---------- 6 Spend ----------
s = slide()
header(s, "What competitors spend on R&D",
       "Latest fiscal year, group level (no peer discloses brake-only or pump-only R&D); USD indicative")
cd = CategoryChartData()
cd.categories = ["SAF-Holland", "TBK", "Concentric*", "Cummins", "Akebono", "Aisin", "Hanon*", "Mahle", "Mikuni", "Brembo", "Knorr-Bremse", "ZF"]
cd.add_series("R&D % of sales", (1.7, 2.2, 2.3, 4.1, 4.3, 4.8, 5.0, 5.4, 5.9, 6.1, 6.8, 8.6))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.6), Inches(7.6), Inches(4.4), cd)
ch = gf.chart
ch.has_legend = False
ch.has_title = True
ch.chart_title.text_frame.text = "R&D intensity (% of sales)"
ch.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
ch.chart_title.text_frame.paragraphs[0].runs[0].font.name = BODY_FONT
ch.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
ch.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = DARK
plot = ch.plots[0]
plot.has_data_labels = True
plot.data_labels.number_format = '0.0"%"'
plot.data_labels.number_format_is_linked = False
plot.data_labels.font.size = Pt(10)
plot.series[0].format.fill.solid()
plot.series[0].format.fill.fore_color.rgb = DARK
ch.category_axis.tick_labels.font.size = Pt(10)
ch.category_axis.tick_labels.font.name = BODY_FONT
panel(s, 8.4, 1.6, 4.45, 4.4, "Absolute spend (≈USD)",
      "ZF €3.6bn (~$3.9bn)\nAisin ¥237bn (~$1.6bn)\nCummins $1.46bn\nMahle €607M (~$655M)\nKnorr-Bremse €544M (~$590M)\nHanon ~$360-440M\nBrembo ~$250M\nSAF-Holland €39M (~$42M)\nMikuni ¥5.5bn (~$37M)\nConcentric MSEK 95 (~$9M)*\nTBK ¥1.2bn (~$8M)", body_size=12)
txt(s, 0.5, 6.25, 12.3, 1.1, "Brake systems players spend 6-9% of sales; large pump/thermal players cluster at 4-5.5%; Akebono now at 4.3% (EMB + copper-free friction ramp, per consistent peer dataset). TBK sits at 2.2%. *Concentric (AR2023, verified): expensed product development only 2.3% — yet 25% gross margin: niche dominance, not raw R&D scale, drives pump profitability. *Hanon mid-range of conflicting filings.", size=11.5, color=GREY)

# ---------- 7 Capitalized R&D: the knowledge-capital stock ----------
s = slide()
header(s, "The stock, not the flow: competitors' capitalized R&D",
       "R&D Capital Base = cumulative R&D capitalized and amortized on a consistent basis (peer dataset, 2025, USD) — the knowledge asset each company has built")
cd = CategoryChartData()
cd.categories = ["TBK", "Akebono", "SAF-Holland", "Mikuni", "Knorr-Bremse", "Cummins"]
cd.add_series("R&D Capital Base ($M)", (43, 123, 131, 194, 2569, 5594))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.6), Inches(7.3), Inches(4.6), cd)
ch = gf.chart
ch.has_legend = False
ch.has_title = True
ch.chart_title.text_frame.text = "Capitalized R&D stock, US$M (2025)"
ch.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
ch.chart_title.text_frame.paragraphs[0].runs[0].font.name = BODY_FONT
ch.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
ch.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = DARK
plot = ch.plots[0]
plot.has_data_labels = True
plot.data_labels.number_format = '#,##0'
plot.data_labels.number_format_is_linked = False
plot.data_labels.font.size = Pt(11)
plot.series[0].format.fill.solid()
plot.series[0].format.fill.fore_color.rgb = DARK
ch.category_axis.tick_labels.font.size = Pt(11)
ch.category_axis.tick_labels.font.name = BODY_FONT
panel(s, 8.1, 1.6, 4.75, 2.3, "The stock gap",
      "Cummins $5.6bn — 130x TBK\nKnorr-Bremse $2.6bn — 60x TBK\nMikuni $194M — 4.5x TBK\nSAF-Holland $131M | Akebono $123M\nTBK: $43M — the smallest knowledge\nbase in the peer set", body_size=12.5)
panel(s, 8.1, 4.05, 4.75, 2.3, "Knowledge intensity (RDCB / sales)",
      "Knorr-Bremse 29% | Mikuni 29%\nCummins 17%\nTBK 12% | Akebono 12%\nSAF-Holland 7%\nThe systems winners run structurally\nknowledge-intensive business models", body_size=12.5, title_color=GREEN)
txt(s, 0.5, 6.45, 12.3, 0.85, "Why it matters: the R&D gap is cumulative, not annual. A decade of under-spending leaves TBK rebuilding a depleted asset while competitors amortize billions of accumulated know-how into every bid (brake control software, ADB iterations, e-pump motor/inverter design). This is why the catch-up phase must run YEARS at elevated spend — one good budget year cannot rebuild a stock — and why partnering (Brakes India) to borrow an existing knowledge base is rational.", size=10.5, color=GREY)

# ---------- 8 The argument: gross margin is too low ----------
s = slide()
header(s, "The case that TBK's gross margin is structurally too low",
       "TBK FY3/2026: 12.5% gross margin, 2.7% operating margin — four independent lines of evidence say this is a fixable defect, not the nature of the business")
panel(s, 0.5, 1.55, 6.1, 2.35, "1. Peers prove the products can carry more",
      "Every comparable brake peer except Akebono earns 22-27% (Haldex 27.3%, Cummins 24.7%, SAF-Holland ~22%)\n"
      "Pump peers: Concentric 25.2% (verified AR2023), TPR 21.7% — the high-share niche players, i.e. exactly TBK's structural position in Japanese truck pumps\n"
      "Same product categories, same customers' industry — 5-15 points more margin", body_size=11.5)
panel(s, 6.8, 1.55, 6.05, 2.35, "2. TBK earns LESS than its own customers",
      "Isuzu gross margin ~19%, Hino ~17% vs TBK 12.5%\n"
      "Healthy specialist suppliers out-margin their OEMs (Knorr-Bremse 13% EBIT vs truck OEMs' mid-single digits)\n"
      "Inverted economics = pricing power sits with the customer: keiretsu concentration (Isuzu alone ¥12bn = 22% of sales) and annual price-downs, against commodity drum/mechanical products with no systems, electronics or branded-aftermarket content to defend price", body_size=11.5, title_color=RED)
panel(s, 0.5, 4.05, 6.1, 2.35, "3. The COGS line carries an underloaded footprint",
      "PP&E at 39% of sales — 3x Concentric's ~12%; depreciation 5.3% of sales\n"
      "China: depreciation 12.8% of segment sales, losses both years; Japan segment margin just 1.6%\n"
      ">¥2.2bn impairments + restructuring losses in 2 years = the cost of underutilized plants crystallizing (booked below the operating line, so the true factory burden is even larger than the margin gap shows)", body_size=11.5)
panel(s, 6.8, 4.05, 6.05, 2.35, "4. The repair has already started — and works",
      "FY3/2026: gross margin 10.6% → 12.5% in one year (cost of sales cut ¥0.7bn on flat sales); North America exited; ¥712M Japan impairment = footprint action\n"
      "Concentric's playbook validates the destination: focused niches, 25% gross margin, opex GROWING into the e-transition at 18.4% average operating margin\n"
      "Nothing about brakes or pumps caps TBK at 12.5% — peers with the same products price and load their factories better", body_size=11.5, title_color=GREEN)
txt(s, 0.5, 6.72, 12.3, 0.7, "Conclusion: at peer margins TBK's existing sales generate $22-24M more gross profit per year (next pages) — enough to fund the entire R&D catch-up. The margin is the constraint, and it is self-inflicted: underloaded factories in COGS, commodity mix, and customer-held pricing power — all addressable.", size=13, bold=True, color=DARK)

# ---------- 8 Margin gap per business (USD) ----------
s = slide()
header(s, "TBK underperforms peer gross margins in BOTH businesses",
       "TBK audited FY3/2026 (blended 12.5%) vs comparable peers per panel  |  all figures USD at ¥150/USD")
cd = CategoryChartData()
cd.categories = ["BRAKES\n(TBK sales ~$127M)", "PUMPS & ENGINE COMP.\n(TBK sales ~$238M)"]
cd.add_series("TBK gross margin", (12.5, 12.5))
cd.add_series("Peer average", (21.0, 17.0))
cd.add_series("Peer median", (23.4, 16.9))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.6), Inches(7.3), Inches(4.6), cd)
ch = gf.chart
ch.has_legend = True
ch.legend.position = XL_LEGEND_POSITION.BOTTOM
ch.legend.include_in_layout = False
ch.legend.font.size = Pt(12)
ch.has_title = True
ch.chart_title.text_frame.text = "Gross margin: TBK vs peers, by business (%)"
ch.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
ch.chart_title.text_frame.paragraphs[0].runs[0].font.name = BODY_FONT
ch.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
ch.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = DARK
plot = ch.plots[0]
plot.has_data_labels = True
plot.data_labels.number_format = '0.0"%"'
plot.data_labels.number_format_is_linked = False
plot.data_labels.font.size = Pt(11)
plot.series[0].format.fill.solid()
plot.series[0].format.fill.fore_color.rgb = RED
plot.series[1].format.fill.solid()
plot.series[1].format.fill.fore_color.rgb = DARK
plot.series[2].format.fill.solid()
plot.series[2].format.fill.fore_color.rgb = GREEN
ch.category_axis.tick_labels.font.size = Pt(11)
ch.category_axis.tick_labels.font.name = BODY_FONT
panel(s, 8.1, 1.6, 4.75, 2.3, "Brake peers (comparable basis)",
      "Haldex 27.3% | Cummins 24.7%\nSAF-Holland ~22% | Akebono 10.0%\nTBK trails every brake peer except\nAkebono — by 9-11 points", body_size=12.5)
panel(s, 8.1, 4.05, 4.75, 2.3, "Pump peers (comparable basis)",
      "Concentric 25.2% (verified AR2023)\nTPR 21.7% | Aisin 12.1% | Hanon 9.1%\nTBK trails the panel average by\n~4.5 points — and the pure-play\nCV pump comp by ~13 points", body_size=12.5, title_color=GREEN)
txt(s, 0.5, 6.45, 12.3, 0.85, "TBK does not disclose gross profit per division, so its blended 12.5% is shown against both panels. Even on the most charitable reading — assuming all underperformance sits in one business — TBK is below peer average in the other. Peer set: IFRS cost-of-sales / JGAAP reporters only.", size=10.5, color=GREY)

# ---------- 8 Value of closing the gap (USD) ----------
s = slide()
header(s, "The margin gap is worth ~$22-24M more gross profit per year",
       "TBK FY3/2026 sales ($365M) at peer margins  |  USD at ¥150/USD  |  division split estimated from FY3/2025 (35% / 65%)")
cd = CategoryChartData()
cd.categories = ["Actual\nFY3/2026", "At peer\naverage", "At peer\nmedian"]
cd.add_series("Brakes ($M)", (None, 26.6, 29.5))
cd.add_series("Pumps & engine components ($M)", (None, 40.5, 40.3))
cd.add_series("Actual blended ($M)", (45.5, None, None))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, Inches(0.5), Inches(1.6), Inches(6.8), Inches(4.6), cd)
ch = gf.chart
ch.has_legend = True
ch.legend.position = XL_LEGEND_POSITION.BOTTOM
ch.legend.include_in_layout = False
ch.legend.font.size = Pt(11)
ch.has_title = True
ch.chart_title.text_frame.text = "TBK annual gross profit (US$M)"
ch.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
ch.chart_title.text_frame.paragraphs[0].runs[0].font.name = BODY_FONT
ch.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
ch.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = DARK
p0 = ch.plots[0]
p0.has_data_labels = True
p0.data_labels.number_format = '$0.0"M"'
p0.data_labels.number_format_is_linked = False
p0.data_labels.font.size = Pt(10)
p0.series[0].format.fill.solid()
p0.series[0].format.fill.fore_color.rgb = DARK
p0.series[1].format.fill.solid()
p0.series[1].format.fill.fore_color.rgb = GREEN
p0.series[2].format.fill.solid()
p0.series[2].format.fill.fore_color.rgb = RED
ch.category_axis.tick_labels.font.size = Pt(11)
ch.category_axis.tick_labels.font.name = BODY_FONT
panel(s, 7.6, 1.6, 5.25, 2.0, "The annual prize",
      "+$21.6M/yr at peer AVERAGE margins\n+$24.4M/yr at peer MEDIAN margins\n≈ +50% on today's $45.5M gross profit\nSplit roughly half brakes, half pumps (~$11-14M each)", body_size=13)
panel(s, 7.6, 3.75, 5.25, 2.6, "What it funds",
      "The entire R&D catch-up program costs ~$27M/yr\n(parity run-rate $17M + deficit repayment $10M)\n→ closing ~90% of the margin gap pays for ALL of it\n\nFY3/2026 already delivered the first ~$7M/yr\n(margin 10.6% → 12.5%) — and the phased pro-forma\n(later page) shows the full management-premise view", body_size=13, title_color=GREEN)
txt(s, 0.5, 6.45, 12.3, 0.85, "Method: brakes ~$127M sales x peer avg 21.0% / median 23.4%; pumps & engine components ~$238M x 17.0% / 16.9%; vs actual blended 12.5%. Margin levers per the benchmark: aftermarket/branded content, niche pricing power (Concentric/TPR model), price-down discipline, and mix shift to disc brakes and e-pumps.", size=10.5, color=GREY)

# ---------- Product-level proof: TBKK vs Concentric ----------
s = slide()
header(s, "Product-level proof: same pumps, half the contribution",
       "Internal data both sides: TBKK (Thailand) profit by items Apr-Sep 2025 vs Concentric Group MR Pack Dec-2021 YTD. Contribution = sales minus variable costs")
cd = CategoryChartData()
cd.categories = ["Water pumps", "Oil pumps", "GROUP TOTAL", "Concentric value-added\n(semi-var./electric/ALFDEX)"]
cd.add_series("TBKK marginal profit %", (19.5, 22.6, 22.7, None))
cd.add_series("Concentric Contribution I %", (38.2, 42.9, 44.5, 56.0))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.6), Inches(7.0), Inches(4.6), cd)
ch = gf.chart
ch.has_legend = True
ch.legend.position = XL_LEGEND_POSITION.BOTTOM
ch.legend.include_in_layout = False
ch.legend.font.size = Pt(11)
ch.has_title = True
ch.chart_title.text_frame.text = "Contribution margin by product (%)"
ch.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
ch.chart_title.text_frame.paragraphs[0].runs[0].font.name = BODY_FONT
ch.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
ch.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = DARK
plot = ch.plots[0]
plot.has_data_labels = True
plot.data_labels.number_format = '0.0"%"'
plot.data_labels.number_format_is_linked = False
plot.data_labels.font.size = Pt(10)
plot.series[0].format.fill.solid()
plot.series[0].format.fill.fore_color.rgb = RED
plot.series[1].format.fill.solid()
plot.series[1].format.fill.fore_color.rgb = GREEN
ch.category_axis.tick_labels.font.size = Pt(10)
ch.category_axis.tick_labels.font.name = BODY_FONT
panel(s, 7.8, 1.6, 5.05, 2.45, "Inside TBKK's own book",
      "Pockets of real pricing power: Isuzu brake VD00 (HR) 64.1%, brake 700P 41.7%, gears 30.6%\n"
      "Bleeding lines: GKN case-set/diff -30.2% (negative contribution), Suzuki -1.2%, ISZ W/P ES 1.7%\n"
      "Customer curve: Isuzu 25.8% > MMTH 23.0% > ITT 22.6% > Kubota 17.0%", body_size=11.5)
panel(s, 7.8, 4.2, 5.05, 2.45, "Concentric's ladder = TBK's roadmap",
      "Conventional pump 38-43% -> semi-variable 54% -> electric 55% -> system (ALFDEX) 56%\n"
      "Each evolution step = roughly +15 points of contribution on the SAME base product\n"
      "TBK's e-pump/TCU program targets exactly this ladder — the data proves the rungs exist", body_size=11.5, title_color=GREEN)
txt(s, 0.5, 6.72, 12.3, 0.7, "This is the micro-mechanism behind every macro gap in this deck: gross margin, EBITA, and ROIC differences all begin at contribution per part number. Caveats: TBKK = Thai subsidiary (~26% of group sales); periods/currencies differ — ratios comparable, absolutes not; Concentric 2021 was a strong year. First actions are free: prune negative-contribution lines (GKN case-set), reprice near-zero lines.", size=10.5, color=GREY)

# ---------- Inside TBK Japan: customer gradient & product truths ----------
s = slide()
header(s, "Inside TBK Japan: the keiretsu discount, quantified",
       "TBK parent FY2025 H1 profitability by customer and product (series delivery, K JPY) — contribution margin after direct material")
cd = CategoryChartData()
cd.categories = ["UD Trucks", "MFTBC", "Isuzu", "Hino", "Kubota", "Aftermarket\n(TBK Sales)", "Caterpillar", "Komatsu-\nOyama", "Toho", "Hitachi", "Cummins"]
cd.add_series("Contribution margin %", (8.1, 15.6, 16.8, 21.0, 22.3, 30.2, 30.5, 32.6, 43.5, 44.8, 79.8))
gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.5), Inches(1.6), Inches(6.6), Inches(4.7), cd)
ch = gf.chart
ch.has_legend = False
ch.has_title = True
ch.chart_title.text_frame.text = "Contribution by customer: keiretsu OEMs 17% blended vs non-keiretsu/aftermarket 30.5%"
ch.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
ch.chart_title.text_frame.paragraphs[0].runs[0].font.name = BODY_FONT
ch.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
ch.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = DARK
plot = ch.plots[0]
plot.has_data_labels = True
plot.data_labels.number_format = '0.0"%"'
plot.data_labels.number_format_is_linked = False
plot.data_labels.font.size = Pt(10)
ser = plot.series[0]
ser.format.fill.solid()
ser.format.fill.fore_color.rgb = GREEN
for i in (0, 1, 2, 3):
    try:
        pt = ser.points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = RED
    except Exception:
        pass
ch.category_axis.tick_labels.font.size = Pt(9)
ch.category_axis.tick_labels.font.name = BODY_FONT
panel(s, 7.4, 1.6, 5.45, 2.45, "Product truths (Japan H1 FY2025)",
      "DISC BRAKE: 68.7% contribution — TBK's HIGHEST-margin product (tiny ¥8.3M volume): the disc transition is margin-accretive, not just defensive\n"
      "Precision niches price well: S/camshaft 52.7%, kits 55.5%\n"
      "Negative lines: O/brake -4.9%, B/housing -8.7% (material > sales)\n"
      "Retarder only 10.4% — the electrification-winner product is underpriced\n"
      "ECWP (electric W/P): -245.7% at pre-scale volume — the first rung of Concentric's ladder has an entry toll", body_size=11)
panel(s, 7.4, 4.2, 5.45, 2.45, "Share map & pipeline",
      "Share map: TBK holds Isuzu + Fuso cells across brake sizes and pump segments; mixed at Hino; ZERO UD cells (the all-disc, Knorr-supplied OEM) — the disc flip's cost, mapped\n"
      "New-business pipeline: ¥8M (FY21) → ¥700M/yr (FY25) incl. ELF EV precharge box, e-W/Pump for Denso AC, Cummins/Komatsu wins — real momentum, but ~2.5% of parent sales: must steepen 10x", body_size=11, title_color=GREEN)
txt(s, 0.5, 6.72, 12.3, 0.7, "Read with the previous page: Japan series contribution is 19.6% (vs TBKK 22.7%, Concentric 44.5%) — the gap is group-wide. The two cheapest margin levers need no R&D: mix toward non-keiretsu/aftermarket/export customers who already pay ~2x, and price the differentiated products (disc, retarder) for what the data says they are worth.", size=10.5, color=GREY)

# ---------- The Haldex lesson: profit pools & the disc warning ----------
s = slide()
header(s, "The Haldex lesson: the profit pool is the aftermarket",
       "Haldex internal analysis (Dec 2019 file): segment P&L Sep YTD 2014 + 2009-2019 history, MSEK — the only Western peer with internal segment economics in this study")
cd = CategoryChartData()
cd.categories = ["Truck OE", "Trailer OE", "Aftermarket"]
cd.add_series("Contribution II %", (22.4, 22.2, 48.6))
cd.add_series("EBIT %", (-3.9, -9.9, 31.6))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.6), Inches(6.4), Inches(4.6), cd)
ch = gf.chart
ch.has_legend = True
ch.legend.position = XL_LEGEND_POSITION.BOTTOM
ch.legend.include_in_layout = False
ch.legend.font.size = Pt(11)
ch.has_title = True
ch.chart_title.text_frame.text = "Haldex segment economics (Sep YTD 2014): OE loses, aftermarket earns it all"
ch.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
ch.chart_title.text_frame.paragraphs[0].runs[0].font.name = BODY_FONT
ch.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
ch.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = DARK
plot = ch.plots[0]
plot.has_data_labels = True
plot.data_labels.number_format = '0.0"%"'
plot.data_labels.number_format_is_linked = False
plot.data_labels.font.size = Pt(10)
plot.series[0].format.fill.solid()
plot.series[0].format.fill.fore_color.rgb = DARK
plot.series[1].format.fill.solid()
plot.series[1].format.fill.fore_color.rgb = GREEN
ch.category_axis.tick_labels.font.size = Pt(11)
ch.category_axis.tick_labels.font.name = BODY_FONT
panel(s, 7.2, 1.55, 5.65, 1.62, "1. Aftermarket is the profit pool",
      "OE truck/trailer: CII 22-32% every year 2009-19, EBIT NEGATIVE (Europe truck -27%); aftermarket CII 42-50%, EBIT +27-37% — ALL of Haldex's profit\n"
      "TBK's aftermarket (30.2% CM) is ~10% of parent sales vs Haldex's 42% — the proven fastest route to group margin", body_size=10.5)
panel(s, 7.2, 3.27, 5.65, 1.62, "2. The pruning playbook (executed at Haldex)",
      "Core products earned 40-56% CII / 15-32% EBIT (air treatment 55.8%/31.6%, suspension 47.1%/20.9%, ABA 41.4%/19.7%) while two bleeders dragged: ModulAir -70.4% EBIT, disc brake -6.3%\n"
      "Verified: excluding the two bleeders, May-2014 EBIT margin = 11.2% vs 8.8% reported — pruning took Haldex over 10%", body_size=10.5, title_color=GREEN)
panel(s, 7.2, 4.99, 5.65, 1.62, "3. The subscale OE disc warning",
      "Haldex disc brakes: CII 12-15%, EBIT -13.6% -> -6.3% — the #3 player fighting Knorr/Meritor LOST money per unit (supplier price rises, stalled plant transfer)\n"
      "TBK's 68.7% D/BRAKE is small-batch pricing, not series-OE economics — go niche/aftermarket/partnered (Brakes India), never frontal", body_size=10.5, title_color=RED)
txt(s, 0.5, 6.72, 12.3, 0.7, "Applied to TBK (Japan H1 FY2025 data): pruning the sub-10% non-strategic lines (O/brake, A/brake, B/housing, G/case, gear & processed components = 6.5% of sales) lifts series contribution 19.6% -> 20.9% immediately — and the EBIT gain is larger once attached capacity costs go. TBK's high-margin cohort (¥3.0bn of lines at >=20% CM) already runs at 27.3%: a portfolio focused at the level of its better products + a grown aftermarket is the Haldex >10% playbook, re-run.", size=10.5, bold=True, color=DARK)

# ---------- Brake segment: TBK vs Haldex ----------
s = slide()
header(s, "Brake segment: TBK vs Haldex — product-level contribution",
       "TBK Japan H1 FY2025 (CM after direct material) vs Haldex May-2014 YTD (CII). Concentric makes no brakes. Periods differ — structure, not cycle, is the message")
cd = CategoryChartData()
cd.categories = ["SAW drum\n(flagship)", "Foundation\ndrum (2L)", "Friction /\nlinings", "Brake\ndrum", "Disc\nbrake", "Brake adjusters\n(ABA)", "ABS / EBS\nelectronics", "Aftermarket"]
cd.add_series("TBK", (19.5, 15.8, 25.1, 22.9, 68.7, None, None, 30.2))
cd.add_series("Haldex", (None, None, 23.2, None, 14.0, 41.4, 43.4, 48.2))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.6), Inches(7.4), Inches(4.6), cd)
ch = gf.chart
ch.has_legend = True
ch.legend.position = XL_LEGEND_POSITION.BOTTOM
ch.legend.include_in_layout = False
ch.legend.font.size = Pt(11)
ch.has_title = True
ch.chart_title.text_frame.text = "Contribution margin by brake product (%)"
ch.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
ch.chart_title.text_frame.paragraphs[0].runs[0].font.name = BODY_FONT
ch.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
ch.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = DARK
plot = ch.plots[0]
plot.has_data_labels = True
plot.data_labels.number_format = '0.0"%"'
plot.data_labels.number_format_is_linked = False
plot.data_labels.font.size = Pt(10)
plot.series[0].format.fill.solid()
plot.series[0].format.fill.fore_color.rgb = RED
plot.series[1].format.fill.solid()
plot.series[1].format.fill.fore_color.rgb = DARK
ch.category_axis.tick_labels.font.size = Pt(9)
ch.category_axis.tick_labels.font.name = BODY_FONT
panel(s, 8.2, 1.6, 4.65, 2.45, "Where the gap lives",
      "On overlapping products, near-parity: friction/linings 25.1% vs 23.2%\n"
      "Haldex's premium came from CONTENT TBK doesn't make: brake adjusters 41.4%, ABS 43.4%, air treatment 55.8% — mechatronic niches\n"
      "…and the AFTERMARKET annuity: 48.2% CII vs TBK Sales 30.2%, on 4x the sales weight", body_size=11)
panel(s, 8.2, 4.2, 4.65, 2.45, "The disc paradox, resolved",
      "TBK disc 68.7% = small-batch/replacement pricing on ¥8.3M\n"
      "Haldex disc 14.0% (EBIT -6.3%) = series-OE pricing fighting Knorr at scale\n"
      "Both are true: discs pay in niches and aftermarket, bleed in frontal OE — TBK's ADB entry must choose the first route", body_size=11, title_color=GREEN)
txt(s, 0.5, 6.72, 12.3, 0.7, "Gross-profit read: TBK's brake book (2L/brake 15.8%, A/brake 7.6%, O/brake -4.9%, drum 22.9%, lining 25.1%) clusters at 16-25% contribution; Haldex's brake-adjacent book ran 23-56% because every point of extra margin came from mechatronics and aftermarket, not from better drum pricing. The brake gap is a PRODUCT-CONTENT gap, not a cost gap.", size=10.5, color=GREY)

# ---------- The SAW case: best product, below-market price ----------
s = slide()
header(s, "The SAW case: TBK's best product, sold ~30% below market",
       "SAW = TBK's flagship drum brake — class-leading stopping time, energy efficiency and weight (per management). ¥4.31bn = 36% of parent H1 series sales, at just 19.5% contribution")
cd = CategoryChartData()
cd.categories = ["SAW today\n(keiretsu price)", "SAW at market\n(+30% price)", "SAW at market\n(price / 0.7)", "Haldex ABA/ABS\navg (ref.)", "Concentric conv.\npumps avg (ref.)"]
cd.add_series("Contribution margin %", (19.5, 38.1, 43.6, 42.4, 40.5))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.6), Inches(7.0), Inches(4.6), cd)
ch = gf.chart
ch.has_legend = False
ch.has_title = True
ch.chart_title.text_frame.text = "SAW contribution: actual vs at-market-price scenarios (vs peer core products)"
ch.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
ch.chart_title.text_frame.paragraphs[0].runs[0].font.name = BODY_FONT
ch.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
ch.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = DARK
plot = ch.plots[0]
plot.has_data_labels = True
plot.data_labels.number_format = '0.0"%"'
plot.data_labels.number_format_is_linked = False
plot.data_labels.font.size = Pt(11)
ser = plot.series[0]
ser.format.fill.solid()
ser.format.fill.fore_color.rgb = DARK
try:
    pt = ser.points[0]
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = RED
    for i in (1, 2):
        pt = ser.points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = GREEN
except Exception:
    pass
ch.category_axis.tick_labels.font.size = Pt(9)
ch.category_axis.tick_labels.font.name = BODY_FONT
panel(s, 7.8, 1.6, 5.05, 2.45, "The triangulation",
      "At market price, SAW earns 38-44% contribution — exactly the level of Haldex's core products (ABA 41.4%, ABS 43.4%) and Concentric's conventional pumps (38-43%)\n"
      "The keiretsu discount fully explains the flagship's gap: a world-leading product, peer-level economics, captured by the customer\n"
      "Price flows straight to contribution: no cost change needed", body_size=11.5)
panel(s, 7.8, 4.2, 5.05, 2.45, "What it is worth — and how to get it",
      "SAW repriced to market: +¥2.6-3.7bn contribution PER YEAR; parent series CM 19.6% -> 27-30% from this one product\n"
      "Realistically gradual at home (keiretsu contracts) — so arbitrage it: SELL SAW AT MARKET PRICE where market prices rule: exports, aftermarket, and India via Brakes India, where drum demand persists at scale\n"
      "Class-leading tech + below-market price = the export business case writes itself", body_size=11.5, title_color=GREEN)
txt(s, 0.5, 6.72, 12.3, 0.7, "Note: 'leading drum brake in stopping time, energy save and weight' and the ~30% below-market pricing are management/first-hand information, not externally verified — consistent, however, with the customer-gradient data (non-keiretsu customers pay ~2x) and with TBK's zero-recall record. If the ~30% discount applies beyond SAW (management indicates Japanese customers broadly), the parent-wide repricing prize is correspondingly larger.", size=10.5, color=GREY)

# ---------- Pump segment: TBK vs Concentric ----------
s = slide()
header(s, "Pump segment: TBK vs Concentric — product-level contribution",
       "TBK Japan H1 FY2025 + TBKK Thailand vs Concentric Dec-2021 YTD CB I. Haldex has no pumps — its hydraulics division IS Concentric (spun off 2011)")
cd = CategoryChartData()
cd.categories = ["Water pumps", "Oil pumps", "Semi-variable\nwater pump", "Electric\nwater pump", "Electric\noil pump", "Fuel / transm.\npumps"]
cd.add_series("TBK (Japan)", (18.5, 26.0, None, None, None, None))
cd.add_series("TBKK (Thailand)", (19.5, 22.6, None, None, None, None))
cd.add_series("Concentric", (38.2, 42.9, 54.1, 54.8, 40.5, 55.6))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.6), Inches(7.4), Inches(4.6), cd)
ch = gf.chart
ch.has_legend = True
ch.legend.position = XL_LEGEND_POSITION.BOTTOM
ch.legend.include_in_layout = False
ch.legend.font.size = Pt(11)
ch.has_title = True
ch.chart_title.text_frame.text = "Contribution margin by pump product (%)"
ch.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
ch.chart_title.text_frame.paragraphs[0].runs[0].font.name = BODY_FONT
ch.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
ch.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = DARK
plot = ch.plots[0]
plot.has_data_labels = True
plot.data_labels.number_format = '0.0"%"'
plot.data_labels.number_format_is_linked = False
plot.data_labels.font.size = Pt(10)
plot.series[0].format.fill.solid()
plot.series[0].format.fill.fore_color.rgb = RED
plot.series[1].format.fill.solid()
plot.series[1].format.fill.fore_color.rgb = RGBColor(0x9A, 0x95, 0x8E)
plot.series[2].format.fill.solid()
plot.series[2].format.fill.fore_color.rgb = GREEN
ch.category_axis.tick_labels.font.size = Pt(9)
ch.category_axis.tick_labels.font.name = BODY_FONT
panel(s, 8.2, 1.6, 4.65, 2.45, "The 2x gap on identical products",
      "Same conventional water pump: TBK 18.5-19.5% vs Concentric 38.2% — DOUBLE\n"
      "Same oil pump: 22.6-26.0% vs 42.9%\n"
      "Consistent in two TBK countries → not a factory or FX artifact: it is price and customer mix (keiretsu price-downs vs Concentric's Cummins/CAT niche pricing)", body_size=11)
panel(s, 8.2, 4.2, 4.65, 2.45, "The ladder TBK hasn't climbed",
      "Concentric's evolved variants: semi-variable 54.1%, electric W/P 54.8%, fuel/transmission 55.6% — each rung ≈ +15pts on the same base product\n"
      "TBK's first rung is bought but not yet scaled: ECWP at -245.7% contribution on ¥2.6M pre-scale volume (off-chart)\n"
      "NB: Concentric = the former Haldex hydraulics division — proof the climb is doable from CV-supplier DNA", body_size=11, title_color=GREEN)
txt(s, 0.5, 6.72, 12.3, 0.7, "Gross-profit read: if TBK's pump book (W/P + O/P = ¥4.9bn of H1 parent sales) earned Concentric's conventional-pump margins, contribution would rise ~¥0.9-1.0bn per HALF-YEAR on those two lines alone — before any ladder-climbing. The pump gap is a PRICING-POWER gap on identical hardware; the brake gap (previous page) is a content gap. Different cures: pricing/mix vs product development.", size=10.5, color=GREY)

# ---------- Pro-forma normalized income statement ----------
s = slide()
header(s, "The normalized TBK: right prices, right factories, real R&D",
       "Pro-forma on audited FY3/2026 — pricing corrected (SAW only), factory structure fixed, competitive R&D loaded. Deficit repayment kept separate, per scope")
cd = CategoryChartData()
cd.categories = ["Actual\nFY3/26", "PHASE 1 (realistic)\n+9% reprice, ~450 heads", "PHASE 2 (full realization)\n+15%, 500 heads", "FULL premises\n(leakage meter)"]
cd.add_series("Operating margin %", (2.7, 10.6, 15.0, 26.6))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.6), Inches(6.9), Inches(4.6), cd)
ch = gf.chart
ch.has_legend = False
ch.has_title = True
ch.chart_title.text_frame.text = "Operating margin by phase: 2.7% -> ~10% (phase 1) -> 15% -> 26.6% ceiling"
ch.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
ch.chart_title.text_frame.paragraphs[0].runs[0].font.name = BODY_FONT
ch.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
ch.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = DARK
plot = ch.plots[0]
plot.has_data_labels = True
plot.data_labels.number_format = '0.0"%"'
plot.data_labels.number_format_is_linked = False
plot.data_labels.font.size = Pt(11)
ser = plot.series[0]
ser.format.fill.solid()
ser.format.fill.fore_color.rgb = DARK
try:
    for i, col in ((0, RED), (1, GREEN), (2, RGBColor(0x31, 0x43, 0x59)), (3, GREY)):
        pt = ser.points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = col
except Exception:
    pass
ch.category_axis.tick_labels.font.size = Pt(8)
ch.category_axis.tick_labels.font.name = BODY_FONT
panel(s, 7.7, 1.6, 5.15, 2.45, "PHASE 1 — the realistic target: ~10% OPM",
      "Sales 59,342 | GP 13,211 (22.3%) | R&D 2,789 at 4.7% | OP 6,263 (10.6%) | Net 4,398 = ¥138 EPS (~2.2x mkt cap)\n"
      "Captures only ~1/3 of the stated pricing gap + ~45% of excess heads — via pruning, renewal-cycle repricing, aftermarket/non-keiretsu mix, first consolidations\n"
      "Funds the entire R&D program (incl. deficit memo) in ~18-24 months", body_size=11)
panel(s, 7.7, 4.2, 5.15, 2.45, "Why not more, at first — operator calibration",
      "Concentric peaked at ~22% OPM — on best-in-class FULL price realization, a discipline built over years\n"
      "Haldex could have reached ~30% by shedding Europe for the US — portfolio focus is the second lever\n"
      "Above ~10% in phase 1 is difficult: keiretsu contracts reprice at renewal, factories close on multi-year clocks. Phase 2 (15%) = Concentric-grade realization + grown aftermarket. FULL premises (26.6%) = leakage meter, not a plan", body_size=11, title_color=GREEN)
txt(s, 0.5, 6.72, 12.3, 0.7, "Nothing here requires new products or customers — existing book, honest prices, loaded factories, peer-level R&D (4.7%) absorbed in every phase. MEMO (separate financing): R&D deficit ~¥1.5bn/yr x5, ADB capex ¥4.5-7.5bn, restructuring cash ~¥2-3bn. Premises are management/first-hand information; phasing reflects operator experience at Concentric and Haldex.", size=10.5, bold=True, color=DARK)

# ---------- OPEX comparison (USD) ----------
s = slide()
header(s, "Why TBK looks lean on OPEX: its factory costs hide in COGS",
       "OPEX here = gross margin − operating margin (SG&A + R&D). Under JGAAP, factory overhead, depreciation and idle capacity sit in COST OF SALES  |  USD at ¥150/USD")
cd = CategoryChartData()
cd.categories = ["BRAKES panel", "PUMPS panel"]
cd.add_series("TBK (blended)", (9.7, 9.7))
cd.add_series("Peer average", (13.8, 10.4))
cd.add_series("Peer median", (13.6, 9.3))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.6), Inches(7.3), Inches(4.6), cd)
ch = gf.chart
ch.has_legend = True
ch.legend.position = XL_LEGEND_POSITION.BOTTOM
ch.legend.include_in_layout = False
ch.legend.font.size = Pt(12)
ch.has_title = True
ch.chart_title.text_frame.text = "OPEX as % of sales: TBK vs peers, by business"
ch.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
ch.chart_title.text_frame.paragraphs[0].runs[0].font.name = BODY_FONT
ch.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
ch.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = DARK
plot = ch.plots[0]
plot.has_data_labels = True
plot.data_labels.number_format = '0.0"%"'
plot.data_labels.number_format_is_linked = False
plot.data_labels.font.size = Pt(11)
plot.series[0].format.fill.solid()
plot.series[0].format.fill.fore_color.rgb = RED
plot.series[1].format.fill.solid()
plot.series[1].format.fill.fore_color.rgb = DARK
plot.series[2].format.fill.solid()
plot.series[2].format.fill.fore_color.rgb = GREEN
ch.category_axis.tick_labels.font.size = Pt(11)
ch.category_axis.tick_labels.font.name = BODY_FONT
panel(s, 8.1, 1.6, 4.75, 2.3, "The lean look is an artifact",
      "TBK's 9.7% covers only SG&A + its\nthin ~2.2% R&D — the future-building costs\nit has starved. The factory burden is in COGS:\nPP&E 39% of sales vs Concentric ~12% (3x)\nDepreciation 5.3% of sales; China: 12.8% of\nsegment sales + losses both years", body_size=11.5)
panel(s, 8.1, 4.05, 4.75, 2.3, "Restructuring bypasses BOTH metrics",
      "Impairments ¥459M + ¥712M and\nrestructuring losses ¥775M + ¥297M\n(>¥2.2bn in 2 years) booked as JGAAP\nextraordinary items — below operating\nprofit, invisible to GM and OPEX alike.\nThat IS the underutilized-footprint cost", body_size=11.5, title_color=RED)
txt(s, 0.5, 6.45, 12.3, 0.85, "Two-front program, mapped to the P&L: (1) footprint/utilization repair attacks COGS — the $22-24M/yr gross-margin prize (North America exit and the ¥712M Japan impairment show it has started); (2) R&D catch-up deliberately GROWS opex toward peer levels (peers spend 10-14% incl. 4%+ R&D and earn 2.5x TBK's operating margin — Concentric grew opex 10.3%→11.7% funding its e-pump transition at 18.4% avg OPM). Cut above the line, invest below it.", size=10.5, color=GREY)

# ---------- 10 Deficit ----------
s = slide()
header(s, "TBK's accumulated R&D deficit: ~US$100M over a decade",
       "TBK dataset FY2016-2025; benchmark = blended peer floor for a brake + pump supplier")
cd = CategoryChartData()
cd.categories = [str(y) for y in range(2017, 2026)]
cd.add_series("TBK actual R&D ($M)", (11.1, 11.4, 10.6, 11.4, 11.6, 11.8, 8.9, 8.1, 7.8))
cd.add_series("At 4.5% of sales ($M)", (19.4, 20.9, 21.8, 21.3, 18.7, 20.5, 17.8, 17.7, 16.1))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.6), Inches(7.6), Inches(4.4), cd)
ch = gf.chart
ch.has_legend = True
ch.legend.position = XL_LEGEND_POSITION.BOTTOM
ch.legend.include_in_layout = False
ch.legend.font.size = Pt(11)
ch.has_title = True
ch.chart_title.text_frame.text = "Actual R&D vs peer-floor benchmark (US$M)"
ch.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
ch.chart_title.text_frame.paragraphs[0].runs[0].font.name = BODY_FONT
ch.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
ch.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = DARK
p0 = ch.plots[0]
p0.series[0].format.fill.solid()
p0.series[0].format.fill.fore_color.rgb = RED
p0.series[1].format.fill.solid()
p0.series[1].format.fill.fore_color.rgb = DARK
ch.category_axis.tick_labels.font.size = Pt(10)
ch.category_axis.tick_labels.font.name = BODY_FONT
panel(s, 8.4, 1.6, 4.45, 4.4, "The deficit, quantified",
      "Cumulative sales FY16-25: $4.26bn\nCumulative R&D: $93.5M (2.2%)\n\nGap vs 4.0% benchmark: ~$77M\nGap vs 4.5% benchmark: ~$98M\nGap vs 5.0% benchmark: ~$120M\n\nCentral estimate: ~US$100M\n\nAnd intensity is FALLING:\n2.6% (FY22) → 2.2% (FY25)", body_size=12.5, title_color=RED)
txt(s, 0.5, 6.25, 12.3, 1.0, "Compounding the problem: TBK's marginal return on invested capital has been negative on every horizon (-2% to -10%) — the deficit is in direction as well as amount. Capex fell to $16M by FY25, though FY3/26 shows a turn: capex back up to ¥3.3bn and ¥1.1bn equity raised from Brakes India.", size=12, color=GREY)

# ---------- 11 Plan ----------
s = slide()
header(s, "The reinvestment requirement: two phases",
       "Catch-up to repay the technological deficit, then a permanently higher run-rate to match peers")
cd = CategoryChartData()
cd.categories = ["FY26\nactual est.", "FY27", "FY28", "FY29", "FY30", "FY31", "Steady state\nFY32+"]
cd.add_series("Parity run-rate (~4.5-5% of sales)", (7.8, 17, 17, 17, 17, 17, 17))
cd.add_series("Deficit repayment (catch-up)", (0, 10, 10, 10, 10, 10, 0))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, Inches(0.5), Inches(1.6), Inches(7.6), Inches(4.4), cd)
ch = gf.chart
ch.has_legend = True
ch.legend.position = XL_LEGEND_POSITION.BOTTOM
ch.legend.include_in_layout = False
ch.legend.font.size = Pt(11)
ch.has_title = True
ch.chart_title.text_frame.text = "R&D spend plan (US$M per year)"
ch.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
ch.chart_title.text_frame.paragraphs[0].runs[0].font.name = BODY_FONT
ch.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
ch.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = DARK
p0 = ch.plots[0]
p0.series[0].format.fill.solid()
p0.series[0].format.fill.fore_color.rgb = DARK
p0.series[1].format.fill.solid()
p0.series[1].format.fill.fore_color.rgb = GREEN
ch.category_axis.tick_labels.font.size = Pt(9)
ch.category_axis.tick_labels.font.name = BODY_FONT
panel(s, 8.4, 1.6, 4.45, 4.4, "The numbers",
      "PHASE 1 — Catch-up (5 yrs):\n~$27M/yr ≈ 7.5% of sales\n= parity run-rate ($16-18M)\n+ deficit repayment (~$10M/yr,\n~$50M split between ADB\nindustrialization and e-pump/TCU\nprograms)\nTotal: ~$135M over 5 years\n\nPHASE 2 — Steady state:\n$16-18M/yr ≈ 4.5-5% of sales\n= 2.2x today's spend, forever\n\nPlus one-off ADB line capex\n(~$30-50M, outside R&D line)", body_size=12, title_color=GREEN)
txt(s, 0.5, 6.25, 12.3, 1.1, "Allocation discipline, per domain: (1) brakes — ADB industrialization for the Japanese disc flip, systems/ADAS access via Brakes India, timed to regulation and platform renewals; (2) pumps — e-pump/TCU sockets on Isuzu/Fuso/Hino electrified programs, timed to program kick-offs happening NOW. Funding: the phase-1 pro-forma uplift (+¥4.8bn / ~$32M operating profit vs actual) more than covers the ~$27M/yr program. The Concentric lesson: focused niche engineering at 2.3% intensity earns 25% gross margins — spend must buy sockets, not just scale.", size=11.5, color=GREY)

# ---------- 14 Product evaluation I: quality & service ----------
s = slide()
header(s, "Product evaluation I: external quality & service evidence",
       "Evidence hierarchy: NHTSA recalls (supplier named by law) > OEM-granted awards > independent surveys > patents > self-claims. Japan's MLIT does not name suppliers")
panel(s, 0.5, 1.55, 6.1, 2.5, "The validated leaders",
      "Knorr-Bremse/Bendix: PACCAR supplier-performance award 6 consecutive yrs; 3 plants on PACCAR's <=10 defects-per-million list; ETM 'Best Brand - Brakes' (German fleet/driver vote) 19 YEARS RUNNING; ~85% of NA air-disc wheel ends\n"
      "Meritor/Cummins: Daimler 'Masters of Quality' + 2020 global quality award ('1.3M axles/brakes at low defect rate'); DTNA EXTENDED its ADB contract — revealed preference\n"
      "Brembo: all 10 F1 teams; AMS brake-brand vote 11 of 13 years", body_size=11)
panel(s, 6.8, 1.55, 6.05, 2.5, "The adverse records",
      "Bendix: 2019 ADB22X caliper recall (fire risk); 2024-25 EC-80 ECU recall, 442k units (software fix) — scale produces both install base and recalls\n"
      "Haldex: lost a WILLFUL patent-infringement verdict to Bendix (copied its ADB); two modest component recalls\n"
      "Akebono: the only Japanese supplier with documented adverse events — 2015 GM recall of its brake products + 2019-21 inspection-data falsification (no recall, governance black mark)", body_size=11, title_color=RED)
panel(s, 0.5, 4.2, 6.1, 2.45, "TBK: externally invisible — in BOTH directions",
      "No OEM award findable, no survey presence, no service network beyond Japan/Thailand\n"
      "BUT also: no recall, no scandal, no misconduct report in Japanese-language search — a meaningful null in a scandal-rich industry (cf. Akebono)\n"
      "75 years inside Isuzu/Fuso/Hino = quality is probably fine; 'probably fine, unproven' is the ceiling the evidence supports\n"
      "Dataset moat scores align: Knorr 7.55 top, Akebono 4.95 bottom — TBK was unscored: literally unmeasurable from outside", body_size=11)
panel(s, 6.8, 4.2, 6.05, 2.45, "TCO & service: only proxies exist",
      "Brand-level TCO data is NOT public (fleets don't publish; IMR studies paywalled) — standing must be inferred:\n"
      "Install base: Bendix 20M+ ADBs in field | Contract renewals: Meritor-DTNA extension | Networks: Knorr's pan-EMEA certified workshop network vs TBK's Japan+Thailand footprint\n"
      "On every proxy the systems players lead — and TBK does not appear", body_size=11, title_color=GREEN)
txt(s, 0.5, 6.75, 12.3, 0.6, "Strategic point: TBK's quality reputation is held entirely inside the keiretsu — like its market share, it cannot be sold to new customers. Making quality externally legible (awards, certifications, published data) is a cheap, necessary complement to the R&D program.", size=12, bold=True, color=DARK)

# ---------- 15 Product evaluation II: innovation output ----------
s = slide()
header(s, "Product evaluation II: innovation output",
       "Patents, documented product firsts, and technical-society awards — the externally visible output of R&D")
panel(s, 0.5, 1.55, 6.1, 2.5, "The patent gap (independent registries)",
      "Knorr-Bremse: ~2,900 patent families / 10,900+ patents (~7,500 in the truck division)\n"
      "Aisin: ~680 JP applications/yr; Clarivate Top-100 Global Innovator 7 consecutive years\n"
      "Akebono: ~321 US-registered patents\n"
      "Mikuni: e-oil-pump fault-detection patents — its BEV-truck win is engineering-backed\n"
      "TBK: 2-3 patent GRANTS PER YEAR — two to three orders of magnitude below the leaders.\nIndependently confirms the capitalized-R&D finding ($43M vs $2.6-5.6bn) from a different data source", body_size=11)
panel(s, 6.8, 1.55, 6.05, 2.5, "Documented product firsts",
      "Knorr: EBS 1989; series air disc brake 1992; Actros/Scania volume 1996\n"
      "Bendix ADB22X: 2005, now ~85% of NA air-disc wheel ends\n"
      "UD Quon 2017: first all-disc Japanese heavy truck\n"
      "Mikuni 2022: e-oil pump in production for a domestic small BEV truck\n"
      "Akebono FY2023: JSME Technology Award — world-first electric parking brake unit for medium/light trucks (innovating through distress)\n"
      "TBK: retarders since the 1960s; 19.5\" ADB exhibited (year unconfirmed: 2019 vs 2022 sources conflict); e-pumps/TCU in development", body_size=11, title_color=GREEN)
panel(s, 0.5, 4.2, 12.35, 1.9, "The pattern — and TBK's fix",
      "Innovation output tracks the moat and margin rankings almost perfectly: the companies with validated quality (Knorr, Cummins) also hold the patent fortresses and the product firsts; the distressed (Akebono) still out-publish TBK; even the similarly-sized peer (Mikuni) protects its e-pump work with patents.\n"
      "TBK's innovation is real but INVISIBLE: 75 years of retarder/brake/pump engineering, an ADB prototype, e-pumps and a thermal control unit in development — none of it patented at scale, published at JSAE, or externally awarded. Invisible innovation cannot win sockets outside the keiretsu.\n"
      "Fix costs little: patent the e-pump/TCU and ADB work (Mikuni shows the value), publish, pursue external validation — a required complement to the R&D spend itself.", body_size=11.5)

# ---------- 16 Moat ranking overview ----------
s = slide()
header(s, "Moat framework ranking: TBK scores 5.38 — 9th of 13",
       "Unified Niche Compounder Methodology v3.2: 10 weighted criteria, domain-blended internal/external evidence. Thresholds: >7.8 Compounder, 6.5-7.8 Watchlist, <6.5 Pass")
cd = CategoryChartData()
cd.categories = ["Akebono", "Mikuni", "ZF", "Mahle", "TBK", "TPR", "Hanon", "Aisin", "SAF-Holland", "Cummins", "Brembo", "Concentric", "Knorr-Bremse"]
cd.add_series("Moat score", (4.95, 5.00, 5.25, 5.35, 5.38, 5.45, 5.88, 6.25, 6.59, 6.73, 6.73, 6.98, 7.55))
gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.5), Inches(1.6), Inches(7.3), Inches(5.0), cd)
ch = gf.chart
ch.has_legend = False
ch.has_title = True
ch.chart_title.text_frame.text = "Final weighted moat scores (TBK in red)"
ch.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
ch.chart_title.text_frame.paragraphs[0].runs[0].font.name = BODY_FONT
ch.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
ch.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = DARK
plot = ch.plots[0]
plot.has_data_labels = True
plot.data_labels.number_format = '0.00'
plot.data_labels.number_format_is_linked = False
plot.data_labels.font.size = Pt(10)
ser = plot.series[0]
ser.format.fill.solid()
ser.format.fill.fore_color.rgb = DARK
try:
    pt = ser.points[4]
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = RED
except Exception:
    pass
ch.category_axis.tick_labels.font.size = Pt(10)
ch.category_axis.tick_labels.font.name = BODY_FONT
panel(s, 8.1, 1.6, 4.75, 2.4, "Bands",
      "Compounder Target (>7.8): NONE —\nno company in this industry qualifies\nWatchlist (6.5-7.8): Knorr 7.55,\nConcentric 6.98, Brembo/Cummins 6.73,\nSAF-Holland 6.59\nPass (<6.5): everyone else —\nincl. TBK 5.38, Mikuni 5.00", body_size=12)
panel(s, 8.1, 4.15, 4.75, 2.45, "TBK's decomposition is unique",
      "The strongest STATIC moat of the\nunscored set (Criticality 7.5 +\nHegemony 8.0 = 2.33 of its 5.38)\n…drained by the worst COMPOUNDING\nengine in the whole group (Capital\nallocation 2.0 + Runway 3.0 = just\n0.50 of a possible 2.00)\nDurable problem, fragile solution", body_size=12, title_color=RED)

# ---------- 17 TBK moat profile: radar ----------
s = slide()
header(s, "TBK's moat profile: where it is strong, where it bleeds",
       "Chapter scores (1-10) — TBK vs Concentric (best newly scored) vs the average of the five newly scored peers (TBK excluded from the average)")
CH_LABELS = ["Criticality", "Standard\npremium", "Market\nhegemony", "Ecosystem", "Lifecycle\ncost adv.", "Substitution\nresistance", "Demand\nhorizon", "Capital\nallocation", "Reinvestment\nrunway", "Reverse-\nLindy"]
TBK_CH = (7.5, 6.0, 8.0, 2.0, 3.0, 3.5, 7.0, 2.0, 3.0, 4.0)
CONC_CH = (7.0, 6.0, 7.5, 2.5, 7.0, 5.0, 7.5, 7.5, 7.0, 5.5)
MIK_CH = (7.0, 5.0, 4.0, 2.0, 4.0, 3.0, 7.0, 2.5, 5.5, 4.5)
AVG_CH = tuple(round(sum(x)/5, 2) for x in zip(
    (7.0, 5.0, 4.0, 2.0, 4.0, 3.0, 7.0, 2.5, 5.5, 4.5),   # Mikuni
    (7.0, 6.0, 7.5, 2.5, 7.0, 5.0, 7.5, 7.5, 7.0, 5.5),   # Concentric
    (7.0, 7.0, 6.0, 5.0, 5.0, 4.0, 7.0, 4.5, 6.0, 4.5),   # Aisin
    (7.5, 6.0, 7.5, 2.0, 6.0, 2.5, 4.0, 4.0, 4.0, 3.5),   # TPR
    (7.0, 5.5, 6.5, 2.5, 3.5, 7.0, 8.0, 2.0, 5.0, 5.0),   # Hanon
))
cd = CategoryChartData()
cd.categories = [c.replace("\n", " ") for c in CH_LABELS]
cd.add_series("TBK", TBK_CH)
cd.add_series("Concentric", CONC_CH)
cd.add_series("Peer average (5 peers)", AVG_CH)
gf = s.shapes.add_chart(XL_CHART_TYPE.RADAR_MARKERS, Inches(0.4), Inches(1.5), Inches(7.0), Inches(5.3), cd)
ch = gf.chart
ch.has_legend = True
ch.legend.position = XL_LEGEND_POSITION.BOTTOM
ch.legend.include_in_layout = False
ch.legend.font.size = Pt(11)
ch.has_title = False
ser_colors = [RED, GREEN, GREY]
for i, sr in enumerate(ch.plots[0].series):
    sr.format.line.color.rgb = ser_colors[i]
    sr.format.line.width = Pt(2.25)
panel(s, 7.7, 1.55, 5.15, 2.5, "TBK's STRONG points",
      "Market hegemony 8.0 — top domestic share in entry-proof niches (best score in the set on this criterion)\n"
      "Criticality 7.5 — life-support product: <1% of truck cost protecting 100% of value\n"
      "Demand horizon 7.0 — stopping & cooling heavy trucks is permanent; content GROWS with electrification", body_size=11.5, title_color=GREEN)
panel(s, 7.7, 4.2, 5.15, 2.5, "TBK's WEAK points",
      "Capital allocation 2.0 — ROIIC negative at every horizon with 100% reinvestment (worst in the group)\n"
      "Ecosystem 2.0 — standalone hardware, no installed-base annuity (vs Knorr's EBS/TruckServices)\n"
      "Lifecycle cost 3.0 — no cost edge: underloaded plants in COGS\n"
      "Reinvestment runway 3.0 — runway exists but unproven; $43M knowledge stock\n"
      "Substitution 3.5 — both products in active displacement", body_size=11.5, title_color=RED)

# ---------- 18 Moat per-criterion charts ----------
s = slide()
header(s, "Moat criteria: TBK vs competitors, criterion by criterion",
       "Newly-scored companies (chapter detail unavailable for the 5 pre-scored names). TBK in red")
def crit_chart(left, cats, tbk, mik, conc, ais, tpr, han, title):
    cd = CategoryChartData()
    cd.categories = cats
    cd.add_series("TBK", tbk)
    cd.add_series("Mikuni", mik)
    cd.add_series("Concentric", conc)
    cd.add_series("Aisin", ais)
    cd.add_series("TPR", tpr)
    cd.add_series("Hanon", han)
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(left), Inches(1.55), Inches(6.2), Inches(4.7), cd)
    ch = gf.chart
    ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.BOTTOM
    ch.legend.include_in_layout = False
    ch.legend.font.size = Pt(9)
    ch.has_title = True
    ch.chart_title.text_frame.text = title
    ch.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
    ch.chart_title.text_frame.paragraphs[0].runs[0].font.name = BODY_FONT
    ch.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
    ch.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = DARK
    colors = [RED, RGBColor(0x9A, 0x95, 0x8E), GREEN, DARK, RGBColor(0xE2, 0xE5, 0xE9), RGBColor(0x33, 0x4E, 0x57)]
    colors[1] = RGBColor(0x64, 0x61, 0x5E)
    for i, sr in enumerate(ch.plots[0].series):
        sr.format.fill.solid()
        sr.format.fill.fore_color.rgb = colors[i]
    ch.category_axis.tick_labels.font.size = Pt(9)
    ch.category_axis.tick_labels.font.name = BODY_FONT
    ch.value_axis.maximum_scale = 10
crit_chart(0.4, ["Criticality", "Standard", "Hegemony", "Ecosystem", "Lifecycle"],
           (7.5, 6.0, 8.0, 2.0, 3.0), (7.0, 5.0, 4.0, 2.0, 4.0), (7.0, 6.0, 7.5, 2.5, 7.0),
           (7.0, 7.0, 6.0, 5.0, 5.0), (7.5, 6.0, 7.5, 2.0, 6.0), (7.0, 5.5, 6.5, 2.5, 3.5),
           "Phase 1 — Structural durability (Ch1-5)")
crit_chart(6.75, ["Substitution", "Demand", "Capital alloc.", "Runway", "Rev-Lindy"],
           (3.5, 7.0, 2.0, 3.0, 4.0), (3.0, 7.0, 2.5, 5.5, 4.5), (5.0, 7.5, 7.5, 7.0, 5.5),
           (4.0, 7.0, 4.5, 6.0, 4.5), (2.5, 4.0, 4.0, 4.0, 3.5), (7.0, 8.0, 2.0, 5.0, 5.0),
           "Phases 2-3 — Endurance & compounding engine (Ch6-10)")
txt(s, 0.5, 6.45, 12.3, 0.85, "Read: in Phase 1 (left) TBK leads or matches on Criticality and Hegemony and collapses on Ecosystem/Lifecycle cost. In Phases 2-3 (right) TBK is last or near-last on every engine criterion — while Concentric (green) leads precisely there. Closing the gap to Concentric is execution (allocation, substitution positioning), not moat-building: the moat already exists.", size=10.5, color=GREY)

# ---------- 19 Research quality notes ----------
s = slide()
header(s, "Research quality & methodology notes",
       "How this material was produced, verified, and where its limits are")
panel(s, 0.5, 1.55, 6.1, 2.6, "Method & confidence tiers",
      "Two deep-research studies (brake market; pump market): 5 parallel search angles each, claims extracted as falsifiable statements, then ADVERSARIAL VERIFICATION — 3 independent voters per claim set, actively attempting refutation\n"
      "Tier 1: peer-reviewed papers, primary regulation, audited filings (TBK tanshin, Concentric AR2023, SEC)\n"
      "Tier 2: government data and multi-source-verified company/industry figures\n"
      "Tier 3: single-source estimates and market research — always labeled in the underlying documents", body_size=11)
panel(s, 6.8, 1.55, 6.05, 2.6, "Corrections made by verification (transparency)",
      "TBK's famous '~70% share' belongs to its WATER PUMPS, not brakes (brakes: 'top share', unaudited)\n"
      "Euro 7 heavy-duty brake-dust limits apply ~2030, not 2028/29 (those are exhaust dates)\n"
      "Concentric R&D = 2.3% of sales per AR2023 — the circulating '~5%' was wrong\n"
      "'85-90% of Asia on drums' could not be corroborated — stated as 'majority' only\n"
      "A claimed TBK-Mikuni alliance is erroneous (conflation with Mikuni-Suzuki)\n"
      "Aggregator's 'TBK FY26 ¥138M profit' contradicted by the audited tanshin (-¥131M)\n"
      "TBK ADB exhibit year unresolved (2019 vs 2022 sources conflict); Concentric MR pack period 2112 = Dec 2021", body_size=11, title_color=RED)
panel(s, 0.5, 4.3, 6.1, 2.6, "Known limitations",
      "Full-text retrieval was blocked in the research environment: paywalled SAE/journal numbers rest on abstracts; flagged where unverified\n"
      "No peer discloses brake-only or pump-only R&D, gross profit, or segment margins — group figures used; ratios are the valid comparison\n"
      "Accounting bases differ: German/Italian nature-of-expense 'gross margins' (Knorr ~50%, Brembo ~50%) are materials-only and EXCLUDED from margin comparisons; JGAAP fully expenses R&D while IFRS peers capitalize part\n"
      "TBK FY3/26 division split estimated from FY3/25 (yuho due June 2026); USD at ¥150/USD", body_size=11)
panel(s, 6.8, 4.3, 6.05, 2.6, "Data hierarchy & framework constructs",
      "Where sources conflict, the audited primary document wins (e.g., TBK tanshin over aggregators)\n"
      "Company-internal documents (TBK/TBKK product P&Ls, Concentric MR pack, Haldex packs) and management statements (~30% below-market pricing, ~1,000 excess heads) are first-hand but NOT externally verified — labeled wherever used\n"
      "Consistent-basis peer metrics (margins, R&D, ROIC) come from the user-provided 10-year/LTM peer dataset — same methodology across all 7 companies\n"
      "ROIIC (ROICm) = return on incremental invested capital over the stated horizon; RR = reinvestment rate; R&D Capital Base = cumulative R&D capitalized/amortized per the dataset framework — a modeling construct, not a balance-sheet line\n"
      "Scenario pages state assumptions on-slide; peer-margin scenarios use comparable reporters only", body_size=11, title_color=GREEN)

# ---------- 20 Sources ----------
s = slide()
header(s, "Sources")
panel(s, 0.5, 1.3, 6.1, 2.85, "Primary & audited",
      "TBK Co., Ltd. — Consolidated Financial Results FY ended 31 Mar 2026 (tanshin, 14 May 2026); securities-report risk factors (有価証券報告書 via kitaishihon/irbank)\n"
      "Concentric AB — Annual Report 2023 (income statement, balance sheet, 5-year key figures)\n"
      "User-provided datasets: TBK 10-year financials; Peers 10-year; Peers LTM (7 companies, consistent methodology); moat methodology v3.2\n"
      "Internal management documents: TBK Japan profitability by product/customer + share map + acquired projects (H1 FY2025); TBKK Thailand profit by items (Apr-Sep 2025); Concentric Group MR Pack (Dec 2021); Haldex analysis file (2019) and May-2014 MR pack\n"
      "SEC filings: Cummins 10-K/10-Q/8-K (FY2024-25); WABCO & Meritor merger 8-Ks\n"
      "Company releases: Knorr-Bremse FY2023-25; ZF AR2024-25; SAF-Holland AR2024; Mahle AR2024-25; Haldex year-end 2021; Rheinmetall FY2025 + AEQUITA divestment (3 Jun 2026); Circle BidCo/A.P. Møller offer documents (2024); Brakes India-TBK alliance (Dec 2025)", body_size=10)
panel(s, 6.8, 1.3, 6.05, 2.85, "Peer-reviewed & academic",
      "Göhring & von Glasner, SAE 902206 (1990); Kajiyama et al. (Isuzu/TBK), SAE 902200 (1990); Day, Harding & Newcomb, Proc IMechE (1984); Hoover et al., SAE 2005-01-3614; Subel & Kienhöfer, Proc IMechE D (2020)\n"
      "Ampadu, Alrejjal & Ksaibati, JSDTL 8(1) 2023 — disc vs drum cost-effectiveness\n"
      "Granitz et al., SN Applied Sciences (2021) — e-coolant pump, 13L HD diesel; Bitsis & Miwa, SAE 2018-01-0980; Cummins SuperTruck 2, SAE 2019-01-0247; MIT S.M. thesis (2012) — HD parasitic losses\n"
      "Zhang et al., J. Hazardous Materials 465 (2024); Hagino, Atmosphere 15:75 (2024, JARI); Environmental Pollution (2023) — regen braking & brake wear\n"
      "Energies 16:4024 (2023), 19:1748 (2026); Sustainability 15:3132 (2023); Energies 16:168 (2023) — BEV/FCEV truck & bus thermal management", body_size=10)
panel(s, 0.5, 4.35, 6.1, 2.85, "Regulatory & legal",
      "NHTSA: FMVSS 121 final rule 74 FR 37122 (2009) + Regulatory Impact Analysis; DOT HS 811 367 (2010, NADS simulator); ESV 07-0242 (2007); heavy-vehicle AEB NPRM (2023)\n"
      "UNECE Regulation No. 13 (braking, Type-II fade tests)\n"
      "EU: Reg 2019/1242 & 2024/1610 (HDV CO2); Reg 2024/1257 (Euro 7); Reg 2017/2400 Annex IX (VECTO auxiliaries) + JRC112015; EPA Phase 2 (40 CFR 1037.520 accessory credits) & Phase 3 (2024)\n"
      "EC merger decisions: M.8102 Valeo/FTE (2017, 90-100% concentration, Raicam remedy); Knorr-Bremse/Haldex Phase II (IP-17-2126, 2017)\n"
      "Japan: MY2025 HDV fuel-economy standards (via ICCT/TransportPolicy); METI supplier-transition studies (2023, 2025); FMCSA Crash Cost Methodology (2025); AAA Foundation (2017)", body_size=10)
panel(s, 6.8, 4.35, 6.05, 2.85, "Industry & financial data (verified, labeled where single-source)",
      "Adoption data: Haldex ADB history; Bendix/Knorr production milestones; FleetOwner, TruckingInfo, Fleet Maintenance, Transport Topics, Fleet Equipment; NACFE confidence reports; NAS 25542; DOE SuperTruck (OSTI)\n"
      "Deal records: Knorr-Bremse-Bendix (2002); ZF-Wabco $7bn (2020); Cummins-Meritor $3.7bn (2022); SAF-Holland-Haldex (2022); Concentric take-private SEK 230 (2024); Pierburg-AEQUITA €350M (2026)\n"
      "Japanese filings via irbank, kitaishihon, buffett-code, kabutan, MarkLines (incl. TBK JSAE 2022 exhibit; Mikuni 2022 BEV-truck e-oil pump disclosure; Yamada e-water pump adoptions)\n"
      "Aggregators for computed ratios: macrotrends, stockanalysis, Simply Wall St — cross-checked against filings; market research (MarketsandMarkets, GVR, IDTechEx) used directionally only", body_size=10)

# Page numbers (skip title slide)
for _i, _s in enumerate(prs.slides, 1):
    if _i == 1:
        continue
    _pn = _s.shapes.add_textbox(Inches(12.75), Inches(7.12), Inches(0.45), Inches(0.3))
    _p = _pn.text_frame.paragraphs[0]
    _p.text = str(_i)
    _p.alignment = PP_ALIGN.RIGHT
    for _r in _p.runs:
        _r.font.size = Pt(9)
        _r.font.name = BODY_FONT
        _r.font.color.rgb = GREY

import warnings
with warnings.catch_warnings():
    warnings.simplefilter("error")
    prs.save("/home/user/Claude/TBK_rd_gap_slides.pptx")
print("saved", len(prs.slides._sldIdLst), "slides cleanly")

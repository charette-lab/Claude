"""Standalone attestation checklist (PDF) — maps each claim in the Athanase
presentation ("what we say") to the corresponding action in the investment
letter ("what we do"). Brand guidelines: white, Times/Arial, monochrome Blue,
the real mark. Two pages: investment process, risk process."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

BLUE1 = RGBColor(0x0E, 0x16, 0x20); BLUE2 = RGBColor(0x15, 0x21, 0x30)
BLUE3 = RGBColor(0x31, 0x43, 0x59); BLUE4 = RGBColor(0x55, 0x6A, 0x83)
BLUE5 = RGBColor(0xE2, 0xE5, 0xE9); BLUE6 = RGBColor(0xF6, 0xF7, 0xF9)
NAVY, NAVY_TX, SLATE, SLATE_LT = BLUE2, BLUE3, BLUE3, BLUE4
HEADERBG, BODY, SUBTLE = BLUE6, BLUE1, BLUE4
WHITE = RGBColor(0xFF, 0xFF, 0xFF); FOOT = BLUE4
SERIF, SANS = "Times New Roman", "Arial"
MARK_DARK = "assets/mark_dark.png"
_MD_AR = (lambda s: s[0] / s[1])(Image.open(MARK_DARK).size)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
_pg = {"n": 0}


def rect(s, x, y, w, h, fill=None, line=None, lw=1.0):
    sp = s.shapes.add_shape(1, x, y, w, h); sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    return sp


def tbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, size, color, bold=False, italic=False, first=False,
         align=PP_ALIGN.LEFT, after=0, lead=1.12, font=SANS):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(after); p.line_spacing = lead
    r = p.add_run(); r.text = text; f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = font
    return p


def page(title, subtitle):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, fill=WHITE)
    s.shapes.add_picture(MARK_DARK, Inches(0.6), Inches(0.5), height=Inches(0.3),
                         width=Emu(int(int(Inches(0.3)) * _MD_AR)))
    para(tbox(s, Inches(1.05), Inches(0.53), Inches(8), Inches(0.3)),
         "ATHANASE · PROCESS ATTESTATION", 11, SLATE_LT, first=True)
    para(tbox(s, Inches(0.6), Inches(1.05), Inches(12.1), Inches(0.7)),
         title, 30, NAVY_TX, first=True, font=SERIF)
    para(tbox(s, Inches(0.62), Inches(1.78), Inches(12.0), Inches(0.5)),
         subtitle, 12.5, SUBTLE, first=True, italic=True, lead=1.2)
    _pg["n"] += 1
    para(tbox(s, Inches(0.6), Inches(7.12), Inches(12.13), Inches(0.3)),
         "Strictly confidential · for professional and intermediary use only", 8,
         FOOT, first=True)
    para(tbox(s, Inches(11.8), Inches(7.12), Inches(0.95), Inches(0.3)),
         str(_pg["n"]), 8, FOOT, first=True, align=PP_ALIGN.RIGHT)
    return s


def checklist(s, rows, top, rh=Inches(0.6)):
    cb_x, say_x, do_x = Inches(0.6), Inches(1.2), Inches(6.35)
    say_w = Emu(int(do_x) - int(say_x) - int(Inches(0.25)))
    do_w = Inches(6.35)
    para(tbox(s, say_x, top, say_w, Inches(0.3)), "WHAT WE SAY  (in the presentation)",
         10, SLATE, first=True, bold=True)
    para(tbox(s, do_x, top, do_w, Inches(0.3)), "WHAT WE DO  (in this letter)",
         10, SLATE, first=True, bold=True)
    y = Emu(int(top) + int(Inches(0.38)))
    rect(s, cb_x, y, Inches(12.13), Pt(1.4), fill=SLATE)
    y = Emu(int(y) + int(Inches(0.1)))
    for say, do in rows:
        rect(s, cb_x, Emu(int(y) + int(Inches(0.07))), Inches(0.24), Inches(0.24),
             line=SLATE_LT, lw=1.2)
        para(tbox(s, say_x, y, say_w, rh, anchor=MSO_ANCHOR.MIDDLE), say, 11.5,
             NAVY_TX, first=True, bold=True, lead=1.12)
        para(tbox(s, do_x, y, do_w, rh, anchor=MSO_ANCHOR.MIDDLE), do, 10.5, BODY,
             first=True, lead=1.16)
        y = Emu(int(y) + int(rh))
        rect(s, cb_x, y, Inches(12.13), Pt(0.7), fill=BLUE5)
    return y


# ===========================================================================
# PAGE 1 · INVESTMENT PROCESS
# ===========================================================================
s = page("Investment process — say what we do",
         "Tick each claim from the presentation against the matching action in "
         "this letter — we say it, the letter does it.")
checklist(s, [
    ("Selective, high-conviction sourcing.",
     "~One new investment a year, from a scored universe of ~20,000 companies; "
     "the sourcing funnel is documented."),
    ("Three tollgates before any capital.",
     "Each idea is minuted through Tollgate 1 (valuation), 2 (plan + board-seat "
     "path) and 3 (IC Go/No-Go), with named decision-makers."),
    ("Buy quality at a discount — never broken turnarounds.",
     "Tollgate 1 requires the core alone to justify the price — a 30–40% margin "
     "of safety versus a PE bidder."),
    ("Agreement before capital.",
     "No sizing to conviction without a secured path to a board seat; if "
     "alignment fails, the position is unwound (walk-away ledger)."),
    ("We create the value from the boardroom.",
     "30+ public board seats; the ownership playbook — refocus, reallocate, "
     "bolt-on, exit — executed per holding."),
    ("Returns from earnings, not re-rating.",
     "Underwritten to a constant exit multiple; portfolio companies grew EBITA "
     "~14.5% a year during our ownership."),
    ("Determined, not forced, sellers.",
     "We exit on a changed thesis, completely and on the facts — never forced by "
     "a fund clock or capital calls."),
], Inches(2.3), rh=Inches(0.62))

# ===========================================================================
# PAGE 2 · RISK PROCESS
# ===========================================================================
s = page("Risk process — do what we say",
         "The same map for risk — each safeguard we describe is a mechanism the "
         "letter commits to and the team applies.")
checklist(s, [
    ("Permanent loss is engineered out before we invest.",
     "Only durable, market-leading cores; we value the rectifiable core, with "
     "zero value for “growth-trap” divisions."),
    ("Every idea clears an entry gauntlet.",
     "Approval requires ≥12% expected IRR and ≤20% probability of a 30% drawdown."),
    ("Diversification by hard limits.",
     "Capped exposures — sector ≤30%, plus economic-cycle, maturity and geography "
     "limits — monitored continuously."),
    ("Mechanised discipline, not emotion.",
     "Automatic −10 / −20 / −30% triggers force a re-underwrite, and an exit where "
     "the thesis is broken."),
    ("The board seat is the kill switch.",
     "Board control freezes poor capital allocation and changes management — risk "
     "falls the moment control is secured."),
    ("No leverage at the fund.",
     "Any leverage is isolated to single investments, never applied at the fund "
     "level."),
    ("Independent checks, segregated duties.",
     "Operations, valuation and risk sit apart from investment; custody SEB, "
     "administration MUFG, audit KPMG."),
    ("Sized so no single mark can dominate.",
     "8–12 concentrated positions within hard limits, on a pre-agreed drawdown "
     "tolerance minuted with the Investment Committee."),
], Inches(2.3), rh=Inches(0.55))

prs.save("Athanase_Process_Checklist.pptx")
print(f"saved ({len(prs.slides._sldIdLst)} pages)")

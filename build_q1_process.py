"""Standalone attestation checklist (PDF) — maps each claim in the Athanase
presentation ("what we say") to the corresponding action in the investment
letter ("what we do"). Brand guidelines: white, Times/Arial, monochrome Blue,
the real mark. Two pages: investment process, risk process."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

BLUE1 = RGBColor(0x0E, 0x16, 0x20); BLUE2 = RGBColor(0x15, 0x21, 0x30)
BLUE3 = RGBColor(0x31, 0x43, 0x59); BLUE4 = RGBColor(0x55, 0x6A, 0x83)
BLUE5 = RGBColor(0xE2, 0xE5, 0xE9); BLUE6 = RGBColor(0xF6, 0xF7, 0xF9)
NAVY, NAVY_TX, SLATE, SLATE_LT = BLUE2, BLUE3, BLUE3, BLUE4
HEADERBG, BODY, SUBTLE = BLUE6, BLUE1, BLUE4
WHITE = RGBColor(0xFF, 0xFF, 0xFF); FOOT = BLUE4
SERIF, SANS = "Times New Roman", "Arial"
MARK_DARK = "assets/mark_dark.png"
_MD_AR = (lambda s: s[0] / s[1])(Image.open(MARK_DARK).size)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
_pg = {"n": 0}


def rect(s, x, y, w, h, fill=None, line=None, lw=1.0):
    sp = s.shapes.add_shape(1, x, y, w, h); sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    return sp


def tbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, size, color, bold=False, italic=False, first=False,
         align=PP_ALIGN.LEFT, after=0, lead=1.12, font=SANS):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(after); p.line_spacing = lead
    r = p.add_run(); r.text = text; f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = font
    return p


def page(title, subtitle):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, fill=WHITE)
    s.shapes.add_picture(MARK_DARK, Inches(0.6), Inches(0.5), height=Inches(0.3),
                         width=Emu(int(int(Inches(0.3)) * _MD_AR)))
    para(tbox(s, Inches(1.05), Inches(0.53), Inches(8), Inches(0.3)),
         "ATHANASE · PROCESS ATTESTATION", 11, SLATE_LT, first=True)
    para(tbox(s, Inches(0.6), Inches(1.05), Inches(12.1), Inches(0.7)),
         title, 30, NAVY_TX, first=True, font=SERIF)
    para(tbox(s, Inches(0.62), Inches(1.78), Inches(12.0), Inches(0.5)),
         subtitle, 12.5, SUBTLE, first=True, italic=True, lead=1.2)
    _pg["n"] += 1
    para(tbox(s, Inches(0.6), Inches(7.12), Inches(12.13), Inches(0.3)),
         "Strictly confidential · for professional and intermediary use only", 8,
         FOOT, first=True)
    para(tbox(s, Inches(11.8), Inches(7.12), Inches(0.95), Inches(0.3)),
         str(_pg["n"]), 8, FOOT, first=True, align=PP_ALIGN.RIGHT)
    return s


def checklist(s, rows, top, rh=Inches(0.62)):
    cb_x, say_x, do_x = Inches(0.6), Inches(1.2), Inches(6.1)
    say_w = Emu(int(do_x) - int(say_x) - int(Inches(0.25)))
    do_w = Inches(6.6)
    para(tbox(s, say_x, top, say_w, Inches(0.3)), "WHAT WE SAY  (the standing claim)",
         10, SLATE, first=True, bold=True)
    para(tbox(s, do_x, top, do_w, Inches(0.3)), "WHAT WE DID  (Q1 2026)",
         10, SLATE, first=True, bold=True)
    y = Emu(int(top) + int(Inches(0.38)))
    rect(s, cb_x, y, Inches(12.13), Pt(1.4), fill=SLATE)
    y = Emu(int(y) + int(Inches(0.1)))
    for say, do in rows:
        rect(s, cb_x, Emu(int(y) + int(Inches(0.08))), Inches(0.2), Inches(0.2), fill=NAVY)
        para(tbox(s, say_x, y, say_w, rh, anchor=MSO_ANCHOR.MIDDLE), say, 11.5,
             NAVY_TX, first=True, bold=True, lead=1.12)
        para(tbox(s, do_x, y, do_w, rh, anchor=MSO_ANCHOR.MIDDLE), do, 10.5, BODY,
             first=True, lead=1.16)
        y = Emu(int(y) + int(rh))
        rect(s, cb_x, y, Inches(12.13), Pt(0.7), fill=BLUE5)
    return y


# ===========================================================================
# PAGE 1 - INVESTMENT PROCESS IN PRACTICE
# ===========================================================================
s = page("Investment process in practice - Q1 2026",
         "Each standing claim from our presentation, mapped to what we actually "
         "did this quarter - so you can verify we do what we say.")
checklist(s, [
    ("Selective; we won't force deployment.",
     "Held ~6.9% cash and flagged difficulty finding value; we are sourcing in "
     "the UK and Japan rather than chase an expensive market."),
    ("Buy quality at a discount.",
     "Core holdings are dominant franchises bought on value - Auto Trader (~10x "
     "its nearest rival), Rightmove (~80% of portal time), Somero (~80% global "
     "share) at a cyclical low, TBK on ~2x normalised earnings."),
    ("Agreement before capital; determined, not forced, sellers.",
     "Tendered Zalaris into Norvestor's NOK 100 cash offer (+37% on the day) - "
     "sold on a fair-value trigger, on the facts, not sentiment."),
    ("We create the value from the boardroom.",
     "Tempest: progressing a Lex Asea spin-off and a buyback / redemption of "
     "cash. TBK: pursuing the operational plan to lift margins toward 10%."),
    ("Returns from earnings, not re-rating.",
     "Operating profit grew across the book - Roko EBITA +5% (SEK 415m), Tempest "
     "operating result up (sales +11%), TBK EBIT JPY 1.50bn (from 0.94bn)."),
    ("Low cadence, high conviction.",
     "17 positions, the ten largest ~84% of NAV; no new platform forced in a "
     "quarter where we judged the market expensive."),
], Inches(2.3), rh=Inches(0.66))

# ===========================================================================
# PAGE 2 - RISK PROCESS IN PRACTICE
# ===========================================================================
s = page("Risk process in practice - Q1 2026",
         "The same map for risk - each safeguard, with the Q1 evidence that we "
         "applied it.")
checklist(s, [
    ("No leverage at the fund.",
     "Gross exposure 85.3% = net exposure 85.3% - zero fund-level leverage, with "
     "~6.9% cash."),
    ("Sized and measured so no mark dominates.",
     "17 positions; downside risk 8.7%, 1-day 95% VaR 1.5%; a full liquidity "
     "profile disclosed (22.5% of NAV within one day)."),
    ("Diversification across exposures.",
     "Spread across Nordic, UK and Japanese listings and across software, "
     "industrials, gaming, media and resources."),
    ("Mechanised discipline, not emotion.",
     "Acted on the Zalaris fair-value trigger; held cash rather than relax the "
     "valuation bar in a market we find expensive."),
    ("Quality cores; avoid permanent loss.",
     "The book is concentrated in durable market leaders; speculative exposure is "
     "limited and disclosed (e.g. NGEx held as exploration optionality)."),
    ("Independent checks and full disclosure.",
     "Performance sourced from StatPro; figures reconciled, with open items "
     "(liquidity placeholders, unit flags) listed explicitly for verification."),
], Inches(2.3), rh=Inches(0.66))

prs.save("Athanase_Q1_2026_Process_In_Practice.pptx")
print(f"saved ({len(prs.slides._sldIdLst)} pages)")

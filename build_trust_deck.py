"""UK INVESTMENT-TRUST pitch — built to the AIP Brand Guidelines (Blue group
only, Times New Roman headlines / Arial body, white content slides, the real
mark, 4:3). Keeps the stronger 20-slide content/structure (asset-class framings
+ trust specifics); launch specifics flagged indicative/proposed.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn as _qn
from pptx.enum.shapes import MSO_SHAPE_TYPE as _MST
from PIL import Image
import openpyxl

# ---- AIP brand palette (Blue group only) / fonts --------------------------
BLUE1 = RGBColor(0x0E, 0x16, 0x20); BLUE2 = RGBColor(0x15, 0x21, 0x30)
BLUE3 = RGBColor(0x31, 0x43, 0x59); BLUE4 = RGBColor(0x55, 0x6A, 0x83)
BLUE5 = RGBColor(0xE2, 0xE5, 0xE9); BLUE6 = RGBColor(0xF6, 0xF7, 0xF9)
NAVY, NAVY_TX, SLATE, SLATE_LT = BLUE2, BLUE3, BLUE3, BLUE4
HEADERBG, BODY, SUBTLE, DIVIDER = BLUE6, BLUE1, BLUE4, BLUE5
WHITE = RGBColor(0xFF, 0xFF, 0xFF); FOOT = BLUE4; LOSS = BLUE4
SERIF, SANS = "Times New Roman", "Arial"
LOGO_WHITE, MARK_DARK = "assets/logo_white.png", "assets/mark_dark.png"
_LW_AR = (lambda s: s[0] / s[1])(Image.open(LOGO_WHITE).size)
_MD_AR = (lambda s: s[0] / s[1])(Image.open(MARK_DARK).size)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
_state = {"n": 0}


def rect(s, x, y, w, h, fill=None, line=None, lw=1.0):
    sp = s.shapes.add_shape(1, x, y, w, h); sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    return sp


def tbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, size, color, bold=False, italic=False, first=False,
         align=PP_ALIGN.LEFT, after=8, lead=1.12, font=SANS, track=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(after); p.line_spacing = lead
    r = p.add_run(); r.text = text; f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = font
    if track is None:
        track = -3 if size >= 24 else 0
    if track:
        r._r.get_or_add_rPr().set("spc", str(int(round(size * track / 100.0 * 100))))
    return p


def runs(p, parts, size, after=8, lead=1.18, font=SANS):
    p.space_after = Pt(after); p.line_spacing = lead
    for txt, col, bold in parts:
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = col; r.font.name = font


def place_mark(s, x, y, h):
    s.shapes.add_picture(MARK_DARK, x, y, height=h, width=Emu(int(int(h) * _MD_AR)))


def place_logo_white(s, x, y, h):
    s.shapes.add_picture(LOGO_WHITE, x, y, height=h, width=Emu(int(int(h) * _LW_AR)))


def footer(s):
    _state["n"] += 1
    para(tbox(s, Inches(8.4), Inches(7.06), Inches(4.3), Inches(0.3)),
         f"Strictly confidential · professional & intermediary use only     "
         f"{_state['n']}", 8, FOOT, first=True, align=PP_ALIGN.RIGHT, after=0)


def content(num, section, title, subtitle=None):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, fill=WHITE)
    place_mark(s, Inches(0.55), Inches(0.24), Inches(0.26))
    para(tbox(s, Inches(0.98), Inches(0.27), Inches(9.0), Inches(0.3)),
         section, 11, SLATE_LT, first=True, after=0)
    band_h = Inches(1.45) if subtitle else Inches(1.05)
    rect(s, 0, Inches(0.62), SW, band_h, fill=HEADERBG)
    tt = tbox(s, Inches(0.6), Inches(0.74), Inches(12.1), Inches(0.85))
    p = tt.paragraphs[0]; p.line_spacing = 1.0
    r0 = p.add_run(); r0.text = f"{num}.  "
    r0.font.size = Pt(30); r0.font.color.rgb = SLATE_LT; r0.font.name = SERIF
    r1 = p.add_run(); r1.text = title
    r1.font.size = Pt(30); r1.font.color.rgb = NAVY_TX; r1.font.name = SERIF
    body_top = Inches(1.95)
    if subtitle:
        para(tbox(s, Inches(0.62), Inches(1.55), Inches(11.9), Inches(0.7)),
             subtitle, 13, SUBTLE, first=True, italic=True, after=0, lead=1.16)
        body_top = Inches(2.4)
    footer(s)
    return s, body_top


def bullets(s, items, top, x=Inches(0.62), w=Inches(11.9), size=13.5, gap=12):
    tf = tbox(s, x, top, w, Inches(4.4))
    for i, it in enumerate(items):
        lead, b = it if isinstance(it, tuple) else ("", it)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.18
        if lead:
            r0 = p.add_run(); r0.text = lead + "  "
            r0.font.size = Pt(size); r0.font.bold = True
            r0.font.color.rgb = NAVY_TX; r0.font.name = SANS
        r1 = p.add_run(); r1.text = b
        r1.font.size = Pt(size); r1.font.color.rgb = BODY; r1.font.name = SANS
    return tf


def card(s, x, y, w, h, title, body, tlead=None, accent=True):
    rect(s, x, y, w, h, fill=HEADERBG)
    if accent:
        rect(s, x, y, Inches(0.07), h, fill=NAVY)
    inx = Emu(int(x) + int(Inches(0.28)))
    para(tbox(s, inx, Emu(int(y) + int(Inches(0.22))), Emu(int(w) - int(Inches(0.5))),
              Inches(0.5)), title, 13.5, SLATE, first=True, bold=True, after=0)
    bt = tbox(s, inx, Emu(int(y) + int(Inches(0.78))), Emu(int(w) - int(Inches(0.54))),
              Emu(int(h) - int(Inches(0.95))))
    if tlead:
        para(bt, tlead, 10.5, NAVY_TX, first=True, bold=True, after=9, lead=1.2)
        para(bt, body, 11, BODY, after=0, lead=1.28)
    else:
        para(bt, body, 11.5, BODY, first=True, after=0, lead=1.3)


def statrow(s, stats, top, h=Inches(1.4)):
    n = len(stats); gap = Inches(0.16)
    cw = Emu(int((int(Inches(12.13)) - (n - 1) * int(gap)) / n))
    for i, (big, lab) in enumerate(stats):
        cx = Emu(int(Inches(0.6)) + i * (int(cw) + int(gap)))
        rect(s, cx, top, cw, h, fill=HEADERBG)
        para(tbox(s, cx, Emu(int(top) + int(Inches(0.18))), cw, Inches(0.6),
                  anchor=MSO_ANCHOR.MIDDLE), big, 33, NAVY_TX, first=True, font=SERIF,
             after=0, align=PP_ALIGN.CENTER)
        para(tbox(s, Emu(int(cx) + int(Inches(0.1))), Emu(int(top) + int(Inches(0.86))),
                  Emu(int(cw) - int(Inches(0.2))), Inches(0.5)), lab, 10.5, SLATE_LT,
             first=True, after=0, align=PP_ALIGN.CENTER, lead=1.12)


def callout(s, text, y=Inches(6.0), h=Inches(0.92), size=12.5):
    rect(s, Inches(0.6), y, Inches(12.13), h, fill=HEADERBG)
    rect(s, Inches(0.6), y, Inches(0.07), h, fill=NAVY)
    para(tbox(s, Inches(0.9), y, Inches(11.6), h, anchor=MSO_ANCHOR.MIDDLE),
         text, size, NAVY_TX, first=True, after=0, lead=1.2)


def closer(s, text, y=Inches(6.46), h=Inches(0.6), size=12.5):
    rect(s, Inches(0.6), y, Inches(12.13), h, fill=NAVY)
    para(tbox(s, Inches(0.8), y, Inches(11.7), h, anchor=MSO_ANCHOR.MIDDLE),
         text, size, WHITE, first=True, italic=True, after=0, track=0)


def dtable(s, headers, rows, top, colx, sizes=(13, 12), rh=Inches(0.56)):
    for j, htext in enumerate(headers):
        para(tbox(s, colx[j], top, Inches(4.6), Inches(0.4)), htext, sizes[0], SLATE,
             first=True, bold=True, after=0)
    ry = Emu(int(top) + int(Inches(0.46)))
    rect(s, colx[0], ry, Emu(int(Inches(12.73)) - int(colx[0])), Pt(1.4), fill=SLATE)
    ry = Emu(int(ry) + int(Inches(0.1)))
    for row in rows:
        for j, cell in enumerate(row):
            bold = isinstance(cell, tuple)
            txt = cell[0] if bold else cell
            col = SLATE if j == 0 else (NAVY_TX if bold else BODY)
            cw = Emu((int(colx[j + 1]) if j + 1 < len(colx) else int(Inches(12.8)))
                     - int(colx[j]) - int(Inches(0.2)))
            para(tbox(s, colx[j], ry, cw, rh, anchor=MSO_ANCHOR.MIDDLE), txt,
                 sizes[1], col, first=True, bold=(j == 0 or bold), after=0, lead=1.1)
        ry = Emu(int(ry) + int(rh))
        rect(s, colx[0], ry, Emu(int(Inches(12.73)) - int(colx[0])), Pt(0.8), fill=BLUE5)
    return ry


# ---- performance data + light cumulative-growth chart ---------------------
import numpy as _np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as _plt
from matplotlib import font_manager as _fmgr
for _ff in ("Arial", "Liberation Sans", "DejaVu Sans"):
    try:
        _fmgr.findfont(_ff, fallback_to_default=False); _plt.rcParams["font.family"] = _ff; break
    except Exception:
        continue


def _load_cmp():
    ws = openpyxl.load_workbook("data/Comparison_returns.xlsx", data_only=True)["Sheet1"]

    def grab(rows):
        out = []
        for c in range(2, ws.max_column + 1):
            for r in rows:
                v = ws.cell(r, c).value
                if isinstance(v, (int, float)):
                    out.append(v)
        return out
    return grab(range(4, 16)), grab(range(22, 34))


_MSCI_X, _ATH_X = _load_cmp()


def _growth_chart(path):
    ga, gm = [1.0], [1.0]
    for r in _ATH_X:
        ga.append(ga[-1] * (1 + r))
    for r in _MSCI_X:
        gm.append(gm[-1] * (1 + r))
    x = _np.linspace(2006, 2026, len(ga))
    fig, ax = _plt.subplots(figsize=(7.4, 4.05))
    ax.plot(x, ga, color="#152130", lw=2.7, label="Athanase (net)", zorder=3)
    ax.plot(x, gm, color="#8A97A6", lw=2.0, label="MSCI World IMI", zorder=2)
    ax.set_yscale("log"); ax.set_yticks([1, 2, 5, 10, 20])
    ax.set_yticklabels(["1×", "2×", "5×", "10×", "20×"])
    ax.set_ylim(0.8, 27); ax.set_xlim(2006, 2027.6)
    ax.set_xticks([2006, 2010, 2014, 2018, 2022, 2026])
    ax.annotate(f"{ga[-1]:.0f}×", (x[-1], ga[-1]), xytext=(7, 0), textcoords="offset points",
                color="#152130", fontsize=16, fontweight="bold", va="center")
    ax.annotate(f"{gm[-1]:.1f}×", (x[-1], gm[-1]), xytext=(7, 0), textcoords="offset points",
                color="#8A97A6", fontsize=12, fontweight="bold", va="center")
    for sp in ("top", "right"):
        ax.spines[sp].set_visible(False)
    for sp in ("left", "bottom"):
        ax.spines[sp].set_color("#AAB2BD")
    ax.tick_params(colors="#556A83", labelsize=10.5)
    ax.grid(axis="y", color="#E2E5E9", lw=0.8); ax.set_axisbelow(True)
    ax.legend(loc="upper left", frameon=False, fontsize=11)
    fig.tight_layout(); fig.savefig(path, dpi=200, bbox_inches="tight", transparent=True)
    _plt.close(fig)


_growth_chart("/tmp/trust_growth.png")


def _pct(v):
    return "n.m." if not isinstance(v, (int, float)) or abs(v) > 3 else f"{v*100:+.0f}%"


def _irr(v):
    return "n.m." if not isinstance(v, (int, float)) or v > 3 else f"{v*100:.0f}%"


def _moic(v):
    return f"{v:.1f}×" if isinstance(v, (int, float)) else str(v)


def load_fund2():
    ws = openpyxl.load_workbook("data/AIP_Trackrecord.xlsx", data_only=True)["Transactions"]
    out = []
    for r in range(4, 21):
        c = ws.cell(r, 3).value
        if not c:
            continue
        out.append(dict(period=str(ws.cell(r, 2).value or "").strip(),
                        company=str(c).strip(), irr=ws.cell(r, 6).value,
                        moic=ws.cell(r, 11).value, outp=ws.cell(r, 8).value))
    return out


IND = "(indicative)"

# ===========================================================================
# 1 · COVER
# ===========================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, fill=NAVY)
place_logo_white(s, Inches(0.6), Inches(0.55), Inches(0.55))
rect(s, Inches(0.62), Inches(3.0), Inches(0.5), Pt(2.4), fill=SLATE_LT)
para(tbox(s, Inches(0.6), Inches(2.45), Inches(11), Inches(0.4)),
     "A PROPOSED UK INVESTMENT TRUST", 13, SLATE_LT, first=True, after=0)
para(tbox(s, Inches(0.6), Inches(3.2), Inches(11.8), Inches(1.3)),
     "Engaged ownership, in a listed trust", 42, WHITE, first=True, after=0, font=SERIF)
para(tbox(s, Inches(0.62), Inches(4.55), Inches(11.2), Inches(0.7)),
     "An institutional, 20-year strategy — daily-dealing access for wealth "
     "portfolios", 16, DIVIDER, first=True, italic=True, after=0, font=SERIF, lead=1.22)
para(tbox(s, Inches(0.62), Inches(6.62), Inches(12), Inches(0.4)),
     "For discussion · June 2026 · professional and intermediary use only", 10.5,
     SLATE_LT, first=True, after=0)

# ===========================================================================
# 2 · AT A GLANCE
# ===========================================================================
s, top = content("1", "Fund overview", "Athanase at a glance",
                 "A proven engaged-ownership strategy — proposed in a listed, "
                 "daily-dealing trust for UK wealth.")
statrow(s, [("38", "investments as\nengaged owners"),
            ("19%", "avg annualised NET,\nany entry month"),
            ("7%", "worst-ever entry,\nstill positive net"),
            ("18×", "growth of capital\nsince 2006")], top)
by = Emu(int(top) + int(Inches(1.7)))
bullets(s, [
    ("~16% net p.a. since 2006", "versus 6.2% for the MSCI World IMI — +10 pts a "
     "year, net of all fees."),
    ("Entry timing barely matters", "the average entry month returned ~19% net; "
     "the worst still compounded at +7%."),
    ("No investor loss", "over any full holding period in 20 years."),
], by, x=Inches(0.62), w=Inches(6.0), size=12.5, gap=11)
card(s, Inches(6.9), by, Inches(5.83), Inches(2.55), "An integrated team",
     "A core team together since 2006 (some since 1996) — 38 companies and 30+ "
     "public board seats. Operations, valuation and risk sit apart from "
     "investment decisions; custody, administration and audit by SEB, MUFG and "
     "KPMG.")

# ===========================================================================
# 3 · TRUST AT A GLANCE (terms)
# ===========================================================================
s, top = content("2", "Fund overview", "The proposed trust at a glance",
                 "Indicative parameters for discussion — final terms to be set "
                 "with the sponsor and prospectus.")
dtable(s, ["Term", "Proposal"], [
    ("Structure", "UK closed-end investment trust, London-listed"),
    ("Listing venue", ("LSE Main Market or Specialist Fund Segment " + IND,)),
    ("Objective", "Long-term capital growth via engaged ownership"),
    ("Portfolio", "8–12 high-conviction listed Nordic / European mid-caps"),
    ("Liquidity", "Daily dealing in the shares; no capital calls"),
    ("Discount control", ("Buyback when shares trade materially below NAV " + IND,)),
    ("Alignment", "Significant team co-investment alongside shareholders"),
    ("Providers", "Custody SEB · Administration MUFG · Audit KPMG"),
], top, colx=[Inches(0.62), Inches(4.1)], rh=Inches(0.44))
para(tbox(s, Inches(0.62), Inches(6.95), Inches(12), Inches(0.3)),
     "Indicative and subject to contract, regulatory approval and final fund "
     "documentation. Not an offer of securities.", 8, FOOT, first=True, italic=True, after=0)

# ===========================================================================
# 4 · WHY NOW — BETA
# ===========================================================================
s, top = content("3", "The opportunity", "Why now — the easy decade for beta is over",
                 "Equities are expensive on almost every measure, and the "
                 "zero-rate tailwind that lifted all assets has reversed.")
para(tbox(s, Inches(0.62), top, Inches(5.9), Inches(0.4)),
     "EXPENSIVE ON MOST MEASURES", 12, SLATE, first=True, bold=True, after=8)
bullets(s, [
    "Shiller CAPE near ~36× — close to the dot-com extreme, roughly double the "
    "long-run ~17×.",
    "Forward P/E ~22× vs a ~16× average; the equity risk premium is at "
    "multi-decade lows.",
    "Most of the 2010s’ return came from re-rating, not earnings growth — that "
    "lever is largely spent.",
], Emu(int(top) + int(Inches(0.45))), x=Inches(0.62), w=Inches(5.85), size=12.5, gap=11)
para(tbox(s, Inches(6.9), top, Inches(5.85), Inches(0.4)),
     "THE MACRO TAILWIND HAS REVERSED", 12, SLATE, first=True, bold=True, after=8)
bullets(s, [
    "The zero-rate, disinflationary backdrop that lifted every asset for a "
    "decade has reversed.",
    "Deglobalisation, deficits and demographics point to structurally higher "
    "inflation and rates.",
    "A higher cost of capital compresses valuations and punishes long-duration, "
    "unprofitable growth.",
], Emu(int(top) + int(Inches(0.45))), x=Inches(6.9), w=Inches(5.85), size=12.5, gap=11)
callout(s, "Whether or not AI lifts productivity, beta is priced for perfection "
           "and the easy tailwind has shifted — so idiosyncratic, "
           "valuation-disciplined returns matter more, not less.")

# ===========================================================================
# 5 · PASSIVE CONCENTRATED
# ===========================================================================
s, top = content("4", "The opportunity", "Passive has become a concentrated bet",
                 "A cap-weighted index puts the most money into the stocks that "
                 "have already risen most — concentration dressed as diversification.")
statrow(s, [("~38%", "of the S&P 500 in its\n10 largest stocks"),
            ("~33%", "in the “Magnificent\nSeven” alone"),
            ("~50 yrs", "since US equity was\nthis concentrated")], top, h=Inches(1.35))
cy = Emu(int(top) + int(Inches(1.6))); cw = Inches(3.92); gx = Inches(0.18)
for i, (t, b) in enumerate([
    ("Cap-weighting buys high", "The more a stock rises, the more you must own — "
     "the opposite of buy-low discipline."),
    ("A concentrated factor bet", "The “diversified” index is now a leveraged "
     "bet on large-cap growth, momentum and a single theme (AI)."),
    ("Breadth is fragile", "A stumble in a handful of mega-caps now moves the "
     "whole index; the rest has been left behind.")]):
    cx = Emu(int(Inches(0.6)) + i * (int(cw) + int(gx)))
    card(s, cx, cy, cw, Inches(2.2), t, b)

# ===========================================================================
# 6 · PASSIVE RISK -> ACTIVE OPPORTUNITY (+AI)
# ===========================================================================
s, top = content("5", "The opportunity", "From passive risk to active opportunity",
                 "An expensive, concentrated, passive-dominated market is what "
                 "active is built for — and as analysis commoditises, the durable "
                 "edge is ownership, not stock-picking.")
bullets(s, [
    ("The setup —", "expensive equities, record concentration and "
     "price-insensitive passive flows that inflate the most crowded names."),
    ("Why active now —", "wider mispricing and dispersion reward valuation "
     "discipline again; the “active underperforms” critique is cyclical, and "
     "weakest in exactly this market."),
    ("AI commoditises analysis —", "quant and specialist edges compete on the "
     "same public data, which commoditises fastest (Kim, Muhn & Nikolaev, 2024; "
     "McLean & Pontiff, 2016)."),
], top, x=Inches(0.62), w=Inches(11.9), size=13.5, gap=13)
callout(s, "The most durable edge is ownership: buying mispriced businesses and "
           "changing the outcome — an edge AI cannot copy.", y=Inches(5.7))

# ===========================================================================
# 7 · VALUATIONS CRACK (+ live PE gating)
# ===========================================================================
s, top = content("6", "The opportunity", "When valuations crack, where will clients hold?",
                 "At record valuations the cycle will turn — and the hedge can "
                 "only be bought before the crack, not after.")
cw = Inches(3.92); gx = Inches(0.18)
for i, (t, b) in enumerate([
    ("Global equities — exposed", "Clients own the most-expensive, "
     "most-concentrated names, with nothing to do when they fall."),
    ("Private equity — the trap", "Marks lag then catch down, while illiquidity "
     "blocks the exit exactly when it is needed."),
    ("This trust — the hedge", "Board control floors the risk, the shares stay "
     "liquid, and the dislocation becomes the entry.")]):
    cx = Emu(int(Inches(0.6)) + i * (int(cw) + int(gx)))
    card(s, cx, top, cw, Inches(2.0), t, b)
rect(s, Inches(0.6), Inches(4.7), Inches(12.13), Inches(0.95), fill=HEADERBG)
rect(s, Inches(0.6), Inches(4.7), Inches(0.07), Inches(0.95), fill=NAVY)
para(tbox(s, Inches(0.9), Inches(4.8), Inches(11.5), Inches(0.3)),
     "LIVE — JUNE 2026 · BLOOMBERG / CNBC", 9.5, SLATE, first=True, bold=True, after=4)
para(tbox(s, Inches(0.9), Inches(5.1), Inches(11.6), Inches(0.5)),
     "Partners Group capped redemptions in its $8.6bn evergreen PE fund at 5% of "
     "NAV after requests hit 9.8% — KKR, Blackstone and Ares fell. The PE exit is "
     "closing, exactly as feared.", 12, NAVY_TX, first=True, after=0, lead=1.18)
closer(s, "While evergreen PE gates redemptions, a listed trust trades every day "
          "the market is open — the liquidity clients keep when it matters.")

# ===========================================================================
# 8 · TWO WAYS TO OWN ACTIVELY
# ===========================================================================
s, top = content("7", "The strategy", "Two ways to own actively: public or private",
                 "The same operational value creation is available in PE and "
                 "public markets — but the public route keeps what private gives up.")
dtable(s, ["Feature", "Private — PE / buyout", "Public — Engaged ownership"], [
    ("How you act", "Buy and control the whole company", ("Board seat and real influence",)),
    ("Entry price", "A ~40% control premium, up front", ("Market price — often at a discount",)),
    ("Liquidity", "10-year lock-up; blind-pool calls", ("Daily liquidity; exit any time",)),
    ("Pricing", "Smoothed quarterly appraisal marks", ("Honest daily market prices",)),
    ("Fees", "2 & 20, plus fees on idle dry powder", ("Lower, on capital deployed",)),
    ("Value engine", "Board-led operational improvement", ("The same — board-led improvement",)),
], top, colx=[Inches(0.62), Inches(3.5), Inches(8.0)], rh=Inches(0.5))

# ===========================================================================
# 9 · THREE HOMES
# ===========================================================================
s, top = content("8", "The strategy", "Where it fits in a portfolio: three homes",
                 "Engaged ownership needs no bucket of its own — it can be funded "
                 "from public equity, alternatives or private equity.")
cw = Inches(3.92); gx = Inches(0.18)
homes = [
    ("Public equity", "A concentrated, active equity sleeve — long-only listed "
     "companies with daily pricing and full transparency.",
     "Adds idiosyncratic, governance-driven alpha to a beta-dominated book."),
    ("Alternatives", "A catalyst strategy — returns from board change, capital "
     "returns and corporate events.",
     "The driver is event-specific, largely independent of rates and growth; it "
     "diversifies."),
    ("Private equity", "PE-style value creation — board influence, operational "
     "improvement, multi-year holds — in public markets.",
     "The same engaged-ownership engine as PE, but liquid, with no J-curve or "
     "lock-up."),
]
for i, (t, lens, why) in enumerate(homes):
    cx = Emu(int(Inches(0.6)) + i * (int(cw) + int(gx)))
    rect(s, cx, top, cw, Inches(3.5), fill=HEADERBG)
    rect(s, cx, top, Inches(0.07), Inches(3.5), fill=NAVY)
    inx = Emu(int(cx) + int(Inches(0.28)))
    para(tbox(s, inx, Emu(int(top) + int(Inches(0.22))), Emu(int(cw) - int(Inches(0.5))),
              Inches(0.4)), t, 14, SLATE, first=True, bold=True, after=0)
    bt = tbox(s, inx, Emu(int(top) + int(Inches(0.8))), Emu(int(cw) - int(Inches(0.54))),
              Inches(2.5))
    para(bt, "THE LENS", 9, SLATE_LT, first=True, bold=True, after=4)
    para(bt, lens, 11, BODY, after=11, lead=1.28)
    para(bt, "WHY FUND IT HERE", 9, SLATE_LT, bold=True, after=4)
    para(bt, why, 11, BODY, after=0, lead=1.28)

# ===========================================================================
# 10 · WHAT WE DO
# ===========================================================================
s, top = content("9", "The strategy", "What we do — we don’t pick stocks, we build companies",
                 "Quality-core constructivism: own good businesses and make them "
                 "better from the boardroom.")
bullets(s, [
    ("Quality at a discount —", "profitable, high-return businesses below "
     "intrinsic value; never broken turnarounds we cannot underwrite."),
    ("Value we create, not discover —", "board control, capital-allocation "
     "discipline and operational change, not leverage or financial engineering."),
    ("Returns from earnings, not re-rating —", "the gain comes from the business "
     "improving under our influence — bankable cash, not a hoped-for multiple."),
    ("Constructive, by agreement —", "board seats won collaboratively, not proxy "
     "wars fought — a Nordic governance model, 20 years honed."),
], top, x=Inches(0.62), w=Inches(11.9), size=14, gap=15)
closer(s, "Private-equity-style value creation — in transparent, listed, "
          "daily-liquid form.")

# ===========================================================================
# 11 · THREE TOLLGATES
# ===========================================================================
s, top = content("10", "Process", "How we invest: three tollgates",
                 "Every idea passes the same three gates — each with an explicit "
                 "bar and a named decision-maker.")
tg = [
    ("Tollgate 1 — Valuation", "“Is it cheap enough even if we change nothing?”",
     "The core business alone justifies the price — a 30–40% margin of safety vs a "
     "PE bidder.  (Investment Team)"),
    ("Tollgate 2 — The plan", "“Is the improvement plan real — and can we deliver it?”",
     "Upside quantified across value levers and validated, with a credible path to "
     "a board seat; includes psychometric evaluation of CEO / Chair.  (Deal Lead + "
     "Head of Research)"),
    ("Tollgate 3 — Go / No-Go", "“Right risk, right size, right fit?”",
     "Within risk and portfolio limits, conviction-sized, with no better use of "
     "the capital.  (Investment Committee)"),
]
for i, (h, q, p) in enumerate(tg):
    cy = Emu(int(top) + i * int(Inches(1.42)))
    rect(s, Inches(0.62), cy, Inches(0.5), Inches(0.5), fill=NAVY)
    para(tbox(s, Inches(0.62), cy, Inches(0.5), Inches(0.5), anchor=MSO_ANCHOR.MIDDLE),
         str(i + 1), 16, WHITE, first=True, bold=True, after=0, align=PP_ALIGN.CENTER, font=SERIF)
    tf = tbox(s, Inches(1.32), cy, Inches(11.4), Inches(1.3))
    para(tf, h, 14.5, SLATE, first=True, bold=True, after=3)
    rp = tf.add_paragraph(); runs(rp, [("The question:  ", NAVY_TX, True), (q, BODY, False)], 11.5, after=2)
    pp = tf.add_paragraph(); runs(pp, [("Pass only if:  ", NAVY_TX, True), (p, BODY, False)], 11.5, after=0)

# ===========================================================================
# 12 · OWNERSHIP PLAYBOOK
# ===========================================================================
s, top = content("11", "Process", "The ownership playbook",
                 "A repeatable sequence, honed over 38 investments, that turns a "
                 "board seat into operational value.")
pb = [
    ("1 · Secure board seat", "Often a precondition to invest — the nomination "
     "committee proposes an AIP director."),
    ("2 · Align the thesis", "Collaborate with management and stakeholders to "
     "agree where the value is."),
    ("3 · Refocus earnings", "Cut overhead, optimise footprint, exit loss-making "
     "“growth-trap” ventures."),
    ("4 · Reallocate to winners", "Redeploy freed-up capital into the profitable "
     "core and its strongest adjacencies."),
    ("5 · Grow by acquisition", "Add complementary products or scale through "
     "value-accretive bolt-ons."),
    ("6 · Exit", "Realise via M&A or the equity market once the core is "
     "refocused and compounding."),
]
cw = Inches(3.92); gx = Inches(0.18); chh = Inches(1.4)
for i, (t, b) in enumerate(pb):
    cx = Emu(int(Inches(0.6)) + (i % 3) * (int(cw) + int(gx)))
    cy = Emu(int(top) + (i // 3) * (int(chh) + int(Inches(0.2))))
    rect(s, cx, cy, cw, chh, fill=HEADERBG)
    rect(s, cx, cy, Inches(0.07), chh, fill=NAVY)
    inx = Emu(int(cx) + int(Inches(0.26)))
    para(tbox(s, inx, Emu(int(cy) + int(Inches(0.18))), Emu(int(cw) - int(Inches(0.5))),
              Inches(0.35)), t, 13, SLATE, first=True, bold=True, after=5)
    para(tbox(s, inx, Emu(int(cy) + int(Inches(0.6))), Emu(int(cw) - int(Inches(0.52))),
              Inches(0.7)), b, 11, BODY, first=True, after=0, lead=1.22)
para(tbox(s, Inches(0.6), Inches(6.4), Inches(12.13), Inches(0.4)),
     "We don’t buy broken companies — we buy hidden cores.", 14, SLATE, first=True,
     bold=True, after=0, italic=True, align=PP_ALIGN.CENTER)

# ===========================================================================
# 13 · RISK, HONESTLY
# ===========================================================================
s, top = content("12", "Process", "Risk, honestly",
                 "A risk system built to engineer out permanent loss — controlled "
                 "before the investment is made.")
para(tbox(s, Inches(0.62), top, Inches(5.9), Inches(0.4)),
     "THREE PILLARS OF PROTECTION", 12, SLATE, first=True, bold=True, after=9)
bullets(s, [
    ("Predictability —", "only durable, market-leading cores in slow-moving "
     "industries where future demand is underwritable."),
    ("Price (a valuation floor) —", "value the rectifiable core, zero for "
     "“growth-trap” divisions — a structural 30–40% margin of safety vs PE."),
    ("Influence —", "a board seat as the “kill switch” — stopping a bad decision "
     "lowers risk below passive or active managers."),
], Emu(int(top) + int(Inches(0.45))), x=Inches(0.62), w=Inches(5.9), size=12, gap=12)
rect(s, Inches(6.9), top, Inches(5.83), Inches(3.7), fill=HEADERBG)
rect(s, Inches(6.9), top, Inches(0.07), Inches(3.7), fill=NAVY)
para(tbox(s, Inches(7.18), Emu(int(top) + int(Inches(0.22))), Inches(5.3), Inches(0.4)),
     "AUTOMATIC GUARDRAILS", 12, SLATE, first=True, bold=True, after=0)
bullets(s, [
    ("Entry gauntlet —", "ideas must clear ≥12% expected IRR and ≤20% probability "
     "of a 30% drawdown."),
    ("Diversification limits —", "capped exposures to sector (≤30%), cycle, "
     "maturity and geography."),
    ("Mechanised discipline —", "automatic −10 / −20 / −30% triggers force a "
     "re-underwrite, and an exit if the thesis is broken."),
    ("No leverage —", "isolated to single investments, never at the fund level."),
], Emu(int(top) + int(Inches(0.78))), x=Inches(7.18), w=Inches(5.35), size=11, gap=9)

# ===========================================================================
# 14 · VERSUS PRIVATE EQUITY
# ===========================================================================
s, top = content("13", "Track record", "Versus private equity",
                 "On true risk, PE and Athanase carry similar volatility — but "
                 "Athanase returns more, with daily liquidity.")
dtable(s, ["Metric", "Private equity (realised)", "Athanase (engaged owner)"], [
    ("Real return", "~11–14% (above PME, below us)", ("~16% net, real, fully invested",)),
    ("Liquidity", "10-year lock-up, blind pools", ("Daily liquidity, listed",)),
    ("Transparency", "Smoothed, model-marked quarterly", ("Daily marks, public scoreboard",)),
    ("True volatility", "~28% (de-smoothed, levered)", ("27% (downside only ~11%)",)),
    ("Downside risk", "~16% chance of 30%+ loss / 3 yrs", ("~2% chance of 30%+ loss / 3 yrs",)),
], top, colx=[Inches(0.62), Inches(3.5), Inches(8.0)], rh=Inches(0.48))
para(tbox(s, Inches(0.62), Inches(6.7), Inches(12), Inches(0.4)),
     "Athanase captures 93% of the market’s upside but only 43% of its downside; "
     "at 0.44 correlation it improves the whole portfolio’s efficient frontier.",
     10, FOOT, first=True, after=0, italic=True)

# ===========================================================================
# 15 · TRACK RECORD CHART
# ===========================================================================
s, top = content("14", "Track record", "£1 has become ~£18, net — vs ~£3 for the index",
                 "Twenty years of compounding from operating improvement — net of "
                 "all fees, independently reconciled.")
s.shapes.add_picture("/tmp/trust_growth.png", Inches(0.5), Inches(2.4), width=Inches(6.9))
rx = Inches(7.7); ry = Inches(2.5)
for big, lab in [("~16%", "net p.a. over 20 years"),
                 ("+10 pts", "a year vs the index, net"),
                 ("0%", "investor loss over a holding period"),
                 ("2.40", "Sortino ratio — downside-aware")]:
    para(tbox(s, rx, ry, Inches(5.0), Inches(0.5)), big, 26, NAVY_TX, first=True, font=SERIF, after=0)
    para(tbox(s, rx, Emu(int(ry) + int(Inches(0.47))), Inches(5.0), Inches(0.32)),
         lab, 11, SLATE_LT, first=True, after=0)
    ry = Emu(int(ry) + int(Inches(0.86)))
para(tbox(s, Inches(0.5), Inches(6.78), Inches(12.2), Inches(0.3)),
     "Cumulative growth of capital, net of fees, 2006–2025 (log scale). Newly "
     "proposed trust; manager’s existing strategy. Past performance is not a "
     "guide to future returns.", 8, FOOT, first=True, italic=True, after=0)

# ===========================================================================
# 16 · FULL FUND II BOOK
# ===========================================================================
s, top = content("15", "Track record", "Every investment in the current fund (AIP Fund II)",
                 "The full book, ranked by money multiple — winners, and the "
                 "discipline on the one that disappointed.")
deals = load_fund2(); deals.sort(key=lambda d: d["moic"] if isinstance(d["moic"], (int, float)) else 0, reverse=True)
heads = ("Company", "Holding", "Gross IRR", "MOIC", "vs index")
colx = [Inches(0.62), Inches(3.6), Inches(6.0), Inches(8.3), Inches(10.6)]
for j, h in enumerate(heads):
    para(tbox(s, colx[j], top, Inches(2.4), Inches(0.3)), h, 11, SLATE, first=True, bold=True,
         after=0, align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT)
ry = Emu(int(top) + int(Inches(0.34)))
rect(s, colx[0], ry, Inches(12.11), Pt(1.2), fill=SLATE)
ry = Emu(int(ry) + int(Inches(0.06)))
rh = Inches(0.247)
for k, d in enumerate(deals):
    if k % 2 == 0:
        rect(s, Inches(0.55), ry, Inches(12.2), rh, fill=HEADERBG)
    loss = isinstance(d["moic"], (int, float)) and d["moic"] < 1.0
    vals = [d["company"], d["period"], _irr(d["irr"]), _moic(d["moic"]), _pct(d["outp"])]
    for j, v in enumerate(vals):
        col = (LOSS if loss else (NAVY_TX if j == 0 else BODY))
        para(tbox(s, colx[j], ry, Inches(2.4) if j else Inches(2.9), rh, anchor=MSO_ANCHOR.MIDDLE),
             v, 10, col, first=True, bold=(j == 0), italic=loss,
             align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT, after=0)
    ry = Emu(int(ry) + int(rh))
para(tbox(s, Inches(0.62), Emu(int(ry) + int(Inches(0.06))), Inches(12.1), Inches(0.3)),
     "Gross deal-level figures, AIP Fund II to 2025; “n.m.” where holding periods "
     "are too short to annualise. Past performance is not a guide to future returns.",
     8, FOOT, first=True, italic=True, after=0)

# ===========================================================================
# 17 · WHY A TRUST
# ===========================================================================
s, top = content("16", "Structure & governance", "Why a listed trust is the ideal home",
                 "A closed-end structure matches long-horizon engaged ownership "
                 "better than any open-ended wrapper.")
bullets(s, [
    ("Permanent capital = patient capital —", "no forced redemptions, so we are "
     "never a forced seller; the multi-year engagement is protected and "
     "dislocations become opportunities."),
    ("Daily liquidity for clients —", "shares trade on the LSE — daily dealing, "
     "no capital calls — while evergreen PE funds gate redemptions (Partners "
     "Group, June 2026)."),
    ("An independent board —", "a majority-independent board oversees the manager "
     "on shareholders’ behalf."),
    ("Tier-one providers —", "operations, valuation and audit segregated from "
     "investment — SEB · MUFG · KPMG."),
], top, x=Inches(0.62), w=Inches(11.9), size=13.5, gap=14)
closer(s, "An institutional strategy clients can buy and sell like any listed "
          "share — without giving up the patient capital it needs.")

# ===========================================================================
# 18 · DISTRIBUTION & DISCOUNT
# ===========================================================================
s, top = content("17", "Structure & governance", "Income and discount control",
                 "A total-return strategy the trust structure can still wrap with "
                 "a predictable income and a tight discount.")
card(s, Inches(0.62), top, Inches(5.96), Inches(3.3), "Distribution policy (proposed)",
     "Capital growth at the core — realised as listed gains. UK trusts can pay a "
     "smoothed annual dividend from realised gains and revenue reserves, so a "
     "steady, quarterly income is available where advisers want it. Target yield "
     "and frequency to be set at launch.")
card(s, Inches(6.77), top, Inches(5.96), Inches(3.3), "Discount control",
     "A buyback policy that triggers when the shares trade materially below NAV — "
     "supported by realised exit gains — to keep any discount tight. Listing on "
     "the LSE Main Market or Specialist Fund Segment, weighed on index "
     "eligibility, investor base and compliance.")
para(tbox(s, Inches(0.62), Inches(6.9), Inches(12), Inches(0.3)),
     "Proposed; distribution policy, yield and discount mechanism to be agreed "
     "with the sponsor and board, subject to reserves and regulation.", 8, FOOT,
     first=True, italic=True, after=0)

# ===========================================================================
# 19 · TIMETABLE
# ===========================================================================
s, top = content("18", "Structure & governance", "An indicative launch timetable",
                 "A standard 3–4 month roadmap from mandate to admission.")
steps = [("MONTH 1", "Structuring, board and provider appointments; mandate and "
          "policies agreed."),
         ("MONTH 2", "Draft prospectus and due diligence; discount and "
          "distribution policies finalised."),
         ("MONTH 3", "Regulatory review; marketing and cornerstone engagement."),
         ("MONTH 4", "Roadshow, book-build and admission to listing.")]
cw = Inches(2.95); gx = Inches(0.16)
for i, (h, b) in enumerate(steps):
    cx = Emu(int(Inches(0.62)) + i * (int(cw) + int(gx)))
    rect(s, cx, top, cw, Inches(2.4), fill=HEADERBG)
    rect(s, cx, top, cw, Inches(0.5), fill=SLATE)
    para(tbox(s, cx, top, cw, Inches(0.5), anchor=MSO_ANCHOR.MIDDLE), h, 13, WHITE,
         first=True, bold=True, after=0, align=PP_ALIGN.CENTER)
    para(tbox(s, Emu(int(cx) + int(Inches(0.24))), Emu(int(top) + int(Inches(0.72))),
              Emu(int(cw) - int(Inches(0.46))), Inches(1.5)), b, 11.5, BODY, first=True,
         after=0, lead=1.3)
para(tbox(s, Inches(0.62), Inches(6.55), Inches(12.1), Inches(0.3)),
     "Illustrative; actual timetable, venue and discount policy subject to the "
     "sponsor, board and regulatory approval.", 8, FOOT, first=True, italic=True, after=0)

# ===========================================================================
# 20 · CLOSE
# ===========================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, fill=NAVY)
place_logo_white(s, Inches(0.6), Inches(0.55), Inches(0.5))
rect(s, Inches(0.62), Inches(2.5), Inches(0.5), Pt(2.4), fill=SLATE_LT)
para(tbox(s, Inches(0.6), Inches(2.7), Inches(12), Inches(1.0)),
     "Bringing institutional engaged ownership to UK wealth", 32, WHITE,
     first=True, after=0, font=SERIF, lead=1.06)
tf = tbox(s, Inches(0.62), Inches(4.0), Inches(12), Inches(2.2))
for ld, b in [
    ("For your clients:", " a liquid, listed, differentiated equity holding — a "
     "20-year net track record, real diversification and downside protection, "
     "daily dealing."),
    ("Why now:", " expensive, concentrated markets, gating private equity and "
     "decaying analytical alpha — the hedge is bought before the crack."),
    ("The proposal:", " partner to launch a UK investment trust on this strategy.")]:
    p = tf.add_paragraph() if tf.paragraphs[0].runs else tf.paragraphs[0]
    runs(p, [(ld, WHITE, True), (b, DIVIDER, False)], 15, after=14, lead=1.25)
para(tbox(s, Inches(0.62), Inches(6.45), Inches(12), Inches(0.5)),
     "Stefan Charette   ·   Athanase Industrial Partner   ·   charette@athanase.se",
     13, SLATE_LT, first=True, after=0)

# ---- final 4:3 rescale (match the AIP brand decks) ------------------------
TARGET_W, TARGET_H = 9753600, 7315200


def _rescale_shape(sh, sx, sy):
    try:
        L, T, W, H = sh.left, sh.top, sh.width, sh.height
    except Exception:
        L = T = W = H = None
    is_pic = sh.shape_type == _MST.PICTURE
    if None not in (L, T, W, H):
        sh.left = int(L * sx); sh.top = int(T * sy)
        sh.width = int(W * sx); sh.height = int(H * sx) if is_pic else int(H * sy)
    if sh.has_text_frame:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size is not None:
                    r.font.size = Pt(round(r.font.size.pt * sx, 1))
            pPr = p._p.find(_qn("a:pPr"))
            if pPr is not None:
                for a in ("marL", "indent"):
                    v = pPr.get(a)
                    if v is not None:
                        pPr.set(a, str(int(round(int(v) * sx))))


_sx = TARGET_W / int(prs.slide_width); _sy = TARGET_H / int(prs.slide_height)
for _sl in prs.slides:
    for _sh in _sl.shapes:
        _rescale_shape(_sh, _sx, _sy)
prs.slide_width = TARGET_W; prs.slide_height = TARGET_H

prs.save("Athanase_Investment_Trust_Pitch.pptx")
print(f"saved ({len(prs.slides._sldIdLst)} slides)")

"""Build a combined Athanase deck:
   Part I  - Why allocate to engaged ownership (activist) as an asset class
   Part II - Why choose Athanase among engaged owners
Styled to match the Athanase Industrial Partner brand.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

# ---- Brand palette (AIP Brand Guidelines v1.0 — BLUE group only) -----------
# Rule: use one base colour group together; never mix groups. Blue is primary.
BLUE1 = RGBColor(0x0E, 0x16, 0x20)
BLUE2 = RGBColor(0x15, 0x21, 0x30)   # primary logo / dark navy
BLUE3 = RGBColor(0x31, 0x43, 0x59)
BLUE4 = RGBColor(0x55, 0x6A, 0x83)
BLUE5 = RGBColor(0xE2, 0xE5, 0xE9)
BLUE6 = RGBColor(0xF6, 0xF7, 0xF9)

NAVY      = BLUE2   # cover / divider background, dark strips, dark bars
NAVY_TX   = BLUE3   # titles on white, table-header fills, lead/emphasis text
SLATE     = BLUE3   # primary accent: table headers, lead bullets, captions
SLATE_LT  = BLUE4   # section labels, secondary accent, lighter bars/lines
HEADERBG  = BLUE6   # light header band / card fill
BODY      = BLUE1   # body text (near-black with blue hint, in-group)
SUBTLE    = BLUE4   # subtitles / secondary
DIVIDER   = BLUE5   # light title text on navy
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
FOOT      = BLUE4   # footer
GOLD      = BLUE4   # accent rule
# Brand palette is monochrome Blue — no red. Losses are de-emphasised in a
# lighter in-group Blue (BLUE4); the negative sign carries the meaning.
LOSS      = BLUE4
HILITE    = BLUE3   # in-group highlight for callouts (was red)

SERIF = "Times New Roman"   # brand headline face
SANS  = "Arial"             # brand body face

import os
import math
from PIL import Image
LOGO_WHITE = "assets/logo_white.png"
MARK_DARK  = "assets/mark_dark.png"
_LW_AR = (lambda s: s[0] / s[1])(Image.open(LOGO_WHITE).size)
_MD_AR = (lambda s: s[0] / s[1])(Image.open(MARK_DARK).size)


def place_logo_white(slide, left, top, height):
    slide.shapes.add_picture(LOGO_WHITE, left, top, height=height,
                             width=Emu(int(int(height) * _LW_AR)))


def place_mark(slide, left, top, height):
    slide.shapes.add_picture(MARK_DARK, left, top, height=height,
                             width=Emu(int(int(height) * _MD_AR)))


# ---- Track-record data (from data/AIP_Trackrecord.xlsx, Transactions tab) ---
import openpyxl


def _fmt_pct(v):
    if not isinstance(v, (int, float)) or abs(v) > 3.0:
        return "n.m."
    return f"{v * 100:+.0f}%"


def _fmt_irr(v):
    if not isinstance(v, (int, float)) or v > 3.0:
        return "n.m."
    return f"{v * 100:.0f}%"


def _fmt_moic(v):
    return f"{v:.1f}x" if isinstance(v, (int, float)) else str(v)


def load_transactions():
    wb = openpyxl.load_workbook("data/AIP_Trackrecord.xlsx", data_only=True)
    ws = wb["Transactions"]

    def grab(r0, r1):
        deals = []
        for r in range(r0, r1 + 1):
            comp = ws.cell(r, 3).value          # C
            if not comp:
                continue
            deals.append(dict(
                period=str(ws.cell(r, 2).value or "").strip(),   # B
                company=str(comp).strip(),
                irr=ws.cell(r, 6).value,         # F
                outp=ws.cell(r, 8).value,        # H
                moic=ws.cell(r, 11).value,       # K
            ))
        return deals

    fund2 = grab(4, 20)         # AIP Fund II (2015-2026)
    hist = grab(27, 48)         # Prior period (2006-2014)
    return fund2, hist


def deal_table(slide, x, y, deals, col_w, font=10.5, rh=Inches(0.265),
               cols=("company", "period", "irr", "moic", "outp"),
               headers=("Company", "Holding", "IRR", "MOIC", "vs Index")):
    """Compact transactions table with header + alternating rows."""
    # header
    cx = x
    for ci, htext in enumerate(headers):
        rect(slide, cx, y, col_w[ci], rh, fill=SLATE)
        htf = tbox(slide, Emu(int(cx) + int(Inches(0.08))), y,
                   Emu(int(col_w[ci]) - int(Inches(0.12))), rh,
                   anchor=MSO_ANCHOR.MIDDLE)
        para(htf, htext, font, WHITE, bold=True, first=True,
             align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT, after=0)
        cx = Emu(int(cx) + int(col_w[ci]))
    yy = Emu(int(y) + int(rh))
    for ri, d in enumerate(deals):
        loss = isinstance(d["moic"], (int, float)) and d["moic"] < 1.0
        fill = HEADERBG if ri % 2 == 0 else WHITE
        cx = x
        for ci, key in enumerate(cols):
            rect(slide, cx, yy, col_w[ci], rh, fill=fill)
            if key == "company":
                txt = d["company"]
            elif key == "period":
                txt = d["period"]
            elif key == "irr":
                txt = _fmt_irr(d["irr"])
            elif key == "moic":
                txt = _fmt_moic(d["moic"])
            else:
                txt = _fmt_pct(d["outp"])
            is_loss_cell = loss and key in ("irr", "moic", "outp")
            col = LOSS if is_loss_cell else (NAVY_TX if ci == 0 else BODY)
            ctf = tbox(slide, Emu(int(cx) + int(Inches(0.08))), yy,
                       Emu(int(col_w[ci]) - int(Inches(0.12))), rh,
                       anchor=MSO_ANCHOR.MIDDLE)
            # In-palette loss treatment: lighter Blue + italic (brand allows
            # italic for emphasis); the negative sign carries the meaning.
            para(ctf, txt, font, col, bold=(ci == 0), italic=is_loss_cell,
                 first=True, align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT,
                 after=0, track=0)
            cx = Emu(int(cx) + int(col_w[ci]))
        yy = Emu(int(yy) + int(rh))
    return yy

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

_state = {"n": 0}


def rect(slide, l, t, w, h, fill=None, line=None, line_w=1.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill is not None:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line is not None:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def tbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    return tf


def para(tf, text, size, color, bold=False, italic=False, first=False,
         align=PP_ALIGN.LEFT, after=8, lead=1.1, font=SANS, track=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(after)
    p.line_spacing = lead
    r = p.add_run(); r.text = text
    f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = font
    # Brand type scale: letter-spacing (track in % of em, e.g. -5 => -0.05em).
    # Auto-apply brand defaults by size if not given: 48pt+ -> -5%, 24-47 -> -3%.
    if track is None:
        track = -5 if size >= 48 else (-3 if size >= 24 else 0)
    if track:
        spc = int(round(size * track / 100.0 * 100))  # EMU-less: points*100
        r._r.get_or_add_rPr().set("spc", str(spc))
    return p


def wordmark(slide, left, top, color=None, scale=1.0):
    """Place the real white Athanase wordmark (for navy backgrounds)."""
    place_logo_white(slide, left, top, Inches(0.52 * scale))


def footer(slide):
    _state["n"] += 1
    tf = tbox(slide, Inches(9.0), Inches(7.05), Inches(4.0), Inches(0.3))
    para(tf, f"Strictly confidential          {_state['n']}", 8.5, FOOT,
         first=True, align=PP_ALIGN.RIGHT, after=0)


# --- research-paper numbering: sections, figures, tables -------------------
_doc = {"part": 0, "sec": 0, "sub": 0, "subsec": 0}
_exhibit = {"fig": 0, "tbl": 0}
TOC = []   # (number, title, ref) for the Contents page


def fig():
    _exhibit["fig"] += 1
    return f'Figure {_exhibit["fig"]}'


def tbl():
    _exhibit["tbl"] += 1
    return f'Table {_exhibit["tbl"]}'


def dlabels(holder, fmt='0.0"%"', size=9, color=None, pos=None):
    """Show numeric data labels on a chart plot or series."""
    holder.has_data_labels = True
    dl = holder.data_labels
    dl.number_format = fmt
    dl.number_format_is_linked = False
    dl.font.size = Pt(size)
    dl.font.color.rgb = color if color is not None else BODY
    if pos is not None:
        dl.position = pos


def divider(title, kicker=None):
    _doc["part"] += 1
    _doc["sec"] = 0
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, fill=NAVY)
    wordmark(s, Inches(0.6), Inches(0.55), WHITE, scale=0.9)
    # brand rule: always left-align (never centre)
    rect(s, Inches(0.62), Inches(3.05), Inches(0.5), Pt(2.2), fill=SLATE_LT)
    if kicker:
        ktf = tbox(s, Inches(0.6), Inches(3.25), Inches(11.7), Inches(0.5))
        para(ktf, kicker.upper(), 13, SLATE_LT, first=True,
             align=PP_ALIGN.LEFT, after=0)
    tf = tbox(s, Inches(0.6), Inches(3.6), Inches(11.9), Inches(1.4))
    para(tf, title, 40, DIVIDER, italic=True, first=True,
         align=PP_ALIGN.LEFT, after=0, font=SERIF)
    return s


_SUBSECTIONS = [
    "Team & operations",
    "The opportunity & our edge",
    "How we invest",
    "Risk, honestly",
    "The 20-year track record",
    "What it means for your portfolio",
    "Versus private equity",
]
_subdiv = {"n": 0}


def subdivider(title):
    """Navy tab divider for a Part II sub-section. Does not affect section
    numbering and is excluded from the TOC. Shows the running tab list with the
    current tab highlighted, so it's easy to flip between subjects."""
    _subdiv["n"] += 1
    cur = _subdiv["n"]
    # advance sub-section numbering: 2.{sub}.{subsec}
    _doc["sub"] += 1
    _doc["subsec"] = 0
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, fill=NAVY)
    wordmark(s, Inches(0.6), Inches(0.55), WHITE, scale=0.9)
    # eyebrow + number
    para(tbox(s, Inches(0.6), Inches(2.45), Inches(11.7), Inches(0.4)),
         f"PART II · WHY ATHANASE   ·   {cur} OF {len(_SUBSECTIONS)}", 13,
         SLATE_LT, first=True, after=0, track=0)
    rect(s, Inches(0.62), Inches(2.95), Inches(0.5), Pt(2.2), fill=SLATE_LT)
    para(tbox(s, Inches(0.6), Inches(3.15), Inches(11.9), Inches(1.0)),
         title, 40, DIVIDER, italic=True, first=True, align=PP_ALIGN.LEFT,
         after=0, font=SERIF)
    # tab list (current highlighted), bottom-left
    ytf = tbox(s, Inches(0.6), Inches(5.0), Inches(12), Inches(2.0))
    for i, name in enumerate(_SUBSECTIONS, 1):
        p = ytf.paragraphs[0] if i == 1 else ytf.add_paragraph()
        p.space_after = Pt(4); p.line_spacing = 1.0
        r0 = p.add_run(); r0.text = f"{i}   "
        r1 = p.add_run(); r1.text = name
        on = (i == cur)
        for rr in (r0, r1):
            rr.font.size = Pt(12.5 if on else 11)
            rr.font.bold = on
            rr.font.color.rgb = WHITE if on else SLATE_LT
            rr.font.name = SANS
    footer(s)
    return s


def content(section, title, subtitle=None, number=True, ref=None):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, fill=WHITE)
    # top section label + real mark
    place_mark(s, Inches(0.55), Inches(0.24), Inches(0.26))
    lt = tbox(s, Inches(0.98), Inches(0.27), Inches(7.0), Inches(0.3))
    para(lt, section, 11, SLATE_LT, first=True, after=0)
    # top-right exhibit reference (Figure N / Table N)
    if ref:
        rtf = tbox(s, Inches(9.9), Inches(0.27), Inches(2.85), Inches(0.3))
        para(rtf, ref, 11, SLATE_LT, first=True, bold=True,
             align=PP_ALIGN.RIGHT, after=0)
    # numbered section header, research-paper style
    if number and _doc["part"] >= 1:
        if _doc["sub"] >= 1:
            _doc["subsec"] += 1
            _num = f'{_doc["sub"]}.{_doc["subsec"]}'
        else:
            _doc["sec"] += 1
            _num = f'{_doc["sec"]}'
        TOC.append((_doc["part"], _num, title, ref))
        title = f'{_num} {title}'
    # grey header band
    band_h = Inches(1.45) if subtitle else Inches(1.05)
    rect(s, 0, Inches(0.62), SW, band_h, fill=HEADERBG)
    tt = tbox(s, Inches(0.6), Inches(0.74), Inches(12.1), Inches(0.85))
    para(tt, title, 30, NAVY_TX, first=True, after=2, font=SERIF)
    body_top = Inches(1.95)
    if subtitle:
        st = tbox(s, Inches(0.62), Inches(1.55), Inches(11.9), Inches(0.5))
        para(st, subtitle, 13, SUBTLE, first=True, italic=True, after=0)
        body_top = Inches(2.35)
    footer(s)
    return s, body_top


def checklist(s, items, top, left=Inches(0.75), width=Inches(11.9),
              size=16, gap=14, lead_bold=True):
    tf = tbox(s, left, top, width, Inches(4.7))
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            lead, body = it
        else:
            lead, body = it, None
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.12
        # Hanging indent: wrapped lines align with the text, not under the marker.
        hang = int(round(size * 1.15 * 12700))   # marker "›  " width in EMU
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(hang))
        pPr.set("indent", str(-hang))
        r0 = p.add_run(); r0.text = "›  "
        r0.font.size = Pt(size); r0.font.bold = True
        r0.font.color.rgb = SLATE; r0.font.name = SANS
        r1 = p.add_run(); r1.text = lead
        r1.font.size = Pt(size); r1.font.bold = lead_bold
        r1.font.color.rgb = NAVY_TX; r1.font.name = SANS
        if body:
            r2 = p.add_run(); r2.text = "  " + body
            r2.font.size = Pt(size); r2.font.color.rgb = BODY; r2.font.name = SANS
    return tf


# ===========================================================================
# II.6h–j  Risk-adjusted analysis vs a global equity fund (MSCI World IMI)
# Data computed live from data/Comparison_returns.xlsx
# ===========================================================================
from statistics import NormalDist as _ND
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager as _fmgr

H_NAVY, H_BLUE3, H_BLUE4 = "#152130", "#314359", "#556A83"
H_BLUE5, H_BODY = "#E2E5E9", "#0E1620"
for _ff in ("Arial", "Liberation Sans", "DejaVu Sans"):
    try:
        _fmgr.findfont(_ff, fallback_to_default=False)
        plt.rcParams["font.family"] = _ff
        break
    except Exception:
        continue


def _load_cmp():
    ws = openpyxl.load_workbook("data/Comparison_returns.xlsx",
                               data_only=True)["Sheet1"]

    def grab(rows):
        out = []
        for c in range(2, ws.max_column + 1):
            for r in rows:
                v = ws.cell(r, c).value
                if isinstance(v, (int, float)):
                    out.append(v)
        return out
    return grab(range(4, 16)), grab(range(22, 34))   # MSCI, Athanase


_MSCI_X, _ATH_X = _load_cmp()
_NM = len(_MSCI_X)
# Annualise over ACTUAL invested time, not observation count. The track record
# is two periods with a gap (team set up AIPFII in between):
#   P1  9 Mar 2006 – 30 Jun 2014  = 3,035 days  (2006 starts 9 Mar; 62.5% is
#                                                 the Mar–Dec 2006 return)
#   gap 30 Jun 2014 – 23 Feb 2015 =   238 days  (EXCLUDED — not invested)
#   P2  23 Feb 2015 – 28 Feb 2026 = 4,023 days
# Total invested = 7,058 days / 365.25 = 19.32 years.
_INVESTED_YEARS = 7058 / 365.25


def _cmp_stats(x, mar=0.0, years=_INVESTED_YEARS):
    n = len(x); mean = sum(x) / n
    dd = math.sqrt(sum(min(r - mar, 0) ** 2 for r in x) / n)
    g = 1.0
    for r in x:
        g *= (1 + r)
    return dict(ann_ret=g ** (1 / years) - 1, cum_ret=g - 1, mult=g,
                ann_dd=dd * math.sqrt(12),
                ann_vol=(math.sqrt(sum((r - mean) ** 2 for r in x) / (n - 1))
                         * math.sqrt(12)))


_MSCI, _ATH = _cmp_stats(_MSCI_X), _cmp_stats(_ATH_X)


def _pair(a, m):
    n = len(a)

    def geo(xs):
        g = 1.0
        for x in xs:
            g *= (1 + x)
        return g ** (1 / len(xs)) - 1 if xs else 0.0
    up = [i for i in range(n) if m[i] > 0]
    dn = [i for i in range(n) if m[i] < 0]
    mb, ab = sum(m) / n, sum(a) / n
    cov = sum((m[i] - mb) * (a[i] - ab) for i in range(n)) / (n - 1)
    vm = sum((x - mb) ** 2 for x in m) / (n - 1)
    va = sum((x - ab) ** 2 for x in a) / (n - 1)

    def cum(xs):
        g = 1.0
        for x in xs:
            g *= (1 + x)
        return g

    def roll(months):
        w = t = 0
        for i in range(0, n - months + 1):
            t += 1
            w += cum(a[i:i + months]) > cum(m[i:i + months])
        return w / t if t else 0.0
    return dict(up_cap=geo([a[i] for i in up]) / geo([m[i] for i in up]),
                dn_cap=geo([a[i] for i in dn]) / geo([m[i] for i in dn]),
                corr=cov / math.sqrt(vm * va), beta=cov / vm,
                roll={y: roll(y * 12) for y in (3, 5, 7)})


_PAIR = _pair(_ATH_X, _MSCI_X)


def _entry_month_stats():
    """Annualised NET return if invested in any single month and held to the end
    — read directly from the workbook's 'invested in the month specified' table
    (AIPFII). Returns (average, worst, n) — the entry-timing-robustness stats."""
    wb = openpyxl.load_workbook("data/AIP_Trackrecord.xlsx",
                                data_only=True)["Monthly Returns"]
    vals = []
    for r in range(35, 44):                 # 2015–2023 entry rows
        for c in range(5, 17):              # Jan–Dec
            v = wb.cell(r, c).value
            if isinstance(v, (int, float)) and v != 0:
                vals.append(v)
    if not vals:
        return 0.0, 0.0, 0
    return sum(vals) / len(vals), min(vals), len(vals)


_ENTRY_AVG, _ENTRY_WORST, _ENTRY_N = _entry_month_stats()   # ~18.8% / ~6.8% / 107


def _project(T, r, dd):
    med = 100 * (1 + r) ** T
    down = 100 * (1 + (r - dd / math.sqrt(T))) ** T
    ploss = _ND().cdf((0 - r) / (dd / math.sqrt(T)))
    return med, down, ploss




# ===========================================================================
# COVER
# ===========================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, fill=NAVY)
wordmark(s, Inches(0.6), Inches(0.55), WHITE, scale=1.0)
rect(s, 0, Inches(4.5), Inches(0.5), Pt(2.2), fill=SLATE_LT)
tf = tbox(s, Inches(0.6), Inches(4.75), Inches(11.8), Inches(2.0))
para(tf, "Engaged Ownership as an Asset Class", 42, WHITE, bold=False,
     first=True, after=6, font=SERIF)
para(tf, "Why allocators should consider it — and why Athanase", 24, DIVIDER,
     italic=True, after=0, font=SERIF)
tf2 = tbox(s, Inches(0.62), Inches(6.6), Inches(8), Inches(0.5))
para(tf2, "An allocator briefing  ·  May 2026", 13, SLATE_LT, first=True, after=0)

# ===========================================================================
# AGENDA
# ===========================================================================
agenda_s, agenda_top = content("Overview", "Contents")
# TOC body rendered at end, once all sections are numbered.

# ===========================================================================
# PART I DIVIDER
# ===========================================================================
divider("Why allocate to engaged ownership", kicker="Part I  ·  The asset class")

# ---- I.0a Market backdrop: expensive equities + sticky rates ---------------
s, top = content("Market Context",
                 "Why now — the easy decade for beta is over",
                 "Equities are expensive on almost every measure, and the "
                 "disinflation and zero-rate tailwind that lifted all assets has "
                 "reversed.")
colL = tbox(s, Inches(0.75), top, Inches(5.85), Inches(4.0))
para(colL, "EXPENSIVE ON MOST MEASURES", 13, SLATE, first=True, bold=True,
     after=8)
for t in ["Shiller CAPE near ~36× — close to the dot-com extreme, and roughly "
          "double the long-run ~17×.",
          "Forward P/E ~22× against a ~16× long-run average; the equity risk "
          "premium is compressed to multi-decade lows.",
          "Most of the 2010s’ return came from re-rating — paying more per "
          "dollar of earnings — not earnings growth. That lever is largely "
          "spent."]:
    para(colL, t, 13.5, BODY, after=10, lead=1.16)
colR = tbox(s, Inches(7.0), top, Inches(5.85), Inches(4.0))
para(colR, "THE MACRO TAILWIND HAS REVERSED", 13, SLATE, first=True, bold=True,
     after=8)
for t in ["The zero-rate, disinflationary backdrop that lifted every asset for "
          "a decade has reversed.",
          "Deglobalisation, fiscal deficits and demographics point to "
          "structurally higher inflation and rates — higher for longer.",
          "A higher cost of capital compresses valuations and punishes "
          "long-duration, unprofitable growth — the opposite of the last "
          "cycle."]:
    para(colR, t, 13.5, BODY, after=10, lead=1.16)
rect(s, Inches(0.6), Inches(6.32), Inches(12.13), Inches(0.6), fill=NAVY)
para(tbox(s, Inches(0.8), Inches(6.32), Inches(11.7), Inches(0.6),
          anchor=MSO_ANCHOR.MIDDLE),
     "When beta is priced for perfection and the macro tailwind has reversed, "
     "index returns are set to be lower and riskier — the moment idiosyncratic, "
     "valuation-disciplined returns matter most.", 13, WHITE, first=True,
     italic=True, after=0)
para(tbox(s, Inches(0.6), Inches(7.12), Inches(8.0), Inches(0.36)),
     "Illustrative, recent data (2024–25); sources: R. Shiller (Yale) CAPE; "
     "S&P Dow Jones Indices; consensus forward earnings. Valuations move — the "
     "regime, not the decimal, is the point.", 7.5, FOOT, first=True, after=0,
     lead=1.05)

# ---- I.0b The passive paradox: concentration & factor risk -----------------
s, top = content("Market Context",
                 "Passive has become a concentrated bet",
                 "A cap-weighted index puts the most money into the stocks that "
                 "have already risen the most — concentration and factor risk "
                 "dressed as diversification.")
cstats = [("~38%", "of the S&P 500 in its\n10 largest stocks*"),
          ("~33%", "in the “Magnificent\nSeven” alone*"),
          ("~50 yrs", "since US equity was\nthis concentrated*")]
bw = Inches(3.9); gpx = Inches(0.2); cx = Inches(0.75)
for big, lab in cstats:
    rect(s, cx, top, bw, Inches(1.5), fill=HEADERBG)
    rect(s, cx, top, bw, Inches(0.06), fill=NAVY)
    ctf = tbox(s, cx, Emu(int(top) + int(Inches(0.22))), bw, Inches(1.2))
    para(ctf, big, 30, NAVY_TX, first=True, align=PP_ALIGN.CENTER, after=3,
         font=SERIF)
    for line in lab.split("\n"):
        para(ctf, line, 11.5, SUBTLE, align=PP_ALIGN.CENTER, after=0)
    cx = Emu(int(cx) + int(bw) + int(gpx))
checklist(s, [
    ("Cap-weighting buys high by design.", "The more a stock rises, the more of "
     "it you must own — the opposite of buy-low discipline."),
    ("A concentrated factor bet.", "The “diversified” index is now a leveraged "
     "position in large-cap growth, momentum and a single theme (AI) — not "
     "diversification."),
    ("Breadth is fragile.", "A stumble in a handful of mega-caps now moves the "
     "entire index; the rest of the market has been left behind."),
], Emu(int(top) + int(Inches(1.72))), size=14, gap=12)
para(tbox(s, Inches(0.6), Inches(7.12), Inches(8.0), Inches(0.36)),
     "*Illustrative, recent data (2024–25); source: S&P Dow Jones Indices. "
     "Index concentration and factor tilts change over time.", 7.5, FOOT,
     first=True, after=0, lead=1.05)

# ---- I.0b2 Reflexivity: passive flows cut both ways ------------------------
def _passive_chart(path):
    yrs = [1995, 2000, 2005, 2010, 2013, 2016, 2019, 2021, 2024]
    shr = [4, 9, 13, 19, 25, 34, 43, 50, 55]
    fig, ax = plt.subplots(figsize=(5.95, 3.5))
    ax.fill_between(yrs, shr, color=H_BLUE4, alpha=0.22, zorder=1)
    ax.plot(yrs, shr, color=H_NAVY, lw=2.6, zorder=3)
    ax.axhline(50, color=H_BLUE4, lw=1.0, ls=(0, (4, 3)), zorder=2)
    ax.text(1996, 51.5, "≈ passive overtakes active", color=H_BLUE3,
            fontsize=9.5, style="italic")
    ax.scatter([2024], [55], color=H_NAVY, s=30, zorder=4)
    ax.annotate("~55%", (2024, 55), textcoords="offset points", xytext=(-4, 7),
                ha="right", color=H_NAVY, fontsize=12, fontweight="bold")
    ax.annotate("~4%", (1995, 4), textcoords="offset points", xytext=(3, 6),
                color=H_BLUE3, fontsize=9.5)
    ax.set_xlim(1994.5, 2025.5); ax.set_ylim(0, 66)
    ax.set_yticks([0, 20, 40, 60]); ax.set_yticklabels(["0%", "20%", "40%", "60%"])
    ax.set_xticks([1995, 2000, 2005, 2010, 2015, 2020, 2024])
    ax.set_title("Passive share of US equity fund assets, 1995–2024",
                 color=H_NAVY, fontsize=12.5, fontweight="bold", loc="left",
                 pad=10)
    for sp in ("top", "right"):
        ax.spines[sp].set_visible(False)
    ax.spines["left"].set_color(H_BLUE5); ax.spines["bottom"].set_color(H_BLUE5)
    ax.tick_params(colors=H_BODY, labelsize=9.5)
    ax.grid(axis="y", color=H_BLUE5, lw=0.6, zorder=0)
    fig.tight_layout()
    fig.savefig(path, dpi=200, transparent=True)
    plt.close(fig)


_passive_chart("/tmp/mc_passive.png")
s, top = content("Market Context",
                 "Passive flows cut both ways",
                 "Index buying has grown from a niche to the majority of US "
                 "equity fund assets — and the same price-insensitive mechanics "
                 "can force selling on the way down.")
s.shapes.add_picture("/tmp/mc_passive.png", Inches(0.6),
                     Emu(int(top) + int(Inches(0.05))), width=Inches(6.15))
colR = tbox(s, Inches(7.05), top, Inches(5.6), Inches(4.2))
para(colR, "WHY IT CUTS BOTH WAYS", 12.5, SLATE, first=True, bold=True, after=8)
for t in ["Markets are inelastic — a $1 flow can move aggregate value by ~$5; "
          "flows increasingly set the level (Gabaix & Koijen, 2021).",
          "ETFs raise the volatility of the stocks they hold (Ben-David et al., "
          "2018); index inclusion inflates prices and comovement (Wurgler, "
          "2011).",
          "The reported passive share understates the true price-insensitive "
          "base (Chinco & Sammon, 2024).",
          "The multiplier runs in reverse: if flows slow, the most-owned "
          "mega-caps face mechanical selling — with few price-sensitive buyers "
          "left to absorb it."]:
    para(colR, t, 12, BODY, after=9, lead=1.16)
rect(s, Inches(0.6), Inches(6.32), Inches(12.13), Inches(0.6), fill=NAVY)
para(tbox(s, Inches(0.8), Inches(6.32), Inches(11.7), Inches(0.6),
          anchor=MSO_ANCHOR.MIDDLE),
     "The concentration that flattered index returns on the way up is a "
     "forced-seller overhang on the way down — a risk you do not carry when "
     "returns come from company-specific change, not from flows.", 13, WHITE,
     first=True, italic=True, after=0)
para(tbox(s, Inches(0.6), Inches(7.12), Inches(12.0), Inches(0.36)),
     "Passive share approximate; sources: ICI, Morningstar (index funds + ETFs "
     "as a share of US equity fund assets). Research citations in references.",
     7.5, FOOT, first=True, after=0, lead=1.05)


# ---- I.0c "Cheap" passive is not cheap -------------------------------------
s, top = content("Market Context",
                 "“Cheap” index exposure is not cheap",
                 "The headline expense ratio is the smallest part of the cost — "
                 "and you pay full price to own concentration and rich "
                 "valuations.")
ccards = [("3–9 bps", "the headline fee —\nall you are quoted"),
          ("20–80 bps", "a year, from index\nreconstitution alone*"),
          ("$1–2 bn", "a year transferred to\narbitrageurs**")]
bw = Inches(3.9); gpx = Inches(0.2); cx = Inches(0.75)
for big, lab in ccards:
    rect(s, cx, top, bw, Inches(1.4), fill=HEADERBG)
    rect(s, cx, top, bw, Inches(0.06), fill=NAVY)
    ctf = tbox(s, cx, Emu(int(top) + int(Inches(0.2))), bw, Inches(1.12))
    para(ctf, big, 29, NAVY_TX, first=True, align=PP_ALIGN.CENTER, after=3,
         font=SERIF)
    for line in lab.split("\n"):
        para(ctf, line, 11.5, SUBTLE, align=PP_ALIGN.CENTER, after=0)
    cx = Emu(int(cx) + int(bw) + int(gpx))
yb = Emu(int(top) + int(Inches(1.64)))
colL = tbox(s, Inches(0.75), yb, Inches(5.85), Inches(2.6))
para(colL, "THE COSTS YOU DON’T SEE", 12.5, SLATE, first=True, bold=True,
     after=7)
for t in ["Index reconstitution forces funds to buy additions high and sell "
          "deletions low, on predictable dates — the ~20–80 bps above "
          "(Petajisto, 2011).",
          "Bid-ask spreads, market impact, premium/discount to NAV and tax drag "
          "all sit on top of the headline fee."]:
    para(colL, t, 12.5, BODY, after=8, lead=1.15)
colR = tbox(s, Inches(7.0), yb, Inches(5.85), Inches(2.6))
para(colR, "THE COST THAT MATTERS MOST", 12.5, SLATE, first=True, bold=True,
     after=7)
for t in ["You buy the most expensive companies at peak weights — no valuation "
          "discipline, and no one improving the businesses.",
          "“Cheap” fees on an expensive, concentrated portfolio are a false "
          "economy: the price you pay is the risk you take."]:
    para(colR, t, 12.5, BODY, after=8, lead=1.15)
rect(s, Inches(0.6), Inches(6.32), Inches(12.13), Inches(0.6), fill=NAVY)
para(tbox(s, Inches(0.8), Inches(6.32), Inches(11.7), Inches(0.6),
          anchor=MSO_ANCHOR.MIDDLE),
     "The default “safe and cheap” choice now embeds maximum valuation and "
     "concentration risk — precisely why a valuation-disciplined, idiosyncratic "
     "return engine earns its place.", 13, WHITE, first=True, italic=True,
     after=0)
para(tbox(s, Inches(0.6), Inches(7.06), Inches(8.7), Inches(0.4)),
     "*Reconstitution cost, lower bound: ~21–28 bps/yr S&P 500, ~38–77 bps/yr "
     "Russell 2000 (Petajisto, 2011). **Combined index-fund losses to arbitrage "
     "(Chen, Noronha & Singal, 2006). See references.", 7.5, FOOT, first=True,
     after=0, lead=1.05)


# ---- I.1 Return & alpha ----
s, top = content("The Proposition",
                 "A catalyst-driven source of return",
                 "Engaged owners change the outcome — they do not merely absorb "
                 "underperformance.")
checklist(s, [
    ("~7% announcement return.", "Average abnormal return around an activist "
     "filing, with no reversal over the following year; success or partial "
     "success in ~two-thirds of campaigns (Brav, Jiang, Partnoy & Thomas, 2008)."),
    ("~10.2% around the 13D,", "plus a further ~11.4% over the next year for "
     "engaged-owner targets (Klein & Zur, 2009)."),
    ("No long-term give-back.", "~2,000 interventions show no negative drift and "
     "no “pump-and-dump” (Bebchuk, Brav & Jiang, 2015)."),
    ("Real, not financial, value.", "Production efficiency improves in the three "
     "years after intervention (Brav, Jiang & Kim, 2015)."),
    ("Idiosyncratic alpha.", "Returns tied to a specific change at a specific "
     "company — scarce in a beta-dominated portfolio."),
], top, size=15.5, gap=12)

# ---- I.2 Horizon & fiduciary ----
s, top = content("Fit", "Aligned with horizon and fiduciary duty")
checklist(s, [
    ("Patient capital wins.", "Board seats, strategic reviews and turnarounds "
     "play out over years — perpetual endowments and long-dated pensions can "
     "hold through to full value realisation."),
    ("Stewardship is the duty.", "Engaged ownership expresses fiduciary "
     "obligation (UK Stewardship Code; UN PRI) rather than conflicting with it."),
    ("Mission alignment.", "Demanding accountability and good governance reflects "
     "institutional values and is defensible to beneficiaries."),
    ("Engagement as risk management.", "Turns hidden governance risk in the "
     "portfolio into a remediable, return-generating exposure."),
], top, size=16.5, gap=16)

# ---- I.3 Diversification ----
s, top = content("Portfolio", "Diversification and the total-portfolio view")
checklist(s, [
    ("A distinct return engine.", "Alpha driven by proxy contests, spin-offs and "
     "capital-return programmes — largely independent of rates, growth and "
     "inflation."),
    ("Diversifies the equity book.", "Small- and mid-cap value and special "
     "situations, typically underweight in institutional portfolios."),
    ("Counter-cyclical dry powder.", "Dislocations create the next targets; "
     "patient capital deploys into weakness rather than selling it."),
    ("PE-style upside, public-market terms.", "Operational engagement with daily "
     "pricing, transparency, lower fees and better liquidity."),
], top, size=16.5, gap=16)

# ---- I.4 Where the allocation fits (which bucket) --------------------------
s, top = content("Portfolio",
                 "Where the allocation fits: three legitimate homes",
                 "Engaged ownership does not need a bucket of its own. It can be "
                 "funded from public equity, event-driven hedge funds or private "
                 "equity — each defensible, depending on what it replaces.")
buckets = [
    ("Public equity", [
        ("The lens", "A concentrated, active equity sleeve — long-only public "
         "companies with daily pricing and full transparency."),
        ("Why fund it here", "Adds idiosyncratic, governance-driven alpha to a "
         "beta-dominated equity book, while keeping liquidity, transparency and "
         "low fees inside an existing equity mandate."),
        ("The trade-off", "Higher active risk than an index sleeve — it is "
         "meant to differ from the benchmark.")]),
    ("Event-driven hedge funds", [
        ("The lens", "A catalyst strategy — returns come from proxy contests, "
         "board change, spin-offs, M&A and capital-return programmes, not "
         "market direction."),
        ("Why fund it here", "The return driver is event-specific and largely "
         "independent of rates and growth; it diversifies an absolute-return "
         "book and matches how the alpha is actually generated."),
        ("The trade-off", "More concentrated and longer-horizon than a typical "
         "multi-strategy fund; positions are owned outright, not hedged.")]),
    ("Private equity", [
        ("The lens", "PE-style value creation — board influence, operational "
         "improvement and multi-year holds — executed in public markets."),
        ("Why fund it here", "The same engaged-ownership return engine as "
         "private equity, but with public-market liquidity, daily transparency, "
         "lower fees and no J-curve or lock-up."),
        ("The trade-off", "Marked to market daily — honest volatility instead "
         "of smoothed quarterly appraisals (see Part II).")]),
]
cw = Inches(3.86); gapx = Inches(0.18); x = Inches(0.62)
hh = Inches(0.5); bodyH = Inches(3.5)
for name, blocks in buckets:
    rect(s, x, top, cw, hh, fill=NAVY)
    htf = tbox(s, x, top, cw, hh, anchor=MSO_ANCHOR.MIDDLE)
    para(htf, name, 15, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER,
         after=0, font=SERIF)
    by = Emu(int(top) + int(hh))
    rect(s, x, by, cw, bodyH, fill=HEADERBG)
    btf = tbox(s, Emu(int(x) + int(Inches(0.18))),
               Emu(int(by) + int(Inches(0.14))),
               Emu(int(cw) - int(Inches(0.36))),
               Emu(int(bodyH) - int(Inches(0.28))))
    first = True
    for label, txt in blocks:
        para(btf, label.upper(), 10.5, SLATE, bold=True, first=first, after=2)
        first = False
        para(btf, txt, 11, BODY, after=8, lead=1.13)
    x = Emu(int(x) + int(cw) + int(gapx))
rect(s, Inches(0.62), Inches(6.45), Inches(11.94), Inches(0.52), fill=NAVY)
para(tbox(s, Inches(0.8), Inches(6.45), Inches(11.6), Inches(0.52),
          anchor=MSO_ANCHOR.MIDDLE),
     "The same strategy, three valid homes — fund it wherever engaged "
     "ownership is the better risk-adjusted alternative to what it replaces.",
     13, WHITE, first=True, italic=True, after=0)

# ---- I.5 Governance / ESG ----
s, top = content("Stewardship", "Governance and the ESG synergy")
checklist(s, [
    ("Applied governance.", "Board independence, pay-for-performance and capital "
     "discipline are the “G” of ESG put into practice."),
    ("Constructive engagement.", "Return and stewardship objectives expressed "
     "through one allocation, disciplined by a value thesis."),
    ("A positive externality.", "Governance precedents improve the quality of the "
     "beta the whole portfolio owns."),
    ("Defensible to stakeholders.", "“We own companies and make them better-"
     "run and more valuable” is an easy story to tell."),
], top, size=16.5, gap=16)

# ---- I.5 The mid-market sweet spot -----------------------------------------
s, top = content("The Opportunity", "Why mid-market is the sweet spot",
                 "Big enough to run well, small enough to change — and to grow. "
                 "The engaged-ownership opportunity is richest in the middle.",
                 ref=tbl())
rows = [
    ("", "Small-cap", "Mid-market", "Large-cap"),
    ("Organisational depth",
     "Often too thin to support a capable team",
     "Enough structure for a real management team and processes",
     "Deep, but layered and slow"),
    ("Growth runway",
     "Long, but fragile and under-resourced",
     "Still long — small enough to compound for years",
     "Largely played out; growth is incremental"),
    ("Speed of change",
     "Fast, but limited resources to execute",
     "Operational change lands within ~1.5 years",
     "Years to turn; entrenched complexity"),
    ("Owner influence",
     "High, but little to work with",
     "A board seat genuinely moves the outcome",
     "Diluted; one owner rarely shifts the agenda"),
]
tl = Inches(0.55)
cw3 = [Inches(2.5), Inches(3.25), Inches(3.5), Inches(2.9)]
rh0 = Inches(0.4); rhd = Inches(0.78); y = top
for ri, row in enumerate(rows):
    head = ri == 0
    rh = rh0 if head else rhd
    x = tl
    for ci, cell in enumerate(row):
        mid = (ci == 2)
        if head:
            fill = SLATE if ci else WHITE
            if mid:
                fill = NAVY
        else:
            fill = (HEADERBG if mid else WHITE) if ri % 2 else \
                   (BLUE5 if mid else HEADERBG)
        rect(s, x, y, cw3[ci], rh, fill=fill)
        ctf = tbox(s, Emu(int(x) + int(Inches(0.12))), y,
                   Emu(int(cw3[ci]) - int(Inches(0.2))), rh,
                   anchor=MSO_ANCHOR.MIDDLE)
        if head:
            col = WHITE
        elif ci == 0:
            col = SLATE
        elif mid:
            col = NAVY_TX
        else:
            col = SUBTLE
        para(ctf, cell, 12 if head else 11, col,
             bold=(head or ci == 0 or mid), first=True, after=0, lead=1.05)
        x = Emu(int(x) + int(cw3[ci]))
    y = Emu(int(y) + int(rh))
para(tbox(s, Inches(0.55), Emu(int(y) + int(Inches(0.16))), Inches(12.2), Inches(0.5)),
     "Mid-market companies are run well enough to be worth owning, yet small "
     "enough that an engaged owner can still change — and grow — them.",
     13, SLATE, first=True, italic=True, after=0)

# ---- I.6 Why mid-cap management listens ------------------------------------
s, top = content("The Opportunity", "Why mid-cap management listens",
                 "The structural reason engaged ownership works here: capable "
                 "operators with predictable blind spots, and the room to act on "
                 "good advice.")
colL = tbox(s, Inches(0.75), top, Inches(5.85), Inches(4.6))
para(colL, "THE BLIND SPOT", 13, SLATE, first=True, bold=True, after=8)
for t in ["Mid-cap CEOs are usually promoted from sales or operations — strong "
          "in their domain, first-time chief executives",
          "They have not run capital allocation, M&A, financing or board "
          "strategy before — predictable, correctable gaps",
          "They know it — so credible, experienced owners are welcomed, not "
          "resisted, unlike at entrenched large-caps"]:
    para(colL, t, 14, BODY, after=11, lead=1.15)
colR = tbox(s, Inches(7.0), top, Inches(5.85), Inches(4.6))
para(colR, "WHAT IT MEANS FOR AN ENGAGED OWNER", 13, SLATE, first=True,
     bold=True, after=8)
for t in ["A board seat is genuinely productive — advice is acted on, not "
          "filtered through layers",
          "Change is collaborative, not a proxy fight — faster, cheaper, less "
          "headline risk",
          "The owner supplies exactly the missing experience: capital "
          "discipline, focus and a wider strategic lens"]:
    para(colR, t, 14, BODY, after=11, lead=1.15)
rect(s, Inches(0.6), Inches(6.35), Inches(12.13), Inches(0.6), fill=NAVY)
para(tbox(s, Inches(0.78), Inches(6.35), Inches(11.8), Inches(0.6),
          anchor=MSO_ANCHOR.MIDDLE),
     "Capable management with a known blind spot, in a company small enough to "
     "change: the conditions under which engaged ownership compounds.",
     13, WHITE, first=True, italic=True, after=0)

# ---- I.5 What to look for (bridge) ----
s, top = content("Selection", "What separates a great engaged owner",
                 "Dispersion is wide — manager selection matters more here "
                 "than in almost any other equity strategy.")
checklist(s, [
    ("A repeatable process", "for sourcing and underwriting correctable, "
     "undervalued companies."),
    ("Genuine operating capability,", "not just financial engineering — the "
     "experience to control the outcome from the boardroom."),
    ("A track record of outcomes,", "board seats won and value delivered across "
     "multiple market cycles."),
    ("Aligned economics", "and a stable, long-tenured team that invests its own "
     "capital."),
    ("Disciplined entry and risk control", "so the strategy delivers without "
     "uncompensated concentration risk."),
], top, size=15.5, gap=12)
# bridge line
bt = tbox(s, Inches(0.75), Inches(6.45), Inches(11.8), Inches(0.5))
para(bt, "Part II shows how Athanase meets each of these criteria.", 14,
     SLATE, first=True, italic=True, after=0)

# ===========================================================================
# PART II DIVIDER
# ===========================================================================
divider("Why Athanase", kicker="Part II  ·  Among engaged owners")

subdivider("Team & operations")
# ---- II.1 At a glance ----
s, top = content("Overview", "Athanase at a glance")
# four big stats
stats = [("38", "investments as\nengaged owners"),
         (f"{_ENTRY_AVG*100:.0f}%", "avg annualised NET,\nany entry month*"),
         (f"{_ENTRY_WORST*100:.0f}%", "worst-ever entry,\nstill positive*"),
         (f"{_ATH['mult']:.0f}×", "growth of capital\nsince 2006")]
bw = Inches(2.95); gapx = Inches(0.18); x = Inches(0.75); y = top
for big, lab in stats:
    card = rect(s, x, y, bw, Inches(1.85), fill=HEADERBG)
    ctf = tbox(s, x, y + Inches(0.28), bw, Inches(1.4), anchor=MSO_ANCHOR.TOP)
    para(ctf, big, 40, NAVY_TX, first=True, align=PP_ALIGN.CENTER, after=4,
         font=SERIF)
    for j, line in enumerate(lab.split("\n")):
        para(ctf, line, 12, SUBTLE, align=PP_ALIGN.CENTER, after=0)
    x = Emu(int(x) + int(bw) + int(gapx))
checklist(s, [
    (f"~{_ENTRY_AVG*100:.0f}% average annualised return*",
     "whatever month you invested and held to today — entry timing has barely "
     "mattered."),
    (f"~{_ENTRY_WORST*100:.0f}% in the single worst entry month",
     "even the unluckiest possible entry still compounded at a mid-single-digit "
     "net return — no investor lost over the holding period."),
    (f"{_ATH['ann_ret']*100:.1f}% time-weighted NET since 2006",
     f"vs MSCI IMI {_MSCI['ann_ret']*100:.1f}% — net of all fees."),
], top + Inches(2.15), size=15, gap=10)
nt = tbox(s, Inches(0.75), Inches(6.75), Inches(11.8), Inches(0.4))
para(nt, "*Average of the annualised NET return earned investing in each "
     f"individual month and holding to the end ({_ENTRY_N} entry months, AIPFII). "
     "Time-weighted figure annualised over actual invested time. Past "
     "performance is not indicative of future results.",
     9, FOOT, first=True, after=0)

# ---- II.2 The team ----
s, top = content("Experience", "An integrated team, 20 years together",
                 "In a world of specialists, a rare unit that has the operating "
                 "experience to control the outcome.")
# heritage timeline visual
nodes = [("1992", "Öresund"), ("1996", "Custos"),
         ("2006", "Creades"), ("2015", "Athanase")]
tl_y = Emu(int(top) + int(Inches(0.42)))           # the line
x0, x1 = Inches(1.6), Inches(11.5)
rect(s, x0, tl_y, Emu(int(x1) - int(x0)), Pt(1.6), fill=SLATE_LT)
n = len(nodes)
for i, (yr, name) in enumerate(nodes):
    cx = Emu(int(x0) + int((int(x1) - int(x0)) * i / (n - 1)))
    # year above
    yt = tbox(s, Emu(int(cx) - int(Inches(1.0))), Emu(int(tl_y) - int(Inches(0.42))),
              Inches(2.0), Inches(0.35))
    para(yt, yr, 15, NAVY_TX, first=True, bold=True, align=PP_ALIGN.CENTER,
         after=0, font=SERIF)
    # node dot
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(cx) - int(Inches(0.09))),
                             Emu(int(tl_y) - int(Inches(0.07))),
                             Inches(0.18), Inches(0.18))
    dot.fill.solid(); dot.fill.fore_color.rgb = NAVY
    dot.line.color.rgb = WHITE; dot.line.width = Pt(1.5); dot.shadow.inherit = False
    # name below
    nt = tbox(s, Emu(int(cx) - int(Inches(1.1))), Emu(int(tl_y) + int(Inches(0.16))),
              Inches(2.2), Inches(0.3))
    para(nt, name, 12, SLATE, first=True, align=PP_ALIGN.CENTER, after=0)
checklist(s, [
    ("38 companies managed,", "with a proven ability to tell a “strong "
     "core” from a “prestige trap”."),
    ("30+ public board seats", "— real experience mandating pivots and "
     "protecting shareholder interests, not writing engagement letters."),
    ("A dual lens:", "a forensic investing lens to find the value and a C-suite "
     "operating lens that understands the “how” of execution."),
    ("Skin in the game.", "The team uses Athanase as its own investment vehicle "
     "and shares in the success economics."),
    ("Institutional memory.", "Multiple market cycles together eliminate "
     "key-person risk and prevent repeating historical mistakes."),
], top + Inches(1.25), size=14.5, gap=10)


# ---- II.2b The people behind the track record (investment + operations) ----
s, top = content("Experience", "The people behind the track record",
                 "An integrated investment team — and an independent operations, "
                 "finance and control function that safeguards investor capital.")


def _person_card(x, y, w, h, nm, role, line):
    rect(s, x, y, w, h, fill=HEADERBG)
    rect(s, x, y, Inches(0.06), h, fill=NAVY)
    para(tbox(s, Emu(int(x) + int(Inches(0.2))), Emu(int(y) + int(Inches(0.13))),
              Emu(int(w) - int(Inches(0.32))), Inches(0.3)),
         nm, 12.5, NAVY_TX, first=True, bold=True, after=0, font=SERIF, track=0)
    para(tbox(s, Emu(int(x) + int(Inches(0.2))), Emu(int(y) + int(Inches(0.42))),
              Emu(int(w) - int(Inches(0.32))), Inches(0.24)),
         role.upper(), 8.5, SLATE, first=True, bold=True, after=0, track=0)
    para(tbox(s, Emu(int(x) + int(Inches(0.2))), Emu(int(y) + int(Inches(0.68))),
              Emu(int(w) - int(Inches(0.34))), Inches(0.7)),
         line, 9, BODY, first=True, after=0, lead=1.1, track=0)


# Band 1 — investment team (five partners, one row)
para(tbox(s, Inches(0.6), Inches(2.18), Inches(8), Inches(0.3)),
     "INVESTMENT TEAM", 11, SLATE, first=True, bold=True, after=0, track=0)
inv = [
    ("Stefan Charette", "CIO & Partner", "CEO/CIO of Custos, Investment AB "
     "Öresund and Creades; CEO of Brokk; investment banker at Lehman Brothers "
     "and Salomon Smith Barney."),
    ("Kenth Eriksson", "Senior PM & Partner", "Ex-CEO of Tradimus; Senior VP, "
     "Electrolux; tech owner-operator."),
    ("Sven Thorén", "PM & Partner", "Ex-PM at Adrigo, Catella, Pan Capital and "
     "Nordea."),
    ("Anders Elsell", "Senior PM", "On the team since 1996; ex-CEO/CIO of HQ "
     "Fonder; Carnegie."),
    ("Daniel Nyhrén", "Partner", "On the team since 2006; ex-CFO of Global "
     "Batterier."),
]
iw = Inches(2.33); igx = Inches(0.12); ix0 = Inches(0.6); iy = Inches(2.5)
for k, (nm, role, line) in enumerate(inv):
    _person_card(Emu(int(ix0) + k * (int(iw) + int(igx))), iy, iw, Inches(1.5),
                 nm, role, line)

# Band 2 — operations, finance & control (four; the safety layer)
para(tbox(s, Inches(0.6), Inches(4.2), Inches(9), Inches(0.3)),
     "OPERATIONS, FINANCE & CONTROL — INDEPENDENT OF THE INVESTMENT TEAM",
     11, SLATE, first=True, bold=True, after=0, track=0)
ops = [
    ("Malin Norrman", "CFO", "Ex-COO/CFO of Captor Fund Mgmt; risk manager at "
     "Carnegie and AB Industrivärden."),
    ("Grant Loon", "Head of Fund Operations", "Ex-COO of Madrague Capital and "
     "Port Capital; Soros, Commerzbank, Morgan Stanley."),
    ("Eva Andersson", "Fund Operations", "MBA, Business School London; sales-"
     "executive background."),
    ("Zetong Liu", "Technology", "MSc applied & computational mathematics, KTH; "
     "projects for the Swiss Finance Institute / EPFL."),
]
ow = Inches(2.96); ogx = Inches(0.13); ox0 = Inches(0.6); oy = Inches(4.52)
for k, (nm, role, line) in enumerate(ops):
    _person_card(Emu(int(ox0) + k * (int(ow) + int(ogx))), oy, ow, Inches(1.45),
                 nm, role, line)

# safety strip — why operations matters to investors
rect(s, Inches(0.6), Inches(6.32), Inches(12.13), Inches(0.62), fill=NAVY)
para(tbox(s, Inches(0.8), Inches(6.32), Inches(11.7), Inches(0.62),
          anchor=MSO_ANCHOR.MIDDLE),
     "Segregation of duties protects investors: operations, valuation and risk "
     "sit apart from investment decisions — with SEB (custody), MUFG "
     "(administration) and KPMG (audit) independently verifying every NAV and fee.",
     12, WHITE, first=True, italic=True, after=0, track=0)




# ---- II.2d Risk system (deepened) ------------------------------------------
s, top = content("Risk Management",
                 "A risk system that engineers out permanent loss",
                 "Risk is controlled before the investment is made — through "
                 "valuation, diversification and influence — not managed after "
                 "the fact.", ref=tbl())
para(tbox(s, Inches(0.7), top, Inches(5.6), Inches(0.3)),
     "BUILT ON THREE PILLARS", 11, SLATE, first=True, bold=True, after=0,
     track=0)
pillars = [
    ("Predictability", "Only durable, market-leading cores in slow-moving "
     "industries — businesses whose future demand we can underwrite."),
    ("Price", "A hard valuation floor: value the rectifiable core, assign zero "
     "to “growth-trap” divisions — a 30–40% margin of safety vs a PE bidder."),
    ("Influence", "A board seat as the “kill switch” — the ability to stop one "
     "bad decision lowers risk below any passive or active manager."),
]
py = top + Inches(0.42)
for ti, bd in pillars:
    rect(s, Inches(0.7), py, Inches(5.65), Inches(1.18), fill=HEADERBG)
    rect(s, Inches(0.7), py, Inches(0.07), Inches(1.18), fill=NAVY)
    para(tbox(s, Inches(0.92), Emu(int(py) + int(Inches(0.14))), Inches(5.3),
              Inches(0.3)), ti, 12.5, NAVY_TX, first=True, bold=True, after=0,
         font=SERIF, track=0)
    para(tbox(s, Inches(0.92), Emu(int(py) + int(Inches(0.46))), Inches(5.3),
              Inches(0.66)), bd, 10, BODY, first=True, after=0, lead=1.13,
         track=0)
    py = Emu(int(py) + int(Inches(1.18)) + int(Inches(0.1)))
para(tbox(s, Inches(7.1), top, Inches(5.6), Inches(0.3)),
     "AUTOMATIC GUARDRAILS", 11, SLATE, first=True, bold=True, after=0, track=0)
gy3 = top + Inches(0.42)
for t in ["Core & satellite: a 30-stock equal-weighted index compounds all "
          "unallocated capital; 8–12 concentrated ideas sit on top.",
          "Entry gauntlet: every idea must clear a ≥12% expected IRR and "
          "≤20% probability of a 30% drawdown before funding.",
          "Eleven standardised risk tags cap any single macro/business factor "
          "at ~20% of the concentrated book.",
          "Hard triggers at −10% / −20% / −30% force thesis review and exit — "
          "removing emotion and PM bias."]:
    para(tbox(s, Inches(7.1), gy3, Inches(5.65), Inches(0.9)),
         t, 12, BODY, first=True, after=0, lead=1.16, track=0)
    gy3 = Emu(int(gy3) + int(Inches(0.92)))
para(tbox(s, Inches(0.6), Inches(7.16), Inches(12.6), Inches(0.4)),
     "The system is designed to make permanent capital loss structurally "
     "unlikely — not merely to react to it.", 9, FOOT, first=True, after=0,
     track=0, lead=1.1)


# ---- II.2e Ownership model (the value-creation playbook) -------------------
s, top = content("Ownership Model",
                 "The ownership playbook, honed over 38 investments",
                 "A repeatable sequence that turns a board seat into operational "
                 "value — refocusing a good core, then growing it.")
steps = [
    ("1", "Secure the board seat",
     "Often a precondition to invest — the nomination committee proposes an AIP "
     "director."),
    ("2", "Align the thesis",
     "Collaborate with management and stakeholders to agree where the value is "
     "and what must change."),
    ("3", "Refocus to drive earnings",
     "Cut overhead and admin, optimise the footprint, exit loss-making "
     "“growth-trap” ventures."),
    ("4", "Reallocate to winners",
     "Redeploy freed-up capital into the profitable core and its strongest "
     "adjacencies; lift R&D where it pays."),
    ("5", "Grow by acquisition",
     "Add complementary products, geographies or scale through value-accretive "
     "bolt-ons."),
    ("6", "Exit",
     "Realise via M&A or in the equity market once the core is refocused and "
     "compounding."),
]
sw = Inches(3.97); sh_ = Inches(1.72); sgx = Inches(0.1); sgy = Inches(0.16)
sx0 = Inches(0.6)
for idx, (num, ti, bd) in enumerate(steps):
    r, c = divmod(idx, 3)
    bx = Emu(int(sx0) + c * (int(sw) + int(sgx)))
    by = Emu(int(top) + r * (int(sh_) + int(sgy)))
    rect(s, bx, by, sw, sh_, fill=HEADERBG)
    bdg = s.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(bx) + int(Inches(0.22))),
                             Emu(int(by) + int(Inches(0.2))), Inches(0.5),
                             Inches(0.5))
    bdg.fill.solid(); bdg.fill.fore_color.rgb = NAVY
    bdg.line.fill.background(); bdg.shadow.inherit = False
    btf = bdg.text_frame; btf.word_wrap = False
    bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
    brn = bp.add_run(); brn.text = num; brn.font.size = Pt(18)
    brn.font.bold = True; brn.font.color.rgb = WHITE; brn.font.name = SERIF
    para(tbox(s, Emu(int(bx) + int(Inches(0.9))), Emu(int(by) + int(Inches(0.24))),
              Emu(int(sw) - int(Inches(1.1))), Inches(0.5)),
         ti, 12.5, NAVY_TX, first=True, bold=True, after=0, font=SERIF, track=0)
    para(tbox(s, Emu(int(bx) + int(Inches(0.25))), Emu(int(by) + int(Inches(0.82))),
              Emu(int(sw) - int(Inches(0.45))), Inches(0.85)),
         bd, 9.5, BODY, first=True, after=0, lead=1.12, track=0)
para(tbox(s, Inches(0.6), Inches(6.5), Inches(12.2), Inches(0.5)),
     "The same playbook across 38 companies — value created by fixing capital "
     "allocation, not by financial engineering or leverage.",
     12.5, SLATE, first=True, italic=True, after=0)


subdivider("The opportunity & our edge")
# ---- II.3 Market opportunity ----
_f = fig()
s, top = content("Market Opportunity", "The inefficiency Athanase harvests",
                 "Corporate decay is structural — it continuously recreates "
                 "low-risk, high-return entry points.", ref=_f)
checklist(s, [
    ("Duration mismatch.", "CEOs run 7–8 year defensive cycles while moats "
     "erode at hyper-speed; an 80%-correct strategy decays to ~50% efficient."),
    ("The prestige trap.", "Empire-building and the sunk-cost fallacy lead "
     "management to double down on unprofitable ventures."),
    ("Shareholder exhaustion.", "As patience runs out, price hits a long-term low "
     "and liquidity dries up — Athanase becomes the “buyer of last "
     "resort” at a discount."),
    ("Low competition.", "Passive buys the most expensive names; active managers "
     "wait for proof; PE needs 10x the capital — leaving these bargains open."),
], top, left=Inches(0.7), width=Inches(6.3), size=13.5, gap=14)
# probability-decay curve (right)
dd = CategoryChartData()
dd.categories = ["Yr 1", "Yr 2", "Yr 3", "Yr 4", "Yr 5", "Yr 6", "Yr 7", "Yr 8"]
dd.add_series("Strategy efficiency", (80, 76, 71, 65, 60, 55, 52, 50))
gf = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(7.15), top + Inches(0.2),
                        Inches(5.7), Inches(3.6), dd)
ch = gf.chart
ch.has_title = True
ch.chart_title.text_frame.text = f"{_f}.  “Probability decay”: strategy efficiency over a CEO cycle (%)"
ch.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
ch.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = SUBTLE
ch.has_legend = False
sr = ch.plots[0].series[0]
sr.format.line.color.rgb = NAVY
sr.format.line.width = Pt(2.5)
sr.smooth = True
dlabels(sr, fmt='0"%"', size=9, color=NAVY_TX, pos=XL_LABEL_POSITION.ABOVE)
ch.category_axis.tick_labels.font.size = Pt(9)
ch.value_axis.tick_labels.font.size = Pt(9)
ch.value_axis.has_major_gridlines = False
ch.value_axis.minimum_scale = 0
ch.value_axis.maximum_scale = 100
para(tbox(s, Inches(7.15), top + Inches(3.85), Inches(5.7), Inches(0.4)),
     "Even an excellent team’s edge erodes as the external environment shifts.",
     10, FOOT, first=True, italic=True, after=0)

# ---- II.3b Negative-selection rebuttal: hidden high-quality core ----
_f2 = fig()
s, top = content("Investment Strategy",
                 "We don’t buy broken companies — we buy hidden cores",
                 "The classic objection to activism is adverse selection. Our "
                 "model inverts it.", ref=_f2)
checklist(s, [
    ("The critique.", "Buy “cheap and troubled” and you select for lemons — "
     "worse than they look. The market discounted them for a reason."),
    ("Our precondition removes it.", "We invest only where the core is a market "
     "leader earning high returns on capital — attractive even if we change "
     "nothing. No structurally-weak turnarounds."),
    ("The real setup.", "Profitable leaders whose management pivoted into weak "
     "adjacencies — chasing growth as the core matured. A capital-allocation "
     "overlay on a good asset, not a bad asset."),
    ("The inversion.", "Blended results make it look mediocre, so it is priced "
     "mediocre. The appearance understates reality."),
], top, left=Inches(0.7), width=Inches(6.4), size=13.5, gap=13)

# --- sum-of-the-parts bar visual (right) ---
ptf = tbox(s, Inches(7.3), top - Inches(0.05), Inches(5.4), Inches(0.3))
para(ptf, f"{_f2}  ·  SUM-OF-THE-PARTS (ILLUSTRATIVE, INDEXED)", 11, SLATE,
     first=True, bold=True, align=PP_ALIGN.CENTER, after=0)
base = Inches(5.5)
unit = 0.017

def _sotp_bar(xc, val, fill, label):
    h = Emu(int(Inches(unit * val)))
    rect(s, Emu(int(xc) - int(Inches(0.8))), Emu(int(base) - int(h)),
         Inches(1.6), h, fill=fill)
    vt = tbox(s, Emu(int(xc) - int(Inches(0.8))),
              Emu(int(base) - int(h) - int(Inches(0.36))), Inches(1.6), Inches(0.32))
    para(vt, str(val), 14, NAVY_TX, first=True, bold=True,
         align=PP_ALIGN.CENTER, after=0)
    lt = tbox(s, Emu(int(xc) - int(Inches(1.2))),
              Emu(int(base) + int(Inches(0.1))), Inches(2.4), Inches(0.6))
    for k, ln in enumerate(label.split("\n")):
        para(lt, ln, 11.5 if k == 0 else 10, BODY, bold=(k == 0),
             first=(k == 0), align=PP_ALIGN.CENTER, after=0, lead=1.0)

_sotp_bar(Inches(8.8), 100, SLATE, "Market price\n(blended optics)")
_sotp_bar(Inches(11.4), 140, NAVY, "Core alone\n(high ROIIC)")
top100 = Emu(int(base) - int(Inches(unit * 100)))
top140 = Emu(int(base) - int(Inches(unit * 140)))
rect(s, Inches(8.0), top100, Inches(3.4), Pt(1.2), fill=SLATE_LT)  # market level
rect(s, Inches(11.95), top140, Pt(2), Emu(int(top100) - int(top140)), fill=HILITE)
gtf = tbox(s, Inches(8.15), Emu(int(top140) + int(Inches(0.04))),
           Inches(3.6), Inches(0.5))
para(gtf, "+40% hidden value", 12, HILITE, first=True, bold=True,
     align=PP_ALIGN.LEFT, after=0)
para(tbox(s, Inches(0.7), Inches(6.45), Inches(12.0), Inches(0.5)),
     "We remove the drag and refocus capital on the core — positive selection on "
     "a concealed asset, not negative selection on a lemon.", 13, SLATE,
     first=True, italic=True, after=0)

# ---- II.3c How Athanase differs: quality-core constructivism vs Elliott ----
s, top = content("Investment Strategy",
                 "Quality-core constructivism — how we differ",
                 "Even Hohn (TCI) and Loeb (Third Point) pivoted to quality — we "
                 "capture quality at a discount, with a catalyst in our control.",
                 ref=tbl())
rows = [
    ("", "Large generalist activists (e.g. Elliott)",
     "Athanase — quality-core constructivism"),
    ("Entry precondition",
     "Wide range of situations — breakups, credit, event-driven; business "
     "quality often secondary",
     "Only a market-leading, high-ROIIC core — “attractive even if we change "
     "nothing”"),
    ("Source of return",
     "Financial / transactional pressure: force a sale, return cash, exit",
     "Operational refocus of a genuinely good business, from a board seat, over "
     "years"),
    ("What gets fixed",
     "The price, the structure, the outcome of a contest",
     "The capital-allocation overlay — the core already works"),
    ("Margin of safety",
     "Winning the campaign",
     "The core’s proven high-ROIIC economics — the hard floor"),
    ("Cadence & style",
     "Many simultaneous campaigns, leverage, often confrontational",
     "~one deal a year, constructive, board-led, mid-cap surgical"),
]
tl = Inches(0.55); cw3 = [Inches(2.35), Inches(4.85), Inches(4.95)]
rh0 = Inches(0.42); rhd = Inches(0.74); y = top
for ri, row in enumerate(rows):
    head = ri == 0
    rh = rh0 if head else rhd
    x = tl
    for ci, cell in enumerate(row):
        if head:
            fill = SLATE if ci else WHITE
        else:
            fill = HEADERBG if ri % 2 else WHITE
        rect(s, x, y, cw3[ci], rh, fill=fill)
        ctf = tbox(s, Emu(int(x) + int(Inches(0.12))), y,
                   Emu(int(cw3[ci]) - int(Inches(0.2))), rh,
                   anchor=MSO_ANCHOR.MIDDLE)
        if head:
            col = WHITE
        elif ci == 0:
            col = SLATE
        elif ci == 2:
            col = NAVY_TX
        else:
            col = BODY
        para(ctf, cell, 12 if head else 11.5, col,
             bold=(head or ci == 0 or ci == 2), first=True, after=0, lead=1.04)
        x = Emu(int(x) + int(cw3[ci]))
    y = Emu(int(y) + int(rh))
para(tbox(s, Inches(0.55), Emu(int(y) + int(Inches(0.14))), Inches(12.2), Inches(0.5)),
     "The megacap activists left mid-cap surgical engagements behind as they "
     "scaled. That hidden-core niche is our lane.", 13, SLATE, first=True,
     italic=True, after=0)

# ---- II.4 Margin for error vs PE (table) ----
s, top = content("Investment Strategy",
                 "We don’t need to be right more than ~50%",
                 "Because we buy at a discount, not at a 40% control premium.",
                 ref=tbl())
# table
rows = [
    ("Metric", "Private Equity (take-private)", "Athanase (engaged owner)"),
    ("Entry basis", "$140  (market + 40% premium)", "$100  (exhaustion price)"),
    ("Holding period", "5 years", "3 years (rerating cycle)"),
    ("Value growth required", "+148% total growth", "+73% total growth"),
    ("Operational success needed", "~85–100% of thesis", "~45–55% of thesis"),
]
tbl_l = Inches(0.75); tbl_t = top; col_w = [Inches(4.0), Inches(4.0), Inches(3.85)]
row_h = Inches(0.62)
y = tbl_t
for ri, row in enumerate(rows):
    x = tbl_l
    head = (ri == 0)
    for ci, cell in enumerate(row):
        fill = SLATE if head else (HEADERBG if ri % 2 == 0 else WHITE)
        cellrect = rect(s, x, y, col_w[ci], row_h, fill=fill)
        ctf = tbox(s, Emu(int(x) + int(Inches(0.15))), y,
                   Emu(int(col_w[ci]) - int(Inches(0.3))), row_h,
                   anchor=MSO_ANCHOR.MIDDLE)
        emph = head or ci == 2
        para(ctf, cell, 13 if not head else 13.5,
             WHITE if head else (NAVY_TX if ci == 2 else BODY),
             bold=emph, first=True, after=0)
        x = Emu(int(x) + int(col_w[ci]))
    y = Emu(int(y) + int(row_h))
para(tbox(s, Inches(0.75), Emu(int(y) + int(Inches(0.2))), Inches(11.8), Inches(0.5)),
     "A margin-for-error strategy: we win by being directionally right where PE "
     "must be precisely perfect.", 14, SLATE, first=True, italic=True, after=0)

# ---- II.4b Structural edge vs mid-market private equity --------------------
s, top = content("Investment Strategy",
                 "The same playbook as mid-market PE — without its costs",
                 "Operational value creation in mid-cap businesses, but in a "
                 "public wrapper that removes private equity’s structural drags.",
                 ref=tbl())
rows = [
    ("", "Mid-market private equity", "Athanase — engaged owner"),
    ("Entry price",
     "Market + ~40% control premium — pays up front for the right to fix it",
     "A discount to the core’s value at “shareholder exhaustion” — gets that "
     "right for free"),
    ("Liquidity",
     "10-year lock-up, blind pool; capital trapped until exit",
     "Listed, daily-priced positions — LPs keep liquidity and can exit"),
    ("Fees & capital efficiency",
     "2&20 on committed capital; fee drag while undeployed dry powder waits",
     "Lower fees; capital is deployed, not sitting in a capital-call queue"),
    ("The J-curve",
     "Opaque, model-marked write-downs before any mark-up",
     "Live and visible — a tradable entry point, not a hidden cost"),
    ("Transparency & control",
     "Annual mark-to-model; thesis underwritten once at close",
     "Daily public scoreboard; board control secures cash flow, capex and pay"),
    ("Scalability for PE",
     "Rational to skip: 10× more capital per buyout, richer fees",
     "Exactly why the niche stays open — little competition for these assets"),
]
tl = Inches(0.55); cw3 = [Inches(2.35), Inches(4.85), Inches(4.95)]
rh0 = Inches(0.36); rhd = Inches(0.64); y = top
for ri, row in enumerate(rows):
    head = ri == 0
    rh = rh0 if head else rhd
    x = tl
    for ci, cell in enumerate(row):
        if head:
            fill = SLATE if ci else WHITE
        else:
            fill = HEADERBG if ri % 2 else WHITE
        rect(s, x, y, cw3[ci], rh, fill=fill)
        ctf = tbox(s, Emu(int(x) + int(Inches(0.12))), y,
                   Emu(int(cw3[ci]) - int(Inches(0.2))), rh,
                   anchor=MSO_ANCHOR.MIDDLE)
        if head:
            col = WHITE
        elif ci == 0:
            col = SLATE
        elif ci == 2:
            col = NAVY_TX
        else:
            col = BODY
        para(ctf, cell, 12 if head else 11, col,
             bold=(head or ci == 0 or ci == 2), first=True, after=0, lead=1.02)
        x = Emu(int(x) + int(cw3[ci]))
    y = Emu(int(y) + int(rh))
para(tbox(s, Inches(0.55), Emu(int(y) + int(Inches(0.18))), Inches(12.2), Inches(0.5)),
     "Private-equity-style operational turnarounds, imported into public markets "
     "— keeping the liquidity, transparency and lower fees PE gives up.",
     13, SLATE, first=True, italic=True, after=0)

# ---- II.4f Specialist vs household name ------------------------------------
s, top = content("Investment Strategy",
                 "Why the specialist, not the household name",
                 "Recommending Athanase over a brand name is a choice for "
                 "structural alpha over institutional inertia.", ref=tbl())
rows = [
    ("", "The household name", "Athanase — the specialist"),
    ("Scale & influence",
     "$10bn+ AUM forces them into large- and mega-caps, where one owner rarely "
     "shifts the agenda and growth is incremental.",
     "Mid-market (€50m–€3bn), where a single board seat changes the outcome — "
     "the high-return, low-competition lane the giants left behind."),
    ("Return source",
     "Large buyouts are leveraged beta — ~0.89 correlated to small/mid equities "
     "once de-smoothed; driven by GDP, rates and multiples.",
     "Idiosyncratic, event-driven alpha at 0.44 correlation — a genuine "
     "diversifier, not a duplicate of your equity book."),
    ("Entry & structure",
     "A ~40% control premium up front, a 10-year blind-pool lock-up, and fees "
     "on undeployed dry powder.",
     "Entry at “shareholder exhaustion” — the turnaround effectively for free — "
     "with daily liquidity and zero capital-call liability."),
    ("What you’re buying",
     "A brand — but brand recognition does not compound capital.",
     "Asymmetric capture (93% up / 43% down) and a validated 16.3% net, "
     "+7% worst-entry record — institutional process, not a logo."),
]
tl = Inches(0.55); cwx = [Inches(2.0), Inches(5.05), Inches(5.05)]
y = top
for ri, row in enumerate(rows):
    head = ri == 0
    rh = Inches(0.42) if head else Inches(0.9)
    x = tl
    for ci, cell in enumerate(row):
        ath = (ci == 2)
        if head:
            fill = NAVY if ath else (SLATE if ci == 1 else WHITE)
        elif ath:
            fill = BLUE5 if ri % 2 == 0 else HEADERBG
        elif ci == 1:
            fill = HEADERBG if ri % 2 == 0 else WHITE
        else:
            fill = WHITE
        rect(s, x, y, cwx[ci], rh, fill=fill)
        ctf = tbox(s, Emu(int(x) + int(Inches(0.14))), y,
                   Emu(int(cwx[ci]) - int(Inches(0.26))), rh,
                   anchor=MSO_ANCHOR.MIDDLE)
        if head:
            col = WHITE
        elif ci == 0:
            col = SLATE
        elif ath:
            col = NAVY_TX
        else:
            col = SUBTLE
        para(ctf, cell, 11.5 if head else (11 if ci == 0 else 10.5), col,
             bold=(head or ci == 0 or ath), first=True, after=0, lead=1.05)
        x = Emu(int(x) + int(cwx[ci]))
    y = Emu(int(y) + int(rh))
rect(s, Inches(0.55), Emu(int(y) + int(Inches(0.12))), Inches(12.1), Inches(0.5),
     fill=NAVY)
para(tbox(s, Inches(0.75), Emu(int(y) + int(Inches(0.12))), Inches(11.7),
          Inches(0.5), anchor=MSO_ANCHOR.MIDDLE),
     "Brand recognition does not compound capital — asymmetric capture does. "
     "The edge is the very lane the household names had to abandon as they "
     "scaled.", 12.5, WHITE, first=True, italic=True, after=0, track=0)
para(tbox(s, Inches(0.55), Inches(7.18), Inches(12.6), Inches(0.32)),
     "De-smoothed large-buyout correlation to small/mid-cap equities ~0.89 "
     "(Two Sigma / Venn, 2024; Couts, Gonçalves & Rossi, 2020; Boyer et al., "
     "2018). Athanase correlation 0.44 (real). See references.", 7.5, FOOT,
     first=True, after=0, track=0, lead=1.05)

# ===========================================================================
# SUB-SECTION: HOW WE INVEST (find -> filter -> execute)
# ===========================================================================
subdivider("How we invest")

# ---- How we invest 0: HOW DECISIONS ARE MADE (gates, owners, governance) ----
s, top = content("How we invest · Decisions",
                 "How decisions are made: three tollgates, clear owners",
                 "Every idea passes through the same three gates — each with an "
                 "explicit bar and a named decision-maker. The process is "
                 "documented, governed and repeatable.")
gates = [
    ("Tollgate 1 · Valuation", "After sourcing & screening",
     "Is it cheap enough even if we change nothing?",
     "The core business alone justifies the price — a 30–40% margin of safety "
     "versus a private-equity bidder.",
     "Investment team"),
    ("Tollgate 2 · The plan", "After research & due diligence",
     "Is the improvement plan real — and can we deliver it?",
     "Upside quantified across the six value levers and validated with "
     "management and industry experts, with a credible path to a board seat.",
     "Deal lead + Head of Research"),
    ("Tollgate 3 · Go / No-Go", "Investment Committee",
     "Right risk, right size, right fit?",
     "Within risk and portfolio limits, conviction-sized, with no better use "
     "of the capital available.",
     "Investment Committee — CIO + PMs; risk independent"),
]
cw = Inches(3.86); gx = Inches(0.18); x = Inches(0.62); hh = Inches(0.5)
bodyH = Inches(3.18)
for name, when, q, bar, who in gates:
    rect(s, x, top, cw, hh, fill=NAVY)
    htf = tbox(s, x, top, cw, hh, anchor=MSO_ANCHOR.MIDDLE)
    para(htf, name, 13.5, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER,
         after=0, font=SERIF, track=0)
    by = Emu(int(top) + int(hh))
    rect(s, x, by, cw, bodyH, fill=HEADERBG)
    tf = tbox(s, Emu(int(x) + int(Inches(0.18))),
              Emu(int(by) + int(Inches(0.14))),
              Emu(int(cw) - int(Inches(0.36))), Emu(int(bodyH) - int(Inches(0.26))))
    para(tf, when.upper(), 8.5, SLATE_LT, first=True, bold=True, after=6, track=0)
    para(tf, "THE QUESTION", 8.5, SLATE, bold=True, after=2, track=0)
    para(tf, "“" + q + "”", 11, NAVY_TX, italic=True, after=7, lead=1.1,
         font=SERIF, track=0)
    para(tf, "PASS ONLY IF", 8.5, SLATE, bold=True, after=2, track=0)
    para(tf, bar, 9.5, BODY, after=7, lead=1.12, track=0)
    para(tf, "DECIDED BY", 8.5, SLATE, bold=True, after=2, track=0)
    para(tf, who, 9.5, NAVY_TX, bold=True, after=0, lead=1.1, track=0)
    x = Emu(int(x) + int(cw) + int(gx))
rect(s, Inches(0.62), Inches(6.2), Inches(11.94), Inches(0.62), fill=NAVY)
para(tbox(s, Inches(0.8), Inches(6.2), Inches(11.6), Inches(0.62),
          anchor=MSO_ANCHOR.MIDDLE),
     "Repeatable by design — the same gates the team has run since 2006, a "
     "written investment case, minuted decisions, and every holding "
     "re-underwritten each quarter. About 1 in 1,000 screened names is bought.",
     12, WHITE, first=True, italic=True, after=0, track=0)
para(tbox(s, Inches(0.62), Inches(7.0), Inches(12.0), Inches(0.4)),
     "A “no” at any gate ends the process — most ideas are rejected at "
     "Tollgate 1 or 2, long before capital is at risk.", 9, FOOT, first=True,
     after=0, track=0, lead=1.1)

# ---- How we invest 1: SOURCE & SCREEN --------------------------------------
_f_src = fig()
s, top = content("How we invest · Source",
                 "Sourcing & screening: from ~7,200 names to a shortlist",
                 "Five sourcing channels feed a funnel that narrows the whole "
                 "European listed market to a handful of candidates — before "
                 "a single krona is committed.", ref=_f_src)
para(tbox(s, Inches(0.7), top, Inches(8), Inches(0.3)),
     "FIVE SOURCING CHANNELS", 11, SLATE, first=True, bold=True, after=0, track=0)
chips = ["Network", "Own research", "Interviews", "Company visits",
         "External experts"]
ccw = Inches(2.36); cgx = Inches(0.105); cx = Inches(0.7); cy = top + Inches(0.36)
for ch in chips:
    rect(s, cx, cy, ccw, Inches(0.46), fill=NAVY)
    para(tbox(s, cx, cy, ccw, Inches(0.46), anchor=MSO_ANCHOR.MIDDLE), ch, 11,
         WHITE, first=True, bold=True, align=PP_ALIGN.CENTER, after=0, track=0)
    cx = Emu(int(cx) + int(ccw) + int(cgx))
fnl = [("≈ 7,200", "listed companies in selected European markets", 7.2),
       ("≈ 2,800", "with a market cap between €50m and €3bn", 6.0),
       ("≈ 2,500", "in our selected sectors", 4.8),
       ("Shortlist", "suitable shareholder base for constructive ownership", 3.6)]
fy = cy + Inches(0.72)
for big, lab, win in fnl:
    fw = Inches(win)
    fx = Emu(int(Inches(0.7)) + int((int(Inches(12.13)) - int(fw)) / 2))
    rect(s, fx, fy, fw, Inches(0.52), fill=HEADERBG)
    rect(s, fx, fy, Inches(0.07), Inches(0.52), fill=NAVY)
    bt = tbox(s, fx, fy, fw, Inches(0.52), anchor=MSO_ANCHOR.MIDDLE)
    p = bt.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.line_spacing = 1.0
    r = p.add_run(); r.text = big + "   "
    r.font.size = Pt(13); r.font.bold = True
    r.font.color.rgb = NAVY_TX; r.font.name = SERIF
    r2 = p.add_run(); r2.text = lab
    r2.font.size = Pt(11); r2.font.color.rgb = BODY; r2.font.name = SANS
    fy = Emu(int(fy) + int(Inches(0.62)))
rect(s, Inches(0.6), Emu(int(fy) + int(Inches(0.06))), Inches(12.13),
     Inches(0.58), fill=NAVY)
para(tbox(s, Inches(0.8), Emu(int(fy) + int(Inches(0.06))), Inches(11.7),
          Inches(0.58), anchor=MSO_ANCHOR.MIDDLE),
     "First tollgate — is the valuation attractive even if we change "
     "nothing? We only proceed when the core alone justifies the price.",
     12.5, WHITE, first=True, italic=True, after=0, track=0)


def _step_row(s, y, h, num, name, body):
    rect(s, Inches(0.6), y, Inches(2.7), h, fill=NAVY)
    lt = tbox(s, Inches(0.76), y, Inches(2.4), h, anchor=MSO_ANCHOR.MIDDLE)
    para(lt, f"STEP {num}", 9, DIVIDER, first=True, bold=True, after=2, track=0)
    para(lt, name, 12.5, WHITE, bold=True, after=0, font=SERIF, track=0,
         lead=1.04)
    bt = tbox(s, Inches(3.5), y, Inches(9.23), h, anchor=MSO_ANCHOR.MIDDLE)
    para(bt, body, 10.5, BODY, first=True, after=0, lead=1.16, track=0)


# ---- How we invest 2: RESEARCH (steps 1-4) ---------------------------------
s, top = content("How we invest · Research",
                 "Research & due diligence: four analytical steps")
steps2 = [
    ("1", "Idea generation",
     "Internal research, our industry network and external advisors — plus "
     "meetings with company management, board members or former employees, and "
     "dialogue with other investors."),
    ("2", "Financial analysis",
     "In-depth analysis of public financials; comparison with industry "
     "standards and similar businesses; interviews with senior and subsidiary "
     "management and with competitors; and an estimate of the optimal capital "
     "structure and M&A headroom."),
    ("3", "Operational analysis",
     "Interviews with management and selected board members; discussions across "
     "our network of industry leaders; visits to competitors, customers and "
     "suppliers; and interviews with former employees and board members."),
    ("4", "Validation of the operating plan",
     "Value the potential operational and structural improvements; pressure-"
     "test the plan with management and industry experts and secure their "
     "commitment; and select potential board members (and, if needed, a new "
     "management team)."),
]
rh = Inches(1.12); gap = Inches(0.1); y = top + Inches(0.02)
for num, name, body in steps2:
    _step_row(s, y, rh, num, name, body)
    y = Emu(int(y) + int(rh) + int(gap))
para(tbox(s, Inches(0.6), Inches(7.14), Inches(12.6), Inches(0.4)),
     "Second tollgate — value the improvement programme across six levers: "
     "organic growth, M&A, cost, structural, operational and financial.",
     9, FOOT, first=True, after=0, track=0, lead=1.1)


# ---- How we invest 3: DECISION & EXECUTION (steps 5-8) ---------------------
s, top = content("How we invest · Execute",
                 "Decision, execution and board engagement")
steps3 = [
    ("5", "Investment team meeting",
     "Discuss the analysis, the recommended size of the investment, and the "
     "risk and portfolio impact."),
    ("6", "Investment committee",
     "Review the recommendation, the stake size and the risk and portfolio "
     "impact, and challenge the team — the third tollgate and the final "
     "Go / No-Go decision."),
    ("7", "Trade execution",
     "Work with our prime brokers to build the stake in the most cost- and "
     "time-efficient way, tracking any new information that could change the "
     "case."),
    ("8", "Inform company, board & stakeholders",
     "Brief the Chairman and CEO, present the case to co-investors, secure at "
     "least one board seat for the Athanase team, and align subsidiary managers "
     "behind the plan."),
]
rh = Inches(1.12); gap = Inches(0.1); y = top + Inches(0.02)
for num, name, body in steps3:
    _step_row(s, y, rh, num, name, body)
    y = Emu(int(y) + int(rh) + int(gap))
para(tbox(s, Inches(0.6), Inches(7.14), Inches(12.6), Inches(0.4)),
     "Roughly one investment a year clears all three tollgates. The board seat "
     "is the engine of both value creation and risk control.",
     9, FOOT, first=True, after=0, track=0, lead=1.1)


subdivider("Risk, honestly")
# ---- II.5 Risk system ----
s, top = content("Risk Management", "A risk system built to avoid permanent loss",
                 "Predictability, price and influence combined to lower risk "
                 "systematically.")
# two columns
colL = tbox(s, Inches(0.75), top, Inches(5.85), Inches(4.5))
para(colL, "THE VALUATION FLOOR", 13, SLATE, first=True, bold=True, after=6)
for t in ["Enter only at “shareholder surrender” — dried-up "
          "liquidity and long-term price lows",
          "Value only the rectifiable core; assign zero to “growth trap” "
          "divisions",
          "Creates a structural 30–40% margin of safety vs a PE bidder"]:
    para(colL, t, 14, BODY, after=9, lead=1.12)
colR = tbox(s, Inches(7.0), top, Inches(5.8), Inches(4.5))
para(colR, "THE BOARD-LEVEL KILL SWITCH", 13, SLATE, first=True, bold=True, after=6)
for t in ["Secure board seats and chairmanships as a prerequisite — not "
          "proxy threats",
          "Mandated rectification: freeze prestige projects the moment the macro "
          "shifts",
          "Automatic risk metrics (10/20/30% triggers) remove portfolio-manager "
          "bias"]:
    para(colR, t, 14, BODY, after=9, lead=1.12)
# core & satellite line
para(tbox(s, Inches(0.75), Inches(6.35), Inches(11.8), Inches(0.6)),
     "Core & satellite: a 30-stock equal-weighted core index compounds capital; "
     "8–12 concentrated ideas must clear a 12% IRR hurdle and a ≤20% "
     "downside gate.", 13, SUBTLE, first=True, italic=True, after=0)

# ---- II.5b Strategic reframing: the discomfort is the moat ----
s, top = content("Strategic Case", "The visible risks are the engine — not the cost",
                 "From “corporate raider” to “industrial constructivist”: "
                 "reputational friction and the public J-curve source the alpha "
                 "and protect it from arbitrage.")
checklist(s, [
    ("Friction is the catalyst.", "Quiet engagement does not reprice a stock; "
     "public engagement forces the market to re-rate — and a 20-year reputation "
     "pulls sophisticated capital in behind us (the coattail effect)."),
    ("The discomfort is the moat.", "Few managers have the industrial background "
     "and stamina to wage a multi-year campaign — so headline-shy capital leaves "
     "a persistent pricing inefficiency for those who will."),
    ("A board seat de-risks the asset.", "We trade optical, mark-to-market "
     "volatility for direct control of cash flow, capex and pay. Fundamental risk "
     "falls the moment control is secured — whatever the share price does."),
    ("The J-curve is tradable, not painful.", "Unlike PE’s accounting J-curve, "
     "ours is live and observable — an engineered entry point (and a co-invest "
     "“free look”) at de-risked governance at distressed prices."),
    ("Public scoreboard, daily.", "There is nowhere to hide poor capital "
     "allocation — the discipline behind a 92% deal-level hit rate over 20 years."),
], top, size=14.5, gap=11)

# ---- II.5c The non-linear path, honestly priced ----
_f3 = fig()
s, top = content("Strategic Case", "A non-linear path — honestly priced",
                 "Public engaged ownership marks to market every day. The "
                 "destination is the multiple; the route is rarely a straight "
                 "line.", ref=_f3)
# explanatory bullets (left)
checklist(s, [
    ("Value moves over the holding period.", "A 5× outcome is realised over "
     "years — the interim mark moves with the market, the cycle and the "
     "campaign. It is a journey, not a coupon."),
    ("PE looks smooth because it is appraised, not traded.", "Infrequent, "
     "model-based marks understate true economic volatility — “volatility "
     "laundering.” Marked daily, a levered buyout would swing as much or more."),
    ("Same economics, honest optics.", "We hold the same kind of assets PE "
     "does — we simply report them at live prices and accept the J-curve as the "
     "toll paid for liquidity."),
], top, left=Inches(0.7), width=Inches(6.2), size=13.5, gap=11)
# illustrative line chart (right)
cd = CategoryChartData()
cd.categories = ["0", "", "", "Yr 1", "", "", "Yr 2", "", "", "Yr 3", "", "", "Exit"]
cd.add_series("Public mark-to-market", (100, 92, 80, 72, 85, 95, 88, 112, 142,
                                        175, 212, 255, 300))
cd.add_series("PE appraisal marks", (100, 104, 110, 118, 128, 140, 155, 172,
                                     192, 215, 245, 272, 300))
gf = s.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(7.05), top - Inches(0.05),
                        Inches(5.9), Inches(3.9), cd)
ch = gf.chart
ch.has_title = True
ch.chart_title.text_frame.text = f"{_f3}.  Same destination, different reported path (indexed to 100)"
ch.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
ch.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = SUBTLE
ch.has_legend = True
ch.legend.position = XL_LEGEND_POSITION.BOTTOM
ch.legend.include_in_layout = False
ch.legend.font.size = Pt(11)
ser = ch.plots[0].series
ser[0].format.line.color.rgb = NAVY
ser[0].format.line.width = Pt(2.5)
ser[1].format.line.color.rgb = SLATE_LT
ser[1].format.line.width = Pt(2.0)
ser[1].format.line.dash_style = 2  # dashed
for sline in ser:
    sline.smooth = True
dlabels(ser[0], fmt='0', size=8, color=NAVY_TX, pos=XL_LABEL_POSITION.BELOW)
ch.category_axis.tick_labels.font.size = Pt(9)
ch.value_axis.tick_labels.font.size = Pt(9)
ch.value_axis.has_major_gridlines = False
para(tbox(s, Inches(7.05), top + Inches(3.95), Inches(5.9), Inches(0.6)),
     "Illustrative. Both reach the same value; only the public line shows the "
     "honest interim drawdown.", 10, FOOT, first=True, italic=True, after=0)

# ---- II.5d Risk reframing summary table (memo Section 8.3) ----
s, top = content("Strategic Case", "Risk reframing for the investment committee",
                 "What the committee sees on the surface — versus what it is "
                 "actually underwriting.", ref=tbl())
rows = [
    ("What the IC sees", "What it is actually underwriting"),
    ("Reputational / headline risk",
     "The catalyst that forces repricing — and a moat that keeps competitors out"),
    ("Public J-curve & mark-to-market drawdown",
     "A live, tradable entry into an asset already de-risked by board control"),
    ("~50% peak-to-trough swings",
     "Optical volatility; two such drawdowns fully recovered (2010, 2020) · Sortino 2.40"),
    ("Concentration in 8–12 names",
     "Deliberate conviction sizing — the skill (85.7% hit on ≥10% names; +0.73)"),
    ("“Catching a falling knife”",
     "Hard-floor valuation gates; entry below the value of the rectifiable core"),
    ("Single-name dependence",
     "One ≥5× MOIC win every ~2 years; competitive even excluding the largest"),
]
tl = Inches(0.6); cw2 = [Inches(4.55), Inches(7.55)]; rh = Inches(0.56); y = top
for ri, (a, b) in enumerate(rows):
    head = ri == 0
    x = tl
    for ci, cell in enumerate((a, b)):
        fill = SLATE if head else (HEADERBG if ri % 2 else WHITE)
        rect(s, x, y, cw2[ci], rh, fill=fill)
        ctf = tbox(s, Emu(int(x) + int(Inches(0.14))), y,
                   Emu(int(cw2[ci]) - int(Inches(0.24))), rh,
                   anchor=MSO_ANCHOR.MIDDLE)
        para(ctf, cell, 12.5 if head else 12,
             WHITE if head else (SUBTLE if ci == 0 else NAVY_TX),
             bold=(head or ci == 1), first=True, after=0, lead=1.05)
        x = Emu(int(x) + int(cw2[ci]))
    y = Emu(int(y) + int(rh))
para(tbox(s, Inches(0.6), Emu(int(y) + int(Inches(0.15))), Inches(12.1), Inches(0.6)),
     "The IC is not paying for safety. It is paying for proven willingness to "
     "take compensated discomfort — +6 percentage points per year, net of every "
     "fee, for 20 years.", 13, SLATE, first=True, italic=True, after=0)

# ---- II.5e Why it is not career-ending (three anxieties disarmed) -----------
s, top = content("Strategic Case",
                 "Why an allocation here is not career-ending",
                 "The three institutional anxieties — mark-to-market "
                 "volatility, headline risk and concentration — each disarmed "
                 "by the mechanics of the strategy.")
anx = [
    ("Mark-to-market volatility", "Isolated by mandate design", [
        ("Sized for immunity", "a 1–3% special-situations sleeve, so no single "
         "mark can dominate the fund or trigger board-level panic."),
        ("Pre-underwritten", "a ~50% drawdown tolerance agreed with the IC up "
         "front; reallocation triggers on manager turnover or style drift — "
         "never the share price."),
        ("No capital-call trap", "unlike PE’s smoothed marks and the "
         "denominator effect, there is zero capital-call liability — you are "
         "never forced to sell liquid assets at the bottom.")]),
    ("Headline risk", "Reframed as fiduciary stewardship", [
        ("Constructive, not confrontational", "quality-core constructivism — "
         "profitable market leaders that welcome capital discipline, not proxy "
         "fights."),
        ("Applied governance", "board independence, pay-for-performance and "
         "capital discipline are the “G” in ESG — turning hidden governance "
         "risk into a return engine."),
        ("Defensible", "aligned with UN PRI fiduciary duty and easy to explain "
         "to stakeholders.")]),
    ("Concentration", "Floored by hardened guardrails", [
        ("A valuation floor", "every idea clears a tollgate where the core "
         "alone justifies the price: a structural 30–40% margin of safety."),
        ("Mechanized triggers", "automatic −10% / −20% / −30% force a thesis "
         "review and exit — removing PM emotion before a position spirals."),
        ("A core-and-satellite buffer", "8–12 conviction ideas sit on a "
         "30-stock equal-weighted core that absorbs idiosyncratic shocks.")]),
]
cw = Inches(3.86); gx = Inches(0.18); x = Inches(0.62); hh = Inches(0.5)
bodyH = Inches(3.4)
for title_a, reframe, bullets in anx:
    rect(s, x, top, cw, hh, fill=NAVY)
    htf = tbox(s, x, top, cw, hh, anchor=MSO_ANCHOR.MIDDLE)
    para(htf, title_a, 13.5, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER,
         after=0, font=SERIF, track=0)
    by = Emu(int(top) + int(hh))
    rect(s, x, by, cw, bodyH, fill=HEADERBG)
    tf = tbox(s, Emu(int(x) + int(Inches(0.18))),
              Emu(int(by) + int(Inches(0.13))),
              Emu(int(cw) - int(Inches(0.36))), Emu(int(bodyH) - int(Inches(0.24))))
    para(tf, reframe.upper(), 8.5, SLATE, first=True, bold=True, after=7, track=0)
    for i, (lead, body) in enumerate(bullets):
        p = tf.add_paragraph()
        p.space_after = Pt(6); p.line_spacing = 1.1
        hang = int(round(9.5 * 1.25 * 12700))
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(hang)); pPr.set("indent", str(-hang))
        r0 = p.add_run(); r0.text = "›  "
        r0.font.size = Pt(9.5); r0.font.bold = True
        r0.font.color.rgb = SLATE_LT; r0.font.name = SANS
        r1 = p.add_run(); r1.text = lead
        r1.font.size = Pt(9.5); r1.font.bold = True
        r1.font.color.rgb = NAVY_TX; r1.font.name = SANS
        r2 = p.add_run(); r2.text = " — " + body
        r2.font.size = Pt(9.5); r2.font.color.rgb = BODY; r2.font.name = SANS
    x = Emu(int(x) + int(cw) + int(gx))
rect(s, Inches(0.62), Inches(6.28), Inches(11.94), Inches(0.62), fill=NAVY)
para(tbox(s, Inches(0.82), Inches(6.28), Inches(11.5), Inches(0.62),
          anchor=MSO_ANCHOR.MIDDLE),
     "Career risk is engineered out — sized so no mark dominates, governed so "
     "headlines read as stewardship, floored so concentration cannot blow up: "
     "a bad quarter is underwritten, not a surprise.",
     12.5, WHITE, first=True, italic=True, after=0, track=0)

subdivider("The 20-year track record")
# ---- II.6 Track record (chart) ----
_f4 = fig()
s, top = content("Track Record", "Proof: outperformance with downside protection",
                 ref=_f4)
chart_data = CategoryChartData()
chart_data.categories = ["All markets", "Up markets", "Down markets"]
chart_data.add_series("Athanase", (1.5, 1.2, 0.2))
chart_data.add_series("MSCI IMI", (0.6, 1.3, -1.7))
cx, cy, cw, ch = Inches(0.7), top, Inches(7.1), Inches(4.2)
gframe = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, cx, cy, cw, ch, chart_data)
chart = gframe.chart
chart.has_title = True
chart.chart_title.text_frame.text = f"{_f4}.  Average monthly return (%)"
chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
chart.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = SUBTLE
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(11)
plot = chart.plots[0]
plot.series[0].format.fill.solid()
plot.series[0].format.fill.fore_color.rgb = NAVY
plot.series[1].format.fill.solid()
plot.series[1].format.fill.fore_color.rgb = SLATE_LT
plot.gap_width = 80
dlabels(plot, fmt='0.0"%"', size=9, color=BODY, pos=XL_LABEL_POSITION.OUTSIDE_END)
for ax in (chart.value_axis, chart.category_axis):
    ax.tick_labels.font.size = Pt(10)
    ax.tick_labels.font.color.rgb = BODY
# right-hand proof points
pr = tbox(s, Inches(8.1), top + Inches(0.1), Inches(4.6), Inches(4.5))
para(pr, "What downside protection means", 16, NAVY_TX, first=True, bold=True,
     after=2, font=SERIF)
para(pr, "In the market’s down months Athanase fell ~1.5% on average versus the "
     "market’s ~3.8% — losing less than half as much when it matters.", 13,
     BODY, after=14, lead=1.15)
para(pr, "Positive when the market wasn’t", 16, NAVY_TX, bold=True, after=2,
     font=SERIF)
para(pr, "Across the six worst years for global equities, Athanase outperformed "
     "in five — and stayed positive in two: +9% in 2015 and +0.3% in 2022, while "
     "the market fell −6% and −20%.", 13, BODY, after=14, lead=1.15)
para(pr, "Why it compounds", 16, NAVY_TX, bold=True, after=2, font=SERIF)
para(pr, f"Losing less in drawdowns leaves more capital working in the recovery "
     f"— the engine behind {_ATH['mult']:.0f}× growth of capital vs "
     f"{_MSCI['mult']:.1f}× for the market since 2006.", 13, BODY,
     after=0, lead=1.15)

# ---- II.6b Transactions: AIP Fund II ----
FUND2, HIST = load_transactions()
# Brokk (2003–2006) ended in 2006 — it belongs in the earlier (pre-2006)
# deals table, not the 2006–2014 prior-period bucket.
HIST = [d for d in HIST if d["company"].strip().lower() != "brokk"]
s, top = content("Track Record", "AIP Fund II — transactions (2015–present)",
                 "Each position shown against its benchmark; invested capital "
                 "weighted MOIC of 2.2x.", ref=tbl())
cw = [Inches(3.3), Inches(2.0), Inches(1.7), Inches(1.7), Inches(2.0)]
endy = deal_table(s, Inches(0.75), top, FUND2, cw, font=10, rh=Inches(0.212))
# summary strip
sy = Emu(int(endy) + int(Inches(0.10)))
rect(s, Inches(0.75), sy, Inches(10.7), Inches(0.34), fill=NAVY)
stf = tbox(s, Inches(0.9), sy, Inches(10.4), Inches(0.34), anchor=MSO_ANCHOR.MIDDLE)
para(stf, "17 deals   ·   weighted MOIC 2.2x   ·   ex best/worst IRR 41% vs "
     "index ~9%   ·   35 of 38 deals profitable since 2006", 10.5, WHITE,
     first=True, after=0)
ny = Emu(int(sy) + int(Inches(0.34)) + int(Inches(0.08)))
nt = tbox(s, Inches(0.75), ny, Inches(11.8), Inches(0.35))
para(nt, "IRR/MOIC net of the position; “n.m.” where short holding periods make "
     "annualised figures not meaningful. Capital in SEK.", 8.5, FOOT, first=True,
     after=0)

# ---- II.6c Transactions: prior period 2006-2014 ----
s, top = content("Track Record", "Prior period — transactions (2006–2014)",
                 "The same strategy, the prior fund — repeatable across two "
                 "decades and multiple cycles.", ref=tbl())
half = 11
cwh = [Inches(2.7), Inches(1.45), Inches(1.45)]
deal_table(s, Inches(0.6), top, HIST[:half], cwh, font=10,
           rh=Inches(0.285), cols=("company", "irr", "moic"),
           headers=("Company", "IRR", "MOIC"))
deal_table(s, Inches(6.95), top, HIST[half:], cwh, font=10,
           rh=Inches(0.285), cols=("company", "irr", "moic"),
           headers=("Company", "IRR", "MOIC"))
sy = top + Inches(3.55)
rect(s, Inches(0.6), sy, Inches(11.05), Inches(0.34), fill=NAVY)
stf = tbox(s, Inches(0.75), sy, Inches(10.7), Inches(0.34), anchor=MSO_ANCHOR.MIDDLE)
para(stf, "21 deals  ·  weighted MOIC 2.6x  ·  capital-weighted IRR 38.4% "
     "(+25 pts vs index)  ·  average IRR 67.7% (+53.5 pts)", 9.5, WHITE,
     first=True, after=0)
nt = tbox(s, Inches(0.6), Inches(6.95), Inches(11.8), Inches(0.4))
para(nt, "Selected larger deals of the investment team. Losses shown in italic. "
     "Capital in SEK; ~$400M invested over the period.", 8.5, FOOT, first=True,
     after=0)

# ---- II.6c-2 Earlier deals — part of the team (1996-2004) ------------------
s, top = content("Track Record",
                 "Earlier deals — part of the team (1996–2006)",
                 "The 38-deal record is the full team since 2006. Earlier, "
                 "members of the team led these larger deals from 1996 — the "
                 "same engaged-ownership playbook, through earlier cycles.",
                 ref=tbl())
EARLY = [
    dict(company="Bilia I", period="1999–03", irr=0.185, moic=2.2, outp=0.170),
    dict(company="Skanska", period="1996–98", irr=0.260, moic=1.6, outp=0.014),
    dict(company="Svedala", period="1997–01", irr=0.111, moic=1.5, outp=0.012),
    dict(company="ASG", period="1997–99", irr=0.668, moic=2.2, outp=0.415),
    dict(company="SCA", period="1996–01", irr=0.159, moic=2.0, outp=-0.270),
    dict(company="Perbio", period="1998–03", irr=1.192, moic=4.2, outp=1.181),
    dict(company="Perstorp", period="1996–01", irr=0.104, moic=1.6, outp=-0.060),
    dict(company="Acando", period="1999–04", irr=-0.023, moic=0.9, outp=-0.068),
    dict(company="Pergo", period="2001–03", irr=-0.224, moic=0.9, outp=-0.151),
    dict(company="Brokk", period="2003–06", irr=0.738, moic=5.25, outp=0.460),
]
cwh = [Inches(2.7), Inches(1.45), Inches(1.45)]
endL = deal_table(s, Inches(0.6), top, EARLY[:5], cwh, font=10,
                  rh=Inches(0.285), cols=("company", "irr", "moic"),
                  headers=("Company", "IRR", "MOIC"))
endR = deal_table(s, Inches(6.95), top, EARLY[5:], cwh, font=10,
                  rh=Inches(0.285), cols=("company", "irr", "moic"),
                  headers=("Company", "IRR", "MOIC"))
sy = Emu(int(max(int(endL), int(endR))) + int(Inches(0.14)))
rect(s, Inches(0.6), sy, Inches(11.05), Inches(0.34), fill=NAVY)
stf = tbox(s, Inches(0.75), sy, Inches(10.7), Inches(0.34),
           anchor=MSO_ANCHOR.MIDDLE)
para(stf, "10 deals  ·  weighted MOIC 2.0x  ·  capital-weighted IRR 22% "
     "(+9.9 pts vs index)  ·  average IRR 31.7% (+19.6 pts)", 9.5, WHITE,
     first=True, after=0)
nt = tbox(s, Inches(0.6), Emu(int(sy) + int(Inches(0.46))), Inches(11.8),
          Inches(0.4))
para(nt, "Selected larger deals led by members of the current team from 1996, "
     "before the full team formed in 2006. Losses shown in lighter type. "
     "Capital in SEK; SEK 7.2bn invested → SEK 14.3bn returned.", 8.5, FOOT,
     first=True, after=0, lead=1.1)


# ---- II.6d Independent validation: 20-year NET record ----
s, top = content("Independent Review",
                 "Independently validated: 20 years, net of every fee",
                 "A prospective LP’s investment committee reconciled the "
                 "security-level (StatPro) data against the LP-level NET return "
                 "series.")
stats = [(f"+{_ATH['cum_ret']*100:,.0f}%", "cumulative NET return\nover 19.3 years"),
         (f"+{_ATH['ann_ret']*100:.1f}%", "annualised NET\n(net of all fees)"),
         (f"+{(_ATH['ann_ret']-_MSCI['ann_ret'])*100:.0f} pts",
          "annualised alpha,\nnet, vs the market"),
         (f"{_ATH['mult']/_MSCI['mult']:.1f}×", "the relevant equity\nbenchmark")]
bw = Inches(2.95); gapx = Inches(0.18); x = Inches(0.75); y = top
for big, lab in stats:
    rect(s, x, y, bw, Inches(1.8), fill=HEADERBG)
    ctf = tbox(s, x, y + Inches(0.26), bw, Inches(1.4))
    para(ctf, big, 34, NAVY_TX, first=True, align=PP_ALIGN.CENTER, after=4,
         font=SERIF)
    for line in lab.split("\n"):
        para(ctf, line, 11.5, SUBTLE, align=PP_ALIGN.CENTER, after=0)
    x = Emu(int(x) + int(bw) + int(gapx))
checklist(s, [
    ("Two full cycles, survived.", "Drawdowns of −52.8% (2008, recovered by "
     "2010) and −45.9% (2017–2020, recovered by Sep 2020) — navigated live with "
     "the same philosophy, not back-tested."),
    ("Top-tier downside-adjusted return.", "20-year Sortino 1.54; AIPFII Sortino "
     "2.40, with upside volatility 3.9× the downside."),
    ("Realised cash, not marks.", "The largest contributors are realised LP cash "
     "— removing the mark-to-model objection."),
], top + Inches(2.05), size=14, gap=10)
nt = tbox(s, Inches(0.75), Inches(6.95), Inches(11.8), Inches(0.4))
para(nt, "Source: independent institutional due-diligence review, May 2026, on "
     "StatPro security-level data and LP-level NET monthly returns (2006–2026).",
     8.5, FOOT, first=True, after=0)

# ---- II.6e Repeatability & conviction sizing ----
s, top = content("Independent Review",
                 "Repeatable and conviction-sized — not lucky",
                 "The review tested the “you can’t rely on the big winners” "
                 "objection directly against 38 deals over 20 years.")
colL = tbox(s, Inches(0.75), top, Inches(5.85), Inches(4.6))
para(colL, "REPEATABLE", 13, SLATE, first=True, bold=True, after=7)
for t in ["92% deal-level hit rate (35 of 38 profitable); losses ~8% of "
          "cumulative invested capital",
          "9 of 38 deals returned ≥5× MOIC — one such win every ~2 years",
          "Two decade-definers, not one: Klarna 28× (2007) and Athanase Tech "
          "5.4× (2021)"]:
    para(colL, t, 14, BODY, after=11, lead=1.13)
colR = tbox(s, Inches(7.0), top, Inches(5.85), Inches(4.6))
para(colR, "CONVICTION-SIZED (THE SKILL SIGNAL)", 13, SLATE, first=True,
     bold=True, after=7)
for t in ["The 7 positions ever sized ≥10% of NAV hit at 85.7% and contributed "
          "+209 pts — ~77% of the positive book",
          "+0.73 rank correlation between position size and outcome within "
          "winners — they size their best ideas biggest",
          "Win/loss asymmetry: average winner 2.74× the average loser; "
          "+333 pts of winners vs −61 pts of losers"]:
    para(colR, t, 14, BODY, after=9, lead=1.13)
para(tbox(s, Inches(0.75), Inches(6.35), Inches(11.8), Inches(0.55)),
     "“The team has done this twice in 20 years, not once — the ‘every 10–20 "
     "years’ objection is empirically unsupportable.”", 13, SLATE, first=True,
     italic=True, after=0)

# ---- II.6f Governing the allocation: why the path won't get you fired ----
s, top = content("Allocation", "Why a non-linear path won’t get you fired",
                 "Career risk is managed by mandate design — not by hoping the "
                 "line goes straight.")
colL = tbox(s, Inches(0.75), top, Inches(5.85), Inches(4.6))
para(colL, "SIZE AND STRUCTURE IT RIGHT", 13, SLATE, first=True, bold=True, after=7)
for t in ["Size at 1–3% of total assets in a special-situations sleeve — no "
          "single mark can dominate the portfolio or force action",
          "Pre-agree a ~50% mark-to-market drawdown tolerance with the IC at "
          "commitment; align reporting to the multi-year holding cycle, not "
          "monthly marks",
          "Define reallocation triggers in advance — manager turnover, style "
          "drift, weaker sizing discipline — never the share price"]:
    para(colL, t, 14, BODY, after=10, lead=1.13)
colR = tbox(s, Inches(7.0), top, Inches(5.85), Inches(4.6))
para(colR, "JUDGE THE RIGHT THING", 13, SLATE, first=True, bold=True, after=7)
for t in ["Underwrite whether business risk is falling — board seat secured, "
          "cash-flow control — not whether the news flow looks calm",
          "The big winners are realised LP cash, not marks: the value is "
          "bankable, not modelled",
          "Two ~50% drawdowns have already been recovered with the same "
          "philosophy — the path has been bumpy before and still paid +6 pts/yr "
          "net"]:
    para(colR, t, 14, BODY, after=10, lead=1.13)
sy = Inches(6.25)
rect(s, Inches(0.75), sy, Inches(11.85), Inches(0.62), fill=NAVY)
para(tbox(s, Inches(0.95), sy, Inches(11.45), Inches(0.62), anchor=MSO_ANCHOR.MIDDLE),
     "A pre-underwritten, correctly-sized, fundamentally-monitored allocation "
     "turns a bad quarter into “as expected” — not a surprise that invites "
     "blame.", 13.5, WHITE, first=True, italic=True, after=0)

subdivider("What it means for your portfolio")
# ---- II.6g Portfolio impact: Athanase vs long-only / other activist ----
_fp = fig()
s, top = content("Portfolio Impact",
                 "What it does to an allocator’s portfolio",
                 "Illustrative: redirecting a sleeve into Athanase rather than a "
                 "long-only global fund or another activist.", ref=_fp)
# --- left: 10-year growth of a 100 sleeve, through a down-market year (yr 6) ---
# Same 10-yr endpoints as a smooth path, but routed through a market drawdown
# in year 6 to make the downside protection VISIBLE: long-only plunges, the
# other activist dips, Athanase holds (cf. p.17 average monthly returns).
_cw_chart = Inches(6.7)
chx, chy = Inches(0.6), top + Inches(0.05)
gd = CategoryChartData()
gd.categories = [str(y) for y in range(0, 11)]
gd.add_series("Athanase (18%)",
              (100, 118, 140, 166, 197, 233, 236, 296, 366, 442, 523))
gd.add_series("Other activist (10.5%)",
              (100, 111, 123, 137, 152, 168, 150, 182, 214, 244, 271))
gd.add_series("Long-only global (7.6%)",
              (100, 108, 117, 126, 136, 147, 119, 142, 168, 190, 208))
gf = s.shapes.add_chart(XL_CHART_TYPE.LINE, chx, chy, _cw_chart, Inches(4.2), gd)
ch = gf.chart
ch.has_title = True
ch.chart_title.text_frame.text = f"{_fp}.  Growth of a 100 sleeve through a market cycle (indexed)"
ch.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
ch.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = SUBTLE
ch.has_legend = True
ch.legend.position = XL_LEGEND_POSITION.BOTTOM
ch.legend.include_in_layout = False
ch.legend.font.size = Pt(10)
_ser = ch.plots[0].series
for _sr, _c, _w in ((_ser[0], NAVY, 3.0), (_ser[1], SLATE_LT, 2.0),
                    (_ser[2], BLUE5, 2.0)):
    _sr.format.line.color.rgb = _c
    _sr.format.line.width = Pt(_w)
    _sr.smooth = True
_ser[2].format.line.dash_style = 2
ch.category_axis.tick_labels.font.size = Pt(9)
ch.value_axis.tick_labels.font.size = Pt(9)
ch.value_axis.has_major_gridlines = False
ch.value_axis.minimum_scale = 0
ch.value_axis.maximum_scale = 600
# --- down-market marker over year 6 (geometry rescales with the chart) ---
# plot area estimate inside the chart frame (tuned against the render)
_pl = Emu(int(chx) + int(Inches(0.62)))          # plot left
_pr = Emu(int(chx) + int(_cw_chart) - int(Inches(0.08)))  # plot right
_pt = top + Inches(0.45)                          # plot top
_pb = top + Inches(3.62)                           # plot bottom
_xv = Emu(int(_pl) + int((int(_pr) - int(_pl)) * 6 / 10))  # year-6 x
_mt = top + Inches(0.62)                          # marker top (below title)
rect(s, _xv, _mt, Pt(1.2), Emu(int(_pb) - int(_mt)), fill=SLATE_LT)  # ref line
mtf = tbox(s, Emu(int(_xv) - int(Inches(1.0))), Emu(int(_mt) - int(Inches(0.26))),
           Inches(2.0), Inches(0.3))
para(mtf, "Down-market year", 10, SLATE, first=True, bold=True,
     align=PP_ALIGN.CENTER, after=0, track=0)
# divergence callouts at year 6
def _yfor(v):
    return Emu(int(_pb) - int((int(_pb) - int(_pt)) * v / 600))
ca = tbox(s, Emu(int(_xv) + int(Inches(0.08))), Emu(int(_yfor(236)) - int(Inches(0.22))),
          Inches(1.6), Inches(0.26))
para(ca, "Athanase: holds", 9.5, NAVY_TX, first=True, bold=True, after=0, track=0)
cl = tbox(s, Emu(int(_xv) + int(Inches(0.08))), Emu(int(_yfor(119)) + int(Inches(0.04))),
          Inches(1.6), Inches(0.26))
para(cl, "Long-only: −19%", 9.5, SLATE_LT, first=True, after=0, track=0)
# end-value callouts
para(tbox(s, Inches(6.55), top + Inches(0.55), Inches(1.3), Inches(0.3)),
     "523", 12, NAVY_TX, first=True, bold=True, after=0)
para(tbox(s, Inches(6.55), top + Inches(2.35), Inches(1.3), Inches(0.3)),
     "271", 11, SLATE_LT, first=True, bold=True, after=0)
para(tbox(s, Inches(6.55), top + Inches(2.95), Inches(1.3), Inches(0.3)),
     "208", 11, SLATE_LT, first=True, after=0)

# --- right: annual uplift to total-portfolio return (bps) ---
_tp = tbl()
para(tbox(s, Inches(8.05), top + Inches(0.05), Inches(4.6), Inches(0.55)),
     f"{_tp}.  Annual uplift to TOTAL-portfolio return (bps), by sleeve size",
     11, SLATE, first=True, bold=True, after=0, lead=1.1, track=0)
rows = [
    ("Sleeve", "vs Long-only", "vs Other activist"),
    ("3%", "+31 bps", "+23 bps"),
    ("5%", "+52 bps", "+38 bps"),
    ("8%", "+83 bps", "+60 bps"),
]
tl = Inches(8.05); cw = [Inches(1.5), Inches(1.55), Inches(1.95)]
rh = Inches(0.46); yy = top + Inches(0.62)
for ri, row in enumerate(rows):
    head = ri == 0
    cx = tl
    for ci, cell in enumerate(row):
        fill = SLATE if head else (HEADERBG if ri % 2 else WHITE)
        rect(s, cx, yy, cw[ci], rh, fill=fill)
        ctf = tbox(s, Emu(int(cx) + int(Inches(0.1))), yy,
                   Emu(int(cw[ci]) - int(Inches(0.16))), rh,
                   anchor=MSO_ANCHOR.MIDDLE)
        para(ctf, cell, 12.5 if head else 12,
             WHITE if head else (NAVY_TX if ci == 0 else BODY),
             bold=(head or ci == 0), first=True,
             align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT,
             after=0, track=0)
        cx = Emu(int(cx) + int(cw[ci]))
    yy = Emu(int(yy) + int(rh))
# how the path is earned (ties to Fig. 4 average monthly returns, p.17)
hy = Emu(int(yy) + int(Inches(0.16)))
para(tbox(s, Inches(8.05), hy, Inches(4.6), Inches(0.3)),
     "HOW THE PATH IS EARNED", 11, SLATE, first=True, bold=True, after=0, track=0)
hp = tbox(s, Inches(8.05), Emu(int(hy) + int(Inches(0.3))), Inches(4.6), Inches(2.2))
para(hp, "The smooth line is not won by chasing rallies. On average monthly "
     "returns (p.17), Athanase roughly matches the market up (1.2% vs 1.3%) but "
     "protects sharply down (+0.2% vs −1.7%).", 11.5, BODY,
     first=True, after=7, lead=1.15, track=0)
para(hp, "Losing less in the down months — not winning more up — is what "
     "compounds the lines apart: a steadier path, not just a higher one.",
     11.5, BODY, after=0, lead=1.15, track=0)
# assumptions footnote
para(tbox(s, Inches(0.6), Inches(7.15), Inches(8.9), Inches(0.4)),
     "Illustrative only; not a projection. Athanase 18% (headline net IRR); "
     "other activist 10.5%; long-only 7.6% (MSCI IMI). Uplift = sleeve × return "
     "spread; growth = annual compounding. Past performance ≠ future results.",
     8.5, FOOT, first=True, after=0, track=0, lead=1.15)

# ---- II.6h  Outcome table by holding period --------------------------------
s, top = content("Risk-Adjusted Outcomes",
                 "Lower downside risk compounds into better outcomes",
                 "Athanase carries higher total volatility than the market — but "
                 "its downside volatility is lower, which is what protects "
                 "compounding.", ref=tbl())
# three risk-stat cards
metrics = [("Annualised return", _MSCI["ann_ret"], _ATH["ann_ret"]),
           ("Total volatility", _MSCI["ann_vol"], _ATH["ann_vol"]),
           ("Downside volatility", _MSCI["ann_dd"], _ATH["ann_dd"])]
bx = Inches(0.6); bw = Inches(3.95); gap = Inches(0.12)
for label, mv, av in metrics:
    rect(s, bx, top, bw, Inches(1.28), fill=HEADERBG)
    para(tbox(s, Emu(int(bx) + int(Inches(0.2))), top + Inches(0.14),
              Emu(int(bw) - int(Inches(0.4))), Inches(0.3)),
         label.upper(), 10.5, SLATE, first=True, bold=True, after=0, track=0)
    vt = tbox(s, Emu(int(bx) + int(Inches(0.2))), top + Inches(0.46),
              Emu(int(bw) - int(Inches(0.4))), Inches(0.8))
    p = vt.paragraphs[0]; p.space_after = Pt(2)
    r1 = p.add_run(); r1.text = "MSCI  "; r1.font.size = Pt(11)
    r1.font.color.rgb = SUBTLE; r1.font.name = SANS
    r2 = p.add_run(); r2.text = f"{mv*100:.1f}%"; r2.font.size = Pt(19)
    r2.font.bold = True; r2.font.color.rgb = NAVY_TX; r2.font.name = SERIF
    p2 = vt.add_paragraph()
    r3 = p2.add_run(); r3.text = "AIP   "; r3.font.size = Pt(11)
    r3.font.color.rgb = SUBTLE; r3.font.name = SANS
    r4 = p2.add_run(); r4.text = f"{av*100:.1f}%"; r4.font.size = Pt(19)
    r4.font.bold = True; r4.font.color.rgb = NAVY; r4.font.name = SERIF
    bx = Emu(int(bx) + int(bw) + int(gap))
para(tbox(s, Inches(9.27), top + Inches(1.32), Inches(3.95), Inches(0.4)),
     "↑ Athanase downside volatility is LOWER than the market’s.",
     10, NAVY, first=True, bold=True, after=0, track=0, lead=1.05)
# projection table
HZ = (3, 5, 7, 10)
ty = top + Inches(1.95)
para(tbox(s, Inches(0.6), ty - Inches(0.3), Inches(12), Inches(0.3)),
     "EXPECTED VALUE OF 100 INVESTED, BY HOLDING PERIOD", 11, SLATE,
     first=True, bold=True, after=0, track=0)
col0 = Inches(3.5); hw = (SW - Inches(1.2) - col0) / len(HZ); rh = Inches(0.35)
trows = [
    ("Median outcome — MSCI",
     [f"{_project(T, _MSCI['ann_ret'], _MSCI['ann_dd'])[0]:.0f}" for T in HZ], "sub"),
    ("Median outcome — Athanase",
     [f"{_project(T, _ATH['ann_ret'], _ATH['ann_dd'])[0]:.0f}" for T in HZ], "navy"),
    ("Adverse case* — MSCI",
     [f"{_project(T, _MSCI['ann_ret'], _MSCI['ann_dd'])[1]:.0f}" for T in HZ], "sub"),
    ("Adverse case* — Athanase",
     [f"{_project(T, _ATH['ann_ret'], _ATH['ann_dd'])[1]:.0f}" for T in HZ], "navy"),
    ("Chance of a loss — MSCI",
     [f"{_project(T, _MSCI['ann_ret'], _MSCI['ann_dd'])[2]*100:.0f}%" for T in HZ], "sub"),
    ("Chance of a loss — Athanase",
     [f"{_project(T, _ATH['ann_ret'], _ATH['ann_dd'])[2]*100:.0f}%" for T in HZ], "navy"),
]
cx = Inches(0.6)
for ci, htext in enumerate(["", "3 years", "5 years", "7 years", "10 years"]):
    w = col0 if ci == 0 else hw
    rect(s, cx, ty, w, rh, fill=SLATE)
    para(tbox(s, Emu(int(cx) + int(Inches(0.12))), ty,
              Emu(int(w) - int(Inches(0.2))), rh, anchor=MSO_ANCHOR.MIDDLE),
         htext, 12, WHITE, bold=True, first=True,
         align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT, after=0, track=0)
    cx = Emu(int(cx) + int(w))
yy = Emu(int(ty) + int(rh))
for ri, (label, vals, kind) in enumerate(trows):
    fill = HEADERBG if ri % 2 == 0 else WHITE
    cx = Inches(0.6)
    rect(s, cx, yy, col0, rh, fill=fill)
    para(tbox(s, Emu(int(cx) + int(Inches(0.12))), yy,
              Emu(int(col0) - int(Inches(0.2))), rh, anchor=MSO_ANCHOR.MIDDLE),
         label, 11.5, NAVY if kind == "navy" else SUBTLE,
         bold=(kind == "navy"), first=True, after=0, track=0)
    cx = Emu(int(cx) + int(col0))
    for v in vals:
        rect(s, cx, yy, hw, rh, fill=fill)
        para(tbox(s, Emu(int(cx) + int(Inches(0.1))), yy,
                  Emu(int(hw) - int(Inches(0.2))), rh, anchor=MSO_ANCHOR.MIDDLE),
             v, 12.5, NAVY_TX if kind == "navy" else SUBTLE,
             bold=(kind == "navy"), first=True, align=PP_ALIGN.RIGHT,
             after=0, track=0)
        cx = Emu(int(cx) + int(hw))
    yy = Emu(int(yy) + int(rh))
tk = Emu(int(yy) + int(Inches(0.08)))
rect(s, Inches(0.6), tk, Inches(12.13), Inches(0.38), fill=NAVY)
para(tbox(s, Inches(0.78), tk, Inches(11.8), Inches(0.38), anchor=MSO_ANCHOR.MIDDLE),
     "Even in the adverse case, Athanase’s lower downside volatility leaves an "
     "investor ahead at every horizon — and the gap widens with time.",
     11.5, WHITE, first=True, italic=True, after=0, track=0)
para(tbox(s, Inches(0.6), Emu(int(tk) + int(Inches(0.46))), Inches(12.6), Inches(0.4)),
     f"Source: monthly net returns, {_NM} months (2006–2025). *Adverse case = "
     "one downside-deviation outcome (≈1-in-6), narrowing with horizon as "
     "dd/√T. Illustrative; past performance is not indicative of future results.",
     7.5, FOOT, first=True, after=0, track=0, lead=1.1)


# ---- II.6i  Up/down capture + diversification ------------------------------
def _capture_chart(path_png):
    fig, ax = plt.subplots(figsize=(5.1, 4.3), dpi=200)
    xm = [0, 1.6]; w = 0.62
    avals = [_PAIR["up_cap"] * 100, _PAIR["dn_cap"] * 100]
    ax.bar([x - w / 2 for x in xm], [100, 100], w, color=H_BLUE4, alpha=0.55,
           label="MSCI World IMI (= market)")
    ax.bar([x + w / 2 for x in xm], avals, w, color=H_NAVY, label="Athanase")
    ax.axhline(100, color=H_BLUE5, lw=1.0, zorder=0)
    for x, v in zip([x + w / 2 for x in xm], avals):
        ax.annotate(f"{v:.0f}%", (x, v), ha="center", va="bottom", fontsize=15,
                    fontweight="bold", color=H_NAVY, xytext=(0, 3),
                    textcoords="offset points")
    for x in [x - w / 2 for x in xm]:
        ax.annotate("100%", (x, 100), ha="center", va="bottom", fontsize=11,
                    color=H_BLUE4, xytext=(0, 3), textcoords="offset points")
    ax.set_xticks(xm)
    ax.set_xticklabels(["Up markets\n(rallies)", "Down markets\n(selloffs)"],
                       fontsize=11, color=H_BODY)
    ax.set_ylabel("Share of the market’s move captured", color=H_BODY, fontsize=10.5)
    ax.set_ylim(0, 120); ax.set_xlim(-0.7, 2.3); ax.set_yticks([0, 50, 100])
    ax.set_yticklabels(["0%", "50%", "100%"], fontsize=10, color=H_BODY)
    ax.tick_params(colors=H_BODY)
    for sp in ("top", "right"):
        ax.spines[sp].set_visible(False)
    for sp in ("left", "bottom"):
        ax.spines[sp].set_color(H_BLUE5)
    ax.grid(True, axis="y", color=H_BLUE5, lw=0.6, alpha=0.6)
    ax.legend(loc="upper center", bbox_to_anchor=(0.5, -0.16), ncol=2,
              fontsize=10, frameon=False, labelcolor=H_BODY, handlelength=1.4,
              columnspacing=1.6)
    ax.set_title("Up- vs down-market capture", color=H_NAVY, fontsize=13,
                 fontweight="bold", pad=10)
    fig.tight_layout(); fig.savefig(path_png, dpi=200, bbox_inches="tight",
                                    facecolor="white"); plt.close(fig)


_capture_chart("/tmp/mc_capture.png")
_f_cap = fig()
s, top = content("Risk-Adjusted Outcomes",
                 "A complement to global equities, not a duplicate",
                 "What a global-equity holder asks of a new manager: do you keep "
                 "up in rallies, lose less in selloffs, and add something the "
                 "index can’t?", ref=_f_cap)
s.shapes.add_picture("/tmp/mc_capture.png", Inches(0.55), Inches(2.25),
                     height=Inches(4.1))
cardx = Inches(6.55); cardw = Inches(6.2); cardh = Inches(1.2); cy = Inches(2.3)
for title, big, body in [
    ("CAPTURE ASYMMETRY",
     f"{_PAIR['up_cap']*100:.0f}% up  /  {_PAIR['dn_cap']*100:.0f}% down",
     "Keeps pace with rallies but takes less than half of selloffs — the "
     "essence of the edge."),
    ("DIVERSIFICATION",
     f"{_PAIR['corr']:.2f} correlation  ·  {_PAIR['beta']:.2f} beta",
     "Low correlation to your existing equity book — a genuine diversifier, "
     "not levered beta."),
    ("RELIABILITY OVER TIME",
     f"Beat MSCI in {_PAIR['roll'][5]*100:.0f}% of 5-yr and "
     f"{_PAIR['roll'][7]*100:.0f}% of 7-yr windows",
     "Over every rolling 7-year holding period in 20 years, Athanase "
     "out-returned the index.")]:
    rect(s, cardx, cy, cardw, cardh, fill=HEADERBG)
    para(tbox(s, Emu(int(cardx) + int(Inches(0.22))), cy + Inches(0.12),
              Emu(int(cardw) - int(Inches(0.44))), Inches(0.25)),
         title, 10.5, SLATE, first=True, bold=True, after=0, track=0)
    para(tbox(s, Emu(int(cardx) + int(Inches(0.22))), cy + Inches(0.36),
              Emu(int(cardw) - int(Inches(0.44))), Inches(0.35)),
         big, 17, NAVY_TX, first=True, bold=True, after=0, font=SERIF, track=0)
    para(tbox(s, Emu(int(cardx) + int(Inches(0.22))), cy + Inches(0.74),
              Emu(int(cardw) - int(Inches(0.44))), Inches(0.45)),
         body, 10.5, BODY, first=True, after=0, lead=1.12, track=0)
    cy = Emu(int(cy) + int(cardh) + int(Inches(0.12)))
rect(s, Inches(0.6), Inches(6.62), Inches(12.13), Inches(0.46), fill=NAVY)
para(tbox(s, Inches(0.78), Inches(6.62), Inches(11.8), Inches(0.46),
          anchor=MSO_ANCHOR.MIDDLE),
     f"Athanase captures {_PAIR['up_cap']*100:.0f}% of the market’s upside but "
     f"only {_PAIR['dn_cap']*100:.0f}% of its downside — at {_PAIR['corr']:.2f} "
     "correlation, it improves the whole portfolio, not just one line of it.",
     12, WHITE, first=True, italic=True, after=0, track=0)
para(tbox(s, Inches(0.6), Inches(7.16), Inches(12.6), Inches(0.4)),
     f"Source: monthly net returns, {_NM} months (2006–2025). Capture = "
     "compounded Athanase return in the market’s up/down months ÷ the market’s. "
     "Past performance is not indicative of future results.",
     7.5, FOOT, first=True, after=0, track=0, lead=1.1)


# ---- II.6j  Efficient frontier ---------------------------------------------
def _ann_b(xs):
    g = 1.0
    for x in xs:
        g *= (1 + x)
    return g ** (12 / len(xs)) - 1


def _dvol_b(xs):
    return math.sqrt(sum(min(x, 0) ** 2 for x in xs) / len(xs)) * math.sqrt(12)


_FRONT = []
for _w in [i / 20 for i in range(21)]:
    _bl = [_w * _ATH_X[i] + (1 - _w) * _MSCI_X[i] for i in range(_NM)]
    _FRONT.append((_w, _dvol_b(_bl) * 100, _ann_b(_bl) * 100))
_b0 = _FRONT[0]
_b8 = min(_FRONT, key=lambda z: abs(z[0] - 0.08))
_fmin = min(_FRONT, key=lambda z: z[1])


def _frontier_chart(path_png):
    fig, ax = plt.subplots(figsize=(8.6, 4.35), dpi=200)
    ax.plot([p[1] for p in _FRONT], [p[2] for p in _FRONT], color=H_BLUE4,
            lw=2.0, zorder=3)

    def mark(wt, label, col, dy=8, dx=8, big=False):
        p = min(_FRONT, key=lambda z: abs(z[0] - wt))
        ax.scatter([p[1]], [p[2]], s=120 if big else 80, color=col, zorder=5,
                   edgecolor="white", linewidth=1.2)
        ax.annotate(label, (p[1], p[2]), color=col, fontsize=10.5,
                    fontweight="bold", xytext=(dx, dy), textcoords="offset points")
    mark(0.0, "100% Global equities\n(MSCI World IMI)", H_BLUE4, dy=-30, dx=-30)
    mark(1.0, "100% Athanase", H_NAVY, dy=6, dx=-150, big=True)
    ax.scatter([_fmin[1]], [_fmin[2]], s=150, color=H_NAVY, zorder=6,
               edgecolor="white", linewidth=1.5, marker="D")
    ax.annotate(f"Minimum downside risk\n≈ {_fmin[0]*100:.0f}% Athanase  "
                f"({_fmin[2]:.1f}% return)", (_fmin[1], _fmin[2]), color=H_NAVY,
                fontsize=10, fontweight="bold", xytext=(12, -4),
                textcoords="offset points", va="center")
    seg = [z for z in _FRONT if 0.03 <= z[0] <= 0.08]
    ax.plot([z[1] for z in seg], [z[2] for z in seg], color=H_NAVY, lw=6,
            alpha=0.95, solid_capstyle="round", zorder=6)
    pmid = min(_FRONT, key=lambda z: abs(z[0] - 0.055))
    ax.annotate("prudent 3–8% zone", (pmid[1], pmid[2]), color=H_NAVY,
                fontsize=10.5, fontweight="bold", xytext=(10.55, 5.0),
                textcoords="data", va="center", ha="center",
                arrowprops=dict(arrowstyle="-", color=H_NAVY, lw=1.0,
                                connectionstyle="arc3,rad=0.2"))
    ax.set_xlabel("Downside volatility (annualised)", color=H_BODY, fontsize=11)
    ax.set_ylabel("Annualised return", color=H_BODY, fontsize=11)
    ax.set_xlim(9.8, 12.2); ax.set_ylim(4, 18)
    ax.set_xticks([10, 10.5, 11, 11.5, 12])
    ax.xaxis.set_major_formatter(lambda v, _: f"{v:.1f}%")
    ax.yaxis.set_major_formatter(lambda v, _: f"{v:.0f}%")
    ax.tick_params(colors=H_BODY, labelsize=10)
    for sp in ("top", "right"):
        ax.spines[sp].set_visible(False)
    for sp in ("left", "bottom"):
        ax.spines[sp].set_color(H_BLUE5)
    ax.grid(True, color=H_BLUE5, lw=0.6, alpha=0.6)
    ax.set_title("Blended portfolio: return vs downside risk", color=H_NAVY,
                 fontsize=13, fontweight="bold", pad=10)
    fig.tight_layout(); fig.savefig(path_png, dpi=200, bbox_inches="tight",
                                    facecolor="white"); plt.close(fig)


_frontier_chart("/tmp/mc_frontier.png")
_f_fr = fig()
s, top = content("Risk-Adjusted Outcomes",
                 "What a blend does to the whole portfolio",
                 "Adding Athanase to a global-equity book moves the portfolio up "
                 "and to the left — more return for less downside risk, thanks to "
                 "low correlation.", ref=_f_fr)
s.shapes.add_picture("/tmp/mc_frontier.png", Inches(0.55), Inches(2.25),
                     height=Inches(4.1))
cx3 = Inches(8.9); cw3 = Inches(3.85); ch3 = Inches(1.2); cy3 = Inches(2.32)
for title, b1, b2, sub in [
    ("ALL GLOBAL EQUITIES", f"{_b0[2]:.1f}% return", f"{_b0[1]:.1f}% downside risk",
     "100% MSCI World IMI"),
    ("ADD A PRUDENT 8% SLEEVE", f"{_b8[2]:.1f}% return",
     f"{_b8[1]:.1f}% downside risk", "92% MSCI  /  8% Athanase")]:
    rect(s, cx3, cy3, cw3, ch3, fill=HEADERBG)
    para(tbox(s, Emu(int(cx3) + int(Inches(0.22))), cy3 + Inches(0.12),
              Emu(int(cw3) - int(Inches(0.44))), Inches(0.25)),
         title, 10.5, SLATE, first=True, bold=True, after=0, track=0)
    rt = tbox(s, Emu(int(cx3) + int(Inches(0.22))), cy3 + Inches(0.38),
              Emu(int(cw3) - int(Inches(0.44))), Inches(0.45))
    p = rt.paragraphs[0]; p.space_after = Pt(0)
    r1 = p.add_run(); r1.text = b1; r1.font.size = Pt(16); r1.font.bold = True
    r1.font.color.rgb = NAVY_TX; r1.font.name = SERIF
    r2 = p.add_run(); r2.text = "   ·   " + b2; r2.font.size = Pt(12)
    r2.font.color.rgb = SLATE; r2.font.name = SANS
    para(tbox(s, Emu(int(cx3) + int(Inches(0.22))), cy3 + Inches(0.82),
              Emu(int(cw3) - int(Inches(0.44))), Inches(0.3)),
         sub, 10.5, BODY, first=True, after=0, track=0)
    cy3 = Emu(int(cy3) + int(ch3) + int(Inches(0.14)))
para(tbox(s, cx3, cy3 + Inches(0.02), cw3, Inches(0.3)),
     "HOW TO READ IT", 10.5, SLATE, first=True, bold=True, after=0, track=0)
para(tbox(s, cx3, cy3 + Inches(0.3), cw3, Inches(1.9)),
     "Every blend on the curve beats holding equities alone, and the benefit "
     f"keeps building until past half the portfolio (downside risk bottoms near "
     f"{_fmin[0]*100:.0f}%). That is a directional case to size up — within "
     "prudent limits — not a target: a 59% position would breach concentration, "
     "liquidity and the 3–8% sizing discipline. Even an 8% sleeve already moves "
     "the portfolio the right way.", 10.5, BODY, first=True, after=0,
     lead=1.16, track=0)
rect(s, Inches(0.6), Inches(6.62), Inches(12.13), Inches(0.46), fill=NAVY)
para(tbox(s, Inches(0.78), Inches(6.62), Inches(11.8), Inches(0.46),
          anchor=MSO_ANCHOR.MIDDLE),
     "The blend is not a trade-off: every allocation raises return and lowers "
     "downside risk — so to a committee this is a mathematical risk-reduction "
     "tool, not a risky bet, sized within prudent 3–8% limits.",
     12, WHITE, first=True, italic=True, after=0, track=0)
para(tbox(s, Inches(0.6), Inches(7.16), Inches(12.6), Inches(0.4)),
     f"Source: monthly net returns, {_NM} months (2006–2025), monthly rebalanced "
     "blends. The minimum-risk point is in-sample and shown for shape, not as "
     "advice. Illustrative; past performance is not indicative of future results.",
     7.5, FOOT, first=True, after=0, track=0, lead=1.1)


subdivider("Versus private equity")
# ===========================================================================
# II.6k-o  Athanase vs mid-market private equity (real, de-smoothed data)
# PE figures from published research; Athanase computed from real returns.
# ===========================================================================
ATH_RET = _ATH["ann_ret"]; ATH_DD = _ATH["ann_dd"]; ATH_VOL = _ATH["ann_vol"]
ATH_CORR = _PAIR["corr"]
PE_HEADLINE = 0.20            # top-quartile mid-market buyout headline net IRR
PE_REAL = 0.114              # Cliffwater realised PE (US state pensions, 2000-22)
PE_VOL_REPORTED = 0.10; PE_VOL_TRUE = 0.28      # reported vs de-smoothed total vol
PE_CORR_REPORTED = 0.75; PE_CORR_TRUE = 0.89    # reported vs de-smoothed corr
PE_EVIDENCE = [
    ("Cliffwater\nUS state pensions\n(2000–2022)", 11.4, 5.8),
    ("BVCA / Capital Dynamics\nUK & Europe\n(2001–2023)", 14.1, 7.7),
]


def _cum_pe(xs):
    g = 1.0
    for x in xs:
        g *= (1 + x)
    return g - 1


_roll36 = [_cum_pe(_ATH_X[i:i + 36]) for i in range(0, _NM - 36 + 1)]
ATH_P30 = sum(1 for r in _roll36 if r < -0.30) / len(_roll36)
PE_P30 = 0.155


# ---- II.6k  What allocators actually earn ----------------------------------
def _pe_realised(path):
    fig_, ax = plt.subplots(figsize=(8.7, 4.35), dpi=200)
    g = [e[0] for e in PE_EVIDENCE]; pe = [e[1] for e in PE_EVIDENCE]
    pme = [e[2] for e in PE_EVIDENCE]; xx = np.arange(len(g)); w = 0.38
    ax.bar(xx - w / 2, pe, w, color=H_BLUE4, label="PE realised (net)")
    ax.bar(xx + w / 2, pme, w, color=H_BLUE5, edgecolor=H_BLUE4, linewidth=1,
           label="Public-market equivalent")
    for xi, (p, q) in enumerate(zip(pe, pme)):
        ax.annotate(f"{p:.1f}%", (xi - w / 2, p), ha="center", va="bottom",
                    fontsize=12, fontweight="bold", color=H_BLUE3, xytext=(0, 3),
                    textcoords="offset points")
        ax.annotate(f"{q:.1f}%", (xi + w / 2, q), ha="center", va="bottom",
                    fontsize=11, color=H_BLUE4, xytext=(0, 3),
                    textcoords="offset points")
        ax.annotate(f"+{p-q:.1f} pts\npremium", (xi, (p + q) / 2), ha="center",
                    va="center", fontsize=9.5, color=H_BLUE3, fontweight="bold")
    ax.axhline(ATH_RET * 100, color=H_NAVY, lw=1.8, ls=(0, (5, 3)), zorder=5)
    ax.annotate(f"Athanase {ATH_RET*100:.0f}%  (real, net, fully invested)",
                (len(g) - 0.5, ATH_RET * 100), ha="right", va="bottom",
                fontsize=10.5, color=H_NAVY, fontweight="bold")
    ax.set_xticks(xx); ax.set_xticklabels(g, fontsize=9.5, color=H_BODY)
    ax.set_ylabel("Annualised return to investors", color=H_BODY, fontsize=11)
    ax.set_ylim(0, 19); ax.set_xlim(-0.6, len(g) - 0.4)
    ax.yaxis.set_major_formatter(lambda v, _: f"{v:.0f}%")
    ax.tick_params(colors=H_BODY, labelsize=10)
    for sp in ("top", "right"):
        ax.spines[sp].set_visible(False)
    for sp in ("left", "bottom"):
        ax.spines[sp].set_color(H_BLUE5)
    ax.grid(True, axis="y", color=H_BLUE5, lw=0.6, alpha=0.6)
    ax.legend(loc="upper center", bbox_to_anchor=(0.5, -0.16), ncol=2,
              fontsize=10, frameon=False, labelcolor=H_BODY)
    ax.set_title("What allocators actually realise from private equity",
                 color=H_NAVY, fontsize=13, fontweight="bold", pad=10)
    fig_.tight_layout(); fig_.savefig(path, dpi=200, bbox_inches="tight",
                                      facecolor="white"); plt.close(fig_)


_pe_realised("/tmp/d_pe_realised.png")
s, top = content("Private Equity",
                 "What allocators actually earn from private equity",
                 "Stripping out headline-IRR marketing, cash-matched (PME) data "
                 "shows PE earns a real ~3–6 pt premium over public markets — but "
                 "lands below Athanase.", ref=tbl())
s.shapes.add_picture("/tmp/d_pe_realised.png", Inches(0.55), Inches(2.3),
                     height=Inches(3.85))
cx = Inches(8.95); cw = Inches(3.8); chh = Inches(1.16); cy = Inches(2.35)
for ti, bd in [
    ("HEADLINE ≠ REALISED",
     "Top-quartile IRR is shown as ~18–22%, struck on drawn capital. Cash-matched "
     "to public markets, the real figure is far lower."),
    ("THE REAL PREMIUM IS GENUINE",
     "PME studies agree: ~1.2× wealth multiple, a real 3–5 pt premium "
     "(Cambridge; Kaplan-Harris-Jenkinson / Burgiss; Cliffwater; BVCA)."),
    ("BUT BELOW ATHANASE",
     "Allocators realise ~11–14% from PE. Athanase has delivered ~16% — net, "
     "real, and fully invested from day one.")]:
    rect(s, cx, cy, cw, chh, fill=HEADERBG)
    para(tbox(s, Emu(int(cx) + int(Inches(0.2))), cy + Inches(0.1),
              Emu(int(cw) - int(Inches(0.4))), Inches(0.25)), ti, 10.5, SLATE,
         first=True, bold=True, after=0, track=0)
    para(tbox(s, Emu(int(cx) + int(Inches(0.2))), cy + Inches(0.36),
              Emu(int(cw) - int(Inches(0.4))), Inches(0.78)), bd, 10.5, BODY,
         first=True, after=0, lead=1.12, track=0)
    cy = Emu(int(cy) + int(chh) + int(Inches(0.14)))
rect(s, Inches(0.6), Inches(6.62), Inches(12.13), Inches(0.46), fill=NAVY)
para(tbox(s, Inches(0.78), Inches(6.62), Inches(11.8), Inches(0.46),
          anchor=MSO_ANCHOR.MIDDLE),
     "PE’s premium over public markets is real — but Athanase’s ~16% sits above "
     "even what the best allocators realise, with daily liquidity.",
     12, WHITE, first=True, italic=True, after=0, track=0)
para(tbox(s, Inches(0.6), Inches(7.12), Inches(12.6), Inches(0.5)),
     "Sources: Cliffwater (US state pensions, 2000–22) PE 11.4% vs PME 5.8%; "
     "BVCA / Capital Dynamics (UK & Europe, 2001–23) 14.1% vs 7.7%; corroborated "
     "by Cambridge Associates mPME 1.15–1.25× and Kaplan-Harris-Jenkinson KS-PME "
     "1.20–1.27×. Athanase real, net, 2006–2025. See references.",
     7.5, FOOT, first=True, after=0, track=0, lead=1.1)


# ---- II.6l  The switch: real return vs true volatility ---------------------
def _pe_scatter(path):
    fig_, ax = plt.subplots(figsize=(6.0, 4.25), dpi=200)
    ax.annotate("", (PE_VOL_TRUE * 100, PE_REAL * 100),
                (PE_VOL_REPORTED * 100, PE_HEADLINE * 100),
                arrowprops=dict(arrowstyle="->", color=H_BLUE4, lw=1.8,
                                linestyle=(0, (3, 2))))
    ax.scatter([PE_VOL_REPORTED * 100], [PE_HEADLINE * 100], s=150, color="white",
               edgecolor=H_BLUE4, linewidth=2.2, zorder=5)
    ax.annotate("PE as marketed\n(headline IRR,\nsmoothed ~10% vol)",
                (PE_VOL_REPORTED * 100, PE_HEADLINE * 100), color=H_BLUE4,
                fontsize=9.5, fontweight="bold", xytext=(6, 6),
                textcoords="offset points", ha="left")
    ax.scatter([PE_VOL_TRUE * 100], [PE_REAL * 100], s=150, color=H_BLUE4,
               zorder=5, edgecolor="white", linewidth=1.5)
    ax.annotate("PE as realised\n(PME return,\nde-smoothed ~28% vol)",
                (PE_VOL_TRUE * 100, PE_REAL * 100), color=H_BLUE4, fontsize=9.5,
                xytext=(-8, -10), textcoords="offset points", ha="right", va="top")
    ax.scatter([ATH_VOL * 100], [ATH_RET * 100], s=210, color=H_NAVY, zorder=6,
               edgecolor="white", linewidth=1.5)
    ax.annotate(f"Athanase\n{ATH_RET*100:.0f}% return", (ATH_VOL * 100, ATH_RET * 100),
                color=H_NAVY, fontsize=11, fontweight="bold", xytext=(-10, 12),
                textcoords="offset points", va="center", ha="right")
    ax.set_xlabel("Total volatility (annualised)", color=H_BODY, fontsize=11)
    ax.set_ylabel("Real net return to investors", color=H_BODY, fontsize=11)
    ax.set_xlim(4, 34); ax.set_ylim(6, 23)
    ax.xaxis.set_major_formatter(lambda v, _: f"{v:.0f}%")
    ax.yaxis.set_major_formatter(lambda v, _: f"{v:.0f}%")
    ax.tick_params(colors=H_BODY, labelsize=10)
    for sp in ("top", "right"):
        ax.spines[sp].set_visible(False)
    for sp in ("left", "bottom"):
        ax.spines[sp].set_color(H_BLUE5)
    ax.grid(True, color=H_BLUE5, lw=0.6, alpha=0.6)
    ax.set_title("Real return vs true (de-smoothed) volatility", color=H_NAVY,
                 fontsize=13, fontweight="bold", pad=10)
    fig_.tight_layout(); fig_.savefig(path, dpi=200, bbox_inches="tight",
                                      facecolor="white"); plt.close(fig_)


_pe_scatter("/tmp/d_pe_scatter.png")
_f_sw = fig()
s, top = content("Private Equity",
                 "Switching a private-equity sleeve into Athanase",
                 "On true risk, PE and Athanase carry similar volatility — but "
                 "Athanase’s is mostly upside, and it returns more, with daily "
                 "liquidity.", ref=_f_sw)
s.shapes.add_picture("/tmp/d_pe_scatter.png", Inches(0.55), Inches(2.3),
                     height=Inches(3.85))
para(tbox(s, Inches(0.7), Inches(6.25), Inches(6.2), Inches(0.4)),
     "De-smoothed, PE’s ~10% reported vol is really ~28% (PIMCO; AQR beta "
     "1.5–1.6) — about Athanase’s 27%, but levered downside, not upside.",
     10, SUBTLE, first=True, italic=True, after=0, lead=1.12)
cardx = Inches(7.05); cardw = Inches(5.7); cardh = Inches(0.78); cy = Inches(2.28)
for ti, big, bd in [
    ("RETURN", "~16% vs ~11–14% realised",
     "Above what even the best allocators realise from PE (Cliffwater; BVCA)."),
    ("LIQUIDITY", "Daily vs 10-yr lock-up",
     "Listed positions you can exit; no blind-pool capital calls."),
    ("FEES & CAPITAL", "Lower drag, fully invested",
     "No fees on undeployed dry powder waiting in a queue."),
    ("TRANSPARENCY", "Daily marks, real control",
     "A public scoreboard plus board control of cash flow, capex and pay."),
    ("RISK, HONESTLY", "Similar vol, opposite shape",
     "PE’s true ~28% vol is levered downside; Athanase’s 27% is mostly upside "
     "(downside only ~11%).")]:
    rect(s, cardx, cy, cardw, cardh, fill=HEADERBG)
    para(tbox(s, Emu(int(cardx) + int(Inches(0.2))), cy + Inches(0.09),
              Inches(2.0), Inches(0.25)), ti, 9.5, SLATE, first=True, bold=True,
         after=0, track=0)
    para(tbox(s, Emu(int(cardx) + int(Inches(2.05))), cy + Inches(0.08),
              Inches(3.4), Inches(0.3)), big, 12.5, NAVY_TX, first=True,
         bold=True, after=0, font=SERIF, track=0)
    para(tbox(s, Emu(int(cardx) + int(Inches(0.2))), cy + Inches(0.36),
              Emu(int(cardw) - int(Inches(0.4))), Inches(0.4)), bd, 9.5, BODY,
         first=True, after=0, lead=1.05, track=0)
    cy = Emu(int(cy) + int(cardh) + int(Inches(0.06)))
rect(s, Inches(0.6), Inches(6.62), Inches(12.13), Inches(0.46), fill=NAVY)
para(tbox(s, Inches(0.78), Inches(6.62), Inches(11.8), Inches(0.46),
          anchor=MSO_ANCHOR.MIDDLE),
     "Once PE’s smoothing is removed, the two carry similar volatility — but "
     "Athanase returns more, its risk is upside not downside, and it stays "
     "liquid and transparent.", 12, WHITE, first=True, italic=True, after=0,
     track=0)
para(tbox(s, Inches(0.6), Inches(7.12), Inches(12.6), Inches(0.5)),
     f"Athanase: net monthly returns, {_NM} months (2006–2025), real (total vol "
     f"{ATH_VOL*100:.0f}%). PE realised ~11–14% vs headline IRR ~18–22%; reported "
     "~10% vol de-smoothed to ~25–30% (PIMCO; AQR beta 1.5–1.6). See references.",
     7.5, FOOT, first=True, after=0, track=0, lead=1.1)


# ---- II.6m  The volatility gap is settled in the literature ----------------
_EVID = [
    ("Cash-flow NPV model", "11% vol", "25% vol", 11, 25),
    ("Secondary-market prices", "beta < 0.5", "beta > 2.0", 11, 28),
    ("Econometric unsmoothing", "smoothed", "serial-corr. removed", 11, 24),
    ("3-step unsmoothing", "low vol", "vol ↑, market ↑", 11, 27),
]


def _pe_evidence(path):
    fig_, ax = plt.subplots(figsize=(7.4, 4.5), dpi=200)
    ys = np.arange(len(_EVID))[::-1]; h = 0.34
    ax.barh(ys + h / 2, [e[3] for e in _EVID], h, color=H_BLUE5,
            edgecolor=H_BLUE4, linewidth=1, label="Reported (appraisal-smoothed)")
    ax.barh(ys - h / 2, [e[4] for e in _EVID], h, color=H_NAVY,
            label="True economic (research)")
    for y, e in zip(ys, _EVID):
        ax.annotate(e[1], (e[3] + 0.5, y + h / 2), va="center", ha="left",
                    fontsize=9, color=H_BLUE4)
        ax.annotate(e[2], (e[4] + 0.5, y - h / 2), va="center", ha="left",
                    fontsize=9, color=H_NAVY, fontweight="bold")
    ax.set_yticks(ys); ax.set_yticklabels([e[0] for e in _EVID], fontsize=10.5,
                                          color=H_BODY)
    ax.set_xlabel("Volatility / risk loading (de-smoothed as vol-equivalent)",
                  color=H_BODY, fontsize=10)
    ax.set_xlim(0, 36); ax.set_ylim(-0.7, len(_EVID) - 0.3)
    ax.xaxis.set_major_formatter(lambda v, _: f"{v:.0f}%")
    ax.tick_params(colors=H_BODY, labelsize=9.5)
    for sp in ("top", "right"):
        ax.spines[sp].set_visible(False)
    for sp in ("left", "bottom"):
        ax.spines[sp].set_color(H_BLUE5)
    ax.grid(True, axis="x", color=H_BLUE5, lw=0.6, alpha=0.6)
    ax.legend(loc="upper center", bbox_to_anchor=(0.5, -0.13), ncol=2,
              fontsize=9.5, frameon=False, labelcolor=H_BODY)
    ax.set_title("Four independent methods, one conclusion", color=H_NAVY,
                 fontsize=13, fontweight="bold", pad=10)
    fig_.tight_layout(); fig_.savefig(path, dpi=200, bbox_inches="tight",
                                      facecolor="white"); plt.close(fig_)


_pe_evidence("/tmp/d_pe_evidence.png")
s, top = content("Private Equity",
                 "The volatility gap is settled in the literature",
                 "Cash-flow models, secondary-market prices and econometric "
                 "unsmoothing all converge: PE’s true risk is roughly 2–3× its "
                 "reported figure.", ref=fig())
s.shapes.add_picture("/tmp/d_pe_evidence.png", Inches(0.5), Inches(2.3),
                     height=Inches(3.95))
cx = Inches(8.5); cw = Inches(4.25); chh = Inches(0.9); cy = Inches(2.35)
for ti, bd in [
    ("CASH-FLOW NPV MODEL",
     "Buyout vol 25% (vs 11% index), beta > 1 — Ang, Chen, Goetzmann & "
     "Phalippou (2018), J. Finance."),
    ("SECONDARY-MARKET PRICES",
     "Traded-stake beta > 2.0 vs < 0.5 on NAV — Boyer, Nadauld, Vorkink & "
     "Weisbach (2018)."),
    ("ECONOMETRIC UNSMOOTHING",
     "Illiquidity creates spurious smoothness — Getmansky, Lo & Makarov (2004), "
     "JFE (cited 1,300+)."),
    ("CATERING / “LAUNDERING”",
     "GPs smooth because LPs value the “phony happiness” — Jackson, Ling & "
     "Naranjo (2022); Couts et al. (2020).")]:
    rect(s, cx, cy, cw, chh, fill=HEADERBG)
    para(tbox(s, Emu(int(cx) + int(Inches(0.2))), cy + Inches(0.1),
              Emu(int(cw) - int(Inches(0.4))), Inches(0.25)), ti, 10, SLATE,
         first=True, bold=True, after=0, track=0)
    para(tbox(s, Emu(int(cx) + int(Inches(0.2))), cy + Inches(0.33),
              Emu(int(cw) - int(Inches(0.4))), Inches(0.55)), bd, 9.5, BODY,
         first=True, after=0, lead=1.1, track=0)
    cy = Emu(int(cy) + int(chh) + int(Inches(0.12)))
rect(s, Inches(0.6), Inches(6.62), Inches(12.13), Inches(0.46), fill=NAVY)
para(tbox(s, Inches(0.78), Inches(6.62), Inches(11.8), Inches(0.46),
          anchor=MSO_ANCHOR.MIDDLE),
     "Independent methods, the same answer: PE is a highly-levered equity "
     "portfolio whose reported smoothness is an artefact.",
     12, WHITE, first=True, italic=True, after=0, track=0)
para(tbox(s, Inches(0.6), Inches(7.12), Inches(12.6), Inches(0.5)),
     "Sources: Ang, Chen, Goetzmann & Phalippou (2018); Boyer, Nadauld, Vorkink "
     "& Weisbach (2018); Getmansky, Lo & Makarov (2004); Couts, Gonçalves & Rossi "
     "(2020); Jackson, Ling & Naranjo (2022); PIMCO (2022); AQR. Bars show "
     "vol-equivalent risk; beta rows scaled for comparison. See references.",
     7, FOOT, first=True, after=0, track=0, lead=1.1)


# ---- II.6n  Same headline volatility, far lower real risk ------------------
def _pe_realrisk(path):
    fig_, ax = plt.subplots(figsize=(6.3, 4.5), dpi=200)
    x = np.arange(2); w = 0.38
    pe = [PE_VOL_TRUE * 100, PE_VOL_TRUE * 0.72 * 100]
    ath = [ATH_VOL * 100, ATH_DD * 100]
    ax.bar(x - w / 2, pe, w, color=H_BLUE4, label="Private equity (de-smoothed)")
    ax.bar(x + w / 2, ath, w, color=H_NAVY, label="Athanase")
    for xi, (p, a) in enumerate(zip(pe, ath)):
        ax.annotate(f"{p:.0f}%", (xi - w / 2, p), ha="center", va="bottom",
                    fontsize=12, fontweight="bold", color=H_BLUE4, xytext=(0, 3),
                    textcoords="offset points")
        ax.annotate(f"{a:.0f}%", (xi + w / 2, a), ha="center", va="bottom",
                    fontsize=12, fontweight="bold", color=H_NAVY, xytext=(0, 3),
                    textcoords="offset points")
    ax.set_xticks(x); ax.set_xticklabels(["Total\nvolatility", "Downside\nvolatility"],
                                         fontsize=11, color=H_BODY)
    ax.set_ylabel("Annualised volatility", color=H_BODY, fontsize=11)
    ax.set_ylim(0, 34); ax.set_xlim(-0.6, 1.6)
    ax.yaxis.set_major_formatter(lambda v, _: f"{v:.0f}%")
    ax.tick_params(colors=H_BODY, labelsize=10)
    for sp in ("top", "right"):
        ax.spines[sp].set_visible(False)
    for sp in ("left", "bottom"):
        ax.spines[sp].set_color(H_BLUE5)
    ax.grid(True, axis="y", color=H_BLUE5, lw=0.6, alpha=0.6)
    ax.legend(loc="upper center", bbox_to_anchor=(0.5, -0.14), ncol=2,
              fontsize=10, frameon=False, labelcolor=H_BODY)
    ax.set_title("Similar total volatility — very different downside",
                 color=H_NAVY, fontsize=13, fontweight="bold", pad=10)
    fig_.tight_layout(); fig_.savefig(path, dpi=200, bbox_inches="tight",
                                      facecolor="white"); plt.close(fig_)


_pe_realrisk("/tmp/d_pe_realrisk.png")
s, top = content("Private Equity",
                 "Same headline volatility — far lower real risk",
                 "Volatility is not symmetric. PE’s true risk is levered "
                 "downside; Athanase’s is mostly upside, so its real loss risk is "
                 "a fraction of PE’s.", ref=fig())
s.shapes.add_picture("/tmp/d_pe_realrisk.png", Inches(0.55), Inches(2.3),
                     height=Inches(3.85))
bx = Inches(7.05); bw = Inches(5.7)
rect(s, bx, Inches(2.35), bw, Inches(1.75), fill=HEADERBG)
para(tbox(s, Emu(int(bx) + int(Inches(0.25))), Inches(2.5), Inches(5.2), Inches(0.3)),
     "CHANCE OF A 30%+ LOSS OVER THREE YEARS", 11, SLATE, first=True, bold=True,
     after=0, track=0)
spx = tbox(s, Emu(int(bx) + int(Inches(0.25))), Inches(2.86), Inches(5.2), Inches(1.1))
p = spx.paragraphs[0]; p.space_after = Pt(2)
r1 = p.add_run(); r1.text = f"~{PE_P30*100:.0f}%"; r1.font.size = Pt(32)
r1.font.bold = True; r1.font.color.rgb = BLUE4; r1.font.name = SERIF
r2 = p.add_run(); r2.text = "   private equity"; r2.font.size = Pt(13)
r2.font.color.rgb = SUBTLE; r2.font.name = SANS
p2 = spx.add_paragraph()
r3 = p2.add_run(); r3.text = f"~{ATH_P30*100:.0f}%"; r3.font.size = Pt(32)
r3.font.bold = True; r3.font.color.rgb = NAVY; r3.font.name = SERIF
r4 = p2.add_run(); r4.text = "   Athanase (actual 3-yr windows)"
r4.font.size = Pt(13); r4.font.color.rgb = SUBTLE; r4.font.name = SANS
cy = Inches(4.32); chh = Inches(0.96)
for ti, bd in [
    ("VOLATILITY ≠ RISK OF LOSS",
     "PE’s ~28% true vol is dominated by downside (levered, left-skewed). "
     "Athanase’s 27% is mostly upside — downside deviation only ~11%."),
    ("THE RISK THAT MATTERS",
     "On large permanent loss — what allocators actually fear — Athanase’s real "
     "risk is roughly one-eighth of private equity’s.")]:
    rect(s, bx, cy, bw, chh, fill=HEADERBG)
    para(tbox(s, Emu(int(bx) + int(Inches(0.25))), cy + Inches(0.12),
              Emu(int(bw) - int(Inches(0.5))), Inches(0.25)), ti, 10.5, SLATE,
         first=True, bold=True, after=0, track=0)
    para(tbox(s, Emu(int(bx) + int(Inches(0.25))), cy + Inches(0.38),
              Emu(int(bw) - int(Inches(0.5))), Inches(0.55)), bd, 10.5, BODY,
         first=True, after=0, lead=1.12, track=0)
    cy = Emu(int(cy) + int(chh) + int(Inches(0.12)))
rect(s, Inches(0.6), Inches(6.62), Inches(12.13), Inches(0.46), fill=NAVY)
para(tbox(s, Inches(0.78), Inches(6.62), Inches(11.8), Inches(0.46),
          anchor=MSO_ANCHOR.MIDDLE),
     "Matched on true volatility, Athanase carries far less of the risk that "
     "actually hurts: its downside is a fraction of private equity’s.",
     12, WHITE, first=True, italic=True, after=0, track=0)
para(tbox(s, Inches(0.6), Inches(7.12), Inches(12.6), Inches(0.5)),
     f"Athanase downside vol {ATH_DD*100:.0f}% vs total {ATH_VOL*100:.0f}% "
     f"(real); 30%+ loss in {ATH_P30*100:.0f}% of actual rolling 3-yr windows. "
     "PE de-smoothed total vol ~28%, ~15–16% chance of a 30% drawdown over 3 "
     "years (PIMCO, 2022). See references.",
     7.5, FOOT, first=True, after=0, track=0, lead=1.1)


# ---- II.6o  PE diversification is an artefact ------------------------------
def _pe_corr(path):
    fig_, ax = plt.subplots(figsize=(6.6, 4.5), dpi=200)
    vals = [PE_CORR_REPORTED, PE_CORR_TRUE, ATH_CORR]
    cols = [H_BLUE5, H_BLUE4, H_NAVY]; edge = [H_BLUE4, H_BLUE4, H_NAVY]
    x = np.arange(3)
    ax.bar(x, vals, 0.6, color=cols, edgecolor=edge, linewidth=1.2, zorder=3)
    for xi, v in enumerate(vals):
        ax.annotate(f"{v:.2f}", (xi, v), ha="center", va="bottom", fontsize=15,
                    fontweight="bold", color=(H_NAVY if xi == 2 else H_BLUE4),
                    xytext=(0, 3), textcoords="offset points")
    ax.axhspan(0.85, 1.0, color=H_BLUE4, alpha=0.08, zorder=0)
    ax.annotate("≈ no real diversification", (0.5, 0.92), ha="center",
                va="center", fontsize=9.5, color=H_BLUE4, style="italic")
    ax.set_xticks(x)
    ax.set_xticklabels(["PE\nas reported", "PE\nde-smoothed", "Athanase\n(real)"],
                       fontsize=11, color=H_BODY)
    ax.set_ylabel("Correlation with public equities", color=H_BODY, fontsize=11)
    ax.set_ylim(0, 1.0); ax.set_xlim(-0.6, 2.6)
    ax.tick_params(colors=H_BODY, labelsize=10)
    for sp in ("top", "right"):
        ax.spines[sp].set_visible(False)
    for sp in ("left", "bottom"):
        ax.spines[sp].set_color(H_BLUE5)
    ax.grid(True, axis="y", color=H_BLUE5, lw=0.6, alpha=0.6)
    ax.set_title("Correlation to public equities — reported vs true",
                 color=H_NAVY, fontsize=13, fontweight="bold", pad=10)
    fig_.tight_layout(); fig_.savefig(path, dpi=200, bbox_inches="tight",
                                      facecolor="white"); plt.close(fig_)


_pe_corr("/tmp/d_pe_corr.png")
s, top = content("Private Equity",
                 "PE’s “diversification” is an artefact — Athanase’s is real",
                 "De-smoothed, PE is ~0.9 correlated to public equities: the same "
                 "risk factors, just reported late. Athanase is genuinely "
                 "uncorrelated.", ref=fig())
s.shapes.add_picture("/tmp/d_pe_corr.png", Inches(0.55), Inches(2.3),
                     height=Inches(3.85))
cx = Inches(7.2); cw = Inches(5.55); chh = Inches(1.14); cy = Inches(2.35)
for ti, bd in [
    ("THE SMOOTHING ILLUSION",
     "Quarterly appraisals lag the market, so reported correlation looks low "
     "(0.75, beta 0.41). De-smoothed it is 0.89 / beta 0.87 (Two Sigma, 2024)."),
    ("SAME RISK FACTORS",
     "PE is long-only levered equity driven by GDP, rates and earnings; over "
     "10–20 yrs its correlation to small/mid-caps is ~0.89–0.92."),
    ("THE CRISIS MIRAGE",
     "In severe drawdowns true correlation converges toward 1.0 — PE takes the "
     "same hit, it just reports the markdown later.")]:
    rect(s, cx, cy, cw, chh, fill=HEADERBG)
    para(tbox(s, Emu(int(cx) + int(Inches(0.22))), cy + Inches(0.12),
              Emu(int(cw) - int(Inches(0.44))), Inches(0.25)), ti, 10.5, SLATE,
         first=True, bold=True, after=0, track=0)
    para(tbox(s, Emu(int(cx) + int(Inches(0.22))), cy + Inches(0.38),
              Emu(int(cw) - int(Inches(0.44))), Inches(0.72)), bd, 10.5, BODY,
         first=True, after=0, lead=1.12, track=0)
    cy = Emu(int(cy) + int(chh) + int(Inches(0.14)))
rect(s, Inches(0.6), Inches(6.62), Inches(12.13), Inches(0.46), fill=NAVY)
para(tbox(s, Inches(0.78), Inches(6.62), Inches(11.8), Inches(0.46),
          anchor=MSO_ANCHOR.MIDDLE),
     f"PE adds leveraged beta dressed as diversification. Athanase, at "
     f"{ATH_CORR:.2f} correlation, is the genuine diversifier in a public-equity "
     "portfolio.", 12, WHITE, first=True, italic=True, after=0, track=0)
para(tbox(s, Inches(0.6), Inches(7.12), Inches(12.6), Inches(0.5)),
     f"Athanase correlation {ATH_CORR:.2f} / beta 0.73 to MSCI World IMI (real). "
     "PE reported 0.75 / beta 0.41 de-smoothed to 0.89 / beta 0.87 (Two Sigma / "
     "Venn, 2024); long-horizon PE–small/mid-cap correlation ~0.89–0.92. See "
     "references.", 7.5, FOOT, first=True, after=0, track=0, lead=1.1)


# ---- II.6p  The drawdown trap: PE smoothing as career risk -----------------
s, top = content("Private Equity",
                 "When the tide goes out: the drawdown trap",
                 "PE’s smooth marks feel safe in calm markets — but in a shock "
                 "they trigger the exact sequence that gets a CIO fired. "
                 "Engaged ownership is the fiduciary safe harbour.")
# left column — the trap (three stages)
para(tbox(s, Inches(0.7), top, Inches(5.9), Inches(0.3)),
     "THE TRAP — HELD AT FAKE NAVs IN A CRASH", 11.5, SLATE, first=True,
     bold=True, after=0, track=0)
trap = [
    ("1 · The denominator effect",
     "Public equities fall 25%; PE is marked flat. The PE weight jumps from "
     "15% to 22%+ — an instant breach of board concentration limits."),
    ("2 · The liquidity death spiral",
     "Distributions stop but capital calls keep coming. To meet them, the "
     "allocator is forced to sell liquid equities at the bottom."),
    ("3 · The secondary-market humiliation",
     "To cure the breach, PE stakes are sold at 15–25% haircuts — publicly "
     "crystallising losses and proving the “low volatility” was a fiction."),
]
ty = top + Inches(0.42)
for t, b in trap:
    rect(s, Inches(0.7), ty, Inches(5.9), Inches(1.12), fill=HEADERBG)
    para(tbox(s, Inches(0.88), ty + Inches(0.13), Inches(5.55), Inches(0.3)),
         t, 12, NAVY_TX, first=True, bold=True, after=0, font=SERIF, track=0)
    para(tbox(s, Inches(0.88), ty + Inches(0.46), Inches(5.55), Inches(0.66)),
         b, 10.5, BODY, first=True, after=0, lead=1.12, track=0)
    ty = Emu(int(ty) + int(Inches(1.12)) + int(Inches(0.1)))
# right column — the safe harbour
para(tbox(s, Inches(6.95), top, Inches(5.85), Inches(0.3)),
     "THE SAFE HARBOUR — ENGAGED OWNERSHIP", 11.5, SLATE, first=True,
     bold=True, after=0, track=0)
safe = [
    ("Honest pricing keeps weights balanced",
     "Daily marks mean the sleeve falls with the market — so the denominator "
     "effect never fires and no threshold is breached."),
    ("Zero capital-call liability",
     "No blind pools, no surprise calls. The allocator is never forced to "
     "become a distressed seller of other assets."),
    ("Alpha from turnarounds, not leverage",
     "Returns come from board-led operational change in high-ROIIC mid-caps — "
     "not the leverage and multiple-expansion that break in a downturn."),
]
ty = top + Inches(0.42)
for t, b in safe:
    rect(s, Inches(6.95), ty, Inches(5.85), Inches(1.12), fill=HEADERBG)
    para(tbox(s, Inches(7.13), ty + Inches(0.13), Inches(5.5), Inches(0.3)),
         t, 12, NAVY, first=True, bold=True, after=0, font=SERIF, track=0)
    para(tbox(s, Inches(7.13), ty + Inches(0.46), Inches(5.5), Inches(0.66)),
         b, 10.5, BODY, first=True, after=0, lead=1.12, track=0)
    ty = Emu(int(ty) + int(Inches(1.12)) + int(Inches(0.1)))
# takeaway strip
rect(s, Inches(0.6), Inches(6.62), Inches(12.13), Inches(0.46), fill=NAVY)
para(tbox(s, Inches(0.78), Inches(6.62), Inches(11.8), Inches(0.46),
          anchor=MSO_ANCHOR.MIDDLE),
     "When the tide goes out, artificially-priced illiquid assets are "
     "indefensible to a board. Engaged ownership delivers PE-like alpha while "
     "protecting portfolio balance, liquidity — and the allocator’s seat.",
     12, WHITE, first=True, italic=True, after=0, track=0)
para(tbox(s, Inches(0.6), Inches(7.16), Inches(12.6), Inches(0.4)),
     "Mechanism described follows from PE’s appraisal-based marking and "
     "capital-call structure (see references). Illustrative; for professional "
     "investors, not investment advice.", 7.5, FOOT, first=True, after=0,
     track=0, lead=1.1)


# ---- II.7 ESG / responsible ownership ----
s, top = content("Ownership Model", "Responsible ownership, by construction")
checklist(s, [
    ("Governance is the lever.", "Value is created through board representation "
     "and active promotion of strong corporate governance across all holdings."),
    ("Five improvement dimensions.", "Operational, financial, structural, "
     "environmental and social improvements — each tied to enterprise value."),
    ("PRI signatory since 2019,", "adhering to the Six Principles for Responsible "
     "Investment."),
    ("Credible oversight.", "SEB custodian · MUFG administrator · KPMG auditor · "
     "Swedish FI-registered AIFM."),
], top, size=16, gap=14)

# ---- Objection handling: declines + the arsenal (not numbered / not in TOC) -
s, top = content("The Case",
                 "The objections we hear — and the arsenal to answer them",
                 "The common declines are rarely about the returns. Each "
                 "concern has a concrete, quantified rebuttal.", number=False)
data = [
    ("“It falls between our buckets.”", "Concern: tracking-error fear", [
        ("0.44 correlation, 43% down-capture", "the divergence is portfolio "
         "armour, not a risk — it stops you riding the index to the bottom."),
        ("Alpha from corporate events", "stripping growth traps, reallocating "
         "to ROIIC — measurable, not macro."),
        ("When the market tanks", "that deliberate divergence defends your "
         "short-term numbers.")]),
    ("“Key-person risk at a boutique.”", "Concern: political cover", [
        ("Segregation of duties", "operations, risk and finance sit fully "
         "independent of the investment team."),
        ("SEB · MUFG · KPMG", "custody, administration and audit verify every "
         "move — a heavyweight institutional fortress."),
        ("The same investment team since 2006", "a 20-year playbook — "
         "institutional memory that engineers out permanent loss, not "
         "key-person risk.")]),
    ("“Capacity — only one deal a year.”", "Concern: effort vs. sleeve size", [
        ("A 3–8% sleeve is not a rounding error", "it shifts the whole equity "
         "book’s frontier up-and-to-the-left — real basis points."),
        ("No fees on undeployed dry powder", "deploy once, capital works "
         "immediately."),
        ("No blind-pool capital calls", "none of the administration that comes "
         "with them.")]),
    ("“Duration ties up our risk budget.”", "Concern: comp-cycle mismatch", [
        ("Down months ~1.5% vs the market’s ~3.8%", "immediate downside "
         "preservation, this year."),
        ("−10% / −20% / −30% hard triggers", "cap drawdowns before they spiral "
         "and remove PM bias."),
        ("Worst-ever entry still +7% net", "a statistical safety blanket over "
         "the holding period.")]),
]
tl = Inches(0.55); c1 = Inches(3.5); c2 = Inches(8.3)
y = top
# header
for cx, cw_, htxt in [(tl, c1, "The objection"),
                      (Emu(int(tl) + int(c1)), c2,
                       "The arsenal — how we answer it")]:
    rect(s, cx, y, cw_, Inches(0.4), fill=NAVY if cw_ == c2 else SLATE)
    para(tbox(s, Emu(int(cx) + int(Inches(0.14))), y,
              Emu(int(cw_) - int(Inches(0.24))), Inches(0.4),
              anchor=MSO_ANCHOR.MIDDLE),
         htxt, 11.5, WHITE, first=True, bold=True, after=0, track=0)
y = Emu(int(y) + int(Inches(0.4)))
rh = Inches(1.0)
for ri, (quote, concern, bullets) in enumerate(data):
    alt = ri % 2 == 0
    # col 1
    rect(s, tl, y, c1, rh, fill=HEADERBG if alt else WHITE)
    t1 = tbox(s, Emu(int(tl) + int(Inches(0.16))), y,
              Emu(int(c1) - int(Inches(0.3))), rh, anchor=MSO_ANCHOR.MIDDLE)
    para(t1, quote, 11.5, NAVY_TX, first=True, bold=True, after=3, lead=1.05,
         font=SERIF, italic=True, track=0)
    para(t1, concern, 9, SUBTLE, italic=True, after=0, lead=1.0, track=0)
    # col 2 — arsenal bullets
    x2 = Emu(int(tl) + int(c1))
    rect(s, x2, y, c2, rh, fill=BLUE5 if alt else HEADERBG)
    t2 = tbox(s, Emu(int(x2) + int(Inches(0.18))), y,
              Emu(int(c2) - int(Inches(0.34))), rh, anchor=MSO_ANCHOR.MIDDLE)
    for bi, (lead, rest) in enumerate(bullets):
        p = t2.paragraphs[0] if bi == 0 else t2.add_paragraph()
        p.space_after = Pt(2.5); p.line_spacing = 1.04
        hang = int(round(9.5 * 1.25 * 12700))
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(hang)); pPr.set("indent", str(-hang))
        r0 = p.add_run(); r0.text = "›  "
        r0.font.size = Pt(9.5); r0.font.bold = True
        r0.font.color.rgb = SLATE_LT; r0.font.name = SANS
        r1 = p.add_run(); r1.text = lead
        r1.font.size = Pt(9.5); r1.font.bold = True
        r1.font.color.rgb = NAVY_TX; r1.font.name = SANS
        r2 = p.add_run(); r2.text = " — " + rest
        r2.font.size = Pt(9.5); r2.font.color.rgb = BODY; r2.font.name = SANS
    y = Emu(int(y) + int(rh))
para(tbox(s, Inches(0.55), Emu(int(y) + int(Inches(0.13))), Inches(12.3),
          Inches(0.4)),
     "We meet each objection where it sits — and answer it with structure, not "
     "optimism. The allocation is engineered to be defensible, not just to "
     "perform.", 12, SLATE, first=True, italic=True, after=0, track=0)


# ---- Capstone: the investment-committee pitch (not numbered / not in TOC) ---
s, top = content("The Case",
                 "The case for your investment committee",
                 "The argument in four moves — structure, engine, downside and "
                 "proof.", number=False)
quads = [
    ("1 · The mandate & structure", [
        ("PE-style operational value creation",
         "executed entirely in the public markets."),
        ("Daily transparency and liquidity, lower fees",
         "and no 10-year blind-pool lock-up."),
        ("Applied corporate governance",
         "board independence, capital discipline and pay-for-performance — "
         "fulfilling UN PRI fiduciary duty.")]),
    ("2 · The engine — quality-core constructivism", [
        ("Market-leading European mid-caps",
         "with highly profitable, durable cores — not broken turnarounds."),
        ("We enter at “shareholder exhaustion”",
         "the buyer of last resort after management diluted the core chasing "
         "prestige projects."),
        ("Secure the board seat",
         "strip the growth-trap adjacencies and reallocate capital to the core "
         "to maximise ROIIC.")]),
    ("3 · The moat — engineered downside protection", [
        ("A valuation floor",
         "the core alone justifies the price; zero value to weak divisions — "
         "a 30–40% margin of safety vs a PE bidder."),
        ("The board seat is the kill switch",
         "authority to freeze bad capital deployment the moment the macro "
         "shifts."),
        ("Mechanized discipline",
         "automatic −10% / −20% / −30% triggers force a thesis review and "
         "exit.")]),
    ("4 · The proof — two decades of compounding", [
        ("16.3% net p.a.; capital grown 18×",
         "across multiple market cycles."),
        ("92% deal-level hit rate (38 deals)",
         "a 5× MOIC roughly every two years."),
        ("93% up-capture vs 43% down-capture",
         "losing less than half preserves the compounding base.")]),
]
cw = Inches(5.86); gx = Inches(0.18); gy = Inches(0.1)
hh = Inches(0.46); bodyH = Inches(1.74)
cardH = int(hh) + int(bodyH)
for qi, (qt, bullets) in enumerate(quads):
    x = Inches(0.62) if qi % 2 == 0 else Emu(int(Inches(0.62)) + int(cw) + int(gx))
    y = top if qi < 2 else Emu(int(top) + cardH + int(gy))
    rect(s, x, y, cw, hh, fill=NAVY)
    para(tbox(s, Emu(int(x) + int(Inches(0.16))), y, Emu(int(cw) - int(Inches(0.3))),
              hh, anchor=MSO_ANCHOR.MIDDLE),
         qt, 12.5, WHITE, first=True, bold=True, after=0, font=SERIF, track=0)
    by = Emu(int(y) + int(hh))
    rect(s, x, by, cw, bodyH, fill=HEADERBG)
    tf = tbox(s, Emu(int(x) + int(Inches(0.18))), Emu(int(by) + int(Inches(0.12))),
              Emu(int(cw) - int(Inches(0.36))), Emu(int(bodyH) - int(Inches(0.2))))
    for bi, (lead, body) in enumerate(bullets):
        p = tf.paragraphs[0] if bi == 0 else tf.add_paragraph()
        p.space_after = Pt(4); p.line_spacing = 1.08
        hang = int(round(9 * 1.3 * 12700))
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(hang)); pPr.set("indent", str(-hang))
        r0 = p.add_run(); r0.text = "›  "
        r0.font.size = Pt(9); r0.font.bold = True
        r0.font.color.rgb = SLATE_LT; r0.font.name = SANS
        r1 = p.add_run(); r1.text = lead
        r1.font.size = Pt(9); r1.font.bold = True
        r1.font.color.rgb = NAVY_TX; r1.font.name = SANS
        r2 = p.add_run(); r2.text = " — " + body
        r2.font.size = Pt(9); r2.font.color.rgb = BODY; r2.font.name = SANS

# ---- Capstone 2: what protects the committee (not numbered / not in TOC) ----
s, top = content("The Case",
                 "What protects the committee — not just the capital",
                 "Beyond the returns, four things an investment committee can "
                 "point to — a due-diligence shield, a capital floor, "
                 "peer-reviewed cover and a liquidity escape hatch.",
                 number=False)
quads2 = [
    ("1 · Independent validation", [
        ("Independently validated",
         "the 20-year net record was reconciled from security-level StatPro "
         "data to actual LP-level returns — not a back-test or internal "
         "spreadsheet."),
        ("SEB · MUFG · KPMG",
         "custody, administration and audit by top-tier institutions verify "
         "every NAV and fee."),
        ("An air-tight paper trail",
         "no auditor or regulator can level an accusation of lax oversight.")]),
    ("2 · No investor has lost capital", [
        (f"{_ENTRY_AVG*100:.0f}% average net return",
         "regardless of the month an investor entered the fund."),
        (f"+{_ENTRY_WORST*100:.0f}% in the worst entry month",
         "the single unluckiest entry over 20 years was still positive, net."),
        ("Zero holding-period losses",
         "historically, no investor has lost money over a holding period.")]),
    ("3 · Backed by peer-reviewed science", [
        ("Settled in the literature",
         "13D alpha and PE’s volatility laundering / de-smoothing are backed by "
         "the Journal of Finance and Review of Financial Studies."),
        (f"{_PAIR['corr']:.2f} correlation, by the science",
         "the low correlation and risk/return profile are established, not "
         "asserted."),
        ("Cover in a bad quarter",
         "point to the evidence, not to hope — a fundamentally sound economic "
         "engine.")]),
    ("4 · A liquidity escape hatch", [
        ("Daily pricing and liquidity",
         "the strategy lives entirely in listed public equities."),
        ("No lock-up or capital-call queue",
         "none of PE’s 10-year blind-pool traps or blind extensions."),
        ("Total optionality",
         "if style drifts or the landscape changes, the committee can vote to "
         "pivot out seamlessly.")]),
]
cw = Inches(5.86); gx = Inches(0.18); gy = Inches(0.1)
hh = Inches(0.46); bodyH = Inches(1.74)
cardH = int(hh) + int(bodyH)
for qi, (qt, bullets) in enumerate(quads2):
    x = Inches(0.62) if qi % 2 == 0 else Emu(int(Inches(0.62)) + int(cw) + int(gx))
    y = top if qi < 2 else Emu(int(top) + cardH + int(gy))
    rect(s, x, y, cw, hh, fill=NAVY)
    para(tbox(s, Emu(int(x) + int(Inches(0.16))), y, Emu(int(cw) - int(Inches(0.3))),
              hh, anchor=MSO_ANCHOR.MIDDLE),
         qt, 12.5, WHITE, first=True, bold=True, after=0, font=SERIF, track=0)
    by = Emu(int(y) + int(hh))
    rect(s, x, by, cw, bodyH, fill=HEADERBG)
    tf = tbox(s, Emu(int(x) + int(Inches(0.18))), Emu(int(by) + int(Inches(0.12))),
              Emu(int(cw) - int(Inches(0.36))), Emu(int(bodyH) - int(Inches(0.2))))
    for bi, (lead, body) in enumerate(bullets):
        p = tf.paragraphs[0] if bi == 0 else tf.add_paragraph()
        p.space_after = Pt(4); p.line_spacing = 1.08
        hang = int(round(9 * 1.3 * 12700))
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(hang)); pPr.set("indent", str(-hang))
        r0 = p.add_run(); r0.text = "›  "
        r0.font.size = Pt(9); r0.font.bold = True
        r0.font.color.rgb = SLATE_LT; r0.font.name = SANS
        r1 = p.add_run(); r1.text = lead
        r1.font.size = Pt(9); r1.font.bold = True
        r1.font.color.rgb = NAVY_TX; r1.font.name = SANS
        r2 = p.add_run(); r2.text = " — " + body
        r2.font.size = Pt(9); r2.font.color.rgb = BODY; r2.font.name = SANS

# ===========================================================================
# CLOSING
# ===========================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, fill=NAVY)
wordmark(s, Inches(0.6), Inches(0.55), WHITE, scale=1.0)
rect(s, 0, Inches(2.7), Inches(0.5), Pt(2.2), fill=SLATE_LT)
tf = tbox(s, Inches(0.6), Inches(2.95), Inches(11.0), Inches(1.6))
para(tf, "The asset class rewards patient, governance-minded capital.", 26,
     WHITE, first=True, after=8, font=SERIF)
para(tf, "Among engaged owners, Athanase brings the team, the discipline and the "
     "20-year proof to deliver it.", 18, DIVIDER, italic=True, after=0, font=SERIF)
# contact
ct = tbox(s, Inches(0.6), Inches(5.4), Inches(6), Inches(1.4))
para(ct, "Stefan Charette", 15, WHITE, first=True, bold=True, after=2)
para(ct, "Athanase Industrial Partner", 12, DIVIDER, after=8)
para(ct, "charette@athanase.se   ·   +46 73 994 7079", 12, SLATE_LT, after=0)
ot = tbox(s, Inches(9.0), Inches(5.4), Inches(3.8), Inches(1.5))
para(ot, "Offices", 12, WHITE, first=True, bold=True, after=4)
para(ot, "Birger Jarlsgatan 6, 114 34 Stockholm, Sweden", 11, SLATE_LT, after=4)
para(ot, "Landmark Square, West Bay Road, Grand Cayman", 11, SLATE_LT, after=0)


# ===========================================================================
# REFERENCES — research-standard, full citations (back matter; not in TOC)
# ===========================================================================
def _ref_slide(title_suffix, part_label):
    sl = prs.slides.add_slide(BLANK)
    rect(sl, 0, 0, SW, SH, fill=WHITE)
    place_mark(sl, Inches(0.55), Inches(0.24), Inches(0.26))
    para(tbox(sl, Inches(0.98), Inches(0.27), Inches(7), Inches(0.3)),
         "References", 11, SLATE_LT, first=True, after=0)
    para(tbox(sl, Inches(9.4), Inches(0.27), Inches(3.35), Inches(0.3)),
         part_label, 11, SLATE_LT, first=True, bold=True, align=PP_ALIGN.RIGHT,
         after=0)
    rect(sl, 0, Inches(0.62), SW, Inches(0.92), fill=HEADERBG)
    para(tbox(sl, Inches(0.6), Inches(0.72), Inches(12.1), Inches(0.7)),
         "References & Methodology" + title_suffix, 26, NAVY_TX, first=True,
         after=0, font=SERIF)
    footer(sl)
    return sl


def _refhead(sl, x, y, w, text):
    para(tbox(sl, x, y, w, Inches(0.3)), text, 11.5, SLATE, first=True,
         bold=True, after=0, track=0)


def _refs(sl, x, y, w, items, size=9):
    tf_ = tbox(sl, x, y, w, Inches(5.0))
    for i, parts in enumerate(items):
        p = tf_.paragraphs[0] if i == 0 else tf_.add_paragraph()
        p.space_after = Pt(5.5); p.line_spacing = 1.04
        pPr = p._p.get_or_add_pPr()
        hang = int(Inches(0.26))
        pPr.set("marL", str(hang)); pPr.set("indent", str(-hang))
        for txt, ital in parts:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.italic = ital
            r.font.name = SANS; r.font.color.rgb = BODY


r1 = _ref_slide(" — Peer-reviewed literature", "1 of 2")
_refhead(r1, Inches(0.6), Inches(1.75), Inches(6.0),
         "ENGAGED OWNERSHIP / ACTIVISM (PART I)")
_refs(r1, Inches(0.6), Inches(2.12), Inches(6.0), [
    [("Bebchuk, L. A., Brav, A., & Jiang, W. (2015). The long-term effects of "
      "hedge fund activism. ", False), ("Columbia Law Review, 115", True),
     ("(5), 1085–1156.", False)],
    [("Brav, A., Jiang, W., Partnoy, F., & Thomas, R. (2008). Hedge fund "
      "activism, corporate governance, and firm performance. ", False),
     ("The Journal of Finance, 63", True),
     ("(4), 1729–1775. https://doi.org/10.1111/j.1540-6261.2008.01373.x", False)],
    [("Brav, A., Jiang, W., & Kim, H. (2015). The real effects of hedge fund "
      "activism: Productivity, asset allocation, and labor outcomes. ", False),
     ("The Review of Financial Studies, 28", True),
     ("(10), 2723–2769. https://doi.org/10.1093/rfs/hhv037", False)],
    [("Klein, A., & Zur, E. (2009). Entrepreneurial shareholder activism: Hedge "
      "funds and other private investors. ", False),
     ("The Journal of Finance, 64", True),
     ("(1), 187–229. https://doi.org/10.1111/j.1540-6261.2008.01432.x", False)],
])
_refhead(r1, Inches(0.6), Inches(4.28), Inches(6.0),
         "PASSIVE INVESTING & MARKET STRUCTURE (PART I)")
_refs(r1, Inches(0.6), Inches(4.55), Inches(6.0), [
    [("Ben-David, I., Franzoni, F., & Moussawi, R. (2018). Do ETFs increase "
      "volatility? ", False), ("The Journal of Finance, 73", True),
     ("(6), 2471–2535. https://doi.org/10.1111/jofi.12727", False)],
    [("Bond, P., & García, D. (2022). The equilibrium consequences of "
      "indexing. ", False), ("The Review of Financial Studies, 35", True),
     ("(7), 3175–3230.", False)],
    [("Chen, H., Noronha, G., & Singal, V. (2006). Index changes and losses to "
      "index fund investors. ", False),
     ("Financial Analysts Journal, 62", True),
     ("(4), 31–47. https://doi.org/10.2469/faj.v62.n4.4185", False)],
    [("Chinco, A., & Sammon, M. (2024). The passive-ownership share is double "
      "what you think. ", False), ("SSRN Electronic Journal", True),
     (".", False)],
    [("Gabaix, X., & Koijen, R. S. J. (2021). In search of the origins of "
      "financial fluctuations: The inelastic markets hypothesis. ", False),
     ("NBER Working Paper No. 28967", True),
     (". https://doi.org/10.3386/w28967", False)],
    [("Petajisto, A. (2011). The index premium and its hidden cost for index "
      "funds. ", False), ("Journal of Empirical Finance, 18", True),
     ("(2), 271–288. https://doi.org/10.1016/j.jempfin.2010.10.002", False)],
    [("Wurgler, J. (2011). On the economic consequences of index-linked "
      "investing. ", False), ("NBER Working Paper No. 16376", True),
     (". https://doi.org/10.3386/w16376", False)],
], size=8.5)
_refhead(r1, Inches(6.95), Inches(1.75), Inches(5.85),
         "PRIVATE EQUITY RISK & RETURN (PART II)")
_refs(r1, Inches(6.95), Inches(2.12), Inches(5.85), [
    [("Ang, A., Chen, B., Goetzmann, W. N., & Phalippou, L. (2018). Estimating "
      "private equity returns from limited partner cash flows. ", False),
     ("The Journal of Finance, 73", True),
     ("(4), 1751–1783. https://doi.org/10.1111/jofi.12688", False)],
    [("Boyer, B. H., Nadauld, T. D., Vorkink, K. P., & Weisbach, M. S. (2018). "
      "Private equity indices based on secondary market transactions. ", False),
     ("SSRN Electronic Journal", True),
     (". https://doi.org/10.2139/ssrn.3272357", False)],
    [("Buchner, A., Kaserer, C., & Wagner, N. F. (2010). Private equity funds: "
      "Valuation, systematic risk and illiquidity. ", False),
     ("SSRN Electronic Journal", True),
     (". https://doi.org/10.2139/ssrn.1102471", False)],
    [("Couts, S., Gonçalves, A. S., & Rossi, A. (2020). Unsmoothing returns of "
      "illiquid funds. ", False), ("SSRN Electronic Journal", True),
     (". https://doi.org/10.2139/ssrn.3544854", False)],
    [("Getmansky, M., Lo, A. W., & Makarov, I. (2004). An econometric model of "
      "serial correlation and illiquidity in hedge fund returns. ", False),
     ("Journal of Financial Economics, 74", True),
     ("(3), 529–609. https://doi.org/10.1016/j.jfineco.2004.04.001", False)],
    [("Harris, R. S., Jenkinson, T., & Kaplan, S. N. (2014). Private equity "
      "performance: What do we know? ", False),
     ("The Journal of Finance, 69", True),
     ("(5), 1851–1882. https://doi.org/10.1111/jofi.12154", False)],
    [("Hayley, S., & Sefiloglu, O. (2022). Biases in private equity returns. ",
      False), ("SSRN Electronic Journal", True),
     (". https://doi.org/10.2139/ssrn.4245715", False)],
    [("Jackson, B., Ling, D. C., & Naranjo, A. (2022). Catering and return "
      "manipulation in private equity. ", False),
     ("SSRN Electronic Journal", True),
     (". https://doi.org/10.2139/ssrn.4244467", False)],
    [("Meyer, T. (2020). Hidden in plain sight — the impact of undrawn "
      "commitments. ", False),
     ("The Journal of Alternative Investments, 23", True),
     ("(2), 94–110. https://doi.org/10.3905/jai.2020.1.101", False)],
    [("Phalippou, L. (2020). An inconvenient fact: Private equity returns and "
      "the billionaire factory. ", False),
     ("The Journal of Investing, 30", True),
     ("(1), 11–39. https://doi.org/10.3905/joi.2020.1.153", False)],
], size=8.5)

r2 = _ref_slide(" — Datasets, practitioner research & methodology", "2 of 2")
_refhead(r2, Inches(0.6), Inches(1.75), Inches(6.0),
         "INDUSTRY DATASETS & PRACTITIONER RESEARCH")
_refs(r2, Inches(0.6), Inches(2.12), Inches(6.0), [
    [("Asness, C. S. (2023). The volatility-laundering hidden in private-asset "
      "marks. ", False), ("AQR Capital Management — Cliff’s Perspectives.", True)],
    [("British Private Equity & Venture Capital Association & Capital Dynamics. "
      "(2024). ", False),
     ("BVCA Private Equity Performance Measurement Survey (PME+).", True),
     (" London: BVCA.", False)],
    [("Cambridge Associates. (2024). ", False),
     ("US Private Equity Index and Selected Benchmark Statistics (Q2 2024).",
      True), (" Boston: Cambridge Associates LLC.", False)],
    [("Cliffwater LLC. (2023). ", False),
     ("Long-term private equity performance in US state pension plans "
      "(2000–2022).", True), (" Marina del Rey: Cliffwater LLC.", False)],
    [("Financial Reporting Council. (2020). ", False),
     ("The UK Stewardship Code 2020.", True), (" London: FRC.", False)],
    [("Morningstar. (2023). ", False),
     ("Volatility laundering: How private equity funds understate the risk of "
      "their investments.", True), (" Chicago: Morningstar, Inc.", False)],
    [("PIMCO. (2022). ", False),
     ("The discreet charm of private assets: De-smoothing and the true "
      "volatility of private equity.", True), (" Newport Beach: PIMCO.", False)],
    [("Principles for Responsible Investment. (2021). ", False),
     ("What are the Principles for Responsible Investment?", True),
     (" London: UN PRI.", False)],
    [("Two Sigma Investments. (2024). ", False),
     ("The alternative truth of private equity (Venn by Two Sigma).", True),
     (" New York: Two Sigma.", False)],
], size=8.5)
_refhead(r2, Inches(6.95), Inches(1.75), Inches(5.85), "DATA & METHODOLOGY")
_md = tbox(r2, Inches(6.95), Inches(2.12), Inches(5.85), Inches(4.6))
para(_md, "All Athanase figures are computed from the fund’s audited net "
     "monthly return series across two vehicles (Öresund/Creades Jan 2006–"
     "Jun 2014; AIPFII from 23 Feb 2015):", 9.5, BODY,
     first=True, after=6, lead=1.15, track=0)
for _mt in [
    "Annualised return — cumulative net return compounded, then annualised over "
    "actual invested time (19.32 yrs, excluding the 2014–15 setup gap): 16.3%.",
    "Total volatility — st. dev. of monthly returns × √12 (27.0%).",
    "Downside volatility — deviation of negative months below a 0% MAR × √12 "
    "(10.9%).",
    "Correlation / beta — versus MSCI World IMI over the common period "
    "(0.44 / 0.73).",
    "Large-loss frequency — share of rolling 36-month windows below −30% (≈2%).",
    "Up/down capture — compounded Athanase return in the market’s up/down "
    "months ÷ the market’s.",
    "Efficient frontier — monthly-rebalanced MSCI/Athanase blends; minimum-risk "
    "point is in-sample.",
]:
    pp = _md.add_paragraph(); pp.space_after = Pt(4); pp.line_spacing = 1.12
    pPr = pp._p.get_or_add_pPr()
    pPr.set("marL", str(int(Inches(0.16))))
    pPr.set("indent", str(-int(Inches(0.16))))
    r = pp.add_run(); r.text = "•  " + _mt
    r.font.size = Pt(9); r.font.name = SANS; r.font.color.rgb = BODY
para(_md, "Private-equity figures are taken from the cited research. "
     "De-smoothed / true-economic / committed-capital values reflect that "
     "literature rather than reported (appraisal-based) data — including the "
     "de-smoothed large-buyout correlation to small/mid-cap equities of ~0.89 "
     "(Two Sigma / Venn, 2024; Couts, Gonçalves & Rossi, 2020; Boyer, Nadauld, "
     "Vorkink & Weisbach, 2018). Bridge and illustrative magnitudes reconcile "
     "cited endpoints and are labelled as such. For professional investors; "
     "not investment advice. Past performance is not indicative of future "
     "results.", 8.5, SUBTLE, italic=True, after=0, lead=1.15, track=0)


# ===========================================================================
# Contents (TOC) page
# ===========================================================================
def _toc_col(col_x, col_w, header, entries, groups=None):
    """entries: list of (num, title, ref). groups: optional dict mapping the
    leading sub-section index (str) -> sub-section name, inserted as headers."""
    t0 = Emu(int(agenda_top) - int(Inches(0.2)))
    htf = tbox(agenda_s, col_x, t0, col_w, Inches(0.4))
    para(htf, header, 14, SLATE, first=True, bold=True, after=0, font=SERIF)
    ytf = tbox(agenda_s, col_x, Emu(int(t0) + int(Inches(0.4))),
               col_w, Inches(5.3))
    n_lines = len(entries) + (len(groups) if groups else 0)
    dense = n_lines > 30
    long_col = len(entries) > 12
    fs = 8.0 if dense else (10 if long_col else 11.5)
    sa = 0.6 if dense else (3 if long_col else 7)
    lh = 1.0 if dense else 1.02
    first = True
    seen = set()
    for num, title, ref in entries:
        # sub-section header when the leading group index changes
        if groups is not None:
            gkey = num.split(".")[0]
            if gkey not in seen:
                seen.add(gkey)
                p = ytf.paragraphs[0] if first else ytf.add_paragraph()
                first = False
                p.space_before = Pt(0 if len(seen) == 1 else (1.5 if dense else 5))
                p.space_after = Pt(1 if dense else 2.5); p.line_spacing = 1.0
                gr = p.add_run()
                gr.text = f"{gkey}   {groups.get(gkey, '')}"
                gr.font.size = Pt(fs + 0.5); gr.font.bold = True
                gr.font.color.rgb = SLATE; gr.font.name = SERIF
        p = ytf.paragraphs[0] if first else ytf.add_paragraph()
        first = False
        p.space_after = Pt(sa); p.line_spacing = lh
        indent = "      " if groups is not None else ""
        r0 = p.add_run(); r0.text = indent + num + "   "
        r0.font.size = Pt(fs); r0.font.bold = (groups is None)
        r0.font.color.rgb = SLATE; r0.font.name = SANS
        r1 = p.add_run(); r1.text = title
        r1.font.size = Pt(fs); r1.font.color.rgb = NAVY_TX; r1.font.name = SANS
        if ref:
            r2 = p.add_run(); r2.text = "   (" + ref + ")"
            r2.font.size = Pt(fs - 2); r2.font.italic = True
            r2.font.color.rgb = SUBTLE; r2.font.name = SANS


_p1 = [(n, t, r) for (pt, n, t, r) in TOC if pt == 1]
_p2 = [(n, t, r) for (pt, n, t, r) in TOC if pt == 2]
_p2_groups = {str(i + 1): name for i, name in enumerate(_SUBSECTIONS)}
_cw = Inches(5.85); _xL = Inches(0.7)
_xR = Emu(int(_xL) + int(_cw) + int(Inches(0.5)))
_toc_col(_xL, _cw, "Part I  ·  Why allocate to engaged ownership", _p1)
_toc_col(_xR, _cw, "Part II  ·  Why choose Athanase", _p2, groups=_p2_groups)

# ===========================================================================
# Convert the whole deck from 16:9 to the brand's 4:3 page format
# (10.667" x 8.0" = the official AIP-deck.pptx page size).  Layout was authored
# in 16:9 coordinates, so we rescale every shape's geometry and font.
#   - geometry: x by sx, y by sy (full-bleed rects keep filling the page)
#   - fonts:    by sx (horizontal is the binding constraint -> wrapping preserved)
#   - pictures: uniform sx (preserve logo aspect ratio)
# ===========================================================================
from pptx.oxml.ns import qn as _qn
from pptx.enum.shapes import MSO_SHAPE_TYPE as _MST

TARGET_W = 9753600   # EMU, 10.667"
TARGET_H = 7315200   # EMU, 8.0"


def _rescale_chart_fonts(gframe, f):
    el = gframe.chart.part._element
    for rpr in el.iter():
        tag = rpr.tag.split("}")[-1]
        if tag in ("rPr", "defRPr") and rpr.get("sz") is not None:
            rpr.set("sz", str(max(100, int(round(int(rpr.get("sz")) * f)))))


def _rescale_shape(sh, sx, sy):
    # geometry
    try:
        L, T = sh.left, sh.top
        W, H = sh.width, sh.height
    except Exception:
        L = T = W = H = None
    is_pic = sh.shape_type == _MST.PICTURE
    if None not in (L, T, W, H):
        sh.left = int(L * sx)
        sh.top = int(T * sy)
        if is_pic:                      # uniform scale -> keep aspect ratio
            sh.width = int(W * sx)
            sh.height = int(H * sx)
        else:
            sh.width = int(W * sx)
            sh.height = int(H * sy)
    # fonts (scale by sx so text/box ratio is preserved horizontally)
    if sh.has_text_frame:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size is not None:
                    r.font.size = Pt(round(r.font.size.pt * sx, 1))
            # scale hanging-indent (marL/indent) by sx to match shrunk fonts
            pPr = p._p.find(_qn("a:pPr"))
            if pPr is not None:
                for attr in ("marL", "indent"):
                    v = pPr.get(attr)
                    if v is not None:
                        pPr.set(attr, str(int(round(int(v) * sx))))
    if sh.has_chart:
        _rescale_chart_fonts(sh, sx)


_sx = TARGET_W / int(prs.slide_width)
_sy = TARGET_H / int(prs.slide_height)
for _slide in prs.slides:
    for _sh in _slide.shapes:
        _rescale_shape(_sh, _sx, _sy)
prs.slide_width = TARGET_W
prs.slide_height = TARGET_H

# ---------------------------------------------------------------------------
# Brand the theme itself: replace the default Office colour scheme with the AIP
# Blue group and set theme fonts to Times New Roman / Arial, so no off-brand
# swatch survives anywhere in the file (guideline: stay within the palette).
# ---------------------------------------------------------------------------
def _brand_theme(prs):
    from lxml import etree
    a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    THEME_REL = ("http://schemas.openxmlformats.org/officeDocument/2006/"
                 "relationships/theme")
    scheme = {
        "dk1": "152130", "lt1": "FFFFFF", "dk2": "0E1620", "lt2": "F6F7F9",
        "accent1": "152130", "accent2": "314359", "accent3": "556A83",
        "accent4": "0E1620", "accent5": "E2E5E9", "accent6": "F6F7F9",
        "hlink": "314359", "folHlink": "556A83",
    }
    seen = set()
    for master in prs.slide_masters:
        theme = master.part.part_related_by(THEME_REL)
        if theme.partname in seen:
            continue
        seen.add(theme.partname)
        root = etree.fromstring(theme.blob)
        clr = root.find(f"{a}themeElements/{a}clrScheme")
        for name, hexv in scheme.items():
            el = clr.find(f"{a}{name}")
            if el is None:
                continue
            srgb = el.find(f"{a}srgbClr")
            sysc = el.find(f"{a}sysClr")
            if srgb is not None:
                srgb.set("val", hexv)
            elif sysc is not None:           # convert sysClr -> srgbClr
                el.remove(sysc)
                etree.SubElement(el, f"{a}srgbClr", {"val": hexv})
        fonts = root.find(f"{a}themeElements/{a}fontScheme")
        for grp, face in (("majorFont", "Times New Roman"), ("minorFont", "Arial")):
            latin = fonts.find(f"{a}{grp}/{a}latin")
            if latin is not None:
                latin.set("typeface", face)
        theme._blob = etree.tostring(root, xml_declaration=True,
                                     encoding="UTF-8", standalone=True)
    return prs


_brand_theme(prs)

out = "Athanase_Engaged_Ownership_Allocator_Deck.pptx"
prs.save(out)
print("Saved", out, "with", len(prs.slides._sldIdLst),
      "slides at 4:3 (%.2f x %.2f in)" % (TARGET_W / 914400, TARGET_H / 914400))

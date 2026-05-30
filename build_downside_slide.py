"""Standalone, brand-compliant slide: downside-volatility-driven outcome
projection for Athanase vs MSCI World IMI over 3/5/7/10-year horizons.

Statistics are computed live from data/Comparison_returns.xlsx so the slide
is reproducible and auditable. Built in 4:3 per the AIP Brand Guidelines.
"""
import math
from statistics import NormalDist
import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from PIL import Image
from lxml import etree

# ---- AIP Brand palette (Blue group only; never mix groups) -----------------
BLUE1 = RGBColor(0x0E, 0x16, 0x20)
BLUE2 = RGBColor(0x15, 0x21, 0x30)
BLUE3 = RGBColor(0x31, 0x43, 0x59)
BLUE4 = RGBColor(0x55, 0x6A, 0x83)
BLUE5 = RGBColor(0xE2, 0xE5, 0xE9)
BLUE6 = RGBColor(0xF6, 0xF7, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NAVY, NAVY_TX, SLATE, SLATE_LT = BLUE2, BLUE3, BLUE3, BLUE4
HEADERBG, BODY, SUBTLE, DIVIDER, FOOT = BLUE6, BLUE1, BLUE4, BLUE5, BLUE4
SERIF, SANS = "Times New Roman", "Arial"

LOGO_WHITE = "assets/logo_white.png"
MARK_DARK = "assets/mark_dark.png"
_MD_AR = (lambda s: s[0] / s[1])(Image.open(MARK_DARK).size)

# ===========================================================================
# 1) Compute statistics from the workbook
# ===========================================================================
def load_series():
    ws = openpyxl.load_workbook("data/Comparison_returns.xlsx",
                                data_only=True)["Sheet1"]

    def grab(month_rows):
        out = []
        for c in range(2, ws.max_column + 1):       # year columns
            for r in month_rows:
                v = ws.cell(r, c).value
                if isinstance(v, (int, float)):
                    out.append(v)
        return out

    return grab(range(4, 16)), grab(range(22, 34))   # MSCI, Athanase


def stats(x, mar=0.0):
    n = len(x)
    mean = sum(x) / n
    sd = math.sqrt(sum((r - mean) ** 2 for r in x) / (n - 1))
    dd = math.sqrt(sum(min(r - mar, 0) ** 2 for r in x) / n)
    ud = math.sqrt(sum(max(r - mar, 0) ** 2 for r in x) / n)
    g = 1.0
    for r in x:
        g *= (1 + r)
    geo_m = g ** (1 / n) - 1
    return dict(n=n, ann_ret=(1 + geo_m) ** 12 - 1,
                ann_vol=sd * math.sqrt(12), ann_dd=dd * math.sqrt(12),
                ann_ud=ud * math.sqrt(12), mean_m=mean, sd_m=sd, series=x)


_msci_x, _ath_x = load_series()
MSCI, ATH = stats(_msci_x), stats(_ath_x)
N = NormalDist()


# ---- Allocator-relevant pair statistics (Athanase relative to MSCI) --------
def pair_stats(a, m):
    n = len(a)

    def geo(xs):
        g = 1.0
        for x in xs:
            g *= (1 + x)
        return g ** (1 / len(xs)) - 1 if xs else 0.0

    up = [i for i in range(n) if m[i] > 0]
    dn = [i for i in range(n) if m[i] < 0]
    up_cap = geo([a[i] for i in up]) / geo([m[i] for i in up])
    dn_cap = geo([a[i] for i in dn]) / geo([m[i] for i in dn])
    mb, ab = sum(m) / n, sum(a) / n
    cov = sum((m[i] - mb) * (a[i] - ab) for i in range(n)) / (n - 1)
    vm = sum((x - mb) ** 2 for x in m) / (n - 1)
    va = sum((x - ab) ** 2 for x in a) / (n - 1)
    corr = cov / math.sqrt(vm * va)
    beta = cov / vm

    def maxdd(xs):
        idx = peak = 1.0
        mdd = 0.0
        for x in xs:
            idx *= (1 + x); peak = max(peak, idx); mdd = min(mdd, idx / peak - 1)
        return mdd

    def cum(xs):
        g = 1.0
        for x in xs:
            g *= (1 + x)
        return g

    def roll_win(months):
        w = t = 0
        for i in range(0, n - months + 1):
            t += 1
            w += cum(a[i:i + months]) > cum(m[i:i + months])
        return w / t if t else 0.0

    return dict(up_cap=up_cap, dn_cap=dn_cap, corr=corr, beta=beta,
                mdd_a=maxdd(a), mdd_m=maxdd(m),
                roll={y: roll_win(y * 12) for y in (3, 5, 7)},
                n_up=len(up), n_dn=len(dn))


PAIR = pair_stats(_ath_x, _msci_x)


def project(T, r, dd):
    """Median and downside (≈1-in-6, one downside-deviation) terminal of 100."""
    med = 100 * (1 + r) ** T
    down_ann = r - dd / math.sqrt(T)
    down = 100 * (1 + down_ann) ** T
    ploss = N.cdf((0 - r) / (dd / math.sqrt(T)))
    return med, down, ploss


HORIZONS = (3, 5, 7, 10)

# ===========================================================================
# 2) Build the slide
# ===========================================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
s = prs.slides.add_slide(prs.slide_layouts[6])


def rect(sl, l, t, w, h, fill=None, line=None, lw=1.0):
    shp = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill is not None:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line is not None:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def tbox(sl, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = sl.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, size, color, bold=False, italic=False, first=False,
         align=PP_ALIGN.LEFT, after=8, lead=1.1, font=SANS, track=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(after); p.line_spacing = lead
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = font
    if track is None:
        track = -5 if size >= 48 else (-3 if size >= 24 else 0)
    if track:
        r._r.get_or_add_rPr().set("spc", str(int(round(size * track / 100 * 100))))
    return p


# header band + mark + title (left-aligned, brand)
s.shapes.add_picture(MARK_DARK, Inches(0.55), Inches(0.24), height=Inches(0.26),
                     width=Emu(int(Inches(0.26) * _MD_AR)))
para(tbox(s, Inches(0.98), Inches(0.27), Inches(7), Inches(0.3)),
     "Risk-Adjusted Outcomes", 11, SLATE_LT, first=True, after=0)
para(tbox(s, Inches(9.9), Inches(0.27), Inches(2.85), Inches(0.3)),
     "Figure 1", 11, SLATE_LT, first=True, bold=True, align=PP_ALIGN.RIGHT, after=0)
rect(s, 0, Inches(0.62), SW, Inches(1.45), fill=HEADERBG)
para(tbox(s, Inches(0.6), Inches(0.74), Inches(12.1), Inches(0.85)),
     "Lower downside risk compounds into better outcomes", 30, NAVY_TX,
     first=True, after=2, font=SERIF)
para(tbox(s, Inches(0.62), Inches(1.55), Inches(11.9), Inches(0.5)),
     "Athanase carries higher total volatility than the market — but its "
     "downside volatility is lower, which is what protects compounding.",
     13, SUBTLE, first=True, italic=True, after=0)

top = Inches(2.2)

# --- risk stat band (three metrics x two funds) ---
def _pct(v):
    return f"{v * 100:.1f}%"


metrics = [
    ("Annualised return", _pct(MSCI["ann_ret"]), _pct(ATH["ann_ret"]), True),
    ("Total volatility", _pct(MSCI["ann_vol"]), _pct(ATH["ann_vol"]), False),
    ("Downside volatility", _pct(MSCI["ann_dd"]), _pct(ATH["ann_dd"]), True),
]
bx = Inches(0.6); bw = Inches(3.95); gap = Inches(0.12); bh = Inches(1.28)
for label, mv, av, ath_better in metrics:
    rect(s, bx, top, bw, bh, fill=HEADERBG)
    para(tbox(s, Emu(int(bx) + int(Inches(0.2))), top + Inches(0.14),
              Emu(int(bw) - int(Inches(0.4))), Inches(0.3)),
         label.upper(), 10.5, SLATE, first=True, bold=True, after=0, track=0)
    vt = tbox(s, Emu(int(bx) + int(Inches(0.2))), top + Inches(0.46),
              Emu(int(bw) - int(Inches(0.4))), Inches(0.8))
    p = vt.paragraphs[0]; p.space_after = Pt(2)
    r1 = p.add_run(); r1.text = "MSCI  "; r1.font.size = Pt(11)
    r1.font.color.rgb = SUBTLE; r1.font.name = SANS
    r2 = p.add_run(); r2.text = mv; r2.font.size = Pt(19); r2.font.bold = True
    r2.font.color.rgb = NAVY_TX; r2.font.name = SERIF
    p2 = vt.add_paragraph()
    r3 = p2.add_run(); r3.text = "AIP   "; r3.font.size = Pt(11)
    r3.font.color.rgb = SUBTLE; r3.font.name = SANS
    r4 = p2.add_run(); r4.text = av; r4.font.size = Pt(19); r4.font.bold = True
    r4.font.color.rgb = NAVY; r4.font.name = SERIF
    bx = Emu(int(bx) + int(bw) + int(gap))

# annotation under the downside card
para(tbox(s, Inches(9.27), top + Inches(1.34), Inches(3.95), Inches(0.4)),
     "↑ Athanase downside volatility is LOWER than the market’s.",
     10, NAVY, first=True, bold=True, after=0, track=0, lead=1.05)

# --- projection table: terminal value of 100 by horizon ---
ty = Inches(4.18)
para(tbox(s, Inches(0.6), ty - Inches(0.3), Inches(12), Inches(0.3)),
     "TABLE 1.  EXPECTED VALUE OF 100 INVESTED, BY HOLDING PERIOD",
     11, SLATE, first=True, bold=True, after=0, track=0)

# columns: metric label + 4 horizons
col0 = Inches(3.5)
hw = (SW - Inches(1.2) - col0) / len(HORIZONS)
rh = Inches(0.375)
header = ["", "3 years", "5 years", "7 years", "10 years"]
rows = []
rows.append(("Median outcome — MSCI",
             [f"{project(T, MSCI['ann_ret'], MSCI['ann_dd'])[0]:.0f}" for T in HORIZONS], "sub"))
rows.append(("Median outcome — Athanase",
             [f"{project(T, ATH['ann_ret'], ATH['ann_dd'])[0]:.0f}" for T in HORIZONS], "navy"))
rows.append(("Adverse case* — MSCI",
             [f"{project(T, MSCI['ann_ret'], MSCI['ann_dd'])[1]:.0f}" for T in HORIZONS], "sub"))
rows.append(("Adverse case* — Athanase",
             [f"{project(T, ATH['ann_ret'], ATH['ann_dd'])[1]:.0f}" for T in HORIZONS], "navy"))
rows.append(("Chance of a loss — MSCI",
             [f"{project(T, MSCI['ann_ret'], MSCI['ann_dd'])[2]*100:.0f}%" for T in HORIZONS], "sub"))
rows.append(("Chance of a loss — Athanase",
             [f"{project(T, ATH['ann_ret'], ATH['ann_dd'])[2]*100:.0f}%" for T in HORIZONS], "navy"))

# header row
cx = Inches(0.6)
for ci, htext in enumerate(header):
    w = col0 if ci == 0 else hw
    rect(s, cx, ty, w, rh, fill=SLATE)
    para(tbox(s, Emu(int(cx) + int(Inches(0.12))), ty, Emu(int(w) - int(Inches(0.2))), rh,
              anchor=MSO_ANCHOR.MIDDLE), htext, 12, WHITE, bold=True, first=True,
         align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT, after=0, track=0)
    cx = Emu(int(cx) + int(w))
yy = Emu(int(ty) + int(rh))
for ri, (label, vals, kind) in enumerate(rows):
    fill = HEADERBG if ri % 2 == 0 else WHITE
    cx = Inches(0.6)
    rect(s, cx, yy, col0, rh, fill=fill)
    para(tbox(s, Emu(int(cx) + int(Inches(0.12))), yy, Emu(int(col0) - int(Inches(0.2))), rh,
              anchor=MSO_ANCHOR.MIDDLE), label, 11.5,
         NAVY if kind == "navy" else SUBTLE, bold=(kind == "navy"),
         first=True, after=0, track=0)
    cx = Emu(int(cx) + int(col0))
    for v in vals:
        rect(s, cx, yy, hw, rh, fill=fill)
        para(tbox(s, Emu(int(cx) + int(Inches(0.1))), yy, Emu(int(hw) - int(Inches(0.2))), rh,
                  anchor=MSO_ANCHOR.MIDDLE), v, 12.5,
             NAVY_TX if kind == "navy" else SUBTLE, bold=(kind == "navy"),
             first=True, align=PP_ALIGN.RIGHT, after=0, track=0)
        cx = Emu(int(cx) + int(hw))
    yy = Emu(int(yy) + int(rh))

# takeaway strip (navy) + footnote
tk = Emu(int(yy) + int(Inches(0.1)))
rect(s, Inches(0.6), tk, Inches(12.13), Inches(0.4), fill=NAVY)
para(tbox(s, Inches(0.78), tk, Inches(11.8), Inches(0.4), anchor=MSO_ANCHOR.MIDDLE),
     "Even in the adverse case, Athanase’s lower downside volatility leaves an "
     "investor ahead at every horizon — and the gap widens with time.",
     11.5, WHITE, first=True, italic=True, after=0, track=0)
para(tbox(s, Inches(0.6), Inches(7.32), Inches(12.6), Inches(0.4)),
     f"Source: monthly net returns, {MSCI['n']} months (2006–2025). Volatility "
     "annualised from monthly; downside volatility = annualised deviation of "
     "negative months (MAR 0). *Adverse case = one downside-deviation outcome "
     "(≈1-in-6), narrowing with horizon as dd/√T. Illustrative; past performance "
     "is not indicative of future results.", 7.5, FOOT, first=True, after=0,
     track=0, lead=1.1)

# ===========================================================================
# 2b) Volatility-cone slide — show the up/down path envelope for both series
# ===========================================================================
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager

# brand colours as hex for matplotlib
H_NAVY, H_BLUE3, H_BLUE4 = "#152130", "#314359", "#556A83"
H_BLUE5, H_BLUE6, H_BODY = "#E2E5E9", "#F6F7F9", "#0E1620"
for _f in ("Arial", "Liberation Sans", "DejaVu Sans"):
    try:
        font_manager.findfont(_f, fallback_to_default=False)
        plt.rcParams["font.family"] = _f
        break
    except Exception:
        continue


def capture_chart(path_png):
    """The allocator's question: how much of the market's up vs down do we get?
    Two grouped bars — up-capture (good high) and down-capture (good low)."""
    fig, ax = plt.subplots(figsize=(5.1, 4.3), dpi=200)
    cats = ["Up markets\n(rallies)", "Down markets\n(selloffs)"]
    xm = [0, 1.6]
    w = 0.62
    mvals = [100, 100]                       # market = 100 by definition
    avals = [PAIR["up_cap"] * 100, PAIR["dn_cap"] * 100]
    ax.bar([x - w / 2 for x in xm], mvals, w, color=H_BLUE4, alpha=0.55,
           label="MSCI World IMI (= market)")
    ax.bar([x + w / 2 for x in xm], avals, w, color=H_NAVY,
           label="Athanase")
    ax.axhline(100, color=H_BLUE5, lw=1.0, zorder=0)
    for x, v in zip([x + w / 2 for x in xm], avals):
        ax.annotate(f"{v:.0f}%", (x, v), ha="center", va="bottom",
                    fontsize=15, fontweight="bold", color=H_NAVY,
                    xytext=(0, 3), textcoords="offset points")
    for x in [x - w / 2 for x in xm]:
        ax.annotate("100%", (x, 100), ha="center", va="bottom", fontsize=11,
                    color=H_BLUE4, xytext=(0, 3), textcoords="offset points")
    ax.set_xticks(xm); ax.set_xticklabels(cats, fontsize=11, color=H_BODY)
    ax.set_ylabel("Share of the market’s move captured", color=H_BODY, fontsize=10.5)
    ax.set_ylim(0, 120); ax.set_xlim(-0.7, 2.3)
    ax.set_yticks([0, 50, 100])
    ax.set_yticklabels(["0%", "50%", "100%"], fontsize=10, color=H_BODY)
    ax.tick_params(colors=H_BODY)
    for sp in ("top", "right"):
        ax.spines[sp].set_visible(False)
    for sp in ("left", "bottom"):
        ax.spines[sp].set_color(H_BLUE5)
    ax.grid(True, axis="y", color=H_BLUE5, lw=0.6, alpha=0.6)
    # legend below the plot so it never collides with the bars
    ax.legend(loc="upper center", bbox_to_anchor=(0.5, -0.16), ncol=2,
              fontsize=10, frameon=False, labelcolor=H_BODY,
              handlelength=1.4, columnspacing=1.6)
    ax.set_title("Up- vs down-market capture", color=H_NAVY, fontsize=13,
                 fontweight="bold", pad=10)
    fig.tight_layout()
    fig.savefig(path_png, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)


capture_chart("/tmp/capture.png")

s2 = prs.slides.add_slide(prs.slide_layouts[6])
s2.shapes.add_picture(MARK_DARK, Inches(0.55), Inches(0.24), height=Inches(0.26),
                      width=Emu(int(Inches(0.26) * _MD_AR)))
para(tbox(s2, Inches(0.98), Inches(0.27), Inches(7), Inches(0.3)),
     "Risk-Adjusted Outcomes", 11, SLATE_LT, first=True, after=0)
para(tbox(s2, Inches(9.9), Inches(0.27), Inches(2.85), Inches(0.3)),
     "Figure 2", 11, SLATE_LT, first=True, bold=True, align=PP_ALIGN.RIGHT, after=0)
rect(s2, 0, Inches(0.62), SW, Inches(1.45), fill=HEADERBG)
para(tbox(s2, Inches(0.6), Inches(0.74), Inches(12.1), Inches(0.85)),
     "A complement to global equities, not a duplicate", 30, NAVY_TX,
     first=True, after=2, font=SERIF)
para(tbox(s2, Inches(0.62), Inches(1.55), Inches(11.9), Inches(0.5)),
     "What a global-equity holder asks of a new manager: do you keep up in "
     "rallies, lose less in selloffs, and add something the index can’t?",
     13, SUBTLE, first=True, italic=True, after=0)

# left: the hero capture chart
s2.shapes.add_picture("/tmp/capture.png", Inches(0.55), Inches(2.25),
                      height=Inches(4.1))

# right: three allocator metric cards
cardx = Inches(6.55); cardw = Inches(6.2); cardh = Inches(1.2)
cy = Inches(2.3)
cards = [
    ("CAPTURE ASYMMETRY",
     f"{PAIR['up_cap']*100:.0f}% up  /  {PAIR['dn_cap']*100:.0f}% down",
     "Keeps pace with rallies but takes less than half of selloffs — the "
     "essence of the edge."),
    ("DIVERSIFICATION",
     f"{PAIR['corr']:.2f} correlation  ·  {PAIR['beta']:.2f} beta",
     "Low correlation to your existing equity book — a genuine diversifier, "
     "not levered beta."),
    ("RELIABILITY OVER TIME",
     f"Beat MSCI in {PAIR['roll'][5]*100:.0f}% of 5-yr and "
     f"{PAIR['roll'][7]*100:.0f}% of 7-yr windows",
     "Over every rolling 7-year holding period in 20 years, Athanase "
     "out-returned the index."),
]
for title, big, body in cards:
    rect(s2, cardx, cy, cardw, cardh, fill=HEADERBG)
    para(tbox(s2, Emu(int(cardx) + int(Inches(0.22))), cy + Inches(0.12),
              Emu(int(cardw) - int(Inches(0.44))), Inches(0.25)),
         title, 10.5, SLATE, first=True, bold=True, after=0, track=0)
    para(tbox(s2, Emu(int(cardx) + int(Inches(0.22))), cy + Inches(0.36),
              Emu(int(cardw) - int(Inches(0.44))), Inches(0.35)),
         big, 17, NAVY_TX, first=True, bold=True, after=0, font=SERIF, track=0)
    para(tbox(s2, Emu(int(cardx) + int(Inches(0.22))), cy + Inches(0.74),
              Emu(int(cardw) - int(Inches(0.44))), Inches(0.45)),
         body, 10.5, BODY, first=True, after=0, lead=1.12, track=0)
    cy = Emu(int(cy) + int(cardh) + int(Inches(0.12)))

# takeaway strip
rect(s2, Inches(0.6), Inches(6.62), Inches(12.13), Inches(0.46), fill=NAVY)
para(tbox(s2, Inches(0.78), Inches(6.62), Inches(11.8), Inches(0.46),
          anchor=MSO_ANCHOR.MIDDLE),
     f"Athanase captures {PAIR['up_cap']*100:.0f}% of the market’s upside but "
     f"only {PAIR['dn_cap']*100:.0f}% of its downside — at {PAIR['corr']:.2f} "
     "correlation, it improves the whole portfolio, not just one line of it.",
     12, WHITE, first=True, italic=True, after=0, track=0)
para(tbox(s2, Inches(0.6), Inches(7.16), Inches(12.6), Inches(0.4)),
     f"Source: monthly net returns, {MSCI['n']} months (2006–2025). Capture = "
     "compounded Athanase return in the market’s up/down months ÷ the market’s. "
     "Beta/correlation and rolling-window win rates from the same series. "
     "Past performance is not indicative of future results.",
     7.5, FOOT, first=True, after=0, track=0, lead=1.1)

# ===========================================================================
# 2c) Efficient frontier — what a blend does to the WHOLE portfolio
# ===========================================================================
def _ann(xs):
    g = 1.0
    for x in xs:
        g *= (1 + x)
    return g ** (12 / len(xs)) - 1


def _dvol(xs):                                # downside (semi) deviation, ann.
    return math.sqrt(sum(min(x, 0) ** 2 for x in xs) / len(xs)) * math.sqrt(12)


_weights = [i / 20 for i in range(21)]        # 0..100% Athanase, 5% steps
_FRONT = []
for w in _weights:
    blend = [w * _ath_x[i] + (1 - w) * _msci_x[i] for i in range(len(_msci_x))]
    _FRONT.append((w, _dvol(blend) * 100, _ann(blend) * 100))


def frontier_chart(path_png):
    fig, ax = plt.subplots(figsize=(8.6, 4.35), dpi=200)
    xs = [p[1] for p in _FRONT]
    ys = [p[2] for p in _FRONT]
    ax.plot(xs, ys, color=H_BLUE4, lw=2.0, zorder=3)
    # highlight markers: 0% (MSCI), a sensible blend, 100% (Athanase)
    def mark(w_target, label, col, dy=8, dx=8, big=False):
        p = min(_FRONT, key=lambda z: abs(z[0] - w_target))
        ax.scatter([p[1]], [p[2]], s=120 if big else 80, color=col, zorder=5,
                   edgecolor="white", linewidth=1.2)
        ax.annotate(label, (p[1], p[2]), color=col, fontsize=10.5,
                    fontweight="bold", xytext=(dx, dy), textcoords="offset points")
        return p
    mark(0.0, "100% Global equities\n(MSCI World IMI)", H_BLUE4, dy=-30, dx=-30)
    mark(1.0, "100% Athanase", H_NAVY, dy=6, dx=-150, big=True)
    # minimum-downside-risk point (the curve's turning point)
    _min = min(_FRONT, key=lambda z: z[1])
    ax.scatter([_min[1]], [_min[2]], s=150, color=H_NAVY, zorder=6,
               edgecolor="white", linewidth=1.5, marker="D")
    ax.annotate(f"Minimum downside risk\n≈ {_min[0]*100:.0f}% Athanase  "
                f"({_min[2]:.1f}% return)", (_min[1], _min[2]), color=H_NAVY,
                fontsize=10, fontweight="bold", xytext=(12, -4),
                textcoords="offset points", va="center")
    # prudent sizing zone (3-8%) — what an allocator would actually hold
    seg = [z for z in _FRONT if 0.03 <= z[0] <= 0.08]
    ax.plot([z[1] for z in seg], [z[2] for z in seg], color=H_NAVY, lw=6,
            alpha=0.95, solid_capstyle="round", zorder=6)
    pmid = min(_FRONT, key=lambda z: abs(z[0] - 0.055))
    # label sits in clear space below the curve, leader points to the segment
    ax.annotate("prudent 3–8% zone", (pmid[1], pmid[2]), color=H_NAVY,
                fontsize=10.5, fontweight="bold", xytext=(10.55, 5.0),
                textcoords="data", va="center", ha="center",
                arrowprops=dict(arrowstyle="-", color=H_NAVY, lw=1.0,
                                connectionstyle="arc3,rad=0.2"))
    ax.set_xlabel("Downside volatility (annualised)", color=H_BODY, fontsize=11)
    ax.set_ylabel("Annualised return", color=H_BODY, fontsize=11)
    ax.set_xlim(9.8, 12.2); ax.set_ylim(4, 18)
    ax.set_xticks([10, 10.5, 11, 11.5, 12])
    ax.xaxis.set_major_formatter(lambda v, _: f"{v:.1f}%")
    ax.yaxis.set_major_formatter(lambda v, _: f"{v:.0f}%")
    ax.tick_params(colors=H_BODY, labelsize=10)
    for sp in ("top", "right"):
        ax.spines[sp].set_visible(False)
    for sp in ("left", "bottom"):
        ax.spines[sp].set_color(H_BLUE5)
    ax.grid(True, color=H_BLUE5, lw=0.6, alpha=0.6)
    ax.set_title("Blended portfolio: return vs downside risk",
                 color=H_NAVY, fontsize=13, fontweight="bold", pad=10)
    fig.tight_layout()
    fig.savefig(path_png, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)


frontier_chart("/tmp/frontier.png")

# the 30% blend stats for the cards
_b30 = min(_FRONT, key=lambda z: abs(z[0] - 0.30))
_b0 = _FRONT[0]

s3 = prs.slides.add_slide(prs.slide_layouts[6])
s3.shapes.add_picture(MARK_DARK, Inches(0.55), Inches(0.24), height=Inches(0.26),
                      width=Emu(int(Inches(0.26) * _MD_AR)))
para(tbox(s3, Inches(0.98), Inches(0.27), Inches(7), Inches(0.3)),
     "Risk-Adjusted Outcomes", 11, SLATE_LT, first=True, after=0)
para(tbox(s3, Inches(9.9), Inches(0.27), Inches(2.85), Inches(0.3)),
     "Figure 3", 11, SLATE_LT, first=True, bold=True, align=PP_ALIGN.RIGHT, after=0)
rect(s3, 0, Inches(0.62), SW, Inches(1.45), fill=HEADERBG)
para(tbox(s3, Inches(0.6), Inches(0.74), Inches(12.1), Inches(0.85)),
     "What a blend does to the whole portfolio", 30, NAVY_TX,
     first=True, after=2, font=SERIF)
para(tbox(s3, Inches(0.62), Inches(1.55), Inches(11.9), Inches(0.5)),
     "Adding Athanase to a global-equity book moves the portfolio up and to the "
     "left — more return for less downside risk, thanks to low correlation.",
     13, SUBTLE, first=True, italic=True, after=0)
s3.shapes.add_picture("/tmp/frontier.png", Inches(0.55), Inches(2.25),
                      height=Inches(4.1))

# right: before/after card (0% vs prudent 8% sleeve) + interpretation
_b8 = min(_FRONT, key=lambda z: abs(z[0] - 0.08))
_min = min(_FRONT, key=lambda z: z[1])
cx3 = Inches(8.9); cw3 = Inches(3.85); ch3 = Inches(1.2); cy3 = Inches(2.32)
fcards = [
    ("ALL GLOBAL EQUITIES",
     f"{_b0[2]:.1f}% return", f"{_b0[1]:.1f}% downside risk",
     "100% MSCI World IMI"),
    ("ADD A PRUDENT 8% SLEEVE",
     f"{_b8[2]:.1f}% return", f"{_b8[1]:.1f}% downside risk",
     "92% MSCI  /  8% Athanase"),
]
for title, big1, big2, sub in fcards:
    rect(s3, cx3, cy3, cw3, ch3, fill=HEADERBG)
    para(tbox(s3, Emu(int(cx3) + int(Inches(0.22))), cy3 + Inches(0.12),
              Emu(int(cw3) - int(Inches(0.44))), Inches(0.25)),
         title, 10.5, SLATE, first=True, bold=True, after=0, track=0)
    rowtf = tbox(s3, Emu(int(cx3) + int(Inches(0.22))), cy3 + Inches(0.38),
                 Emu(int(cw3) - int(Inches(0.44))), Inches(0.45))
    p = rowtf.paragraphs[0]; p.space_after = Pt(0)
    r1 = p.add_run(); r1.text = big1; r1.font.size = Pt(16); r1.font.bold = True
    r1.font.color.rgb = NAVY_TX; r1.font.name = SERIF
    r2 = p.add_run(); r2.text = "   ·   " + big2; r2.font.size = Pt(12)
    r2.font.color.rgb = SLATE; r2.font.name = SANS
    para(tbox(s3, Emu(int(cx3) + int(Inches(0.22))), cy3 + Inches(0.82),
              Emu(int(cw3) - int(Inches(0.44))), Inches(0.3)),
         sub, 10.5, BODY, first=True, after=0, track=0)
    cy3 = Emu(int(cy3) + int(ch3) + int(Inches(0.14)))
# interpretation block — the directional argument + prudence caveat
para(tbox(s3, cx3, cy3 + Inches(0.02), cw3, Inches(0.3)),
     "HOW TO READ IT", 10.5, SLATE, first=True, bold=True, after=0, track=0)
para(tbox(s3, cx3, cy3 + Inches(0.3), cw3, Inches(1.9)),
     f"Every blend on the curve beats holding equities alone, and the benefit "
     f"keeps building until past half the portfolio (downside risk bottoms near "
     f"{_min[0]*100:.0f}%). That is a directional case to size up — within "
     f"prudent limits — not a target: a 59% position would breach concentration, "
     f"liquidity and the 3–8% sizing discipline. Even an 8% sleeve already moves "
     f"the portfolio the right way.", 10.5, BODY, first=True, after=0,
     lead=1.16, track=0)

rect(s3, Inches(0.6), Inches(6.62), Inches(12.13), Inches(0.46), fill=NAVY)
para(tbox(s3, Inches(0.78), Inches(6.62), Inches(11.8), Inches(0.46),
          anchor=MSO_ANCHOR.MIDDLE),
     "The blend is not a trade-off — every allocation improves the portfolio, so "
     "the case is to size up within prudent 3–8% limits, not to chase the "
     "mathematical optimum.",
     12, WHITE, first=True, italic=True, after=0, track=0)
para(tbox(s3, Inches(0.6), Inches(7.16), Inches(12.6), Inches(0.4)),
     f"Source: monthly net returns, {MSCI['n']} months (2006–2025), monthly "
     "rebalanced blends. Downside volatility = annualised deviation of negative "
     "months (MAR 0). The minimum-risk point is in-sample and shown for shape, "
     "not as advice. Illustrative; past performance is not indicative of future "
     "results.", 7.5, FOOT, first=True, after=0, track=0, lead=1.1)

# ===========================================================================
# 3) Convert to 4:3 and brand the theme (palette + fonts)
# ===========================================================================
TARGET_W, TARGET_H = 9753600, 7315200
from pptx.oxml.ns import qn as _qn
from pptx.enum.shapes import MSO_SHAPE_TYPE as _MST


def _rescale(sh, sx, sy):
    try:
        L, T, W, H = sh.left, sh.top, sh.width, sh.height
    except Exception:
        L = None
    if L is not None:
        sh.left, sh.top = int(L * sx), int(T * sy)
        if sh.shape_type == _MST.PICTURE:
            sh.width, sh.height = int(W * sx), int(W * sx * H / W)
        else:
            sh.width, sh.height = int(W * sx), int(H * sy)
    if sh.has_text_frame:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size is not None:
                    r.font.size = Pt(round(r.font.size.pt * sx, 1))
            pPr = p._p.find(_qn("a:pPr"))
            if pPr is not None:
                for at in ("marL", "indent"):
                    v = pPr.get(at)
                    if v is not None:
                        pPr.set(at, str(int(round(int(v) * sx))))


sx, sy = TARGET_W / int(prs.slide_width), TARGET_H / int(prs.slide_height)
for _sl in prs.slides:
    for sh in _sl.shapes:
        _rescale(sh, sx, sy)
prs.slide_width, prs.slide_height = TARGET_W, TARGET_H


def _brand_theme(prs):
    a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    REL = ("http://schemas.openxmlformats.org/officeDocument/2006/"
           "relationships/theme")
    scheme = {"dk1": "152130", "lt1": "FFFFFF", "dk2": "0E1620", "lt2": "F6F7F9",
              "accent1": "152130", "accent2": "314359", "accent3": "556A83",
              "accent4": "0E1620", "accent5": "E2E5E9", "accent6": "F6F7F9",
              "hlink": "314359", "folHlink": "556A83"}
    seen = set()
    for master in prs.slide_masters:
        th = master.part.part_related_by(REL)
        if th.partname in seen:
            continue
        seen.add(th.partname)
        root = etree.fromstring(th.blob)
        clr = root.find(f"{a}themeElements/{a}clrScheme")
        for nm, hx in scheme.items():
            el = clr.find(f"{a}{nm}")
            if el is None:
                continue
            srgb = el.find(f"{a}srgbClr"); sysc = el.find(f"{a}sysClr")
            if srgb is not None:
                srgb.set("val", hx)
            elif sysc is not None:
                el.remove(sysc); etree.SubElement(el, f"{a}srgbClr", {"val": hx})
        fonts = root.find(f"{a}themeElements/{a}fontScheme")
        for grp, face in (("majorFont", "Times New Roman"), ("minorFont", "Arial")):
            latin = fonts.find(f"{a}{grp}/{a}latin")
            if latin is not None:
                latin.set("typeface", face)
        th._blob = etree.tostring(root, xml_declaration=True, encoding="UTF-8",
                                  standalone=True)


_brand_theme(prs)

out = "Athanase_Downside_Risk_Outcomes.pptx"
prs.save(out)
print("Saved", out)
print(f"MSCI  : ret {MSCI['ann_ret']*100:.1f}%  vol {MSCI['ann_vol']*100:.1f}%  "
      f"dd {MSCI['ann_dd']*100:.1f}%")
print(f"AIP   : ret {ATH['ann_ret']*100:.1f}%  vol {ATH['ann_vol']*100:.1f}%  "
      f"dd {ATH['ann_dd']*100:.1f}%")

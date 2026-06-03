"""Standalone single page: DistIT case study (acting on a thesis shift; the
origin of the determined-seller framework). NOT part of the combined deck —
built for review, styled to match build_combined_deck.py exactly (incl. the
final 4:3 rescale)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn as _qn
from pptx.enum.shapes import MSO_SHAPE_TYPE as _MST
from PIL import Image

# ---- palette / fonts (AIP Blue group) -------------------------------------
BLUE1 = RGBColor(0x0E, 0x16, 0x20); BLUE2 = RGBColor(0x15, 0x21, 0x30)
BLUE3 = RGBColor(0x31, 0x43, 0x59); BLUE4 = RGBColor(0x55, 0x6A, 0x83)
BLUE5 = RGBColor(0xE2, 0xE5, 0xE9); BLUE6 = RGBColor(0xF6, 0xF7, 0xF9)
NAVY, NAVY_TX, SLATE, SLATE_LT = BLUE2, BLUE3, BLUE3, BLUE4
HEADERBG, BODY, SUBTLE, DIVIDER = BLUE6, BLUE1, BLUE4, BLUE5
WHITE = RGBColor(0xFF, 0xFF, 0xFF); FOOT = BLUE4
SERIF, SANS = "Times New Roman", "Arial"
MARK_DARK = "assets/mark_dark.png"
_MD_AR = (lambda s: s[0] / s[1])(Image.open(MARK_DARK).size)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def rect(slide, l, t, w, h, fill=None, line=None, line_w=1.0):
    sp = slide.shapes.add_shape(1, l, t, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    return sp


def tbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, size, color, bold=False, italic=False, first=False,
         align=PP_ALIGN.LEFT, after=8, lead=1.1, font=SANS, track=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(after); p.line_spacing = lead
    r = p.add_run(); r.text = text; f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = font
    if track is None:
        track = -5 if size >= 48 else (-3 if size >= 24 else 0)
    if track:
        r._r.get_or_add_rPr().set("spc", str(int(round(size * track / 100.0 * 100))))
    return p


def place_mark(slide, left, top, height):
    slide.shapes.add_picture(MARK_DARK, left, top, height=height,
                             width=Emu(int(int(height) * _MD_AR)))


def content(section, title, subtitle=None):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, fill=WHITE)
    place_mark(s, Inches(0.55), Inches(0.24), Inches(0.26))
    para(tbox(s, Inches(0.98), Inches(0.27), Inches(7.0), Inches(0.3)),
         section, 11, SLATE_LT, first=True, after=0)
    band_h = Inches(1.45) if subtitle else Inches(1.05)
    rect(s, 0, Inches(0.62), SW, band_h, fill=HEADERBG)
    para(tbox(s, Inches(0.6), Inches(0.74), Inches(12.1), Inches(0.85)),
         title, 30, NAVY_TX, first=True, after=2, font=SERIF)
    body_top = Inches(1.95)
    if subtitle:
        para(tbox(s, Inches(0.62), Inches(1.55), Inches(11.9), Inches(0.7)),
             subtitle, 13, SUBTLE, first=True, italic=True, after=0, lead=1.16)
        body_top = Inches(2.4)
    para(tbox(s, Inches(9.0), Inches(7.05), Inches(4.0), Inches(0.3)),
         "Strictly confidential          proposed", 8.5, FOOT, first=True,
         align=PP_ALIGN.RIGHT, after=0)
    return s, body_top


# ---------------------------------------------------------------------------
s, top = content(
    "Risk system · In practice",
    "The risk ledger: our system in practice — and how it developed",
    "Every engagement tests the risk system, and our mistakes refine it. Each "
    "row shows what the system flagged, the disciplined action it forced, the "
    "outcome — and what it taught.")

# ---- ledger table (rows = companies; add more rows as the page grows) -------
cols = ["Company & position", "Thesis & engagement", "The risk trigger",
        "The disciplined action", "Outcome", "What it taught the system"]
cw = [Inches(2.05), Inches(1.9), Inches(2.0), Inches(2.35), Inches(1.7),
      Inches(2.1)]
rows = [
    ["DistIT — Nasdaq Stockholm. Bought from 2015; built it with a new "
     "management team and board.",
     "Active owner with a board seat — compound a Nordic distribution platform.",
     "2021: the EFUEL acquisition we had warned against went ahead; we judged "
     "the company overvalued — thesis broken.",
     "From the board, forced the price down ~40% (SEK 300m → 180m + earn-out, "
     "which never paid). Sold across two blocks — into a rising price, not a "
     "falling one.",
     "Profitable overall — but we sold too gradually instead of exiting in full.",
     "Origin of the determined-seller framework: when the thesis breaks, exit "
     "completely — on the facts, not the price."],
    ["Robit Plc — Nasdaq Helsinki. Invested in stages from May 2015; a niche "
     "drilling-consumables franchise below intrinsic value.",
     "Board-level plan to fix the cost base and capital allocation.",
     "From 2019 we sought a board seat to formalise the plan — and could not "
     "secure board alignment.",
     "Walked away — unwound rather than stay without influence.",
     "Exited at ~14% IRR (>2× the index); the shares have since roughly halved.",
     "Reinforced agreement-before-capital: no board alignment, no position."],
]
lx = Inches(0.6); ly = top
# header
hh = Inches(0.55)
x = lx
for ci, c in enumerate(cols):
    rect(s, x, ly, cw[ci], hh, fill=SLATE)
    para(tbox(s, Emu(int(x) + int(Inches(0.1))), ly,
              Emu(int(cw[ci]) - int(Inches(0.18))), hh, anchor=MSO_ANCHOR.MIDDLE),
         c, 9, WHITE, first=True, bold=True, after=0, lead=1.04)
    x = Emu(int(x) + int(cw[ci]))
ly = Emu(int(ly) + int(hh))
# data rows
drh = Inches(1.55)
for ri, row in enumerate(rows):
    x = lx
    fill = HEADERBG if ri % 2 == 0 else WHITE
    for ci, cell in enumerate(row):
        rect(s, x, ly, cw[ci], drh, fill=fill)
        if ci == 0:
            rect(s, x, ly, Inches(0.06), drh, fill=NAVY)
        para(tbox(s, Emu(int(x) + int(Inches(0.13))), Emu(int(ly) + int(Inches(0.1))),
                  Emu(int(cw[ci]) - int(Inches(0.22))), Emu(int(drh) - int(Inches(0.18)))),
             cell, 8.5, NAVY_TX if ci == 0 else BODY, bold=(ci == 0), first=True,
             after=0, lead=1.12)
        x = Emu(int(x) + int(cw[ci]))
    ly = Emu(int(ly) + int(drh))
para(tbox(s, Inches(0.6), Emu(int(ly) + int(Inches(0.08))), Inches(12.1), Inches(0.3)),
     "Two entries shown; further cases to be added. The system is not static — "
     "Robit shaped “agreement before capital,” DistIT built the determined-seller "
     "exit.", 8, FOOT, first=True, italic=True, after=0, lead=1.1)
rect(s, Inches(0.6), Inches(6.4), Inches(12.16), Inches(0.6), fill=NAVY)
para(tbox(s, Inches(0.8), Inches(6.4), Inches(11.8), Inches(0.6),
          anchor=MSO_ANCHOR.MIDDLE),
     "A risk system we both use and keep building — every engagement tests it, "
     "and our own mistakes are written back into it.", 12, WHITE, first=True,
     italic=True, after=0, track=0)

# ---- final 4:3 rescale (match the combined deck output) --------------------
TARGET_W, TARGET_H = 9753600, 7315200


def _rescale_shape(sh, sx, sy):
    try:
        L, T, W, H = sh.left, sh.top, sh.width, sh.height
    except Exception:
        L = T = W = H = None
    is_pic = sh.shape_type == _MST.PICTURE
    if None not in (L, T, W, H):
        sh.left = int(L * sx); sh.top = int(T * sy)
        sh.width = int(W * sx)
        sh.height = int(H * sx) if is_pic else int(H * sy)
    if sh.has_text_frame:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size is not None:
                    r.font.size = Pt(round(r.font.size.pt * sx, 1))
            pPr = p._p.find(_qn("a:pPr"))
            if pPr is not None:
                for attr in ("marL", "indent"):
                    v = pPr.get(attr)
                    if v is not None:
                        pPr.set(attr, str(int(round(int(v) * sx))))


_sx = TARGET_W / int(prs.slide_width); _sy = TARGET_H / int(prs.slide_height)
for _sl in prs.slides:
    for _sh in _sl.shapes:
        _rescale_shape(_sh, _sx, _sy)
prs.slide_width = TARGET_W; prs.slide_height = TARGET_H

prs.save("Athanase_DistIT_page.pptx")
print("saved Athanase_DistIT_page.pptx")
